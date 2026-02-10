import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BLUE = RGBColor(0, 32, 96)
MEDIUM_BLUE = RGBColor(0, 70, 140)
LIGHT_BLUE = RGBColor(0, 112, 192)
ACCENT_GOLD = RGBColor(255, 192, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(240, 240, 240)

os.makedirs("figures4", exist_ok=True)

def add_slide_background(slide, color=DARK_BLUE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=None):
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = ACCENT_GOLD
        p2.alignment = PP_ALIGN.LEFT

def add_content_box(slide, left, top, width, height, items, font_size=14, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, size, bold, color = item
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
        elif item == "":
            p.text = ""
            p.font.size = Pt(6)
        else:
            prefix = "\u2022 " if bullet else ""
            p.text = f"{prefix}{item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY

def add_table(slide, left, top, width, height, data, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            if i == 0:
                from pptx.oxml.ns import qn
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '002060'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

def add_accent_bar(slide, top=Inches(1.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()

# ============================================================
# FIGURE 1: Demographic Distribution (Table 1)
# ============================================================
fig, axes = plt.subplots(1, 3, figsize=(14, 5))

age_groups = ['Below 20', '21-30', '31-40', '41-50', 'Above 50']
age_pcts = [15.8, 29.9, 25.3, 17.4, 11.6]
colors_age = ['#3498DB', '#2E86C1', '#1F618D', '#154360', '#0B2545']
axes[0].barh(age_groups, age_pcts, color=colors_age, edgecolor='black', linewidth=0.5)
for i, v in enumerate(age_pcts):
    axes[0].text(v + 0.5, i, f'{v}%', va='center', fontsize=10, fontweight='bold')
axes[0].set_xlabel('Percentage (%)', fontsize=11, fontweight='bold')
axes[0].set_title('Age Distribution', fontsize=13, fontweight='bold')
axes[0].set_xlim(0, 38)

edu_labels = ['Illiterate', 'Primary', 'Secondary', '12th', 'Graduate', 'PG+']
edu_pcts = [12.4, 18.7, 22.0, 15.8, 20.3, 10.8]
colors_edu = ['#E74C3C', '#E67E22', '#F1C40F', '#2ECC71', '#3498DB', '#9B59B6']
axes[1].barh(edu_labels, edu_pcts, color=colors_edu, edgecolor='black', linewidth=0.5)
for i, v in enumerate(edu_pcts):
    axes[1].text(v + 0.5, i, f'{v}%', va='center', fontsize=10, fontweight='bold')
axes[1].set_xlabel('Percentage (%)', fontsize=11, fontweight='bold')
axes[1].set_title('Education Level', fontsize=13, fontweight='bold')
axes[1].set_xlim(0, 30)

occ_labels = ['Farmer', 'Student', 'Self-Emp.', 'Housewife', 'Salaried', 'Other']
occ_pcts = [25.7, 20.7, 19.9, 17.4, 12.0, 4.3]
colors_occ = ['#27AE60', '#2980B9', '#E67E22', '#8E44AD', '#E74C3C', '#95A5A6']
axes[2].barh(occ_labels, occ_pcts, color=colors_occ, edgecolor='black', linewidth=0.5)
for i, v in enumerate(occ_pcts):
    axes[2].text(v + 0.5, i, f'{v}%', va='center', fontsize=10, fontweight='bold')
axes[2].set_xlabel('Percentage (%)', fontsize=11, fontweight='bold')
axes[2].set_title('Occupation', fontsize=13, fontweight='bold')
axes[2].set_xlim(0, 33)

plt.suptitle('Demographic Profile of Respondents (n=241)\n(Ahmad et al., 2026 \u2014 Table 1)',
             fontsize=14, fontweight='bold', y=1.02)
plt.tight_layout()
plt.savefig('figures4/demographics.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 2: KMO Adequacy Scale (Table 3)
# ============================================================
fig, ax = plt.subplots(figsize=(10, 5))
kmo_thresholds = [0.0, 0.5, 0.6, 0.7, 0.8, 0.9, 1.0]
kmo_labels = ['Unacceptable', 'Miserable', 'Mediocre', 'Middling', 'Meritorious', 'Marvellous']
kmo_colors = ['#E74C3C', '#E67E22', '#F1C40F', '#2ECC71', '#27AE60', '#1ABC9C']

for i in range(len(kmo_labels)):
    ax.barh(0, kmo_thresholds[i+1] - kmo_thresholds[i], left=kmo_thresholds[i],
            height=0.5, color=kmo_colors[i], edgecolor='black', linewidth=0.5)
    mid = (kmo_thresholds[i] + kmo_thresholds[i+1]) / 2
    ax.text(mid, 0, kmo_labels[i], ha='center', va='center', fontsize=9, fontweight='bold', color='white')

ax.axvline(x=0.90, color='red', linestyle='--', linewidth=3)
ax.annotate('KMO = 0.90\n(Marvellous)', xy=(0.90, 0.3), xytext=(0.90, 0.55),
            fontsize=14, fontweight='bold', color='red', ha='center',
            arrowprops=dict(arrowstyle='->', color='red', lw=2))
ax.set_xlim(0, 1.0)
ax.set_ylim(-0.5, 0.8)
ax.set_xlabel('KMO Value', fontsize=13, fontweight='bold')
ax.set_title('Kaiser-Meyer-Olkin (KMO) Sampling Adequacy\n(Ahmad et al., 2026 \u2014 Table 3)',
             fontsize=14, fontweight='bold')
ax.set_yticks([])
ax.grid(True, alpha=0.3, axis='x')
plt.tight_layout()
plt.savefig('figures4/kmo_scale.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 3: Cronbach's Alpha Reliability (Table 2)
# ============================================================
factors = ['Content &\nInformativeness', 'Trust &\nCredibility', 'Peer &\nCommunity',
           'Emotional\nEngagement', 'Convenience', 'Overall']
alphas = [0.84, 0.88, 0.86, 0.85, 0.87, 0.92]
colors_alpha = ['#3498DB', '#2E86C1', '#1F618D', '#E67E22', '#27AE60', '#E74C3C']

fig, ax = plt.subplots(figsize=(10, 6))
bars = ax.bar(factors, alphas, color=colors_alpha, edgecolor='black', linewidth=0.8, width=0.6)
ax.axhline(y=0.70, color='red', linestyle='--', linewidth=1.5, label='Acceptable Threshold (0.70)')
for bar, alpha in zip(bars, alphas):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01,
            f'{alpha:.2f}', ha='center', va='bottom', fontsize=12, fontweight='bold')
ax.set_ylabel("Cronbach's Alpha", fontsize=13, fontweight='bold')
ax.set_title("Reliability Analysis: Cronbach's Alpha by Construct\n(Ahmad et al., 2026 \u2014 Table 2, n=241)",
             fontsize=14, fontweight='bold')
ax.set_ylim(0, 1.05)
ax.legend(fontsize=11, loc='lower right')
ax.grid(True, alpha=0.3, axis='y')
plt.tight_layout()
plt.savefig('figures4/cronbach_alpha.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 4: Factor Loading Ranges (Table 4)
# ============================================================
factor_names = ['Content &\nInformativeness', 'Trust &\nCredibility', 'Peer\nInfluence',
                'Emotional\nEngagement', 'Convenience']
loading_min = [0.74, 0.81, 0.76, 0.79, 0.79]
loading_max = [0.81, 0.85, 0.80, 0.85, 0.84]
loading_mid = [(mn + mx) / 2 for mn, mx in zip(loading_min, loading_max)]
loading_err = [(mx - mn) / 2 for mn, mx in zip(loading_min, loading_max)]

fig, ax = plt.subplots(figsize=(10, 6))
colors_load = ['#3498DB', '#E74C3C', '#2ECC71', '#E67E22', '#9B59B6']
y_pos = range(len(factor_names))
bars = ax.barh(y_pos, loading_mid, xerr=loading_err, color=colors_load,
               edgecolor='black', linewidth=0.8, height=0.5, capsize=8, ecolor='black')
for i in range(len(factor_names)):
    ax.text(loading_max[i] + 0.01, i, f'{loading_min[i]:.2f}\u2013{loading_max[i]:.2f}',
            va='center', fontsize=11, fontweight='bold')
ax.axvline(x=0.50, color='red', linestyle='--', linewidth=1.5, label='Minimum Threshold (0.50)')
ax.set_yticks(y_pos)
ax.set_yticklabels(factor_names, fontsize=11)
ax.set_xlabel('Factor Loading', fontsize=13, fontweight='bold')
ax.set_title('Factor Loading Ranges by Construct (EFA Results)\n(Ahmad et al., 2026 \u2014 Table 4, All loadings > 0.74)',
             fontsize=14, fontweight='bold')
ax.set_xlim(0.4, 0.95)
ax.legend(fontsize=10, loc='lower right')
ax.grid(True, alpha=0.3, axis='x')
plt.tight_layout()
plt.savefig('figures4/factor_loadings.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 5: CFA Model Fit Indices (Table 5)
# ============================================================
fit_indices = ['CFI', 'TLI', 'RMSEA', 'SRMR']
fit_values = [0.95, 0.93, 0.05, 0.06]
fit_thresholds = [0.90, 0.90, 0.08, 0.08]
fit_labels = ['\u2265 0.90', '\u2265 0.90', '\u2264 0.08', '\u2264 0.08']

fig, axes = plt.subplots(1, 4, figsize=(14, 5))
colors_fit = ['#27AE60', '#2ECC71', '#3498DB', '#2E86C1']

for i, ax in enumerate(axes):
    ax.barh([0], [fit_values[i]], height=0.4, color=colors_fit[i],
            edgecolor='black', linewidth=0.8)
    if i < 2:
        ax.axvline(x=fit_thresholds[i], color='red', linestyle='--', linewidth=2,
                   label=f'Threshold: {fit_labels[i]}')
        ax.set_xlim(0, 1.1)
    else:
        ax.axvline(x=fit_thresholds[i], color='red', linestyle='--', linewidth=2,
                   label=f'Threshold: {fit_labels[i]}')
        ax.set_xlim(0, 0.12)
    ax.text(fit_values[i]/2, 0, f'{fit_values[i]:.2f}', ha='center', va='center',
            fontsize=16, fontweight='bold', color='white')
    ax.set_title(fit_indices[i], fontsize=14, fontweight='bold')
    ax.set_yticks([])
    ax.legend(fontsize=9, loc='upper right')
    verdict = 'PASS' if (i < 2 and fit_values[i] >= fit_thresholds[i]) or \
                        (i >= 2 and fit_values[i] <= fit_thresholds[i]) else 'FAIL'
    color_v = '#27AE60' if verdict == 'PASS' else '#E74C3C'
    ax.text(0.5, -0.15, f'\u2713 {verdict}', transform=ax.transAxes, ha='center',
            fontsize=13, fontweight='bold', color=color_v)

plt.suptitle('Confirmatory Factor Analysis: Model Fit Indices\n(Ahmad et al., 2026 \u2014 Table 5)',
             fontsize=14, fontweight='bold', y=1.05)
plt.tight_layout()
plt.savefig('figures4/cfa_fit.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 6: SEM Path Coefficients (Hypothesis Testing)
# ============================================================
hypotheses = ['H1: Content &\nInformativeness', 'H2: Trust &\nCredibility',
              'H3: Peer\nInfluence', 'H4: Emotional\nEngagement', 'H5: Convenience']
beta_values = [0.67, 0.72, 0.65, 0.61, 0.68]
colors_sem = ['#3498DB', '#E74C3C', '#2ECC71', '#E67E22', '#9B59B6']

fig, ax = plt.subplots(figsize=(10, 6))
bars = ax.barh(hypotheses, beta_values, color=colors_sem, edgecolor='black', linewidth=0.8, height=0.5)
for bar, beta in zip(bars, beta_values):
    ax.text(bar.get_width() + 0.01, bar.get_y() + bar.get_height()/2,
            f'\u03B2 = {beta:.2f}***', va='center', fontsize=12, fontweight='bold')
ax.set_xlabel('Standardized Path Coefficient (\u03B2)', fontsize=13, fontweight='bold')
ax.set_title('SEM Results: Factors Influencing Purchase Decisions\n(Ahmad et al., 2026 \u2014 All P < 0.001)',
             fontsize=14, fontweight='bold')
ax.set_xlim(0, 0.90)
ax.grid(True, alpha=0.3, axis='x')
ax.text(0.5, -0.12, '*** P < 0.001 (All hypotheses supported)',
        transform=ax.transAxes, ha='center', fontsize=11, fontstyle='italic', color='gray')
plt.tight_layout()
plt.savefig('figures4/sem_paths.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 7: Factor Importance Ranking
# ============================================================
factor_rank = ['Trust &\nCredibility', 'Convenience', 'Content &\nInformativeness',
               'Peer\nInfluence', 'Emotional\nEngagement']
beta_rank = [0.72, 0.68, 0.67, 0.65, 0.61]
colors_rank = ['#E74C3C', '#9B59B6', '#3498DB', '#2ECC71', '#E67E22']

fig, ax = plt.subplots(figsize=(10, 6))
bars = ax.barh(factor_rank, beta_rank, color=colors_rank, edgecolor='black', linewidth=0.8, height=0.5)
for i, (bar, beta) in enumerate(zip(bars, beta_rank)):
    ax.text(bar.get_width() + 0.01, bar.get_y() + bar.get_height()/2,
            f'\u03B2 = {beta:.2f}  (Rank #{len(beta_rank)-i})', va='center', fontsize=11, fontweight='bold')
ax.set_xlabel('Standardized Path Coefficient (\u03B2)', fontsize=13, fontweight='bold')
ax.set_title('Factor Importance Ranking: Strongest to Weakest Influence\n(Ahmad et al., 2026 \u2014 SEM Results)',
             fontsize=14, fontweight='bold')
ax.set_xlim(0, 0.92)
ax.grid(True, alpha=0.3, axis='x')
plt.tight_layout()
plt.savefig('figures4/factor_ranking.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 8: Digital Adoption Profile (Table 1)
# ============================================================
fig, axes = plt.subplots(1, 3, figsize=(14, 5))

phone_labels = ['Feature Phone', 'Smartphone']
phone_pcts = [37.8, 62.2]
axes[0].pie(phone_pcts, labels=phone_labels, autopct='%1.1f%%', startangle=90,
            colors=['#E67E22', '#2E86C1'], textprops={'fontsize': 12, 'fontweight': 'bold'},
            wedgeprops={'edgecolor': 'black', 'linewidth': 0.8})
axes[0].set_title('Mobile Device Type', fontsize=13, fontweight='bold')

usage_labels = ['<1 hour', '1-2 hours', '2-4 hours', '>4 hours']
usage_pcts = [19.5, 34.4, 27.8, 18.3]
axes[1].pie(usage_pcts, labels=usage_labels, autopct='%1.1f%%', startangle=90,
            colors=['#3498DB', '#2ECC71', '#E67E22', '#E74C3C'],
            textprops={'fontsize': 11, 'fontweight': 'bold'},
            wedgeprops={'edgecolor': 'black', 'linewidth': 0.8})
axes[1].set_title('Daily Internet Usage', fontsize=13, fontweight='bold')

purchase_labels = ['Yes', 'No']
purchase_pcts = [67.2, 32.8]
axes[2].pie(purchase_pcts, labels=purchase_labels, autopct='%1.1f%%', startangle=90,
            colors=['#27AE60', '#E74C3C'], textprops={'fontsize': 12, 'fontweight': 'bold'},
            wedgeprops={'edgecolor': 'black', 'linewidth': 0.8})
axes[2].set_title('Purchased via Social Media', fontsize=13, fontweight='bold')

plt.suptitle('Digital Adoption Profile of Rural Consumers (n=241)\n(Ahmad et al., 2026 \u2014 Table 1)',
             fontsize=14, fontweight='bold', y=1.02)
plt.tight_layout()
plt.savefig('figures4/digital_adoption.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 9: Bartlett's Test Visualization (Table 3)
# ============================================================
fig, ax = plt.subplots(figsize=(10, 5))
test_labels = ["Bartlett's\nChi-Square", 'Degrees of\nFreedom']
test_values = [2890.7, 300]
colors_bart = ['#2E86C1', '#E67E22']
bars = ax.bar(test_labels, test_values, color=colors_bart, edgecolor='black',
              linewidth=0.8, width=0.4)
for bar, val in zip(bars, test_values):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 30,
            f'{val:,.1f}' if isinstance(val, float) else f'{val}',
            ha='center', va='bottom', fontsize=16, fontweight='bold')
ax.set_ylabel('Value', fontsize=13, fontweight='bold')
ax.set_title("Bartlett's Test of Sphericity\n(\u03C7\u00B2 = 2890.7, df = 300, P < 0.001)",
             fontsize=14, fontweight='bold')
ax.grid(True, alpha=0.3, axis='y')

ax.text(0.95, 0.95, 'P < 0.001\n\u2713 SIGNIFICANT', transform=ax.transAxes,
        ha='right', va='top', fontsize=14, fontweight='bold', color='#27AE60',
        bbox=dict(boxstyle='round,pad=0.5', facecolor='#D5F5E3', edgecolor='#27AE60'))
plt.tight_layout()
plt.savefig('figures4/bartlett_test.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 10: Reliability & SEM Path Coefficients Radar
# ============================================================
categories = ['Content &\nInform.', 'Trust &\nCredib.', 'Peer &\nCommunity',
              'Emotional\nEngag.', 'Convenience']
alpha_vals = [0.84, 0.88, 0.86, 0.85, 0.87]
beta_vals = [0.67, 0.72, 0.65, 0.61, 0.68]

N = len(categories)
angles = [n / float(N) * 2 * np.pi for n in range(N)]
angles += angles[:1]
alpha_vals_r = alpha_vals + [alpha_vals[0]]
beta_vals_r = beta_vals + [beta_vals[0]]

fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))
ax.plot(angles, alpha_vals_r, 'o-', linewidth=2.5, color='#E74C3C', label="Cronbach's Alpha (Table 2)", markersize=8)
ax.fill(angles, alpha_vals_r, alpha=0.15, color='#E74C3C')
ax.plot(angles, beta_vals_r, 's-', linewidth=2.5, color='#2E86C1', label='SEM Path Coeff. \u03B2 (SEM Results)', markersize=8)
ax.fill(angles, beta_vals_r, alpha=0.15, color='#2E86C1')

ax.set_xticks(angles[:-1])
ax.set_xticklabels(categories, fontsize=10, fontweight='bold')
ax.set_ylim(0.4, 1.0)
ax.set_title('Reliability (\u03B1) vs. Influence Strength (\u03B2) by Construct\n(Ahmad et al., 2026 \u2014 Tables 2 & SEM Results)',
             fontsize=14, fontweight='bold', pad=20)
ax.legend(loc='lower right', fontsize=11, bbox_to_anchor=(1.2, -0.05))
plt.tight_layout()
plt.savefig('figures4/radar_summary.png', dpi=200, bbox_inches='tight')
plt.close()

print("All 10 figures generated successfully in figures4/\n")

# ============================================================
# CREATE PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SLIDE 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()
add_content_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(1.5), [
    ("SSIE-605: Applied Multivariate Data Analysis", 18, False, ACCENT_GOLD)
])
add_content_box(slide, Inches(0.5), Inches(2.0), Inches(9), Inches(2.5), [
    ("Factor Analysis:", 36, True, WHITE),
    ("", 6, False, WHITE),
    ("Social Media Influence on Consumer", 32, True, WHITE),
    ("Purchase Decisions", 32, True, WHITE)
])
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.7), Inches(3), Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GOLD
bar.line.fill.background()
add_content_box(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(2.0), [
    ("A Factor Analytical Study of Rural Consumers in Uttar Pradesh", 16, False, RGBColor(180, 198, 231)),
    ("Based on: Ahmad et al. (2026) \u2014 Intl. Review of Management & Marketing", 14, False, RGBColor(180, 198, 231)),
    ("", 6, False, WHITE),
    ("Lana Jalal Gidan", 18, True, WHITE),
    ("Binghamton University | Professor Susan Lu", 14, False, RGBColor(180, 198, 231))
])

# SLIDE 2: Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Presentation Agenda")
add_accent_bar(slide)
agenda_items = [
    ("1. Introduction to Factor Analysis", 15, True, DARK_BLUE),
    ("2. Why Factor Analysis in Marketing?", 15, False, DARK_GRAY),
    ("3. Study Overview: Social Media & Purchase Decisions", 15, True, DARK_BLUE),
    ("4. Research Methodology & Likert Scale Design", 15, False, DARK_GRAY),
    ("5. Demographic Profile of Respondents", 15, True, DARK_BLUE),
    ("6. Step 1: Data Adequacy \u2014 KMO & Bartlett's Test", 15, False, DARK_GRAY),
    ("7. Step 2: Reliability Analysis (Cronbach's Alpha)", 15, True, DARK_BLUE),
    ("8. Step 3: Exploratory Factor Analysis (EFA)", 15, False, DARK_GRAY),
    ("9. Step 4: Confirmatory Factor Analysis (CFA)", 15, True, DARK_BLUE),
    ("10. Step 5: Structural Equation Modeling (SEM)", 15, False, DARK_GRAY),
    ("11. Key Findings & Marketing Implications", 15, True, DARK_BLUE),
    ("12. Strengths, Limitations & Conclusion", 15, False, DARK_GRAY),
]
add_content_box(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5), agenda_items)

# SLIDE 3: What is Factor Analysis?
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "What is Factor Analysis?", "Simplifying Complex Survey Data")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Definition:", 16, True, DARK_BLUE),
    ("Factor Analysis is a statistical method that reduces many survey questions into a", 14, False, DARK_GRAY),
    ("smaller number of meaningful groups (called 'factors' or 'dimensions').", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("In Simple Terms:", 16, True, DARK_BLUE),
    ("\u2022 Imagine you ask consumers 25 questions about social media shopping", 14, False, DARK_GRAY),
    ("\u2022 Factor Analysis finds that these 25 questions really measure just 5 big themes", 14, False, DARK_GRAY),
    ("\u2022 Each theme groups together questions that consumers answer in similar ways", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Two Main Types:", 16, True, DARK_BLUE),
    ("\u2022 Exploratory Factor Analysis (EFA) \u2014 discovers hidden patterns in data", 14, False, DARK_GRAY),
    ("\u2022 Confirmatory Factor Analysis (CFA) \u2014 tests if those patterns hold up", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Key Steps:", 16, True, DARK_BLUE),
    ("\u2022 Check if data is suitable (KMO & Bartlett's Test)", 14, False, DARK_GRAY),
    ("\u2022 Extract factors (using eigenvalue > 1 criterion)", 14, False, DARK_GRAY),
    ("\u2022 Rotate factors for clearer interpretation (Varimax rotation)", 14, False, DARK_GRAY),
    ("\u2022 Validate the structure (CFA / SEM)", 14, False, DARK_GRAY),
])

# SLIDE 4: Why FA in Marketing?
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Why Factor Analysis in Marketing?", "Applications in Consumer Analytics")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Marketing researchers use Factor Analysis to:", 15, True, DARK_BLUE),
    ("", 6, False, DARK_GRAY),
    ("\u2022 Identify what really drives consumer purchase decisions", 14, False, DARK_GRAY),
    ("\u2022 Understand brand perception dimensions", 14, False, DARK_GRAY),
    ("\u2022 Measure customer satisfaction across multiple touchpoints", 14, False, DARK_GRAY),
    ("\u2022 Evaluate the effectiveness of marketing channels", 14, False, DARK_GRAY),
    ("\u2022 Build reliable measurement scales for surveys", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Why Questionnaires & Likert Scales?", 15, True, DARK_BLUE),
    ("", 6, False, DARK_GRAY),
    ("\u2022 Likert scales (e.g., 1 = Strongly Disagree to 5 = Strongly Agree) capture", 14, False, DARK_GRAY),
    ("  consumer attitudes in a structured, measurable way", 14, False, DARK_GRAY),
    ("\u2022 Factor Analysis groups correlated Likert items into meaningful dimensions", 14, False, DARK_GRAY),
    ("\u2022 This helps marketers focus on key drivers rather than individual questions", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Today's Study: How does social media influence rural consumer purchase", 15, True, MEDIUM_BLUE),
    ("decisions? Which social media roles matter most?", 15, True, MEDIUM_BLUE),
])

# SLIDE 5: Study Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Study Overview", "Ahmad et al. (2026) \u2014 Intl. Review of Management & Marketing")
add_accent_bar(slide)
data = [
    ["Aspect", "Details"],
    ["Title", "Identifying Determinants of Social Media Influence\non Purchase Decisions"],
    ["Authors", "Ahmad, A., Singh, N.K., Ahmad, S., Fatima, R.,\nAlam, A. & Ghani, U."],
    ["Journal", "International Review of Management & Marketing,\nVol. 16, No. 2 (2026)"],
    ["DOI", "10.32479/irmm.21497"],
    ["Sample", "n = 241 rural consumers across 5 districts\nof Uttar Pradesh, India"],
    ["Instrument", "25-item structured questionnaire,\n5-point Likert scale (1=Strongly Disagree to 5=Strongly Agree)"],
    ["Methods", "EFA (PCA + Varimax), CFA, SEM"],
    ["Software", "SPSS and AMOS"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5), data, col_widths=[2.0, 7.4])

# SLIDE 6: Research Context
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Research Context", "Social Media and Rural Consumer Markets")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Why Study Social Media's Role in Rural Markets?", 16, True, DARK_BLUE),
    ("", 6, False, DARK_GRAY),
    ("\u2022 Rural India represents ~65% of the total population (Census of India, 2011)", 14, False, DARK_GRAY),
    ("\u2022 Rapid smartphone adoption: 62.2% of rural respondents now use smartphones", 14, False, DARK_GRAY),
    ("\u2022 67.2% of respondents have already purchased products via social media", 14, False, DARK_GRAY),
    ("\u2022 Social media is replacing traditional word-of-mouth in rural areas", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Research Gap:", 16, True, DARK_BLUE),
    ("\u2022 Most studies focus on urban consumers; rural behavior is under-researched", 14, False, DARK_GRAY),
    ("\u2022 The multidimensional nature of social media influence (trust, emotions, peer", 14, False, DARK_GRAY),
    ("  recommendations, content quality) is rarely explored in rural contexts", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Research Objective:", 16, True, DARK_BLUE),
    ("\u2022 Identify the key factors (dimensions) of social media that influence rural", 14, False, DARK_GRAY),
    ("  consumers' purchase decisions using Exploratory Factor Analysis", 14, False, DARK_GRAY),
    ("\u2022 Validate the factor structure using CFA and test causal relationships via SEM", 14, False, DARK_GRAY),
])

# SLIDE 7: Five Constructs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Five Constructs Measured", "25 Likert-Scale Items Across 5 Dimensions")
add_accent_bar(slide)
data = [
    ["Factor", "Items", "\u03B1", "What It Measures"],
    ["Content &\nInformativeness", "5", "0.84", "Product info quality, updates,\ncomparisons on social media"],
    ["Trust &\nCredibility", "5", "0.88", "Influencer recommendations,\nverified accounts, real reviews"],
    ["Peer &\nCommunity", "5", "0.86", "Family/friend recommendations,\nshares, community feedback"],
    ["Emotional\nEngagement", "5", "0.85", "Entertainment, emotional storytelling,\nregional content appeal"],
    ["Convenience", "5", "0.87", "Ease of buying, click-to-buy,\nreduced decision time"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.0), data, col_widths=[2.0, 1.0, 1.0, 5.4])
add_content_box(slide, Inches(0.5), Inches(6.6), Inches(9), Inches(0.5), [
    ("Overall Cronbach's Alpha = 0.92 | All constructs adapted from validated scales (Table 2)", 12, False, DARK_GRAY)
])

# SLIDE 8: Likert Scale Design
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Questionnaire Design", "5-Point Likert Scale Structure")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("5-Point Likert Scale Used:", 16, True, DARK_BLUE),
    ("", 6, False, DARK_GRAY),
    ("  1 = Strongly Disagree    2 = Disagree    3 = Neutral    4 = Agree    5 = Strongly Agree", 15, True, MEDIUM_BLUE),
    ("", 6, False, DARK_GRAY),
    ("Questionnaire Structure:", 16, True, DARK_BLUE),
    ("\u2022 Part 1: Demographic information (age, gender, education, occupation, income,", 14, False, DARK_GRAY),
    ("  mobile type, internet usage, district)", 14, False, DARK_GRAY),
    ("\u2022 Part 2: 25 scale-based items measuring five social media influence constructs", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Sampling Design:", 16, True, DARK_BLUE),
    ("\u2022 Stratified random sampling + purposive sampling", 14, False, DARK_GRAY),
    ("\u2022 5 districts: Lucknow, Kanpur, Agra, Jhansi, Gorakhpur", 14, False, DARK_GRAY),
    ("\u2022 241 valid responses from rural consumers with social media experience", 14, False, DARK_GRAY),
    ("\u2022 Both offline paper-based and online survey methods used", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Items adapted from validated scales:", 14, True, DARK_BLUE),
    ("\u2022 Bhatt & Joshi (2021), Kapoor & Dwivedi (2022), and others", 14, False, DARK_GRAY),
])

# SLIDE 9: Demographics
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Demographic Profile", "Sample Characteristics (Table 1, n=241)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/demographics.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(5.0))
add_content_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.5), [
    ("Majority: Age 21-30 (29.9%) | Secondary+ education (58.1%) | Diverse occupations (Table 1)", 12, False, DARK_GRAY)
])

# SLIDE 10: Digital Adoption
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Digital Adoption Profile", "Technology Usage Among Rural Consumers (Table 1)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/digital_adoption.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(5.0))
add_content_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.5), [
    ("62.2% use smartphones | 80%+ spend >1 hour online daily | 67.2% have purchased via social media", 12, False, DARK_GRAY)
])

# SLIDE 11: KMO Test
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 1: KMO Sampling Adequacy", "Is the Data Suitable for Factor Analysis? (Table 3)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/kmo_scale.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(4.5))
add_content_box(slide, Inches(0.5), Inches(6.0), Inches(9), Inches(1.2), [
    ("KMO = 0.90 (Classified as 'Marvellous' by Hair et al., 2019)", 14, True, DARK_BLUE),
    ("\u2022 Well above the 0.60 minimum \u2192 Excellent sampling adequacy for factor analysis", 13, False, DARK_GRAY),
    ("\u2022 This means the survey items share enough common patterns to extract meaningful factors", 13, False, DARK_GRAY),
])

# SLIDE 12: Bartlett's Test
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 1: Bartlett's Test of Sphericity", "Testing for Sufficient Correlations (Table 3)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/bartlett_test.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(4.5))
add_content_box(slide, Inches(0.5), Inches(6.0), Inches(9), Inches(1.2), [
    ("Bartlett's Test: \u03C7\u00B2 = 2890.7, df = 300, P < 0.001 (Table 3)", 14, True, DARK_BLUE),
    ("\u2022 Highly significant \u2192 The variables are sufficiently correlated for factor analysis", 13, False, DARK_GRAY),
    ("\u2022 Rejects the null hypothesis that the correlation matrix is an identity matrix", 13, False, DARK_GRAY),
])

# SLIDE 13: KMO + Bartlett Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Data Adequacy Summary", "Both Tests Confirm Suitability for Factor Analysis")
add_accent_bar(slide)
data = [
    ["Test", "Statistic", "Value", "Threshold", "Result"],
    ["KMO", "Sampling Adequacy", "0.90", "\u2265 0.60", "\u2713 Marvellous"],
    ["Bartlett's", "Chi-Square (\u03C7\u00B2)", "2890.7", "P < 0.05", "\u2713 P < 0.001"],
    ["Bartlett's", "Degrees of Freedom", "300", "\u2014", "\u2014"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.5), data, col_widths=[1.5, 2.2, 1.5, 1.8, 2.0])
add_content_box(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(3.0), [
    ("Interpretation:", 16, True, DARK_BLUE),
    ("", 6, False, DARK_GRAY),
    ("\u2022 KMO = 0.90 is classified as 'Marvellous' (Hair et al., 2019)", 14, False, DARK_GRAY),
    ("\u2022 This is the highest rating on the Kaiser adequacy scale", 14, False, DARK_GRAY),
    ("\u2022 Bartlett's test confirms that the 25 survey items are sufficiently inter-correlated", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("\u2022 Both tests strongly support proceeding with Factor Analysis", 14, True, MEDIUM_BLUE),
    ("\u2022 The data patterns are meaningful, not random noise", 14, False, DARK_GRAY),
])

# SLIDE 14: Cronbach's Alpha
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 2: Reliability Analysis", "Cronbach's Alpha by Construct (Table 2)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/cronbach_alpha.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(5.0))
add_content_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.5), [
    ("All constructs exceed 0.70 threshold (Nunnally, 1978) | Overall \u03B1 = 0.92 (Excellent)", 12, False, DARK_GRAY)
])

# SLIDE 15: Reliability Interpretation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Reliability Interpretation", "What Cronbach's Alpha Tells Us (Table 2)")
add_accent_bar(slide)
data_rel = [
    ["Construct", "Items", "\u03B1", "Rating"],
    ["Content & Informativeness", "5", "0.84", "Good"],
    ["Trust & Credibility", "5", "0.88", "Very Good"],
    ["Peer & Community", "5", "0.86", "Good"],
    ["Emotional Engagement", "5", "0.85", "Good"],
    ["Convenience", "5", "0.87", "Very Good"],
    ["Overall Scale", "25", "0.92", "Excellent"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.8), data_rel, col_widths=[3.5, 1.2, 1.2, 3.1])
add_content_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(1.5), [
    ("What This Means:", 15, True, DARK_BLUE),
    ("\u2022 Reliability = how consistently the questions measure the same concept", 13, False, DARK_GRAY),
    ("\u2022 \u03B1 \u2265 0.70 = acceptable | \u03B1 \u2265 0.80 = good | \u03B1 \u2265 0.90 = excellent", 13, False, DARK_GRAY),
    ("\u2022 All five factors and the overall scale pass reliability checks \u2713", 13, False, DARK_GRAY),
])

# SLIDE 16: EFA Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 3: Exploratory Factor Analysis", "Extracting the Underlying Factors (Table 4)")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("EFA Method:", 16, True, DARK_BLUE),
    ("\u2022 Extraction: Principal Component Analysis (PCA)", 14, False, DARK_GRAY),
    ("\u2022 Rotation: Varimax with Kaiser Normalization", 14, False, DARK_GRAY),
    ("\u2022 Criterion: Eigenvalue > 1 (Kaiser Criterion)", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Results:", 16, True, DARK_BLUE),
    ("\u2022 Five latent factors were extracted from 25 questionnaire items", 14, False, DARK_GRAY),
    ("\u2022 All factor loadings are above 0.74 (well above the 0.50 threshold)", 14, False, DARK_GRAY),
    ("\u2022 Strong construct validity confirmed", 14, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Factor Loading Ranges (from Table 4):", 16, True, DARK_BLUE),
    ("\u2022 Content & Informativeness: 0.74 \u2013 0.81", 14, False, DARK_GRAY),
    ("\u2022 Trust & Credibility: 0.81 \u2013 0.85 (strongest loadings)", 14, False, DARK_GRAY),
    ("\u2022 Peer Influence: 0.76 \u2013 0.80", 14, False, DARK_GRAY),
    ("\u2022 Emotional Engagement: 0.79 \u2013 0.85", 14, False, DARK_GRAY),
    ("\u2022 Convenience: 0.79 \u2013 0.84", 14, False, DARK_GRAY),
])

# SLIDE 17: Factor Loading Visualization
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Loading Ranges", "EFA Rotated Component Matrix Summary (Table 4)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/factor_loadings.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(5.0))
add_content_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.5), [
    ("All 25 items load above 0.74 | Well above 0.50 minimum | Trust & Credibility has the strongest loadings", 12, False, DARK_GRAY)
])

# SLIDE 18: EFA Interpretation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Interpreting the Five Factors", "What Each Factor Represents")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Factor 1: Content & Informativeness (Loadings: 0.74\u20130.81)", 15, True, RGBColor(52, 152, 219)),
    ("\u2022 Product details, price comparisons, new launch updates on social media", 13, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Factor 2: Trust & Credibility (Loadings: 0.81\u20130.85)", 15, True, RGBColor(231, 76, 60)),
    ("\u2022 Verified accounts, influencer recommendations, authentic user reviews", 13, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Factor 3: Peer Influence (Loadings: 0.76\u20130.80)", 15, True, RGBColor(46, 204, 113)),
    ("\u2022 Family/friend opinions, shares, likes, community-based recommendations", 13, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Factor 4: Emotional Engagement (Loadings: 0.79\u20130.85)", 15, True, RGBColor(230, 126, 34)),
    ("\u2022 Entertaining ads, emotional storytelling, culturally relevant regional content", 13, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("Factor 5: Convenience (Loadings: 0.79\u20130.84)", 15, True, RGBColor(155, 89, 182)),
    ("\u2022 Easy access, click-to-buy options, reduced shopping time, direct purchase links", 13, False, DARK_GRAY),
])

# SLIDE 19: CFA Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 4: Confirmatory Factor Analysis", "Validating the Five-Factor Structure (Table 5)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/cfa_fit.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(4.5))
add_content_box(slide, Inches(0.5), Inches(6.0), Inches(9), Inches(1.2), [
    ("All four fit indices meet or exceed recommended thresholds (Byrne, 2016)", 13, True, DARK_BLUE),
    ("\u2022 CFI=0.95, TLI=0.93 (\u22650.90) | RMSEA=0.05, SRMR=0.06 (\u22640.08) \u2192 Excellent model fit", 13, False, DARK_GRAY),
])

# SLIDE 20: CFA Interpretation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "CFA Model Fit Interpretation", "Confirming Convergent & Discriminant Validity")
add_accent_bar(slide)
data_cfa = [
    ["Index", "Value", "Threshold", "Interpretation"],
    ["CFI", "0.95", "\u2265 0.90", "Excellent fit"],
    ["TLI", "0.93", "\u2265 0.90", "Good fit"],
    ["RMSEA", "0.05", "\u2264 0.08", "Close fit (minimal error)"],
    ["SRMR", "0.06", "\u2264 0.08", "Acceptable residuals"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.0), data_cfa, col_widths=[1.5, 1.5, 2.0, 4.0])
add_content_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(2.2), [
    ("What This Means:", 16, True, DARK_BLUE),
    ("\u2022 The five-factor structure identified by EFA is statistically confirmed", 14, False, DARK_GRAY),
    ("\u2022 Convergent validity: Items within each factor measure the same concept", 14, False, DARK_GRAY),
    ("\u2022 Discriminant validity: The five factors are distinct from each other", 14, False, DARK_GRAY),
    ("\u2022 The model accurately represents how social media influences purchase decisions", 14, False, DARK_GRAY),
])

# SLIDE 21: SEM Path Coefficients
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 5: Structural Equation Modeling", "Testing Causal Relationships (SEM Results)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/sem_paths.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(5.0))
add_content_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.5), [
    ("All five hypotheses supported | All paths significant at P < 0.001 | Trust has strongest effect (\u03B2=0.72)", 12, False, DARK_GRAY)
])

# SLIDE 22: Hypothesis Testing
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Hypothesis Testing Results", "All Five Hypotheses Supported")
add_accent_bar(slide)
data_hyp = [
    ["Hypothesis", "Path", "\u03B2", "P-value", "Result"],
    ["H1", "Content & Inform. \u2192 Purchase", "0.67", "< 0.001", "\u2713 Supported"],
    ["H2", "Trust & Credibility \u2192 Purchase", "0.72", "< 0.001", "\u2713 Supported"],
    ["H3", "Peer Influence \u2192 Purchase", "0.65", "< 0.001", "\u2713 Supported"],
    ["H4", "Emotional Engage. \u2192 Purchase", "0.61", "< 0.001", "\u2713 Supported"],
    ["H5", "Convenience \u2192 Purchase", "0.68", "< 0.001", "\u2713 Supported"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.5), data_hyp, col_widths=[1.0, 3.4, 1.0, 1.5, 2.5])
add_content_box(slide, Inches(0.5), Inches(5.3), Inches(9), Inches(1.8), [
    ("Key Takeaway:", 15, True, DARK_BLUE),
    ("\u2022 All five social media dimensions significantly influence purchase decisions", 14, False, DARK_GRAY),
    ("\u2022 Trust & Credibility is the strongest predictor (\u03B2 = 0.72)", 14, False, DARK_GRAY),
    ("\u2022 Emotional Engagement, while significant, has the smallest effect (\u03B2 = 0.61)", 14, False, DARK_GRAY),
    ("\u2022 These findings support the Theory of Reasoned Action (Ajzen, 1991)", 14, False, DARK_GRAY),
])

# SLIDE 23: Factor Ranking
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Importance Ranking", "Which Social Media Roles Matter Most?")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/factor_ranking.png',
                         Inches(0.2), Inches(1.4), Inches(9.6), Inches(5.0))
add_content_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.5), [
    ("Trust (\u03B2=0.72) > Convenience (\u03B2=0.68) > Content (\u03B2=0.67) > Peer (\u03B2=0.65) > Emotional (\u03B2=0.61)", 12, False, DARK_GRAY)
])

# SLIDE 24: Radar Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Reliability & Validity Overview", "Combined View of All Constructs (Tables 2 & 4)")
add_accent_bar(slide)
slide.shapes.add_picture('figures4/radar_summary.png',
                         Inches(1.0), Inches(1.3), Inches(8.0), Inches(5.8))
add_content_box(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.3), [
    ("All constructs show strong reliability (\u03B1 = 0.84\u20130.88) and significant SEM influence (\u03B2 = 0.61\u20130.72)", 11, False, DARK_GRAY)
])

# SLIDE 25: Marketing Implications
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Marketing Implications", "Practical Takeaways for Digital Strategists")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Based on the factor analysis results, marketers should:", 15, True, DARK_BLUE),
    ("", 6, False, DARK_GRAY),
    ("1. Build Trust First (\u03B2 = 0.72 \u2014 Strongest Factor)", 15, True, RGBColor(231, 76, 60)),
    ("\u2022 Use verified accounts, authentic reviews, and credible influencers", 13, False, DARK_GRAY),
    ("\u2022 Trust is the single most important driver of purchase decisions", 13, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("2. Maximize Convenience (\u03B2 = 0.68)", 15, True, RGBColor(155, 89, 182)),
    ("\u2022 Simplify the path from social media content to actual purchase", 13, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("3. Invest in Quality Content (\u03B2 = 0.67)", 15, True, RGBColor(52, 152, 219)),
    ("\u2022 Provide clear product information, comparisons, and updates", 13, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("4. Leverage Peer Networks (\u03B2 = 0.65)", 15, True, RGBColor(46, 204, 113)),
    ("\u2022 Encourage sharing, community reviews, and family-oriented campaigns", 13, False, DARK_GRAY),
    ("", 6, False, DARK_GRAY),
    ("5. Create Emotional Connections (\u03B2 = 0.61)", 15, True, RGBColor(230, 126, 34)),
    ("\u2022 Use culturally localized, emotionally engaging content and regional language", 13, False, DARK_GRAY),
])

# SLIDE 26: Strengths & Limitations
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Strengths & Limitations", "Evaluating the Study")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), [
    ("Strengths:", 16, True, RGBColor(39, 174, 96)),
    ("", 6, False, DARK_GRAY),
    ("\u2022 Strong KMO (0.90) and highly", 13, False, DARK_GRAY),
    ("  significant Bartlett's test", 13, False, DARK_GRAY),
    ("\u2022 Excellent overall reliability", 13, False, DARK_GRAY),
    ("  (Cronbach's \u03B1 = 0.92)", 13, False, DARK_GRAY),
    ("\u2022 Both EFA and CFA used", 13, False, DARK_GRAY),
    ("  (comprehensive validation)", 13, False, DARK_GRAY),
    ("\u2022 All CFA fit indices pass", 13, False, DARK_GRAY),
    ("\u2022 SEM confirms all 5 hypotheses", 13, False, DARK_GRAY),
    ("\u2022 Addresses under-researched", 13, False, DARK_GRAY),
    ("  rural consumer segment", 13, False, DARK_GRAY),
])
add_content_box(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5), [
    ("Limitations:", 16, True, RGBColor(231, 76, 60)),
    ("", 6, False, DARK_GRAY),
    ("\u2022 Geographically limited to 5", 13, False, DARK_GRAY),
    ("  districts of Uttar Pradesh", 13, False, DARK_GRAY),
    ("\u2022 Cross-sectional design (single", 13, False, DARK_GRAY),
    ("  time point, no trends)", 13, False, DARK_GRAY),
    ("\u2022 Moderate sample size (n=241)", 13, False, DARK_GRAY),
    ("\u2022 Self-reported data may have", 13, False, DARK_GRAY),
    ("  social desirability bias", 13, False, DARK_GRAY),
    ("\u2022 Individual eigenvalues and", 13, False, DARK_GRAY),
    ("  total variance not reported", 13, False, DARK_GRAY),
    ("  in the published tables", 13, False, DARK_GRAY),
])

# SLIDE 27: References & Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()
add_content_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8), [
    ("References", 28, True, WHITE)
])
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(3), Inches(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GOLD
bar.line.fill.background()
add_content_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(4.5), [
    ("1. Ahmad, A., Singh, N.K., Ahmad, S., Fatima, R., Alam, A. & Ghani, U. (2026).", 12, False, RGBColor(180, 198, 231)),
    ("   Identifying Determinants of Social Media Influence on Purchase Decisions:", 12, False, RGBColor(180, 198, 231)),
    ("   A Factor Analytical Study of Rural Consumers. IRMM, 16(2), 111-119.", 12, False, RGBColor(180, 198, 231)),
    ("   DOI: 10.32479/irmm.21497", 12, True, ACCENT_GOLD),
    ("", 6, False, WHITE),
    ("2. Hair, J.F., Black, W.C., Babin, B.J. & Anderson, R.E. (2019).", 12, False, RGBColor(180, 198, 231)),
    ("   Multivariate Data Analysis (8th Ed.). Cengage Learning.", 12, False, RGBColor(180, 198, 231)),
    ("", 6, False, WHITE),
    ("3. Nunnally, J.C. (1978). Psychometric Theory (2nd Ed.).", 12, False, RGBColor(180, 198, 231)),
    ("   McGraw-Hill, New York.", 12, False, RGBColor(180, 198, 231)),
])
add_content_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(1.8), [
    ("Thank You!", 32, True, ACCENT_GOLD),
    ("", 6, False, WHITE),
    ("Lana Jalal Gidan", 18, True, WHITE),
    ("Binghamton University | SSIE-605 | Professor Susan Lu", 14, False, RGBColor(180, 198, 231)),
])

# Save
prs.save('Factor_Analysis_SocialMedia_Presentation.pptx')
print(f"\nPresentation saved as: Factor_Analysis_SocialMedia_Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
print("All data sourced from: Ahmad et al. (2026), DOI: 10.32479/irmm.21497")
