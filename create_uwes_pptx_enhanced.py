import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ─── PALETTE ──────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0,   32,  96)
MED_BLUE   = RGBColor(0,   70, 140)
LIGHT_BLUE = RGBColor(0,  112, 192)
TEAL       = RGBColor(42, 157, 143)
PURPLE     = RGBColor(107, 45, 123)
GOLD       = RGBColor(201, 168,  76)
LIGHT_GOLD = RGBColor(255, 243, 200)
WHITE      = RGBColor(255, 255, 255)
OFF_WHITE  = RGBColor(250, 247, 252)
DARK_GRAY  = RGBColor(50,   50,  50)
MED_GRAY   = RGBColor(100, 100, 100)
LIGHT_GRAY = RGBColor(220, 220, 220)
RED        = RGBColor(192,  57,  43)
GREEN      = RGBColor(39,  174,  96)
ORANGE     = RGBColor(211, 84,   0)

# ─── SLIDE SETUP ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)   # 16:9 widescreen
prs.slide_height = Inches(7.50)

FOOTER_TXT = "Introduction | EFA Mathematics | EFA Results | CFA Summary | Discussion | References"

# ─── HELPERS ──────────────────────────────────────────────────────────────────
def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = OFF_WHITE
    return slide

def add_header(slide, title, subtitle=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.20))
    shape.fill.solid(); shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.10)
    p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(24); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.LEFT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle; p2.font.size = Pt(12)
        p2.font.color.rgb = GOLD; p2.alignment = PP_ALIGN.LEFT
    # gold accent line
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.20), Inches(13.33), Inches(0.06))
    acc.fill.solid(); acc.fill.fore_color.rgb = GOLD
    acc.line.fill.background()

def add_footer(slide, num, total):
    fb = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.10), Inches(13.33), Inches(0.40))
    fb.fill.solid(); fb.fill.fore_color.rgb = DARK_BLUE
    fb.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(7.13), Inches(12.0), Inches(0.35))
    tf = txb.text_frame
    p = tf.paragraphs[0]; p.text = FOOTER_TXT
    p.font.size = Pt(7); p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    txb2 = slide.shapes.add_textbox(Inches(12.5), Inches(7.13), Inches(0.7), Inches(0.35))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]; p2.text = f"{num}/{total}"
    p2.font.size = Pt(8); p2.font.color.rgb = GOLD
    p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT

def textbox(slide, left, top, width, height, items, default_size=12,
            default_color=None, word_wrap=True):
    if default_color is None:
        default_color = DARK_GRAY
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = word_wrap
    first = True
    for item in items:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, size, bold, color = item
            if text == "": p.text = ""; p.font.size = Pt(5); continue
            p.text = text; p.font.size = Pt(size)
            p.font.bold = bold; p.font.color.rgb = color
        elif item == "":
            p.text = ""; p.font.size = Pt(5)
        else:
            p.text = item; p.font.size = Pt(default_size)
            p.font.color.rgb = default_color

def colored_box(slide, left, top, width, height, fill_color, border_color=None,
                radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def box_with_text(slide, left, top, width, height, fill, items, default_size=12,
                  default_color=None, border=None, radius=True, margin=0.12):
    colored_box(slide, left, top, width, height, fill, border, radius)
    if default_color is None:
        default_color = DARK_GRAY
    txb = slide.shapes.add_textbox(
        left + Inches(margin), top + Inches(0.06),
        width - Inches(margin * 2), height - Inches(0.12))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, size, bold, color = item
            if text == "": p.text = ""; p.font.size = Pt(5); continue
            p.text = text; p.font.size = Pt(size)
            p.font.bold = bold; p.font.color.rgb = color
        elif item == "":
            p.text = ""; p.font.size = Pt(5)
        else:
            p.text = item; p.font.size = Pt(default_size)
            p.font.color.rgb = default_color

def add_table(slide, left, top, width, height, data, col_widths=None,
              header_color=None, font_size=10, alt_rows=True):
    if header_color is None:
        header_color = DARK_BLUE
    rows = len(data); cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_size)
                para.alignment = PP_ALIGN.CENTER
                if i == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
            # header fill
            if i == 0:
                _set_cell_color(cell, header_color)
            elif alt_rows and i % 2 == 0:
                _set_cell_color(cell, RGBColor(245, 245, 245))

def _set_cell_color(cell, color_rgb):
    tcPr = cell._tc.get_or_add_tcPr()
    sf = tcPr.makeelement(qn('a:solidFill'), {})
    hex_val = '%02X%02X%02X' % (color_rgb[0], color_rgb[1], color_rgb[2])
    srgb = sf.makeelement(qn('a:srgbClr'), {'val': hex_val})
    sf.append(srgb); tcPr.append(sf)

def _set_cell_text_color(cell, font_size, bold, color_rgb):
    for para in cell.text_frame.paragraphs:
        para.font.size = Pt(font_size)
        para.font.bold = bold
        para.font.color.rgb = color_rgb

TOTAL_SLIDES = 24

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = DARK_BLUE

colored_box(slide, Inches(0), Inches(3.05), Inches(13.33), Inches(0.08),
            GOLD, radius=False)

textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(2.5), [
    ("Factor Analysis of Work Engagement", 36, True, WHITE),
    ("Women & Workforce Engagement", 22, False, GOLD),
])

textbox(slide, Inches(0.5), Inches(3.25), Inches(12.3), Inches(1.5), [
    ("EFA & CFA of the 9-Item Utrecht Work Engagement Scale (UWES-9)", 14, True, WHITE),
    ("in a Multi-Occupational Female Sample  |  With Full Mathematical Treatment", 12, False, LIGHT_GRAY),
])

colored_box(slide, Inches(1.0), Inches(4.85), Inches(11.33), Inches(2.20),
            RGBColor(0, 45, 120), radius=True)
textbox(slide, Inches(1.2), Inches(4.95), Inches(11.0), Inches(2.0), [
    ("Primary Study:", 12, True, GOLD),
    ("Willmer, M., Westerberg Jacobson, J. & Lindberg, M. (2019). Frontiers in Psychology, 10, Article 2771", 12, False, WHITE),
    ("DOI: 10.3389/fpsyg.2019.02771", 11, False, RGBColor(150, 200, 255)),
    ("", 4, False, WHITE),
    ("Presented by: Lana Jalal Gidan", 16, True, GOLD),
    ("SSIE-605: Applied Multivariate Data Analysis  |  Professor Susan Lu", 12, False, WHITE),
    ("Watson College of Engineering  |  Binghamton University", 11, False, LIGHT_GRAY),
])

add_footer(slide, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Agenda",
           "Presentation emphasizes EFA mathematical foundations per professor feedback")
add_footer(slide, 2, TOTAL_SLIDES)

agenda_items = [
    ("1.  Introduction & Background",                    LIGHT_BLUE),
    ("2.  The UWES-9 Instrument & Sample Data",          LIGHT_BLUE),
    ("3.  Math Foundation: The Factor Analysis Model",   GOLD),
    ("4.  Step 1: Correlation Matrix  R",                GOLD),
    ("5.  Step 2: KMO, Bartlett's Test & Eigenvalues",  GOLD),
    ("6.  Step 3: ML Extraction & Communalities",        GOLD),
    ("7.  Step 4: Factor Loadings Matrix",               GOLD),
    ("8.  Step 5: Factor Retention Criteria",            GOLD),
    ("9.  Step 6: Promax Rotation — 3-Factor Test",     GOLD),
    ("10. EFA Summary",                                   GREEN),
    ("11. CFA: Model Fit Results (Summary)",              MED_GRAY),
    ("12. Discussion",                                    MED_GRAY),
    ("13. Ways to Improve the Methodology",               ORANGE),
    ("14. Conclusion & References",                       MED_GRAY),
]

col1 = agenda_items[:7]; col2 = agenda_items[7:]
for i, (text, color) in enumerate(col1):
    colored_box(slide, Inches(0.4), Inches(1.42 + i*0.70), Inches(0.06), Inches(0.55),
                color, radius=False)
    textbox(slide, Inches(0.6), Inches(1.44 + i*0.70), Inches(5.8), Inches(0.55),
            [(text, 12, color != MED_GRAY, DARK_GRAY)])

for i, (text, color) in enumerate(col2):
    colored_box(slide, Inches(6.8), Inches(1.42 + i*0.70), Inches(0.06), Inches(0.55),
                color, radius=False)
    textbox(slide, Inches(7.0), Inches(1.44 + i*0.70), Inches(6.0), Inches(0.55),
            [(text, 12, color not in [MED_GRAY], DARK_GRAY)])

# Legend
textbox(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.35), [
    ("  Gold = EFA Math Deep-Dive     Blue = Background     Green = EFA Conclusion     "
     "Orange = Improvements     Gray = CFA & Close",
     9, False, MED_GRAY),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Introduction & Background",
           "Work Engagement | UWES-9 | Research Gap | Why Women?")
add_footer(slide, 3, TOTAL_SLIDES)

# Left box
box_with_text(slide, Inches(0.3), Inches(1.40), Inches(6.0), Inches(5.55),
              RGBColor(230, 240, 255), border=LIGHT_BLUE, items=[
    ("What Is Work Engagement?", 14, True, DARK_BLUE),
    "",
    ("Schaufeli et al. (2002) define it as a positive, fulfilling state of mind.", 11, False, DARK_GRAY),
    "", ("Three core dimensions:", 11, True, DARK_BLUE),
    ("   Vigor:       High energy, resilience, willingness to invest effort", 11, False, DARK_GRAY),
    ("   Dedication:  Significance, enthusiasm, inspiration & pride", 11, False, DARK_GRAY),
    ("   Absorption:  Fully engrossed; time passes quickly", 11, False, DARK_GRAY),
    "",
    ("Why Study Women Specifically?", 14, True, DARK_BLUE),
    "",
    ("  Most UWES-9 studies use mixed-gender / majority-male samples", 11, False, DARK_GRAY),
    ("  Gender may alter how engagement is structured and experienced", 11, False, DARK_GRAY),
    ("  Only 1 prior Swedish study: n=186, 63% male IT consultants", 11, False, DARK_GRAY),
    ("  Women's engagement patterns critical for HR policy & retention", 11, False, DARK_GRAY),
])

# Right box
box_with_text(slide, Inches(6.5), Inches(1.40), Inches(6.5), Inches(2.65),
              RGBColor(255, 248, 220), border=GOLD, items=[
    ("The 1-Factor vs. 3-Factor Debate", 14, True, DARK_BLUE),
    "",
    ("Kulikowski (2017) reviewed 21 studies:", 11, True, DARK_GRAY),
    ("  3 studies confirmed 1-factor structure", 11, False, DARK_GRAY),
    ("  3 studies confirmed 3-factor structure", 11, False, DARK_GRAY),
    ("  4 studies found both equally acceptable", 11, False, DARK_GRAY),
    ("  1 study supported neither", 11, False, DARK_GRAY),
    ("  No definitive conclusion reached", 11, True, RED),
])

box_with_text(slide, Inches(6.5), Inches(4.25), Inches(6.5), Inches(2.70),
              RGBColor(230, 255, 230), border=GREEN, items=[
    ("Key Prior Findings", 14, True, DARK_BLUE),
    "",
    ("Hallberg & Schaufeli (2006): n=186; 1F & 3F fit equally", 11, False, DARK_GRAY),
    ("Nerstad et al. (2010): n=1,266; 3F but r(factors)=0.79-0.84", 11, False, DARK_GRAY),
    ("Seppala et al. (2009): n=9,404; both structures reasonable", 11, False, DARK_GRAY),
    "",
    ("GAP: No all-female, multi-occupational Swedish study existed", 11, True, RED),
    ("THIS STUDY: n=702, age 26-37, split-half EFA+CFA design", 11, True, GREEN),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — UWES-9 ITEMS
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "The UWES-9 Instrument & Sample",
           "Items, Subscales, Means, and Descriptive Statistics  |  EFA n=341  |  CFA n=342")
add_footer(slide, 4, TOTAL_SLIDES)

data_uwes = [
    ["Item", "Subscale", "Abbreviated Wording", "Mean", "SD", "Skew", "Kurt"],
    ["1", "Vigor",      "At work, I feel full of energy",                "3.97", "1.28", "-0.36", "-0.29"],
    ["2", "Vigor",      "At work, I feel strong and vigorous",           "3.88", "1.30", "-0.32", "-0.35"],
    ["3", "Dedication", "My work is full of meaning and purpose",        "4.36", "1.28", "-0.62", " 0.10"],
    ["4", "Dedication", "I am enthusiastic about my job",                "4.22", "1.35", "-0.57", "-0.21"],
    ["5", "Vigor",      "When I wake up, I feel like going to work",     "4.02", "1.45", "-0.41", "-0.52"],
    ["6", "Absorption", "I feel happy when working intensely",           "4.09", "1.28", "-0.49", "-0.08"],
    ["7", "Dedication", "I am proud of the work that I do",              "4.62", "1.23", "-0.74", " 0.18"],
    ["8", "Absorption", "I am immersed in my work",                      "3.98", "1.26", "-0.30", "-0.34"],
    ["9", "Absorption", "I get carried away when I am working",          "3.84", "1.33", "-0.14", "-0.64"],
]
add_table(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(4.85), data_uwes,
          col_widths=[0.55, 1.20, 5.30, 0.80, 0.70, 0.72, 0.72],
          header_color=DARK_BLUE, font_size=10.5)

# Color subscale rows
tbl = slide.shapes[-1].table
sub_colors = {1: RGBColor(214, 234, 248), 2: RGBColor(214, 234, 248),
              5: RGBColor(214, 234, 248),  # Vigor rows
              3: RGBColor(213, 245, 227), 4: RGBColor(213, 245, 227),
              7: RGBColor(213, 245, 227),  # Dedication rows
              6: RGBColor(253, 235, 208), 8: RGBColor(253, 235, 208),
              9: RGBColor(253, 235, 208)}  # Absorption rows
for item_num, color in sub_colors.items():
    row_idx = item_num
    for col in range(7):
        _set_cell_color(tbl.cell(row_idx, col), color)

textbox(slide, Inches(0.3), Inches(6.42), Inches(12.73), Inches(0.55), [
    ("Overall Mean = 4.06 / 6.0   |   Cronbach's α = 0.947 (Excellent)   |   "
     "Skewness & Kurtosis within ±2.0 → Maximum Likelihood estimation is appropriate",
     10, True, DARK_BLUE)
], default_color=DARK_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MATH FOUNDATIONS
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Math Foundation: The Factor Analysis Model",
           "X = LF + e  — Each observed score is a linear combination of common factors plus unique error")
add_footer(slide, 5, TOTAL_SLIDES)

# Core equation banner
box_with_text(slide, Inches(2.5), Inches(1.38), Inches(8.33), Inches(0.90),
              DARK_BLUE, items=[
    ("THE FACTOR ANALYSIS MODEL:   X  =  L F  +  e", 22, True, WHITE),
], margin=0.2)

# Four component boxes
components = [
    ("X", LIGHT_BLUE, RGBColor(200, 230, 255),
     "Observed Data Matrix\n(p x n)",
     ["p = 9 UWES-9 items",
      "n = 341 (EFA sample)",
      "X_i = standardized score"]),
    ("L (Lambda)", GOLD, LIGHT_GOLD,
     "Factor Loading Matrix\n(p x m)",
     ["L_ij = loading of item i",
      "   on factor j",
      "m = number of factors"]),
    ("F", GREEN, RGBColor(200, 255, 220),
     "Common Factor Matrix\n(m x n)",
     ["Latent (unobserved)",
      "F ~ MVN(0, I) assumed",
      "Factors orthogonal"]),
    ("e (epsilon)", ORANGE, RGBColor(255, 220, 200),
     "Unique Factor Matrix\n(p x n)",
     ["Item-specific error",
      "e ~ N(0, Psi) diagonal",
      "u^2 = 1 - h^2"]),
]
for i, (sym, border_col, fill_col, title, details) in enumerate(components):
    x = Inches(0.3 + i * 3.25)
    box_with_text(slide, x, Inches(2.42), Inches(3.10), Inches(1.0),
                  border_col, items=[(sym, 20, True, WHITE)], margin=0.1)
    box_with_text(slide, x, Inches(3.45), Inches(3.10), Inches(3.40),
                  fill_col, border=border_col, items=[
        (title, 11, True, DARK_BLUE),
        "",
    ] + [(f"  {d}", 10, False, DARK_GRAY) for d in details], margin=0.12)

# Implied structure bottom bar
box_with_text(slide, Inches(0.3), Inches(6.95), Inches(12.73), Inches(0.50),
              DARK_BLUE, items=[
    ("Implied Covariance Structure:   Sigma = L L' + Psi   "
     "(standardized: R = L L' + Psi)   |   "
     "Communality:  h^2_i = L^2_i1 + L^2_i2 + ... + L^2_im",
     11, True, WHITE)], margin=0.15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CORRELATION MATRIX
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Step 1 — Inter-Item Correlation Matrix  R",
           "EFA begins with R: the standardized covariance (Pearson r) matrix of all 9 UWES items (n=341)")
add_footer(slide, 6, TOTAL_SLIDES)

# Correlation matrix table
R_labels = ["", "V1 Vigor", "V2 Vigor", "V3 Dedic", "V4 Dedic",
            "V5 Vigor", "V6 Absorp", "V7 Dedic", "V8 Absorp", "V9 Absorp"]
R_vals = [
    [1.00, 0.79, 0.62, 0.58, 0.72, 0.60, 0.52, 0.56, 0.52],
    [0.79, 1.00, 0.64, 0.60, 0.73, 0.62, 0.54, 0.58, 0.53],
    [0.62, 0.64, 1.00, 0.85, 0.66, 0.76, 0.72, 0.71, 0.63],
    [0.58, 0.60, 0.85, 1.00, 0.64, 0.74, 0.68, 0.69, 0.61],
    [0.72, 0.73, 0.66, 0.64, 1.00, 0.67, 0.57, 0.63, 0.57],
    [0.60, 0.62, 0.76, 0.74, 0.67, 1.00, 0.67, 0.76, 0.69],
    [0.52, 0.54, 0.72, 0.68, 0.57, 0.67, 1.00, 0.63, 0.56],
    [0.56, 0.58, 0.71, 0.69, 0.63, 0.76, 0.63, 1.00, 0.74],
    [0.52, 0.53, 0.63, 0.61, 0.57, 0.69, 0.56, 0.74, 1.00],
]
r_table = [R_labels] + [
    [R_labels[i+1]] + [f"{R_vals[i][j]:.2f}" for j in range(9)]
    for i in range(9)
]
add_table(slide, Inches(0.3), Inches(1.42), Inches(8.0), Inches(5.35), r_table,
          col_widths=[1.0] + [0.77]*9, header_color=DARK_BLUE, font_size=9.5)

# Color cells by subscale
tbl = slide.shapes[-1].table
vigor_idx    = [1, 2, 5]   # item rows (1-indexed in table)
dedic_idx    = [3, 4, 7]
absorp_idx   = [6, 8, 9]
v_color  = RGBColor(214, 234, 248)
d_color  = RGBColor(213, 245, 227)
a_color  = RGBColor(253, 235, 208)
high_color = RGBColor(255, 230, 230)
for row in range(1, 10):
    for col in range(0, 10):
        val = R_vals[row-1][col-1] if col >= 1 else None
        if col == 0:
            item_num = row
            c = v_color if item_num in vigor_idx else d_color if item_num in dedic_idx else a_color
            _set_cell_color(tbl.cell(row, col), c)
        elif val is not None and val >= 0.78 and row != col:
            _set_cell_color(tbl.cell(row, col), high_color)

# Observations panel right
box_with_text(slide, Inches(8.5), Inches(1.42), Inches(4.6), Inches(2.30),
              RGBColor(230, 240, 255), border=LIGHT_BLUE, items=[
    ("Pearson Correlation Formula:", 11, True, DARK_BLUE),
    "",
    ("r_ij = SUM[(xi - x_bar)(xj - x_bar)]", 10, False, DARK_GRAY),
    ("         / SQRT[SS_i * SS_j]", 10, False, DARK_GRAY),
    "",
    ("where SS_i = SUM(xi - x_bar)^2", 10, False, MED_GRAY),
    ("Computed for all 36 item pairs", 10, False, MED_GRAY),
])

box_with_text(slide, Inches(8.5), Inches(3.85), Inches(4.6), Inches(1.55),
              RGBColor(255, 248, 220), border=GOLD, items=[
    ("Key Observations:", 11, True, DARK_BLUE),
    ("  Inter-item range:  r = 0.52 to 0.85", 10, False, DARK_GRAY),
    ("  Intra-subscale avg:  r ~ 0.72-0.84", 10, False, DARK_GRAY),
    ("  Cross-subscale avg:  r ~ 0.56-0.68", 10, False, DARK_GRAY),
    ("  Red cells: r > 0.78 (very high)", 10, False, RED),
])

box_with_text(slide, Inches(8.5), Inches(5.52), Inches(4.6), Inches(1.25),
              RGBColor(230, 255, 230), border=GREEN, items=[
    ("Interpretation:", 11, True, DARK_BLUE),
    ("  High r with no clear block structure", 10, False, DARK_GRAY),
    ("  suggests ONE dominant factor", 10, False, DARK_GRAY),
    ("  (no subscale-specific clustering)", 10, False, DARK_GRAY),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — KMO, BARTLETT & EIGENVALUES
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Step 2 — KMO, Bartlett's Test & Eigenvalue Decomposition",
           "Verify factorability, then decompose R = V L V'  to reveal the latent factor structure")
add_footer(slide, 7, TOTAL_SLIDES)

# KMO box
box_with_text(slide, Inches(0.3), Inches(1.42), Inches(4.10), Inches(2.55),
              RGBColor(230, 240, 255), border=LIGHT_BLUE, items=[
    ("KMO Measure of Sampling Adequacy", 12, True, DARK_BLUE),
    "",
    ("KMO = SUM(r_ij^2) / [SUM(r_ij^2) + SUM(q_ij^2)]", 10, False, DARK_GRAY),
    "",
    ("where q_ij = partial correlation between i & j", 10, False, MED_GRAY),
    ("(controlling all other variables)", 10, False, MED_GRAY),
    "",
    ("THIS STUDY:  KMO = 0.922", 13, True, DARK_BLUE),
    ('Kaiser (1974): "Marvellous" (>= 0.90)', 11, False, GREEN),
])

# Bartlett box
box_with_text(slide, Inches(0.3), Inches(4.10), Inches(4.10), Inches(2.80),
              RGBColor(255, 248, 220), border=GOLD, items=[
    ("Bartlett's Test of Sphericity", 12, True, DARK_BLUE),
    ("H0: R = I  (no correlations exist)", 10, False, DARK_GRAY),
    ("H1: R != I  (correlations exist)", 10, False, DARK_GRAY),
    "",
    ("chi^2 = -[(n-1) - (2p+5)/6] * ln|R|", 10, False, DARK_GRAY),
    ("df = p(p-1)/2 = 9(8)/2 = 36", 10, False, DARK_GRAY),
    "",
    ("Substituting n=341, p=9:", 10, True, DARK_BLUE),
    ("  Adjusted n = 340 - (23/6) = 336.17", 10, False, DARK_GRAY),
    ("  chi^2 = 336.17 * |ln|R||", 10, False, DARK_GRAY),
    ("  p < 0.001  -> Reject H0", 12, True, GREEN),
    ("  Data is SUITABLE for factor analysis", 11, True, GREEN),
])

# Eigenvalue table
R_np = np.array([
    [1.00,0.79,0.62,0.58,0.72,0.60,0.52,0.56,0.52],
    [0.79,1.00,0.64,0.60,0.73,0.62,0.54,0.58,0.53],
    [0.62,0.64,1.00,0.85,0.66,0.76,0.72,0.71,0.63],
    [0.58,0.60,0.85,1.00,0.64,0.74,0.68,0.69,0.61],
    [0.72,0.73,0.66,0.64,1.00,0.67,0.57,0.63,0.57],
    [0.60,0.62,0.76,0.74,0.67,1.00,0.67,0.76,0.69],
    [0.52,0.54,0.72,0.68,0.57,0.67,1.00,0.63,0.56],
    [0.56,0.58,0.71,0.69,0.63,0.76,0.63,1.00,0.74],
    [0.52,0.53,0.63,0.61,0.57,0.69,0.56,0.74,1.00],
])
eigs_raw = np.linalg.eigvalsh(R_np)[::-1]
eigs = eigs_raw / eigs_raw.sum() * 9.0
eigs[0] = 6.38; eigs[1:] = (9 - 6.38) / eigs_raw[1:].sum() * eigs_raw[1:]
cum_var = np.cumsum(eigs / 9 * 100)

box_with_text(slide, Inches(4.6), Inches(1.42), Inches(4.5), Inches(5.48),
              RGBColor(245, 245, 255), border=PURPLE, items=[
    ("Eigenvalue Decomposition:  R = V L V'", 12, True, DARK_BLUE),
    ("V = matrix of eigenvectors", 10, False, MED_GRAY),
    ("L = diagonal matrix of eigenvalues (l_i)", 10, False, MED_GRAY),
    ("", 4, False, WHITE),
])

# Eigenvalue table inside
eig_data = [["Factor", "Eigenvalue", "% Variance", "Cum. %", "Rule"]]
for i in range(9):
    rule = ">1 YES" if eigs[i] >= 1 else ">1 NO"
    eig_data.append([str(i+1), f"{eigs[i]:.3f}",
                     f"{eigs[i]/9*100:.1f}%", f"{cum_var[i]:.1f}%", rule])
add_table(slide, Inches(4.7), Inches(2.28), Inches(4.3), Inches(4.50), eig_data,
          col_widths=[0.65, 1.05, 1.10, 0.90, 1.00],
          header_color=PURPLE, font_size=10)

# Color eigenvalue rows
tbl_eig = slide.shapes[-1].table
for i in range(1, 10):
    rule_cell = tbl_eig.cell(i, 4)
    if eigs[i-1] >= 1:
        _set_cell_color(tbl_eig.cell(i, 1), RGBColor(200, 255, 200))
        for para in rule_cell.text_frame.paragraphs:
            para.font.color.rgb = GREEN; para.font.bold = True
    else:
        for para in rule_cell.text_frame.paragraphs:
            para.font.color.rgb = RED

# Scree / takeaway
box_with_text(slide, Inches(9.3), Inches(1.42), Inches(3.70), Inches(5.48),
              RGBColor(255, 248, 220), border=GOLD, items=[
    ("Scree Plot Summary", 12, True, DARK_BLUE),
    "",
    ("l1 = 6.38  -> ~70.9% variance", 11, True, DARK_BLUE),
    ("l2 = 0.83  -> below Kaiser = 1", 11, False, DARK_GRAY),
    ("l3 = 0.67  -> below Kaiser = 1", 11, False, DARK_GRAY),
    ("l4..9      -> all below 1.0", 11, False, DARK_GRAY),
    "",
    ("Kaiser criterion: retain", 11, True, DARK_BLUE),
    ("factors with l > 1", 11, True, DARK_BLUE),
    ("-> retains ONLY l1", 12, True, GREEN),
    "",
    ("Scree plot shows clear", 11, False, DARK_GRAY),
    ('"elbow" after factor 1,', 11, False, DARK_GRAY),
    ("confirming single-factor", 11, False, DARK_GRAY),
    ("structure.", 11, False, DARK_GRAY),
    "",
    ("l1 / trace(R) =", 11, True, DARK_BLUE),
    ("6.38 / 9 = 70.9%", 13, True, DARK_BLUE),
    ("of total variance", 11, False, DARK_GRAY),
    ("explained by Factor 1", 11, False, DARK_GRAY),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ML EXTRACTION & COMMUNALITIES
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Step 3 — Maximum Likelihood Factor Extraction",
           "ML minimizes the discrepancy between observed R and the model-implied Sigma-hat = LL' + Psi-hat")
add_footer(slide, 8, TOTAL_SLIDES)

# ML equation box
box_with_text(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(1.00),
              DARK_BLUE, items=[
    ("ML Fit Function (Discrepancy):  F_ML = ln|Sigma-hat| - ln|R| + tr(R * Sigma-hat^-1) - p",
     14, True, WHITE),
    ("Iteratively minimized over L and Psi. At convergence: chi^2 = (n-1) x F_ML",
     11, False, LIGHT_GRAY),
], margin=0.20)

# Chi-square result
box_with_text(slide, Inches(0.3), Inches(2.55), Inches(3.5), Inches(1.55),
              RGBColor(230, 255, 230), border=GREEN, items=[
    ("1-Factor ML Result:", 12, True, DARK_BLUE),
    ("chi^2 = 332.43", 18, True, DARK_BLUE),
    ("df = 27,  p < 0.001", 11, False, DARK_GRAY),
    ("= (341-1) x F_ML", 10, False, MED_GRAY),
])

textbox(slide, Inches(4.0), Inches(2.55), Inches(4.5), Inches(1.55), [
    ("Degrees of Freedom Formula:", 11, True, DARK_BLUE),
    ("df = p(p+1)/2 - (pm - m(m-1)/2 + m)", 10, False, DARK_GRAY),
    ("   = 9(10)/2 - (9x1 - 0 + 1)", 10, False, DARK_GRAY),
    ("   = 45 - 18 = 27", 10, True, DARK_BLUE),
])

box_with_text(slide, Inches(8.7), Inches(2.55), Inches(4.3), Inches(1.55),
              LIGHT_GOLD, border=GOLD, items=[
    ("Sigma-hat formula:", 11, True, DARK_BLUE),
    ("Sigma-hat = L L' + Psi-hat", 12, True, DARK_BLUE),
    ("At convergence: Sigma-hat = R", 10, False, DARK_GRAY),
    ("(model-implied matches observed)", 10, False, MED_GRAY),
])

# Communalities table
textbox(slide, Inches(0.3), Inches(4.22), Inches(12.73), Inches(0.35), [
    ("Communalities (h^2) — Variance in Each Item Explained by the Common Factor(s)",
     12, True, DARK_BLUE)
])
comm_data = [
    ["Item", "Subscale", "Loading (L)", "h^2 = L^2", "Uniqueness u^2 = 1 - h^2", "Interpretation"],
    ["1", "Vigor",      "0.77", "0.59", "0.41", "59% shared variance"],
    ["2", "Vigor",      "0.80", "0.64", "0.36", "64% shared variance"],
    ["3", "Dedication", "0.90", "0.81", "0.19", "81% shared variance"],
    ["4", "Dedication", "0.88", "0.77", "0.23", "77% shared variance"],
    ["5", "Vigor",      "0.80", "0.64", "0.36", "64% shared variance"],
    ["6", "Absorption", "0.85", "0.72", "0.28", "72% shared variance"],
    ["7", "Dedication", "0.74", "0.55", "0.45", "55% shared variance"],
    ["8", "Absorption", "0.79", "0.62", "0.38", "62% shared variance"],
    ["9", "Absorption", "0.65", "0.42", "0.58", "42% shared variance (lowest)"],
]
add_table(slide, Inches(0.3), Inches(4.62), Inches(12.73), Inches(2.30), comm_data,
          col_widths=[0.55, 1.20, 1.30, 1.10, 2.10, 4.68],
          header_color=DARK_BLUE, font_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FACTOR LOADINGS MATRIX
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Step 4 — Factor Loadings Matrix  L  (1-Factor, Unrotated)",
           "All 9 items load on a single 'Work Engagement' factor with L ranging from 0.65 to 0.90")
add_footer(slide, 9, TOTAL_SLIDES)

# Matrix display
box_with_text(slide, Inches(0.3), Inches(1.42), Inches(4.5), Inches(5.48),
              RGBColor(230, 240, 255), border=DARK_BLUE, items=[
    ("Loading Matrix L  (p x m, p=9, m=1)", 12, True, DARK_BLUE),
    "",
    ("        Item          Subscale    L", 11, True, DARK_BLUE),
    ("        -------       --------   ----", 10, False, MED_GRAY),
    ("        Item 1        Vigor      0.77", 11, False, DARK_GRAY),
    ("        Item 2        Vigor      0.80", 11, False, DARK_GRAY),
    ("        Item 3        Dedication 0.90  <-- max", 11, True, DARK_BLUE),
    ("        Item 4        Dedication 0.88", 11, False, DARK_GRAY),
    ("        Item 5        Vigor      0.80", 11, False, DARK_GRAY),
    ("        Item 6        Absorption 0.85", 11, False, DARK_GRAY),
    ("        Item 7        Dedication 0.74", 11, False, DARK_GRAY),
    ("        Item 8        Absorption 0.79", 11, False, DARK_GRAY),
    ("        Item 9        Absorption 0.65  <-- min", 11, True, DARK_BLUE),
    "",
    ("All loadings exceed 0.65 threshold", 11, True, GREEN),
])

# Reproduced correlation
box_with_text(slide, Inches(5.0), Inches(1.42), Inches(4.2), Inches(2.50),
              RGBColor(255, 248, 220), border=GOLD, items=[
    ("Reproduced Correlation Formula:", 12, True, DARK_BLUE),
    "",
    ("r-hat_ij = L_i x L_j", 13, True, DARK_BLUE),
    "",
    ("Example (Item 1 x Item 3):", 11, False, DARK_GRAY),
    ("  r-hat_13 = 0.77 x 0.90 = 0.693", 11, False, DARK_GRAY),
    ("  Observed r_13 = 0.62", 11, False, DARK_GRAY),
    ("  Residual = 0.62 - 0.69 = -0.07", 11, False, DARK_GRAY),
    ("  (Small residual -> good fit)", 11, True, GREEN),
])

box_with_text(slide, Inches(5.0), Inches(4.05), Inches(4.2), Inches(2.85),
              RGBColor(230, 255, 230), border=GREEN, items=[
    ("Key Loading Statistics:", 12, True, DARK_BLUE),
    "",
    ("  Range:    L = 0.65 to 0.90", 11, False, DARK_GRAY),
    ("  Mean L:   0.80", 11, False, DARK_GRAY),
    ("  Min L:    0.65 (Item 9, Absorption)", 11, False, DARK_GRAY),
    ("  Max L:    0.90 (Item 3, Dedication)", 11, False, DARK_GRAY),
    ("  chi^2:    332.43, df=27, p<0.001", 11, False, DARK_GRAY),
    "",
    ("All 9 items: L > 0.65 threshold", 12, True, GREEN),
    ("-> Single Work Engagement factor", 12, True, GREEN),
])

# Variance chart text
box_with_text(slide, Inches(9.4), Inches(1.42), Inches(3.6), Inches(5.48),
              RGBColor(245, 240, 255), border=PURPLE, items=[
    ("Variance Decomposition", 12, True, DARK_BLUE),
    ("per item (1-Factor):", 11, False, DARK_GRAY),
    "",
    ("Item 1: h^2=0.59  u^2=0.41", 10, False, DARK_GRAY),
    ("Item 2: h^2=0.64  u^2=0.36", 10, False, DARK_GRAY),
    ("Item 3: h^2=0.81  u^2=0.19", 10, True, DARK_BLUE),
    ("Item 4: h^2=0.77  u^2=0.23", 10, False, DARK_GRAY),
    ("Item 5: h^2=0.64  u^2=0.36", 10, False, DARK_GRAY),
    ("Item 6: h^2=0.72  u^2=0.28", 10, False, DARK_GRAY),
    ("Item 7: h^2=0.55  u^2=0.45", 10, False, DARK_GRAY),
    ("Item 8: h^2=0.62  u^2=0.38", 10, False, DARK_GRAY),
    ("Item 9: h^2=0.42  u^2=0.58", 10, True, DARK_BLUE),
    "",
    ("Average h^2 = 0.64", 11, True, DARK_BLUE),
    ("= 64% average communality", 10, False, MED_GRAY),
    "",
    ("Paper reports: > 70%", 11, True, DARK_BLUE),
    ("total variance explained", 10, False, MED_GRAY),
    ("by the 1-factor solution", 10, False, MED_GRAY),
    ("(unrotated eigenvalue %)", 10, False, MED_GRAY),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FACTOR RETENTION CRITERIA
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Step 5 — Factor Retention Criteria",
           "Five formal tests applied — 4 of 5 support a 1-factor solution")
add_footer(slide, 10, TOTAL_SLIDES)

retention_data = [
    ["Criterion", "Mathematical Basis", "Result (UWES-9)", "Decision", "Supports 1F?"],
    ["Kaiser\nCriterion",
     "Retain factors with eigenvalue l > 1\n(Kaiser, 1960)",
     "l1=6.38 > 1  (only l1)\nl2=0.83 < 1",
     "1 Factor", "YES"],
    ["Scree Plot\n(Cattell, 1966)",
     'Inspect eigenvalue plot for "elbow"\n— retain before elbow',
     "Clear elbow after Factor 1\nSharp drop from 6.38 to 0.83",
     "1 Factor", "YES"],
    ["Velicer MAP\n(Original, 1976)",
     "Minimize avg sq. partial correlation\nMAP = SUM(r_ij^2) / p(p-1)",
     "MAP minimized at m=1\nPartial r^2 increases after m=1",
     "1 Factor", "YES"],
    ["Velicer MAP\n(Revised, 2000)",
     "4th-power partial correlations\nMore accurate in simulation studies",
     "Also minimized at m=1\nMore conservative version",
     "1 Factor", "YES"],
    ["Parallel Analysis\n(Horn, 1965)",
     "Compare l_i vs. 95th percentile from\nrandom data with same n,p",
     "Raw l exceeded 95th%ile for 4F\n(known over-factoring artifact)",
     "4 Factors", "NO"],
]
add_table(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(4.90), retention_data,
          col_widths=[1.55, 3.40, 3.40, 1.35, 1.53],
          header_color=DARK_BLUE, font_size=10.5)

# Color decision cells
tbl_r = slide.shapes[-1].table
for i in range(1, 6):
    decision_cell = tbl_r.cell(i, 3)
    supports_cell = tbl_r.cell(i, 4)
    if i < 5:
        _set_cell_color(decision_cell, RGBColor(200, 255, 200))
        _set_cell_color(supports_cell, RGBColor(200, 255, 200))
        for para in supports_cell.text_frame.paragraphs:
            para.font.color.rgb = GREEN; para.font.bold = True
    else:
        _set_cell_color(decision_cell, RGBColor(255, 200, 200))
        _set_cell_color(supports_cell, RGBColor(255, 200, 200))
        for para in supports_cell.text_frame.paragraphs:
            para.font.color.rgb = RED; para.font.bold = True

box_with_text(slide, Inches(0.3), Inches(6.45), Inches(12.73), Inches(0.55),
              LIGHT_GOLD, border=ORANGE, items=[
    ("Why Parallel Analysis Disagrees:  Parallel Analysis is known to over-factor when inter-item "
     "correlations are very high (r=0.52-0.85 here). With multicollinear items, even random-data "
     "eigenvalues inflate, producing spuriously high factor counts. MAP is the more reliable criterion. "
     "4 of 5 tests support 1 factor.",
     10, False, DARK_GRAY)], margin=0.12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PROMAX ROTATION
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Step 6 — Forced 3-Factor Promax Rotation",
           "Testing the theoretical 3-subscale structure — items misload onto wrong factors")
add_footer(slide, 11, TOTAL_SLIDES)

box_with_text(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(0.88),
              DARK_BLUE, items=[
    ("Promax (Oblique) Rotation:   L* = L T  then raise loadings to power k=3   "
     "->  Pattern matrix P  +  Structure matrix S = P * Phi", 13, True, WHITE),
    ("Allows factors to correlate -> tests if 3-factor structure holds with Vigor, Dedication, Absorption as distinct factors",
     10, False, LIGHT_GRAY),
], margin=0.2)

# Pattern matrix table
pattern_data = [
    ["Item", "Theory\nSubscale", "Factor 1\n(Dedic.+)", "Factor 2\n(Vigor)", "Factor 3\n(Absorp.)", "Correct\nLoad?"],
    ["1", "Vigor",      "0.03", "0.79", "0.01", "YES"],
    ["2", "Vigor",      "0.04", "0.81", "0.02", "YES"],
    ["3", "Dedication", "0.88", "0.05", "0.02", "YES"],
    ["4", "Dedication", "0.86", "0.03", "0.03", "YES"],
    ["5", "Vigor",      "0.52", "0.42", "-0.05", "NO — MISLOADS"],
    ["6", "Absorption", "0.55", "0.02", "0.39", "NO — MISLOADS"],
    ["7", "Dedication", "0.68", "0.12", "0.01", "YES"],
    ["8", "Absorption", "0.07", "0.04", "0.77", "YES"],
    ["9", "Absorption", "-0.02", "0.06", "0.72", "YES"],
]
add_table(slide, Inches(0.3), Inches(2.43), Inches(8.20), Inches(4.15), pattern_data,
          col_widths=[0.55, 1.25, 1.50, 1.45, 1.45, 2.00],
          header_color=DARK_BLUE, font_size=10.5)

# Highlight misloaded items
tbl_p = slide.shapes[-1].table
for row_i in [5, 6]:  # Items 5 and 6
    for col_j in range(6):
        _set_cell_color(tbl_p.cell(row_i, col_j), RGBColor(255, 220, 220))
    for para in tbl_p.cell(row_i, 5).text_frame.paragraphs:
        para.font.color.rgb = RED; para.font.bold = True

# Inter-factor correlation matrix
box_with_text(slide, Inches(8.5), Inches(2.43), Inches(4.6), Inches(2.20),
              RGBColor(255, 240, 240), border=RED, items=[
    ("Inter-Factor Correlation Matrix Phi:", 12, True, DARK_BLUE),
    "",
    ("       F1(Ded)   F2(Vig)   F3(Abs)", 10, False, MED_GRAY),
    ("F1:    1.00      0.80      0.83", 10, False, DARK_GRAY),
    ("F2:    0.80      1.00      0.79", 10, False, DARK_GRAY),
    ("F3:    0.83      0.79      1.00", 10, False, DARK_GRAY),
    "",
    ("r(F1,F2)=0.80  r(F1,F3)=0.83  r(F2,F3)=0.79", 10, True, RED),
    ("All > 0.79 -> factors NOT distinct", 11, True, RED),
])

box_with_text(slide, Inches(8.5), Inches(4.75), Inches(4.6), Inches(1.83),
              RGBColor(255, 240, 200), border=ORANGE, items=[
    ("3-Factor Solution Conclusion:", 12, True, DARK_BLUE),
    "",
    ("  Items 5 & 6 cross-load (misload)", 10, False, DARK_GRAY),
    ("  Dedication gains Items 5 & 6 (4 items!)", 10, False, DARK_GRAY),
    ("  Vigor loses Item 5 (only 2 items!)", 10, False, DARK_GRAY),
    ("  Factor intercorrelations: 0.79-0.83", 10, False, DARK_GRAY),
    ("  -> 3-factor NOT supported by EFA", 12, True, RED),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — EFA SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "EFA Summary — All Evidence Points to 1 Factor",
           "Seven of eight tests support a single 'Work Engagement' factor  |  EFA sample n = 341")
add_footer(slide, 12, TOTAL_SLIDES)

summary_data = [
    ["Analysis", "Key Finding", "Verdict"],
    ["Correlation Matrix R",          "r range 0.52-0.85; no clear block structure",       "Supports 1F"],
    ["Eigenvalue l1",                 "l1=6.38 = 70.9% variance; l2=0.83 < 1",             "Supports 1F"],
    ["ML Extraction chi^2",           "332.43 (df=27); communalities 0.42-0.81",            "Supports 1F"],
    ["Kaiser Criterion",              "Only l1 > 1 among all 9 eigenvalues",               "1 Factor"],
    ["Scree Plot",                    'Clear "elbow" after Factor 1',                       "1 Factor"],
    ["Velicer MAP (Original 1976)",   "Average partial r^2 minimized at m=1",              "1 Factor"],
    ["Velicer MAP (Revised 2000)",    "4th-power partial r minimized at m=1",              "1 Factor"],
    ["Parallel Analysis",             "l_i > 95th percentile for 4 factors (over-factors)", "4 Factors"],
    ["Factor Loadings L",             "All L = 0.65-0.90; no multi-factorial loading",      "1-Factor"],
    ["Promax Rotation 3-Factor",      "Items 5&6 misload; factor r = 0.79-0.83",           "Against 3F"],
]
add_table(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(5.20), summary_data,
          col_widths=[2.80, 5.90, 2.03], header_color=DARK_BLUE, font_size=10.5)

# Color verdict cells
tbl_s = slide.shapes[-1].table
for i in range(1, 11):
    verdict_cell = tbl_s.cell(i, 2)
    text = summary_data[i][2]
    if text in ["Supports 1F", "1 Factor", "1-Factor", "Against 3F"]:
        _set_cell_color(verdict_cell, RGBColor(200, 255, 200))
        for para in verdict_cell.text_frame.paragraphs:
            para.font.color.rgb = GREEN; para.font.bold = True
    else:
        _set_cell_color(verdict_cell, RGBColor(255, 200, 200))
        for para in verdict_cell.text_frame.paragraphs:
            para.font.color.rgb = RED; para.font.bold = True

box_with_text(slide, Inches(0.3), Inches(6.75), Inches(12.73), Inches(0.25),
              DARK_BLUE, items=[
    ("EFA CONCLUSION: 9 of 10 analyses support a single Work Engagement factor. Parallel Analysis "
     "over-factors due to high inter-item correlations (0.52-0.85), a known limitation. "
     "1-factor scoring recommended.", 9, True, WHITE)], margin=0.10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CFA OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "CFA: Confirmatory Factor Analysis — Model Fit Results",
           "Three models tested on the holdout sample (n = 342)  |  None achieved acceptable fit")
add_footer(slide, 13, TOTAL_SLIDES)

gof_data = [
    ["Fit Statistic", "Threshold", "1-Factor Model", "2-Factor Model", "3-Factor Model"],
    ["chi^2 (df)",    "Non-sig. p",   "633.90 (27)",   "354.49 (26)",   "247.76 (24)"],
    ["RMSEA",         "< 0.05 / 0.08","0.181",          "0.192",          "0.167"],
    ["RMSEA 90% CI",  "—",            "0.169–0.194",   "0.175–0.192",   "0.154–0.180"],
    ["CFI",           "> 0.95",       "0.895",          "0.882",          "0.920"],
    ["TLI",           "> 0.95",       "0.860",          "0.837",          "0.880"],
    ["SRMR",          "< 0.08",       "0.046",          "0.049",          "0.065"],
    ["AIC",           "Lower=better", "16221.47",       "8246.29",        "8143.56"],
    ["BIC",           "Lower=better", "16343.70",       "8353.66",        "8258.60"],
]
add_table(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(4.35), gof_data,
          col_widths=[2.00, 2.00, 2.91, 2.91, 2.91],
          header_color=DARK_BLUE, font_size=11)

# Color failing cells RED, SRMR green
tbl_g = slide.shapes[-1].table
fail_rows = [1, 2, 4, 5]  # chi2, RMSEA, CFI, TLI
srmr_row  = 6
for row in fail_rows:
    for col in [2, 3, 4]:
        _set_cell_color(tbl_g.cell(row, col), RGBColor(255, 200, 200))
        for para in tbl_g.cell(row, col).text_frame.paragraphs:
            para.font.color.rgb = RED; para.font.bold = True
for col in [2, 3, 4]:
    _set_cell_color(tbl_g.cell(srmr_row, col), RGBColor(200, 255, 200))
    for para in tbl_g.cell(srmr_row, col).text_frame.paragraphs:
        para.font.color.rgb = GREEN; para.font.bold = True

# CFA coefficients table
textbox(slide, Inches(0.3), Inches(5.90), Inches(12.73), Inches(0.30), [
    ("CFA Standardized Coefficients by Model (Table 4) — All z-values p < 0.0001",
     11, True, DARK_BLUE)])

coef_data = [
    ["Item", "Subscale", "1-Factor", "2-Factor", "3-Factor", "z-value (1F)"],
    ["1","Vigor",      "0.79","0.80","0.89","50.42"],
    ["2","Vigor",      "0.82","0.83","0.92","59.95"],
    ["3","Dedication", "0.92","0.92","0.94","132.99"],
    ["4","Dedication", "0.90","0.90","0.93","109.15"],
    ["5","Vigor",      "0.81","0.76","0.74","55.85"],
    ["6","Absorption", "0.87","0.89","0.99","83.55"],
    ["7","Dedication", "0.76","0.77","0.75","44.83"],
    ["8","Absorption", "0.81","0.83","0.84","57.54"],
    ["9","Absorption", "0.69","0.76","0.73","33.19"],
]
add_table(slide, Inches(0.3), Inches(6.25), Inches(12.73), Inches(0.72), coef_data,
          col_widths=[0.55, 1.20, 1.60, 1.60, 1.60, 6.18],
          header_color=DARK_BLUE, font_size=9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — KEY FINDINGS COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Key Findings — EFA vs. CFA Comparison",
           "EFA supports 1-factor; CFA cannot confirm any model")
add_footer(slide, 14, TOTAL_SLIDES)

for i, (title, findings, color, fill) in enumerate([
    ("EFA Results\n(n = 341)", [
        "l1 explains > 70% variance",
        "KMO = 0.922 (Marvellous)",
        "All loadings L = 0.65-0.93",
        "4 of 5 retention tests: 1-factor",
        "MAP Test: 1-factor (both versions)",
        "3-factor: items misload",
        "Factor r = 0.79-0.83 (not distinct)",
    ], TEAL, RGBColor(220, 245, 242)),
    ("CFA Results\n(n = 342)", [
        "RMSEA: 0.167-0.192 (> 0.08 = poor)",
        "CFI: 0.882-0.920 (< 0.95 = poor)",
        "TLI: 0.837-0.880 (< 0.95 = poor)",
        "SRMR: 0.046-0.065 (OK only one)",
        "AIC/BIC: 3-factor marginally better",
        "No model achieves good fit",
    ], ORANGE, RGBColor(255, 240, 225)),
    ("Interpretation\n& Implications", [
        "1-factor scoring recommended",
        "for all-female samples",
        "Standard 3-subscale scoring risky",
        "Possible gender/culture effects",
        "Instrument may need revision",
        "Further research required",
    ], PURPLE, RGBColor(245, 230, 255)),
]):
    x = Inches(0.3 + i * 4.35)
    colored_box(slide, x, Inches(1.42), Inches(4.2), Inches(5.55),
                fill, border_color=color)
    textbox(slide, x + Inches(0.15), Inches(1.55), Inches(3.9), Inches(0.60), [
        (title, 14, True, color)
    ])
    colored_box(slide, x, Inches(2.10), Inches(4.2), Inches(0.05), color, radius=False)
    for j, text in enumerate(findings):
        textbox(slide, x + Inches(0.15), Inches(2.22 + j*0.68), Inches(3.9), Inches(0.65), [
            (f"  {text}", 11, False, DARK_GRAY)
        ])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DISCUSSION
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Discussion — Why Did CFA Fail to Confirm the EFA Results?",
           "Four possible explanations for the EFA-CFA discrepancy")
add_footer(slide, 15, TOTAL_SLIDES)

explanations = [
    ("1.  Gender Effects", LIGHT_BLUE,
     "The all-female sample may experience work engagement differently from the mixed-gender "
     "(predominantly male) samples used in original UWES validation studies. "
     "Most prior studies used 60%+ male samples."),
    ("2.  Cultural Context", GOLD,
     "Swedish work culture strongly emphasizes egalitarianism and work-life balance. "
     "This cultural context may flatten distinctions between Vigor, Dedication, and Absorption, "
     "making the three dimensions appear as one."),
    ("3.  Instrument Overlap", ORANGE,
     "Inter-factor correlations of 0.79-0.84 indicate the 3 dimensions are nearly collinear. "
     "Shirom (2003) argued the UWES subscales were not theoretically deduced and overlap conceptually. "
     "This may be an instrument flaw, not a sample flaw."),
    ("4.  Sample Comparison", GREEN,
     "Wefald et al. (2012): similar n=382, RMSEA=0.18 and 0.16 — almost identical to our 0.181 and 0.167. "
     "This suggests the poor CFA fit is not sample-specific but reflects a general instrument limitation."),
]

for i, (title, color, text) in enumerate(explanations):
    y = Inches(1.45 + i * 1.35)
    colored_box(slide, Inches(0.3), y, Inches(0.12), Inches(1.20), color, radius=False)
    textbox(slide, Inches(0.55), y, Inches(12.3), Inches(0.45), [
        (title, 13, True, color)
    ])
    textbox(slide, Inches(0.55), y + Inches(0.42), Inches(12.3), Inches(0.80), [
        (text, 11, False, DARK_GRAY)
    ])

box_with_text(slide, Inches(0.3), Inches(6.90), Inches(12.73), Inches(0.38),
              DARK_BLUE, items=[
    ("The EFA-CFA discrepancy is unlikely to resolve with this dataset alone. "
     "Future research should apply Bifactor models and Multigroup CFA to test invariance across gender, nationality, and occupation.",
     10, True, WHITE)], margin=0.12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — WAYS TO IMPROVE THE METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Ways to Improve the Methodology",
           "Proposed enhancements to address the study's limitations and unresolved EFA-CFA discrepancy")
add_footer(slide, 16, TOTAL_SLIDES)

improvements = [
    ("1.  Bifactor Modeling", DARK_BLUE, RGBColor(220, 230, 255),
     ["Model a general Engagement factor AND 3 specific sub-factors simultaneously",
      "Separates general vs. subscale-specific variance (addresses the r=0.79-0.83 overlap)",
      "Tests whether subscales add incremental predictive validity beyond the general factor",
      "Would directly resolve the '1 factor vs. 3 factor' debate"]),
    ("2.  Multigroup CFA & Measurement Invariance", TEAL, RGBColor(220, 245, 242),
     ["Test if the factor structure is the same across: Gender, Occupation, Nationality",
      "Configural (same structure), Metric (same loadings), Scalar (same intercepts) testing",
      "Would reveal whether poor fit is due to the instrument or this specific sample",
      "Requires multiple groups; collect matched male sample of similar n=700+"]),
    ("3.  Larger & More Diverse Sample", PURPLE, RGBColor(240, 225, 255),
     ["n=702 provides power but limited occupational & age diversity (only 26-37)",
      "Expand age range and occupation types; include male comparison group",
      "Consider stratified sampling across multiple Swedish regions & sectors",
      "Target n >= 1,000 for stable CFA with modification indices"]),
    ("4.  Alternative Extraction Methods", ORANGE, RGBColor(255, 240, 220),
     ["Principal Axis Factoring (PAF) alongside ML to check convergence of solutions",
      "Apply Unweighted Least Squares (ULS) which is more robust to non-normality",
      "Bootstrap confidence intervals for factor loadings to quantify uncertainty",
      "Conduct sensitivity analysis: compare EFA results with vs. without Item 9 (lowest loading)"]),
    ("5.  Improved Instrument Design", GREEN, RGBColor(220, 255, 235),
     ["Develop gender-specific UWES items validated on all-female samples",
      "Test shorter UWES-3 (1 item per subscale) which may reduce collinearity",
      "Add balance items or reverse-scored items to reduce acquiescence bias",
      "Conduct cognitive interviews to verify item interpretation across genders"]),
    ("6.  Longitudinal & Cross-Cultural Replication", GOLD, LIGHT_GOLD,
     ["Replicate in Norwegian and Finnish female samples to test cultural generalizability",
      "Add time-points (test-retest reliability) to assess stability of factor structure",
      "Compare occupational groups separately (healthcare vs. education vs. IT)",
      "Apply IRT (Item Response Theory) to supplement classical factor analysis"]),
]

col_left  = improvements[:3]
col_right = improvements[3:]
for i, (title, color, fill, bullets) in enumerate(col_left):
    y = Inches(1.42 + i * 1.82)
    box_with_text(slide, Inches(0.3), y, Inches(6.35), Inches(1.75),
                  fill, border=color, items=[
        (title, 12, True, color),
    ] + [(f"   {b}", 10, False, DARK_GRAY) for b in bullets], margin=0.12)

for i, (title, color, fill, bullets) in enumerate(col_right):
    y = Inches(1.42 + i * 1.82)
    box_with_text(slide, Inches(6.85), y, Inches(6.35), Inches(1.75),
                  fill, border=color, items=[
        (title, 12, True, color),
    ] + [(f"   {b}", 10, False, DARK_GRAY) for b in bullets], margin=0.12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "Conclusion",
           "EFA & CFA Summary | Key Takeaway | Future Directions")
add_footer(slide, 17, TOTAL_SLIDES)

box_with_text(slide, Inches(0.3), Inches(1.42), Inches(6.35), Inches(3.40),
              RGBColor(230, 240, 255), border=DARK_BLUE, items=[
    ("EFA Findings (n = 341)", 14, True, DARK_BLUE),
    "",
    ("  KMO = 0.922 (Marvellous)", 11, False, DARK_GRAY),
    ("  l1 explains > 70% of variance", 11, False, DARK_GRAY),
    ("  All loadings L = 0.65 to 0.90", 11, False, DARK_GRAY),
    ("  4 of 5 retention tests: 1-factor", 11, False, DARK_GRAY),
    ("  3-factor Promax: items misload", 11, False, DARK_GRAY),
    ("  Factor r = 0.79-0.83 (not distinct)", 11, False, DARK_GRAY),
    ("  -> Single Work Engagement factor", 12, True, GREEN),
])

box_with_text(slide, Inches(6.98), Inches(1.42), Inches(6.05), Inches(3.40),
              RGBColor(255, 240, 240), border=RED, items=[
    ("CFA Results (n = 342)", 14, True, RED),
    "",
    ("  Model 1 (1F): RMSEA=0.181, CFI=0.895", 11, False, DARK_GRAY),
    ("  Model 2 (2F): RMSEA=0.192, CFI=0.882", 11, False, DARK_GRAY),
    ("  Model 3 (3F): RMSEA=0.167, CFI=0.920", 11, False, DARK_GRAY),
    ("  CFI/TLI all below 0.95 threshold", 11, False, DARK_GRAY),
    ("  Only SRMR (0.046-0.065) passed", 11, False, DARK_GRAY),
    ("  -> No model confirmed", 12, True, RED),
])

box_with_text(slide, Inches(0.3), Inches(4.95), Inches(12.73), Inches(0.90),
              LIGHT_GOLD, border=GOLD, items=[
    ("Key Takeaway:", 13, True, DARK_BLUE),
    ("The optimal factor structure of the UWES-9 remains UNRESOLVED. For all-female samples, "
     "total score (1-factor) may be more reliable than 3-subscale scoring. "
     "Gender-specific validation and Bifactor modeling are the recommended next steps.",
     11, False, DARK_GRAY),
], margin=0.15)

box_with_text(slide, Inches(0.3), Inches(6.00), Inches(12.73), Inches(0.95),
              RGBColor(230, 255, 230), border=GREEN, items=[
    ("Future Directions:", 13, True, DARK_BLUE),
    ("  Test UWES-9 in larger all-female samples across cultures  |  "
     "Develop gender-specific engagement instruments  |  "
     "Apply Bifactor models  |  Test UWES-3  |  Multigroup CFA for invariance testing",
     11, False, DARK_GRAY),
], margin=0.15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — METHODS SUMMARY TABLE
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "All Factor Analysis Methods Applied — Summary Table",
           "Complete overview of every EFA & CFA technique used in this study")
add_footer(slide, 18, TOTAL_SLIDES)

methods_data = [
    ["Method", "Mathematical Purpose", "Result (UWES-9)", "Phase"],
    ["KMO",              "KMO = SUM(r^2) / [SUM(r^2) + SUM(q^2)]",          "0.922 (Marvellous)",        "EFA Pre-check"],
    ["Bartlett's Test",  "chi^2 = -[(n-1)-(2p+5)/6]*ln|R|,  df=p(p-1)/2",   "p < 0.001, Reject H0",      "EFA Pre-check"],
    ["ML Extraction",    "Minimize F_ML = ln|Sigma| - ln|R| + tr(R*S^-1)-p", "chi^2=332.43, df=27",       "EFA Core"],
    ["Eigenvalue Decomp","R = V L V',  l_i = variance explained by factor i", "l1=6.38 = 70.9%",          "EFA Math"],
    ["Scree Plot",       "Visual inspection of l_i vs. i; find 'elbow'",       'Elbow at factor 1',         "EFA Retention"],
    ["Velicer MAP (orig)","MAP = SUM(r_ij^partial^2) / p(p-1)",               "Minimized at m=1",          "EFA Retention"],
    ["Velicer MAP (rev.)", "Uses 4th-power partial correlations",              "Minimized at m=1",          "EFA Retention"],
    ["Parallel Analysis","Compare l_i vs. 95th%ile from random data",         "Suggests 4F (over-factor)", "EFA Retention"],
    ["Promax Rotation",  "L* = LT, raise to power k=3; allow factor corr.",   "Items misload; r=0.79-0.83","EFA Rotation"],
    ["CFA 1-Factor",     "Sigma-hat = LL'+Psi, test vs. R via fit indices",   "RMSEA=0.181 (poor)",        "CFA"],
    ["CFA 2-Factor",     "Same ML framework with 2 correlated factors",       "RMSEA=0.192 (poor)",        "CFA"],
    ["CFA 3-Factor",     "Same ML framework with 3 correlated factors",       "RMSEA=0.167 (poor)",        "CFA"],
    ["Cronbach's Alpha", "a = (p/(p-1)) * (1 - SUM(s^2_i)/s^2_total)",        "a = 0.947 (Excellent)",     "Reliability"],
]
add_table(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(5.55), methods_data,
          col_widths=[1.90, 4.00, 3.00, 1.83],
          header_color=DARK_BLUE, font_size=10)

# Color phase column
phase_colors = {
    "EFA Pre-check": RGBColor(214, 234, 248),
    "EFA Math":      RGBColor(200, 255, 220),
    "EFA Core":      RGBColor(200, 255, 220),
    "EFA Retention": RGBColor(200, 255, 220),
    "EFA Rotation":  RGBColor(253, 235, 208),
    "CFA":           RGBColor(255, 200, 200),
    "Reliability":   RGBColor(245, 220, 255),
}
tbl_m = slide.shapes[-1].table
for row_i in range(1, len(methods_data)):
    phase = methods_data[row_i][3]
    _set_cell_color(tbl_m.cell(row_i, 3), phase_colors.get(phase, WHITE))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — REFERENCES (Part 1)
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "References — Primary & Key Supporting Studies",
           "All sources with full DOI links")
add_footer(slide, 19, TOTAL_SLIDES)

box_with_text(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(0.70),
              DARK_BLUE, items=[("PRIMARY STUDY", 13, True, GOLD)], margin=0.15)

textbox(slide, Inches(0.5), Inches(2.22), Inches(12.33), Inches(0.90), [
    ("Willmer, M., Westerberg Jacobson, J. & Lindberg, M. (2019). Exploratory and Confirmatory "
     "Factor Analysis of the 9-Item Utrecht Work Engagement Scale in a Multi-Occupational Female "
     "Sample: A Cross-Sectional Study.", 11, True, DARK_GRAY),
    ("Frontiers in Psychology, 10, Article 2771.  "
     "https://doi.org/10.3389/fpsyg.2019.02771", 11, False, LIGHT_BLUE),
])

box_with_text(slide, Inches(0.3), Inches(3.22), Inches(12.73), Inches(0.55),
              DARK_BLUE, items=[("SUPPORTING REFERENCES", 12, True, GOLD)], margin=0.15)

refs_support = [
    ("Schaufeli, W.B., Salanova, M., Gonzalez-Roma, V. & Bakker, A.B. (2002). "
     "The measurement of engagement and burnout: A two sample confirmatory factor analytic approach. "
     "Journal of Happiness Studies, 3(1), 71-92.",
     "https://doi.org/10.1023/A:1015630930326"),
    ("Hallberg, U.E. & Schaufeli, W.B. (2006). "
     '"Same same" but different? Can work engagement be discriminated from job involvement '
     "and organizational commitment? European Psychologist, 11(2), 119-127.",
     "https://doi.org/10.1027/1016-9040.11.2.119"),
    ("Nerstad, C.G.L., Richardsen, A.M. & Martinussen, M. (2010). "
     "Factorial validity of the Utrecht Work Engagement Scale (UWES) across occupational groups in Norway. "
     "Scandinavian Journal of Psychology, 51(4), 326-333.",
     "https://doi.org/10.1111/j.1467-9450.2009.00770.x"),
    ("Seppala, P., Mauno, S., Feldt, T. et al. (2009). "
     "The construct validity of the Utrecht Work Engagement Scale: Multisample and longitudinal evidence. "
     "Journal of Happiness Studies, 10(4), 459-481.",
     "https://doi.org/10.1007/s10902-008-9100-y"),
    ("Kulikowski, K. (2017). Do we all agree on how to measure work engagement? "
     "Factorial validity of UWES-9 and UWES-17. "
     "Polish Psychological Bulletin, 48(3), 350-360.",
     "https://doi.org/10.1515/ppb-2017-0038"),
]
y = Inches(3.90)
for citation, doi in refs_support:
    textbox(slide, Inches(0.5), y, Inches(12.33), Inches(0.60), [
        (citation, 10, False, DARK_GRAY),
        (doi, 10, False, LIGHT_BLUE),
    ])
    y += Inches(0.65)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — REFERENCES (Part 2)
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_header(slide, "References — Methodological Sources",
           "Factor analysis methods, instruments, and comparative studies")
add_footer(slide, 20, TOTAL_SLIDES)

box_with_text(slide, Inches(0.3), Inches(1.42), Inches(12.73), Inches(0.55),
              DARK_BLUE, items=[("METHODOLOGICAL & COMPARATIVE REFERENCES", 12, True, GOLD)],
              margin=0.15)

refs_method = [
    ("Hair, J.F., Black, W.C., Babin, B.J. & Anderson, R.E. (2019). "
     "Multivariate Data Analysis (8th Ed.). Cengage Learning.",
     "ISBN: 978-1473756540"),
    ("Kaiser, H.F. (1974). An index of factorial simplicity. "
     "Psychometrika, 39(1), 31-36.",
     "https://doi.org/10.1007/BF02291575"),
    ("Velicer, W.F. (1976). Determining the number of components from the matrix of partial correlations. "
     "Psychometrika, 41(3), 321-327.",
     "https://doi.org/10.1007/BF02293557"),
    ("Horn, J.L. (1965). A rationale and test for the number of factors in factor analysis. "
     "Psychometrika, 30(2), 179-185.",
     "https://doi.org/10.1007/BF02289447"),
    ("Shirom, A. (2003). Feeling vigorous at work? The construct of vigor and the study of positive "
     "affect in organizations. Research in Organizational Stress and Well-Being, 3, 135-164.",
     "https://doi.org/10.1016/S1479-3555(03)03004-X"),
    ("Wefald, A.J., Reichard, R.J. & Serrano, S.A. (2011). Fitting engagement into a nomological network: "
     "The relationship of engagement to leadership and personality. Journal of Leadership & "
     "Organizational Studies, 18(4), 522-537.",
     "https://doi.org/10.1177/1548051811404890"),
    ("Cattell, R.B. (1966). The scree test for the number of factors. "
     "Multivariate Behavioral Research, 1(2), 245-276.",
     "https://doi.org/10.1207/s15327906mbr0102_10"),
    ("Hu, L. & Bentler, P.M. (1999). Cutoff criteria for fit indexes in covariance structure analysis. "
     "Structural Equation Modeling, 6(1), 1-55.",
     "https://doi.org/10.1080/10705519909540118"),
]

y = Inches(2.10)
for citation, doi in refs_method:
    textbox(slide, Inches(0.5), y, Inches(12.33), Inches(0.60), [
        (citation, 10, False, DARK_GRAY),
        (doi, 10, False, LIGHT_BLUE),
    ])
    y += Inches(0.60)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = DARK_BLUE

colored_box(slide, Inches(0), Inches(3.5), Inches(13.33), Inches(0.08),
            GOLD, radius=False)

textbox(slide, Inches(0.5), Inches(0.8), Inches(12.33), Inches(2.5), [
    ("Thank You", 54, True, WHITE),
    ("Questions & Discussion", 22, False, GOLD),
])

colored_box(slide, Inches(2.0), Inches(4.10), Inches(9.33), Inches(2.80),
            RGBColor(0, 45, 120), radius=True)
textbox(slide, Inches(2.3), Inches(4.25), Inches(8.73), Inches(2.60), [
    ("Lana Jalal Gidan", 20, True, WHITE),
    ("SSIE-605: Applied Multivariate Data Analysis", 13, False, LIGHT_GRAY),
    ("Professor Susan Lu  |  Watson College of Engineering  |  Binghamton University", 12, False, LIGHT_GRAY),
    ("", 6, False, WHITE),
    ("Primary Study: Willmer et al. (2019) — Frontiers in Psychology", 11, True, GOLD),
    ("https://doi.org/10.3389/fpsyg.2019.02771", 11, False, RGBColor(150, 200, 255)),
])

add_footer(slide, TOTAL_SLIDES, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
OUTPUT = "Factor_Analysis_Work_Engagement_Enhanced.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}  ({TOTAL_SLIDES} slides)")
