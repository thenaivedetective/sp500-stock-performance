"""
Build Final_Presentation_Video.mp4 from scratch.
Renders each PPTX slide as 1280x720 PNG then stitches with audio via ffmpeg.
"""
import os, io, subprocess, zipfile, shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

PPTX    = 'Final_Presentation_AutoPlay.pptx'
OUTPUT  = 'Final_Presentation_Video.mp4'
FRAMES  = '/tmp/v2_frames'
AUDIO   = '/tmp/v2_audio'
SEGS    = '/tmp/v2_segs'
for d in (FRAMES, AUDIO, SEGS):
    shutil.rmtree(d, ignore_errors=True)
    os.makedirs(d)

W, H = 1280, 720

AUDIO_MAP = {
     4: ('media5.wav',   85.90),
     5: ('media6.wav',   42.75),
     6: ('media7.wav',   25.59),
     7: ('media8.wav',   56.18),
     8: ('media9.wav',   33.42),
     9: ('media10.wav',  31.94),
    10: ('media11.wav',  68.68),
    11: ('media12.wav',  53.95),
    12: ('media13.wav',  47.26),
    13: ('media14.wav', 101.16),
    14: ('media15.wav', 102.76),
}

# ── Extract WAV audio ─────────────────────────────────────────────────────────
print('Extracting audio...')
with zipfile.ZipFile(PPTX, 'r') as z:
    for _, (wav, _) in AUDIO_MAP.items():
        data = z.read(f'ppt/media/{wav}')
        open(f'{AUDIO}/{wav}', 'wb').write(data)

# ── Font helper ───────────────────────────────────────────────────────────────
_font_cache = {}
def get_font(size):
    size = max(10, size)
    if size in _font_cache:
        return _font_cache[size]
    for p in [
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
        '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
        '/usr/share/fonts/truetype/freefont/FreeSans.ttf',
    ]:
        if os.path.exists(p):
            try:
                f = ImageFont.truetype(p, size)
                _font_cache[size] = f
                return f
            except Exception:
                pass
    f = ImageFont.load_default()
    _font_cache[size] = f
    return f

def ep(emu, emu_dim, px_dim):
    return int(int(emu) / int(emu_dim) * px_dim)

def shape_color(shape):
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.type:
                    rgb = run.font.color.rgb
                    return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return (255, 255, 255)

def shape_fontsize(shape):
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    return max(10, int(run.font.size.pt * 96 / 72))
            if para.font and para.font.size:
                return max(10, int(para.font.size.pt * 96 / 72))
    except Exception:
        pass
    return 18

def slide_bg(slide):
    try:
        fill = slide.background.fill
        if fill.type and str(fill.type) == 'SOLID (1)':
            rgb = fill.fore_color.rgb
            return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    try:
        fill = slide.slide_layout.background.fill
        if fill.type and str(fill.type) == 'SOLID (1)':
            rgb = fill.fore_color.rgb
            return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    try:
        fill = slide.slide_layout.slide_master.background.fill
        if fill.type and str(fill.type) == 'SOLID (1)':
            rgb = fill.fore_color.rgb
            return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return (15, 36, 74)   # default navy

# ── Render slides ─────────────────────────────────────────────────────────────
prs = Presentation(PPTX)
ew = int(prs.slide_width)
eh = int(prs.slide_height)
print(f'Rendering {len(prs.slides)} slides...')

for si, slide in enumerate(prs.slides):
    bg = slide_bg(slide)
    img = Image.new('RGB', (W, H), color=bg)
    draw = ImageDraw.Draw(img)

    # Draw pictures first, then text on top
    pics   = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    texts  = [s for s in slide.shapes if s.shape_type not in
              (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA)]

    for shape in pics:
        try:
            l = ep(shape.left  or 0, ew, W)
            t = ep(shape.top   or 0, eh, H)
            w = ep(shape.width or 0, ew, W)
            h = ep(shape.height or 0, eh, H)
            if w < 1 or h < 1:
                continue
            pi = Image.open(io.BytesIO(shape.image.blob))
            pi = pi.resize((w, h), Image.LANCZOS)
            if pi.mode == 'RGBA':
                img.paste(pi, (l, t), pi)
            else:
                img.paste(pi.convert('RGB'), (l, t))
        except Exception:
            pass

    for shape in texts:
        if not getattr(shape, 'has_text_frame', False):
            continue
        if not shape.has_text_frame:
            continue
        try:
            l  = ep(shape.left   or 0, ew, W)
            t  = ep(shape.top    or 0, eh, H)
            w  = ep(shape.width  or 0, ew, W)
            h  = ep(shape.height or 0, eh, H)
            col = shape_color(shape)
            fsz = shape_fontsize(shape)
            font = get_font(fsz)
            y = t + 2
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if not txt:
                    y += fsz // 2
                    continue
                # per-paragraph font/color override
                try:
                    if para.runs and para.runs[0].font.size:
                        fsz2 = max(10, int(para.runs[0].font.size.pt * 96 / 72))
                        font2 = get_font(fsz2)
                    else:
                        fsz2, font2 = fsz, font
                    if para.runs and para.runs[0].font.color and para.runs[0].font.color.type:
                        r = para.runs[0].font.color.rgb
                        col2 = (r.red, r.green, r.blue)
                    else:
                        col2 = col
                except Exception:
                    fsz2, font2, col2 = fsz, font, col

                # word-wrap
                words = txt.split()
                lines, cur = [], ''
                max_w = max(1, w - 8)
                for word in words:
                    test = (cur + ' ' + word).strip()
                    bb = draw.textbbox((0, 0), test, font=font2)
                    if bb[2] <= max_w:
                        cur = test
                    else:
                        if cur:
                            lines.append(cur)
                        cur = word
                if cur:
                    lines.append(cur)

                for line in lines:
                    if y + fsz2 > t + h + 10:
                        break
                    # shadow
                    draw.text((l + 5, y + 1), line, font=font2, fill=(0, 0, 0))
                    draw.text((l + 4, y),     line, font=font2, fill=col2)
                    y += fsz2 + 2
                y += 3
        except Exception:
            pass

    path = f'{FRAMES}/slide_{si:02d}.png'
    img.save(path, 'PNG')
    non_bg = sum(1 for px in img.getdata() if px != bg)
    print(f'  Slide {si+1:2d} saved  ({non_bg:,} non-bg pixels)')

# ── Build per-slide MP4 segments ──────────────────────────────────────────────
print('\nBuilding segments...')
segs = []
for si in range(len(prs.slides)):
    frame = f'{FRAMES}/slide_{si:02d}.png'
    seg   = f'{SEGS}/seg_{si:02d}.mp4'

    if si in AUDIO_MAP:
        wav, dur = AUDIO_MAP[si]
        cmd = [
            'ffmpeg', '-y',
            '-loop', '1', '-i', frame,
            '-i', f'{AUDIO}/{wav}',
            '-c:v', 'libx264', '-preset', 'ultrafast', '-tune', 'stillimage',
            '-c:a', 'aac', '-b:a', '128k',
            '-pix_fmt', 'yuv420p',
            '-t', str(dur + 0.5), '-shortest', seg,
        ]
        tag = f'{dur:.0f}s audio'
    else:
        cmd = [
            'ffmpeg', '-y',
            '-loop', '1', '-i', frame,
            '-c:v', 'libx264', '-preset', 'ultrafast', '-tune', 'stillimage',
            '-pix_fmt', 'yuv420p',
            '-t', '5', '-an', seg,
        ]
        tag = '5s silent'

    r = subprocess.run(cmd, capture_output=True)
    if r.returncode == 0:
        segs.append(seg)
        print(f'  Seg {si+1:2d} done  ({tag})')
    else:
        print(f'  ERROR seg {si+1}: {r.stderr.decode()[-150:]}')

# ── Concatenate ────────────────────────────────────────────────────────────────
print(f'\nConcatenating {len(segs)} segments → {OUTPUT}')
cl = '/tmp/v2_concat.txt'
open(cl, 'w').write('\n'.join(f"file '{s}'" for s in segs) + '\n')
r = subprocess.run(
    ['ffmpeg', '-y', '-f', 'concat', '-safe', '0', '-i', cl, '-c', 'copy', OUTPUT],
    capture_output=True
)
if r.returncode == 0:
    sz = os.path.getsize(OUTPUT) / 1024 / 1024
    print(f'\nDone!  {OUTPUT}  ({sz:.1f} MB)')
else:
    print('Concat error:', r.stderr.decode()[-300:])
