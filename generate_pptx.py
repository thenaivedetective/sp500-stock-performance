"""
Generates a professional .pptx presentation for the credit risk analysis project.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ── Colours ────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2E)
NAVY2   = RGBColor(0x1A, 0x2F, 0x4A)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
LBLUE   = RGBColor(0x60, 0xA5, 0xFA)
GOLD    = RGBColor(0xF5, 0x9E, 0x0B)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
RED     = RGBColor(0xEF, 0x44, 0x44)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGREY   = RGBColor(0xCB, 0xD5, 0xE1)
DGREY   = RGBColor(0x33, 0x41, 0x55)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helper functions ───────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width     = line_w
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, italic=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf  = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return box

def tag(slide, text, x, y):
    txbox(slide, text.upper(), x, y, 8, 0.35,
          size=9, bold=True, color=LBLUE, font="Calibri")

def heading(slide, text, x, y, w=12, size=32, color=WHITE):
    txbox(slide, text, x, y, w, 1.0, size=size, bold=True, color=color)

def body(slide, text, x, y, w=5.5, h=0.5, size=14, color=LGREY):
    txbox(slide, text, x, y, w, h, size=size, color=color, wrap=True)

def bullet_box(slide, items, x, y, w=5.5, size=13, color=LGREY, dot_color=BLUE):
    bx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5))
    bx.word_wrap = True
    tf = bx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # dot
        r1 = p.add_run(); r1.text = "▸  "
        r1.font.color.rgb = dot_color; r1.font.size = Pt(size); r1.font.name = "Calibri"
        # text
        r2 = p.add_run(); r2.text = item
        r2.font.color.rgb = color; r2.font.size = Pt(size); r2.font.name = "Calibri"
        p.space_after = Pt(6)
    return bx

def stat_card(slide, stat, label, x, y, stat_color=BLUE):
    rect(slide, x, y, 2.8, 1.5, fill=NAVY2)
    txbox(slide, stat,  x+0.1, y+0.1, 2.6, 0.85,
          size=36, bold=True, color=stat_color, align=PP_ALIGN.CENTER)
    txbox(slide, label, x+0.1, y+0.9, 2.6, 0.5,
          size=11, color=LGREY, align=PP_ALIGN.CENTER)

def add_image(slide, path, x, y, w):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

def card(slide, title, body_text, x, y, w=5.5, h=1.5, title_color=LBLUE, fill=NAVY2):
    rect(slide, x, y, w, h, fill=fill)
    txbox(slide, title,     x+0.18, y+0.15, w-0.3, 0.35,
          size=13, bold=True, color=title_color)
    txbox(slide, body_text, x+0.18, y+0.52, w-0.3, h-0.55,
          size=11, color=LGREY, wrap=True)

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)

def accent_bar(slide, color=BLUE):
    rect(slide, 0, 6.9, 13.33, 0.08, fill=color)

def slide_num(slide, n, total=15):
    txbox(slide, f"{n} / {total}", 12.3, 7.15, 1.0, 0.3,
          size=9, color=DGREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
# gradient overlay rectangle
rect(sl, 0, 0, 7, 7.5, fill=RGBColor(0x0F, 0x2D, 0x52))
rect(sl, 0, 6.6, 13.33, 0.12, fill=BLUE)

txbox(sl, "MULTIVARIATE STATISTICS · BINGHAMTON UNIVERSITY",
      0.5, 0.6, 9, 0.4, size=10, bold=True, color=LBLUE)
txbox(sl, "Credit Risk Classification\nUsing ML Techniques",
      0.5, 1.1, 9, 2.4, size=42, bold=True, color=WHITE, font="Calibri")
txbox(sl, "Applying Factor Analysis, Linear Discriminant Analysis, and\n"
          "Logistic Regression to 500,000 real LendingClub loan records.",
      0.5, 3.6, 8, 0.9, size=15, color=LGREY)

# Meta boxes
for i,(label,val) in enumerate([
    ("Author","Lana Jalal Gidan"),
    ("Department","Systems Science & Industrial Engineering"),
    ("Dataset","LendingClub · 500,000 Loans"),
    ("Models","LDA & Logistic Regression"),
]):
    xi = 0.5 + i*3.2
    txbox(sl, label.upper(), xi, 5.0, 3.0, 0.3, size=8, bold=True, color=LBLUE)
    txbox(sl, val,           xi, 5.3, 3.0, 0.5, size=12, bold=True, color=WHITE)

slide_num(sl, 1)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Roadmap", 0.5, 0.3)
heading(sl, "What We'll Cover", 0.5, 0.65, size=30)

items = [
    ("1","The Problem","What is credit risk?"),
    ("2","Dataset","LendingClub · 500k real loans"),
    ("3","Data Cleaning","From messy 151 columns to 76 clean features"),
    ("4","Factor Analysis","Interdependence technique — 6 latent factors"),
    ("5","LDA","Linear Discriminant Analysis — separating classes"),
    ("6","Logistic Regression","Probability-based binary classification"),
    ("7","Model Comparison","ROC · Accuracy · AUC"),
    ("8","Live Predictor","Interactive demo — real-time prediction"),
]
cols = [(0.4,1.6),(6.8,1.6)]
for i, (num, title, sub) in enumerate(items):
    col = i % 2
    row = i // 2
    xi, yi = cols[col][0], cols[col][1] + row * 1.25
    rect(sl, xi, yi, 5.9, 1.1, fill=NAVY2)
    txbox(sl, num,   xi+0.12, yi+0.1,  0.5,  0.9, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, xi+0.5, yi+0.05, 0.06, 1.0, fill=BLUE)
    txbox(sl, title, xi+0.7, yi+0.05, 5.0, 0.5, size=14, bold=True, color=WHITE)
    txbox(sl, sub,   xi+0.7, yi+0.52, 5.0, 0.5, size=11, color=LGREY)
slide_num(sl, 2)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Context", 0.5, 0.3)
heading(sl, "What Is Credit Risk?", 0.5, 0.65, size=30)

card(sl, "💳  The Loan",
     "A bank lends you money today. You promise to repay it over time — with interest. "
     "The bank earns interest; you get capital.", 0.4, 1.5, 5.8, 1.4)
card(sl, "⚠️  The Risk",
     "The borrower might DEFAULT — stop paying. Banks lend billions, so even a small "
     "default rate means massive losses.", 0.4, 3.05, 5.8, 1.4, title_color=GOLD,
     fill=RGBColor(0x1A, 0x17, 0x10))
card(sl, "🎯  The Goal",
     "Before lending, predict: 'Will this person repay?' A good model lets banks "
     "charge fair rates and avoid bad loans.", 0.4, 4.60, 5.8, 1.4, title_color=GREEN,
     fill=RGBColor(0x10, 0x1F, 0x1A))

stat_card(sl, "$82B+",  "US consumer loan\nlosses per year",   6.8, 1.5, stat_color=RED)
stat_card(sl, "~71%",   "Model accuracy using\norigination features", 9.8, 1.5, stat_color=GREEN)
card(sl, "Why Multivariate Statistics?",
     "A FICO score is just one number. Default risk depends on DOZENS of factors "
     "interacting simultaneously. Multivariate methods see the full picture.",
     6.8, 3.2, 6.0, 2.8, title_color=LBLUE)
slide_num(sl, 3)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — DATASET
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Data", 0.5, 0.3)
heading(sl, "The LendingClub Dataset", 0.5, 0.65, size=30)

stat_card(sl, "500K",  "Real loan records",          0.4, 1.6)
stat_card(sl, "151",   "Features per loan",           3.4, 1.6)
stat_card(sl, "20.9%", "Default rate in data",        6.4, 1.6, stat_color=RED)
stat_card(sl, "$15K",  "Avg loan amount",             9.4, 1.6)

bullet_box(sl, [
    "Loan info: amount, term, interest rate, monthly installment",
    "Borrower info: income, employment, debt-to-income ratio",
    "Credit history: FICO score, open accounts, delinquencies",
    "LendingClub grade: A (safest) → G (riskiest)",
    "Outcome: Did they repay or default?",
], 0.4, 3.4, 6.0, size=13)

# Status table
txbox(sl, "Loan Status", 6.8, 3.3, 5.5, 0.4, size=13, bold=True, color=LBLUE)
rows = [
    ("Fully Paid",       "312,340", GREEN),
    ("Charged Off",      " 78,824", RED),
    ("Late 31–120 days", "  2,977", RED),
    ("Current",          "104,240", LGREY),
    ("In Grace Period",  "  1,046", LGREY),
]
for i,(status,count,col) in enumerate(rows):
    y = 3.8 + i*0.54
    fill = NAVY2 if i%2==0 else RGBColor(0x10,0x1C,0x30)
    rect(sl, 6.8, y, 5.9, 0.52, fill=fill)
    txbox(sl, status, 7.0,  y+0.1, 3.5, 0.35, size=12, color=LGREY)
    txbox(sl, count,  10.5, y+0.1, 1.8, 0.35, size=12, bold=True, color=col)
slide_num(sl, 4)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — DATA CLEANING
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Preprocessing", 0.5, 0.3)
heading(sl, "Data Cleaning Pipeline", 0.5, 0.65, size=30)

steps = [
    ("1","Dropped 58 columns","with >40% missing values — hardship/settlement fields"),
    ("2","Dropped 17 columns","irrelevant identifiers and free-text descriptions"),
    ("3","Engineered target","coded loan_status → binary: 0 = paid, 1 = default"),
    ("4","Parsed text fields",'"13.5%" → 13.5, "36 months" → 36, "10+ years" → 10'),
    ("5","Label-encoded","grade, home_ownership, purpose → numeric values"),
    ("6","Removed outliers","cut top/bottom 1% of income, loan amount, revolving balance"),
    ("7","Dropped NaN rows","listwise deletion for complete-case analysis"),
]
for i,(num,title,sub) in enumerate(steps):
    y = 1.5 + i*0.72
    rect(sl, 0.4, y, 0.55, 0.55, fill=BLUE)
    txbox(sl, num,   0.4,  y+0.05, 0.55, 0.45, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, title, 1.1,  y+0.02, 3.5,  0.3,  size=13, bold=True, color=WHITE)
    txbox(sl, sub,   1.1,  y+0.3,  5.8,  0.35, size=11, color=LGREY)

rect(sl, 7.3, 1.5, 5.5, 1.7, fill=NAVY2)
txbox(sl, "Before  →  After", 7.5, 1.55, 5.0, 0.4, size=12, color=LGREY, align=PP_ALIGN.CENTER)
txbox(sl, "500K × 151", 7.5, 1.95, 2.3, 0.6, size=22, bold=True, color=LGREY, align=PP_ALIGN.CENTER)
txbox(sl, "→",          9.9, 1.95, 0.8, 0.6, size=22, bold=True, color=BLUE,  align=PP_ALIGN.CENTER)
txbox(sl, "292K × 76",  10.5,1.95, 2.1, 0.6, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

card(sl, "⚖️  Class Imbalance Handling",
     "79% good loans vs 21% defaults. We balanced training at 2:1 ratio (undersampling) "
     "so the model learns real patterns — not just 'always predict good.'",
     7.3, 3.4, 5.5, 1.4, title_color=GOLD, fill=RGBColor(0x1A,0x17,0x10))
card(sl, "📏  Z-Score Standardization",
     "All features scaled to mean=0, std=1 so dollar amounts don't overpower "
     "year-based variables. Required before LDA and Logistic Regression.",
     7.3, 4.95, 5.5, 1.4, title_color=GREEN, fill=RGBColor(0x10,0x1F,0x1A))
slide_num(sl, 5)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — FACTOR ANALYSIS METHOD
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Interdependence Technique", 0.5, 0.3)
heading(sl, "Factor Analysis", 0.5, 0.65, size=30)

card(sl, "🔍  What It Does",
     "Groups many correlated variables into a smaller number of hidden (latent) FACTORS. "
     "No outcome variable — purely about structure in the data.",
     0.4, 1.5, 6.0, 1.3)
card(sl, "💡  The Analogy",
     "You can't measure 'intelligence' directly — but you see it reflected in test scores "
     "across math, reading, and logic. FA finds that hidden dimension automatically.",
     0.4, 2.95, 6.0, 1.5, title_color=GOLD, fill=RGBColor(0x1A,0x17,0x10))
card(sl, "✅  Bartlett's Test — PASSED",
     "χ² = 657,867   |   p < 0.001\nVariables are significantly correlated. FA is appropriate.",
     0.4, 4.6, 6.0, 1.3, title_color=GREEN, fill=RGBColor(0x10,0x1F,0x1A))

txbox(sl, "The 6 Factors Found", 6.9, 1.4, 5.8, 0.45, size=14, bold=True, color=LBLUE)
factors = [
    ("F1","Loan Size & Cost","14.6%"),
    ("F2","Credit Risk Pricing","13.7%"),
    ("F3","Credit Utilization","10.0%"),
    ("F4","Credit Breadth","9.4%"),
    ("F5","Credit Limit Stress","3.9%"),
    ("F6","Wealth & Assets","4.8%"),
]
txbox(sl, "TOTAL VARIANCE EXPLAINED", 6.9, 6.0, 4.0, 0.35, size=10, bold=True, color=LGREY)
txbox(sl, "56.5%", 6.9, 6.35, 4.0, 0.55, size=30, bold=True, color=GREEN)

for i,(f,name,pct) in enumerate(factors):
    y = 1.9 + i*0.65
    fill = NAVY2 if i%2==0 else RGBColor(0x10,0x1C,0x30)
    rect(sl, 6.9, y, 5.9, 0.6, fill=fill)
    txbox(sl, f,    7.05, y+0.1, 0.6, 0.4, size=13, bold=True, color=BLUE)
    txbox(sl, name, 7.7,  y+0.1, 3.8, 0.4, size=12, color=WHITE)
    txbox(sl, pct,  11.5, y+0.1, 1.2, 0.4, size=13, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
slide_num(sl, 6)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — FACTOR ANALYSIS CHARTS
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Interdependence Technique · Results", 0.5, 0.3)
heading(sl, "Factor Analysis — Visual Results", 0.5, 0.65, size=28)

txbox(sl, "Scree Plot — How Many Factors?",
      0.4, 1.5, 6.0, 0.4, size=13, bold=True, color=LBLUE)
add_image(sl, "results/fa_scree_plot.png",  0.4, 1.95, 6.0)

txbox(sl, "Factor Loadings Heatmap",
      6.9, 1.5, 6.0, 0.4, size=13, bold=True, color=LBLUE)
add_image(sl, "results/fa_loadings_heatmap.png", 6.9, 1.95, 5.8)

txbox(sl, "The 'elbow' at factor 6 confirms we keep 6 factors (eigenvalue > 1).",
      0.4, 6.55, 5.8, 0.4, size=10, italic=True, color=LGREY)
txbox(sl, "Dark blue = strong positive loading. Dark red = strong negative loading.",
      6.9, 6.55, 5.8, 0.4, size=10, italic=True, color=LGREY)
slide_num(sl, 7)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — LDA METHOD
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Dependence Technique 1", 0.5, 0.3)
heading(sl, "Linear Discriminant Analysis (LDA)", 0.5, 0.65, size=28)

card(sl, "🎯  What It Does",
     "Finds the linear combination of variables that MAXIMALLY SEPARATES two groups — "
     "Good Borrowers vs. Defaulters. Produces one 'discriminant score' per borrower.",
     0.4, 1.5, 5.8, 1.4)
card(sl, "💡  The Intuition",
     "All borrowers are dots in a multi-dimensional space. LDA finds the single direction "
     "where the blue (good) and red (default) groups overlap as little as possible.",
     0.4, 3.05, 5.8, 1.4, title_color=GOLD, fill=RGBColor(0x1A,0x17,0x10))

txbox(sl, "Assumptions:", 0.4, 4.6, 5.8, 0.35, size=12, bold=True, color=LBLUE)
bullet_box(sl, [
    "Variables roughly normally distributed within each class",
    "Equal covariance matrices across both classes",
    "Linear boundary between the two groups",
], 0.4, 5.0, 5.8, size=12)

txbox(sl, "LDA Projection — Class Separation",
      6.9, 1.5, 6.0, 0.4, size=13, bold=True, color=LBLUE)
add_image(sl, "results/lda_projection.png", 6.9, 1.95, 6.0)
txbox(sl, "Two peaks = two classes. Less overlap = better model.",
      6.9, 6.5, 5.8, 0.4, size=10, italic=True, color=LGREY)
slide_num(sl, 8)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — LDA RESULTS
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Dependence Technique 1 · Results", 0.5, 0.3)
heading(sl, "LDA — Performance Results", 0.5, 0.65, size=28)

stat_card(sl, "97.3%", "Accuracy",  0.4, 1.5, stat_color=BLUE)
stat_card(sl, "0.997", "AUC-ROC",   3.4, 1.5, stat_color=GREEN)
stat_card(sl, "0.975", "Precision", 6.4, 1.5, stat_color=GOLD)
stat_card(sl, "0.960", "F1-Score",  9.4, 1.5, stat_color=LBLUE)

txbox(sl, "Top Discriminant Variables:", 0.4, 3.4, 6.0, 0.4, size=12, bold=True, color=LBLUE)
bullet_box(sl, [
    "FICO Score — single most powerful separator between classes",
    "Interest Rate — encodes the lender's risk assessment",
    "Debt Settlement Flag — red flag indicating financial distress",
    "Installment Amount — higher monthly burden = more risk",
], 0.4, 3.85, 6.0, size=12)

add_image(sl, "results/lda_confusion_matrix.png", 7.0, 3.2, 5.8)
txbox(sl, "Confusion Matrix — Diagonal = Correct Predictions",
      7.0, 3.0, 5.8, 0.35, size=12, bold=True, color=LBLUE)
slide_num(sl, 9)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — LOGISTIC REGRESSION METHOD
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Dependence Technique 2", 0.5, 0.3)
heading(sl, "Logistic Regression", 0.5, 0.65, size=30)

card(sl, "📊  What It Does",
     "Predicts the PROBABILITY of default — a number between 0% and 100%. Uses the "
     "sigmoid (S-shaped) function. If probability > 50% → predict Default.",
     0.4, 1.5, 6.0, 1.4)

txbox(sl, "LDA  vs  Logistic Regression:", 0.4, 3.1, 6.0, 0.4, size=13, bold=True, color=LBLUE)
rows_lr = [("LDA","Logistic Regression"),
           ("Assumes normality","No distribution assumption"),
           ("Maximizes class separation","Maximizes likelihood"),
           ("Outputs discriminant score","Outputs probability (0–1)")]
for i,(a,b) in enumerate(rows_lr):
    y = 3.6 + i*0.55
    fill = BLUE if i==0 else (NAVY2 if i%2==1 else RGBColor(0x10,0x1C,0x30))
    col  = WHITE if i==0 else LGREY
    rect(sl, 0.4,  y, 2.8, 0.52, fill=fill)
    rect(sl, 3.25, y, 3.1, 0.52, fill=fill)
    txbox(sl, a, 0.5,  y+0.1, 2.6, 0.35, size=11, bold=(i==0), color=col)
    txbox(sl, b, 3.35, y+0.1, 2.8, 0.35, size=11, bold=(i==0), color=col)

card(sl, "✅  Why LR Wins Here",
     "Financial variables are NOT normally distributed — incomes and balances are "
     "right-skewed. LR's weaker assumptions make it more robust and it outperforms LDA.",
     0.4, 5.75, 6.0, 1.1, title_color=GREEN, fill=RGBColor(0x10,0x1F,0x1A))

txbox(sl, "Top Coefficients — What Drives Default?",
      6.9, 1.5, 6.0, 0.4, size=13, bold=True, color=LBLUE)
add_image(sl, "results/lr_coefficients.png", 6.9, 1.95, 6.0)
txbox(sl, "Blue = lowers default risk   |   Red = raises default risk",
      6.9, 6.5, 5.8, 0.4, size=10, italic=True, color=LGREY)
slide_num(sl, 10)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — LR RESULTS
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Dependence Technique 2 · Results", 0.5, 0.3)
heading(sl, "Logistic Regression — Performance Results", 0.5, 0.65, size=26)

stat_card(sl, "99.6%",  "Accuracy",  0.4, 1.5, stat_color=BLUE)
stat_card(sl, "0.9997", "AUC-ROC",   3.4, 1.5, stat_color=GREEN)
stat_card(sl, "0.998",  "Precision", 6.4, 1.5, stat_color=GOLD)
stat_card(sl, "0.995",  "F1-Score",  9.4, 1.5, stat_color=LBLUE)

card(sl, "⭐  LR Outperforms LDA",
     "LR beats LDA by +2.3% accuracy and +0.003 AUC. Consistent with theory: "
     "LDA's normality assumptions are violated by skewed financial data.",
     0.4, 3.3, 5.8, 1.3, title_color=GOLD, fill=RGBColor(0x1A,0x17,0x10))
bullet_box(sl, [
    "Precision 1.00 on defaults — almost zero false alarms",
    "Recall 0.99 on defaults — almost zero missed defaults",
    "F1-Score 0.99 — excellent balance of precision and recall",
    "AUC = 0.9997 — near-perfect discriminative ability",
], 0.4, 4.75, 5.8, size=12)

add_image(sl, "results/lr_confusion_matrix.png", 7.0, 3.2, 5.8)
txbox(sl, "Confusion Matrix",
      7.0, 3.0, 5.8, 0.35, size=12, bold=True, color=LBLUE)
slide_num(sl, 11)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — ROC COMPARISON
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Model Comparison", 0.5, 0.3)
heading(sl, "LDA vs Logistic Regression — ROC Comparison", 0.5, 0.65, size=26)

add_image(sl, "results/roc_comparison.png", 0.4, 1.5, 6.5)

card(sl, "📈  What Is the ROC Curve?",
     "Plots True Positive Rate (catching actual defaults) vs. False Positive Rate "
     "(wrongly flagging good loans) at every possible decision threshold.",
     7.1, 1.5, 5.8, 1.4)

txbox(sl, "Results Summary:", 7.1, 3.1, 5.8, 0.4, size=13, bold=True, color=LBLUE)
rows_roc = [("Model","Accuracy","AUC-ROC"),
            ("LDA","97.3%","0.997"),
            ("Logistic Regression","99.6%","0.9997")]
for i,(m,acc,auc) in enumerate(rows_roc):
    y = 3.6 + i*0.62
    fill = BLUE if i==0 else (NAVY2 if i%2==1 else RGBColor(0x10,0x1C,0x30))
    col  = WHITE if i==0 else LGREY
    rect(sl, 7.1, y, 5.8, 0.58, fill=fill)
    txbox(sl, m,   7.2, y+0.1, 2.8, 0.4, size=12, bold=(i==0), color=col)
    txbox(sl, acc, 9.8, y+0.1, 1.5, 0.4, size=12, bold=(i==0), color=(GREEN if i>0 else col), align=PP_ALIGN.CENTER)
    txbox(sl, auc,11.2, y+0.1, 1.5, 0.4, size=12, bold=(i==0), color=(GREEN if i>0 else col), align=PP_ALIGN.CENTER)

card(sl, "📐  AUC Interpretation",
     "AUC = 1.0 → perfect. AUC = 0.5 → random guessing.\n"
     "Both models are at 0.997+ — exceptional discrimination.",
     7.1, 5.65, 5.8, 1.2, title_color=GREEN, fill=RGBColor(0x10,0x1F,0x1A))
slide_num(sl, 12)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — LIVE PREDICTOR (info slide)
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl, color=GOLD)
tag(sl, "★ Interactive Demo", 0.5, 0.3)
heading(sl, "Live Credit Risk Predictor", 0.5, 0.65, size=30, color=GOLD)

txbox(sl, "Open the live web app to enter real loan values and get instant predictions from the trained model.",
      0.5, 1.4, 12.3, 0.6, size=15, color=LGREY)

# 4 preset cards
presets = [
    ("✅","Ideal Borrower","FICO 790 · Income $120K · DTI 8%\nRevol Util 12% · Grade A",GREEN,"Will Repay"),
    ("⚖️","Average Borrower","FICO 700 · Income $65K · DTI 18%\nRevol Util 45% · Grade C",LBLUE,"Will Repay"),
    ("⚠️","Risky Borrower","FICO 620 · Income $42K · DTI 32%\nRevol Util 78% · Grade E",GOLD,"Borderline"),
    ("🚨","Very High Risk","FICO 560 · Income $28K · DTI 40%\nRevol Util 95% · Grade G",RED,"Default Risk"),
]
for i,(icon,name,details,color,verdict) in enumerate(presets):
    x = 0.4 + i*3.25
    rect(sl, x, 2.2, 3.0, 3.5, fill=NAVY2)
    txbox(sl, icon,    x+1.2, 2.3, 0.8, 0.7, size=28, align=PP_ALIGN.CENTER)
    txbox(sl, name,    x+0.1, 3.1, 2.8, 0.5, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    txbox(sl, details, x+0.1, 3.65,2.8, 0.9, size=11, color=LGREY, align=PP_ALIGN.CENTER)
    rect(sl, x+0.2, 4.7, 2.6, 0.55, fill=color)
    txbox(sl, verdict, x+0.2, 4.75,2.6, 0.45, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txbox(sl, "The live predictor (web app) lets you:",
      0.5, 6.0, 6.0, 0.4, size=12, bold=True, color=LBLUE)
bullet_box(sl, [
    "Switch between Logistic Regression and LDA with one click",
    "Enter any custom loan values across 15 variables",
    "See probability bars, risk tier, and top 3 risk drivers instantly",
], 0.5, 6.4, 12.3, size=12)
slide_num(sl, 13)

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — LIMITATION
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl, color=RED)
tag(sl, "Critical Analysis", 0.5, 0.3)
heading(sl, "Key Limitation: Data Leakage", 0.5, 0.65, size=30)

txbox(sl, "The full dataset models (97–99%) include variables only known AFTER the loan resolves.",
      0.5, 1.4, 12.3, 0.5, size=14, color=LGREY)

leaks = [
    ("total_rec_prncp","Total principal recovered — only exists after the loan ends"),
    ("recoveries","Money recovered post-default — only populated IF there was a default"),
    ("collection_recovery_fee","Fee charged to recover defaulted funds — post-event only"),
    ("last_pymnt_amnt","Last payment made — known only after payment history unfolds"),
]
for i,(var,desc) in enumerate(leaks):
    y = 2.0 + i*0.8
    rect(sl, 0.4, y, 5.8, 0.7, fill=RGBColor(0x1F,0x10,0x10))
    txbox(sl, "🔴  "+var, 0.6, y+0.05, 5.4, 0.3, size=13, bold=True, color=RED)
    txbox(sl, desc,       0.6, y+0.35, 5.4, 0.3, size=11, color=LGREY)

card(sl, "💡  The Analogy",
     "Trying to predict if a student will fail — and using their final exam score as a "
     "predictor. Of course the model is perfect. You're using the answer to predict the answer.",
     6.5, 2.0, 6.4, 1.5, title_color=GOLD, fill=RGBColor(0x1A,0x17,0x10))
card(sl, "✅  Our Live Predictor Fixes This",
     "The interactive predictor uses ONLY 15 origination-time features — variables known "
     "at application time. That is why accuracy is ~71%, not 99%. Honest and realistic.",
     6.5, 3.65, 6.4, 1.5, title_color=GREEN, fill=RGBColor(0x10,0x1F,0x1A))
card(sl, "📚  What Real Banks Do",
     "Only use application-time data · Must comply with ECOA & fair lending laws · "
     "Models are audited for discriminatory patterns before deployment.",
     6.5, 5.3, 6.4, 1.35, title_color=LBLUE)
slide_num(sl, 14)

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSION
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); accent_bar(sl)
tag(sl, "Conclusion", 0.5, 0.3)
heading(sl, "Summary & Takeaways", 0.5, 0.65, size=30)

cards_c = [
    ("🧹","Real Data, Real Cleaning","500K messy records → 292K clean rows through a rigorous 7-step pipeline.",WHITE),
    ("🔬","Factor Analysis","6 latent factors — Loan Cost, Risk Pricing, Utilization, Breadth, Stress, Wealth — explaining 56.5% of variance.",WHITE),
    ("📐","LDA","Separated good from default borrowers: 97.3% accuracy, AUC = 0.997.",BLUE),
    ("📊","Logistic Regression","Outperformed LDA: 99.6% accuracy, AUC = 0.9997. Better fit for skewed financial data.",GREEN),
    ("⚠️","Data Leakage","High full-model accuracy is partly due to post-outcome variables. Origination-only model gives honest ~71%.",RED),
    ("⚡","Live Predictor","Interactive web app: enter loan values and watch the trained model decide in real time.",GOLD),
]
for i,(icon,title,txt,col) in enumerate(cards_c):
    r,c = divmod(i,3)
    x = 0.4 + c*4.3
    y = 1.6 + r*2.5
    rect(sl, x, y, 4.0, 2.3, fill=NAVY2)
    txbox(sl, icon,  x+1.5, y+0.1, 1.0, 0.7, size=24, align=PP_ALIGN.CENTER)
    txbox(sl, title, x+0.1, y+0.8, 3.8, 0.5, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    txbox(sl, txt,   x+0.1, y+1.3, 3.8, 0.9, size=10, color=LGREY, align=PP_ALIGN.CENTER, wrap=True)

# Bottom stats row
for val,lbl,col in [("292K","Clean Records",WHITE),("6","Latent Factors",BLUE),
                     ("56.5%","Variance Explained",GREEN),
                     ("0.997","LDA AUC",LBLUE),("0.9997","LR AUC",GOLD)]:
    pass  # already shown in cards

slide_num(sl, 15)

# ── Save ──────────────────────────────────────────────────────────
OUT = "Gidan_CreditRisk_Presentation.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
