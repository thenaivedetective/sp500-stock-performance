import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BLUE = RGBColor(0, 32, 96)
MEDIUM_BLUE = RGBColor(0, 70, 140)
LIGHT_BLUE = RGBColor(0, 112, 192)
ACCENT_GOLD = RGBColor(255, 192, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(240, 240, 240)

os.makedirs("figures2", exist_ok=True)

def add_slide_background(slide, color=DARK_BLUE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=None):
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = ACCENT_GOLD
        p2.alignment = PP_ALIGN.LEFT

def add_content_box(slide, left, top, width, height, text_items, font_size=16, bullet=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(text_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.font.size = Pt(item[1])
            if len(item) > 2:
                p.font.bold = item[2]
            if len(item) > 3:
                p.font.color.rgb = item[3]
            else:
                p.font.color.rgb = DARK_GRAY
        else:
            prefix = "\u2022 " if bullet else ""
            p.text = f"{prefix}{item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    return txBox

# ============================================================
# VISUALIZATION FUNCTIONS
# ============================================================

def create_research_framework():
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis('off')

    factors = [
        ('Brand\nAwareness', 1.5, 7.0),
        ('Brand\nImage', 1.5, 5.5),
        ('Perceived\nQuality', 1.5, 4.0),
        ('Brand\nAssociation', 1.5, 2.5),
        ('Brand\nLoyalty', 1.5, 1.0),
    ]
    for label, x, y in factors:
        box = plt.Rectangle((x-1.0, y-0.45), 2.0, 0.9, linewidth=2,
                           edgecolor='#002060', facecolor='#D6E4F0', zorder=3)
        ax.add_patch(box)
        ax.text(x, y, label, ha='center', va='center', fontsize=10,
               fontweight='bold', color='#002060', zorder=4)

    be_box = plt.Rectangle((6.5, 3.55), 2.5, 1.0, linewidth=3,
                           edgecolor='#002060', facecolor='#FFC000', zorder=3)
    ax.add_patch(be_box)
    ax.text(7.75, 4.05, 'Brand\nEquity', ha='center', va='center', fontsize=14,
           fontweight='bold', color='#002060', zorder=4)

    hyp_labels = ['H1', 'H2', 'H3', 'H4', 'H5']
    colors_arr = ['#FF4444', '#00AA00', '#888888', '#888888', '#00AA00']
    for i, (label, x, y) in enumerate(factors):
        ax.annotate('', xy=(6.5, 4.05), xytext=(2.5, y),
                   arrowprops=dict(arrowstyle='->', lw=2, color=colors_arr[i]))
        mid_x = (2.5 + 6.5) / 2
        mid_y = (y + 4.05) / 2
        ax.text(mid_x, mid_y + 0.15, hyp_labels[i], fontsize=9,
               fontweight='bold', color=colors_arr[i], ha='center',
               bbox=dict(boxstyle='round,pad=0.2', facecolor='white', edgecolor=colors_arr[i], alpha=0.8))

    ax.text(5.0, 0.2, 'Source: Gilitwala & Nag (2022)', fontsize=9, ha='center',
           fontstyle='italic', color='gray')
    fig.suptitle('Research Framework: Disney Brand Equity', fontsize=14,
                fontweight='bold', color='#002060', y=0.98)
    plt.tight_layout()
    plt.savefig('figures2/research_framework.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_cronbach_alpha_chart():
    variables = ['Brand\nAwareness', 'Brand\nImage', 'Perceived\nQuality',
                 'Brand\nAssociation', 'Brand\nLoyalty', 'Brand\nEquity']
    alphas = [0.909, 0.711, 0.889, 0.943, 0.912, 0.897]
    colors = ['#004080' if a >= 0.7 else '#CC0000' for a in alphas]

    fig, ax = plt.subplots(figsize=(8, 4.5))
    bars = ax.bar(variables, alphas, color=colors, alpha=0.85, edgecolor='#002060', linewidth=1.5)
    ax.axhline(y=0.7, color='red', linestyle='--', linewidth=1.5, label='Threshold = 0.70')
    ax.axhline(y=0.6, color='orange', linestyle=':', linewidth=1, label='Minimum = 0.60')

    for bar, val in zip(bars, alphas):
        ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 0.01,
               f'{val:.3f}', ha='center', va='bottom', fontweight='bold', fontsize=11, color='#002060')

    ax.set_ylabel("Cronbach's Alpha (α)", fontsize=12, fontweight='bold')
    ax.set_title("Reliability Analysis: Cronbach's Alpha Coefficients", fontsize=14,
                fontweight='bold', color='#004080')
    ax.set_ylim(0, 1.1)
    ax.legend(fontsize=10)
    ax.grid(True, axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig('figures2/cronbach_alpha.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_model_fit_chart():
    indices = ['CMIN/df', 'GFI', 'CFI', 'IFI', 'NFI', 'TLI', 'RMSEA']
    values = [2.975, 0.880, 0.970, 0.971, 0.956, 0.965, 0.070]
    thresholds = [3.0, 0.90, 0.90, 0.90, 0.90, 0.90, 0.08]
    passed = [True, False, True, True, True, True, True]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4.5))

    fit_indices = indices[1:6]
    fit_values = values[1:6]
    fit_thresh = thresholds[1:6]
    fit_pass = passed[1:6]
    colors = ['#004080' if p else '#CC6600' for p in fit_pass]
    x = np.arange(len(fit_indices))
    bars = ax1.bar(x, fit_values, color=colors, alpha=0.85, width=0.6, edgecolor='#002060')
    ax1.axhline(y=0.90, color='red', linestyle='--', linewidth=1.5, label='Threshold = 0.90')
    for bar, val in zip(bars, fit_values):
        ax1.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 0.01,
               f'{val:.3f}', ha='center', va='bottom', fontweight='bold', fontsize=10, color='#002060')
    ax1.set_xticks(x)
    ax1.set_xticklabels(fit_indices, fontsize=10)
    ax1.set_ylabel('Value', fontsize=11, fontweight='bold')
    ax1.set_title('Fit Indices (> 0.90)', fontsize=12, fontweight='bold', color='#004080')
    ax1.set_ylim(0, 1.15)
    ax1.legend(fontsize=9)
    ax1.grid(True, axis='y', alpha=0.3)

    special = ['CMIN/df', 'RMSEA']
    sp_values = [2.975, 0.070]
    sp_thresh = [3.0, 0.08]
    sp_colors = ['#004080', '#004080']
    bars2 = ax2.bar(special, sp_values, color=sp_colors, alpha=0.85, width=0.5, edgecolor='#002060')
    ax2.bar(special, sp_thresh, color='none', edgecolor='red', linewidth=2, linestyle='--', width=0.5, label='Threshold')
    for bar, val in zip(bars2, sp_values):
        ax2.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 0.05,
               f'{val:.3f}', ha='center', va='bottom', fontweight='bold', fontsize=12, color='#002060')
    ax2.set_title('CMIN/df (< 3.0) & RMSEA (< 0.08)', fontsize=12, fontweight='bold', color='#004080')
    ax2.legend(fontsize=9)
    ax2.grid(True, axis='y', alpha=0.3)

    fig.suptitle('Model Fit Assessment', fontsize=14, fontweight='bold', color='#002060', y=1.02)
    plt.tight_layout()
    plt.savefig('figures2/model_fit.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_sem_path_diagram():
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    ax.axis('off')

    factors = [
        ('Brand Awareness', 1.8, 6.0, -0.086, 'Negative'),
        ('Brand Image', 1.8, 4.8, 0.412, 'Highest'),
        ('Perceived Quality', 1.8, 3.6, 0.087, 'Low'),
        ('Brand Association', 1.8, 2.4, 0.064, 'Low'),
        ('Brand Loyalty', 1.8, 1.2, 0.291, 'Significant'),
    ]

    ax.add_patch(plt.Rectangle((6.5, 3.1), 2.5, 1.0, linewidth=3,
                               edgecolor='#002060', facecolor='#FFC000', zorder=3))
    ax.text(7.75, 3.6, 'Brand Equity', ha='center', va='center', fontsize=13,
           fontweight='bold', color='#002060', zorder=4)

    for label, x, y, weight, status in factors:
        if weight < 0:
            fc = '#FFCCCC'
            ec = '#CC0000'
        elif status == 'Highest':
            fc = '#CCFFCC'
            ec = '#006600'
        elif status == 'Significant':
            fc = '#CCE5FF'
            ec = '#004080'
        else:
            fc = '#F0F0F0'
            ec = '#888888'

        box = plt.Rectangle((x-1.3, y-0.35), 2.6, 0.7, linewidth=2,
                           edgecolor=ec, facecolor=fc, zorder=3)
        ax.add_patch(box)
        ax.text(x, y, label, ha='center', va='center', fontsize=9,
               fontweight='bold', color='#002060', zorder=4)

        arrow_color = ec
        lw = 3 if status in ['Highest', 'Significant'] else 1.5
        ax.annotate('', xy=(6.5, 3.6), xytext=(3.1, y),
                   arrowprops=dict(arrowstyle='->', lw=lw, color=arrow_color))

        mid_x = (3.1 + 6.5) / 2
        mid_y = (y + 3.6) / 2
        sign = '+' if weight > 0 else ''
        ax.text(mid_x, mid_y + 0.18, f'β = {sign}{weight:.3f}',
               fontsize=9, fontweight='bold', color=arrow_color, ha='center',
               bbox=dict(boxstyle='round,pad=0.2', facecolor='white',
                        edgecolor=arrow_color, alpha=0.9))

    legend_items = [
        ('Significant Positive', '#006600'),
        ('Positive (p < 0.05)', '#004080'),
        ('Insignificant', '#888888'),
        ('Negative', '#CC0000'),
    ]
    for i, (text, color) in enumerate(legend_items):
        ax.plot([], [], 's', color=color, markersize=8)
        ax.text(0.3 + i*2.5, 0.2, text, fontsize=8, color=color, fontweight='bold')

    ax.text(5.0, 0.2, 'β = -0.086 directly reported; others illustrative of relative effect sizes',
           fontsize=7, ha='center', fontstyle='italic', color='gray')
    fig.suptitle('Structural Equation Model: Path Coefficients (Illustrative)',
                fontsize=14, fontweight='bold', color='#002060', y=0.98)
    plt.tight_layout()
    plt.savefig('figures2/sem_path.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_brand_equity_comparison():
    factors = ['Brand\nAwareness', 'Brand\nImage', 'Perceived\nQuality',
               'Brand\nAssociation', 'Brand\nLoyalty']
    weights = [-0.086, 0.412, 0.087, 0.064, 0.291]
    colors = ['#CC0000', '#006600', '#888888', '#888888', '#004080']

    fig, ax = plt.subplots(figsize=(8, 4.5))
    bars = ax.barh(factors, weights, color=colors, alpha=0.85, edgecolor='#002060', height=0.6)
    ax.axvline(x=0, color='black', linewidth=1)
    ax.axvline(x=0.1, color='gray', linestyle=':', linewidth=1, alpha=0.5)

    for bar, val in zip(bars, weights):
        offset = 0.02 if val >= 0 else -0.02
        ha = 'left' if val >= 0 else 'right'
        sign = '+' if val > 0 else ''
        ax.text(val + offset, bar.get_y() + bar.get_height()/2.,
               f'β = {sign}{val:.3f}', ha=ha, va='center', fontweight='bold',
               fontsize=11, color='#002060')

    ax.set_xlabel('Regression Weight (β)', fontsize=12, fontweight='bold')
    ax.set_title('Impact of Brand Equity Dimensions on Disney Brand Equity\n(Illustrative β values based on reported effect sizes)',
                fontsize=12, fontweight='bold', color='#004080')
    ax.grid(True, axis='x', alpha=0.3)
    plt.tight_layout()
    plt.savefig('figures2/brand_equity_impact.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_media_efa_chart():
    stages = ['Awareness\n(AWR)', 'Interest\n(INT)', 'Conviction\n(CON)',
              'Purchase\n(PUR)', 'Post-Purchase\n(PPUR)']
    factor1 = [0.82, 0.79, 0.74, 0.31, 0.28]
    factor2 = [0.25, 0.30, 0.33, 0.85, 0.81]

    x = np.arange(len(stages))
    width = 0.35

    fig, ax = plt.subplots(figsize=(8, 4.5))
    bars1 = ax.bar(x - width/2, factor1, width, label='Factor 1: Pre-Purchase',
                  color='#004080', alpha=0.85, edgecolor='#002060')
    bars2 = ax.bar(x + width/2, factor2, width, label='Factor 2: Purchase',
                  color='#FFC000', alpha=0.85, edgecolor='#CC9900')
    ax.axhline(y=0.5, color='red', linestyle='--', linewidth=1.5, label='Loading Threshold = 0.50')

    ax.set_ylabel('Factor Loading', fontsize=12, fontweight='bold')
    ax.set_title('EFA Results: Consumer Behaviour Stages (Illustrative)\n(Based on Sama, 2019)',
                fontsize=13, fontweight='bold', color='#004080')
    ax.set_xticks(x)
    ax.set_xticklabels(stages, fontsize=10)
    ax.set_ylim(0, 1.1)
    ax.legend(fontsize=10, loc='upper right')
    ax.grid(True, axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig('figures2/media_efa.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_correlation_heatmap():
    labels = ['BA', 'BI', 'PQ', 'BAs', 'BL', 'BE']
    corr = np.array([
        [1.000, 0.421, 0.385, 0.512, 0.398, 0.312],
        [0.421, 1.000, 0.456, 0.489, 0.534, 0.687],
        [0.385, 0.456, 1.000, 0.467, 0.412, 0.398],
        [0.512, 0.489, 0.467, 1.000, 0.478, 0.375],
        [0.398, 0.534, 0.412, 0.478, 1.000, 0.589],
        [0.312, 0.687, 0.398, 0.375, 0.589, 1.000],
    ])

    fig, ax = plt.subplots(figsize=(6, 5))
    mask = np.triu(np.ones_like(corr, dtype=bool), k=1)
    sns.heatmap(corr, mask=mask, annot=True, fmt='.3f', cmap='Blues',
               xticklabels=labels, yticklabels=labels, ax=ax,
               linewidths=2, linecolor='white',
               vmin=0.2, vmax=1.0,
               annot_kws={'size': 11, 'fontweight': 'bold'})
    ax.set_title('Correlation Matrix: Brand Equity Dimensions (Illustrative)',
                fontsize=13, fontweight='bold', color='#004080', pad=15)
    note = 'BA=Awareness, BI=Image, PQ=Perceived Quality,\nBAs=Association, BL=Loyalty, BE=Brand Equity\n(Illustrative values based on reported relationships in Gilitwala & Nag, 2022)'
    ax.text(0.5, -0.18, note, transform=ax.transAxes, fontsize=7,
           ha='center', fontstyle='italic', color='gray')
    plt.tight_layout()
    plt.savefig('figures2/correlation_heatmap.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_fa_procedure_flowchart():
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    ax.axis('off')

    steps = [
        (1.2, 6.2, 'Step 1:\nDefine Research\nProblem'),
        (3.7, 6.2, 'Step 2:\nDesign Survey\n(Likert Scale)'),
        (6.2, 6.2, 'Step 3:\nCollect Data\n(n=400)'),
        (8.7, 6.2, 'Step 4:\nReliability Test\n(Cronbach α)'),
        (1.2, 3.5, 'Step 5:\nCFA Model\n(AMOS)'),
        (3.7, 3.5, 'Step 6:\nValidity Check\n(AVE, CR)'),
        (6.2, 3.5, 'Step 7:\nModel Fit\nAssessment'),
        (8.7, 3.5, 'Step 8:\nSEM & Path\nAnalysis'),
    ]

    colors_top = ['#002060', '#004080', '#006699', '#0088AA']
    colors_bot = ['#006600', '#008800', '#00AA00', '#FFC000']
    all_colors = colors_top + colors_bot

    for i, (x, y, text) in enumerate(steps):
        box = plt.Rectangle((x-1.0, y-0.55), 2.0, 1.1, linewidth=2,
                           edgecolor=all_colors[i], facecolor=all_colors[i],
                           alpha=0.15, zorder=2)
        ax.add_patch(box)
        box2 = plt.Rectangle((x-1.0, y-0.55), 2.0, 1.1, linewidth=2,
                            edgecolor=all_colors[i], facecolor='none', zorder=3)
        ax.add_patch(box2)
        ax.text(x, y, text, ha='center', va='center', fontsize=9,
               fontweight='bold', color=all_colors[i], zorder=4)

    for i in range(3):
        ax.annotate('', xy=(steps[i+1][0]-1.0, steps[i+1][1]),
                   xytext=(steps[i][0]+1.0, steps[i][1]),
                   arrowprops=dict(arrowstyle='->', lw=2, color='#004080'))
    ax.annotate('', xy=(steps[4][0], steps[4][1]+0.55),
               xytext=(steps[3][0], steps[3][1]-0.55),
               arrowprops=dict(arrowstyle='->', lw=2, color='#004080'))
    for i in range(4, 7):
        ax.annotate('', xy=(steps[i+1][0]-1.0, steps[i+1][1]),
                   xytext=(steps[i][0]+1.0, steps[i][1]),
                   arrowprops=dict(arrowstyle='->', lw=2, color='#006600'))

    fig.suptitle('Factor Analysis Procedure (CFA/SEM Approach)',
                fontsize=14, fontweight='bold', color='#002060', y=0.98)
    ax.text(5.0, 1.8, 'Applied to Disney Brand Equity Study (Gilitwala & Nag, 2022)',
           fontsize=10, ha='center', fontstyle='italic', color='gray')
    plt.tight_layout()
    plt.savefig('figures2/fa_procedure.png', dpi=200, bbox_inches='tight')
    plt.close()

print("Generating all visualizations...")
create_research_framework()
create_cronbach_alpha_chart()
create_model_fit_chart()
create_sem_path_diagram()
create_brand_equity_comparison()
create_media_efa_chart()
create_correlation_heatmap()
create_fa_procedure_flowchart()
print("All visualizations created!")

# ============================================================
# BUILD THE POWERPOINT PRESENTATION
# ============================================================

print("\nBuilding PowerPoint presentation...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)

gold_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), Inches(10), Inches(0.08))
gold_bar.fill.solid()
gold_bar.fill.fore_color.rgb = ACCENT_GOLD
gold_bar.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Factor Analysis in Entertainment Industry"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Brand Equity Analysis: Disney & Media Companies"
p2.font.size = Pt(22)
p2.font.color.rgb = ACCENT_GOLD
p2.alignment = PP_ALIGN.CENTER

info_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(3.5))
tf = info_box.text_frame
tf.word_wrap = True
info_lines = [
    ("Lana Jalal Gidan", Pt(22), True, WHITE),
    ("", Pt(10), False, WHITE),
    ("SSIE-605: Applied Multivariate Data Analysis", Pt(18), False, ACCENT_GOLD),
    ("Professor Susan Lu", Pt(16), False, WHITE),
    ("Binghamton University", Pt(16), False, WHITE),
    ("", Pt(10), False, WHITE),
    ("Technical Presentation 2: Factor Analysis Application", Pt(14), False, LIGHT_BLUE),
]
for i, (text, size, bold, color) in enumerate(info_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 2: AGENDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Agenda")

agenda_left = [
    ("1.", "Problem Statement", "Why Factor Analysis in Entertainment?"),
    ("2.", "Article Information", "Research Papers & Sources"),
    ("3.", "Factor Analysis Overview", "EFA vs CFA"),
    ("4.", "Case Study Introduction", "Disney Brand Equity"),
    ("5.", "Research Framework", "5 Dimensions of Brand Equity"),
    ("6.", "Research Methodology", "Data Collection & Approach"),
]
agenda_right = [
    ("7.", "Reliability Analysis", "Cronbach's Alpha Results"),
    ("8.", "CFA & Validity", "Factor Loadings, AVE, CR"),
    ("9.", "Model Fit Assessment", "Goodness-of-Fit Indices"),
    ("10.", "SEM Path Analysis", "Structural Model Results"),
    ("11.", "Supporting Study", "Media Advertising Impact"),
    ("12.", "Discussion & Conclusion", "Findings, Limitations, Improvements"),
]

for col, items, x_start in [(agenda_left, agenda_left, 0.3), (agenda_right, agenda_right, 5.2)]:
    for i, (num, title, sub) in enumerate(items):
        y = 1.5 + i * 0.95
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_start), Inches(y), Inches(0.5), Inches(0.5))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = DARK_BLUE
        num_box.line.fill.background()
        ntf = num_box.text_frame
        ntf.paragraphs[0].text = num
        ntf.paragraphs[0].font.size = Pt(14)
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.color.rgb = WHITE
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.margin_top = Inches(0.08)

        txt = slide.shapes.add_textbox(Inches(x_start + 0.6), Inches(y), Inches(4.0), Inches(0.5))
        ttf = txt.text_frame
        ttf.paragraphs[0].text = title
        ttf.paragraphs[0].font.size = Pt(14)
        ttf.paragraphs[0].font.bold = True
        ttf.paragraphs[0].font.color.rgb = DARK_BLUE
        p2 = ttf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY

# ============================================================
# SLIDE 3: PROBLEM STATEMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Problem Statement")

items = [
    ("Why Factor Analysis in Entertainment?", 18, True, DARK_BLUE),
    "",
    "The entertainment industry (Disney, NBCUniversal, Warner Bros.) manages complex brand portfolios across multiple business segments",
    "",
    "Understanding which brand dimensions drive brand equity is critical for strategic decision-making",
    "",
    ("Key Questions:", 16, True, MEDIUM_BLUE),
    "Which dimensions (awareness, image, quality, association, loyalty) most strongly influence brand equity?",
    "Can factor analysis reveal hidden patterns in consumer brand perception?",
    "How do entertainment brands differ from traditional product brands?",
    "",
    ("Business Impact:", 16, True, RGBColor(0, 128, 0)),
    "Disney's brand value: $57.0 billion (Brand Finance, 2022)",
    "Multi-segment strategy: Media, Parks, Studio, Streaming",
    "Factor Analysis helps identify which dimensions to invest in",
]
add_content_box(slide, Inches(0.5), Inches(1.4), Inches(9.0), Inches(5.5), items, font_size=14, bullet=True)

# ============================================================
# SLIDE 4: ARTICLE INFORMATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Article Information")

papers = [
    ("Paper 1 (Primary Case Study):", 15, True, DARK_BLUE),
    ('Gilitwala, B. & Nag, A.K. (2022). "Understanding Effective Factors Affecting Brand Equity."', 12, False, DARK_GRAY),
    ("Cogent Business & Management, 9(1), 2104431.", 12, False, DARK_GRAY),
    ("DOI: 10.1080/23311975.2022.2104431", 11, True, LIGHT_BLUE),
    ("Method: CFA + SEM | Sample: 400 | Application: Disney Brand Equity", 11, False, RGBColor(0, 128, 0)),
    "",
    ("Paper 2 (Supporting Study):", 15, True, DARK_BLUE),
    ('Sama, R. (2019). "Impact of Media Advertisements on Consumer Behaviour."', 12, False, DARK_GRAY),
    ("Journal of Creative Communications, 14(1), 54-68.", 12, False, DARK_GRAY),
    ("DOI: 10.1177/0973258618822624", 11, True, LIGHT_BLUE),
    ("Method: EFA + PCA | Sample: 529 | Application: Media Advertising", 11, False, RGBColor(0, 128, 0)),
    "",
    ("Paper 3 (Methodology Reference):", 15, True, DARK_BLUE),
    ('Shrestha, N. (2021). "Factor Analysis as a Tool for Survey Analysis."', 12, False, DARK_GRAY),
    ("American Journal of Applied Mathematics and Statistics, 9(1), 4-11.", 12, False, DARK_GRAY),
    ("DOI: 10.12691/ajams-9-1-2", 11, True, LIGHT_BLUE),
    ("Method: EFA | Application: FA Methodology & Best Practices", 11, False, RGBColor(0, 128, 0)),
]
add_content_box(slide, Inches(0.5), Inches(1.4), Inches(9.0), Inches(5.8), papers, font_size=12, bullet=False)

# ============================================================
# SLIDE 5: WHAT IS FACTOR ANALYSIS?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "What is Factor Analysis?", "Brief Recap")

items_left = [
    ("Definition:", 16, True, DARK_BLUE),
    "A statistical method that reduces many observed variables into a smaller set of underlying factors",
    "",
    ("Mathematical Model:", 14, True, MEDIUM_BLUE),
    ("Xi = ai1F1 + ai2F2 + ... + aimFm + ei", 13, True, RGBColor(128, 0, 0)),
    "",
    "Xi = observed variable i",
    "F1...Fm = common factors",
    "ai = factor loading coefficients",
    "ei = unique/error factor",
    "",
    ("In This Study:", 14, True, RGBColor(0, 128, 0)),
    "Brand Equity = f(Awareness, Image, Quality, Association, Loyalty)",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(4.8), Inches(5.5), items_left, font_size=13, bullet=False)

items_right = [
    ("Two Main Types:", 16, True, DARK_BLUE),
    "",
    ("EFA (Exploratory):", 14, True, MEDIUM_BLUE),
    "Discovers factor structure",
    "No prior assumptions",
    "Uses eigenvalues & scree plots",
    "Rotation: Varimax/Oblimin",
    "",
    ("CFA (Confirmatory):", 14, True, MEDIUM_BLUE),
    "Tests hypothesized structure",
    "Requires theoretical model",
    "Uses goodness-of-fit indices",
    "Part of SEM framework",
    "",
    ("This study uses CFA + SEM", 13, True, RGBColor(0, 128, 0)),
]
add_content_box(slide, Inches(5.3), Inches(1.4), Inches(4.5), Inches(5.5), items_right, font_size=13, bullet=False)

# ============================================================
# SLIDE 6: FA PROCEDURE FLOWCHART
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Analysis Procedure (CFA/SEM)")

slide.shapes.add_picture('figures2/fa_procedure.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))

# ============================================================
# SLIDE 7: EFA vs CFA COMPARISON
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "EFA vs CFA: Understanding the Difference")

table_data = [
    ['Aspect', 'EFA (Exploratory)', 'CFA (Confirmatory)'],
    ['Purpose', 'Discover factor structure', 'Test hypothesized structure'],
    ['Theory', 'No prior theory needed', 'Theory-driven model'],
    ['Extraction', 'PCA / Principal Axis', 'Maximum Likelihood'],
    ['Rotation', 'Varimax / Oblimin', 'Not applicable'],
    ['Fit Indices', 'Not typically used', 'CFI, RMSEA, TLI, etc.'],
    ['Output', 'Factor loadings, eigenvalues', 'Path coefficients, fit indices'],
    ['Software', 'SPSS, R, Python', 'AMOS, Mplus, lavaan'],
    ['Used By', 'Sama (2019)', 'Gilitwala & Nag (2022)'],
]

table = slide.shapes.add_table(len(table_data), 3, Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.5)).table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(3.5)

for i, row_data in enumerate(table_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            elif j == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = DARK_BLUE
            else:
                paragraph.font.color.rgb = DARK_GRAY
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        elif i == len(table_data) - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(220, 235, 255)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9.0), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "This presentation primarily focuses on CFA + SEM (Gilitwala & Nag, 2022), with EFA comparison from Sama (2019)"
p.font.size = Pt(12)
p.font.color.rgb = MEDIUM_BLUE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 8: CASE STUDY INTRO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Case Study: Disney Brand Equity", "Gilitwala & Nag (2022)")

items_left = [
    ("About Disney:", 16, True, DARK_BLUE),
    "The Walt Disney Company is one of the world's largest entertainment conglomerates",
    "",
    ("Four Business Segments:", 14, True, MEDIUM_BLUE),
    "Media Networks (ABC, ESPN, Disney Channel)",
    "Parks, Experiences & Products",
    "Studio Entertainment (Pixar, Marvel, Lucasfilm)",
    "Direct-to-Consumer (Disney+, Hulu)",
    "",
    ("Competitors:", 14, True, MEDIUM_BLUE),
    "NBCUniversal (Comcast) - Peacock, Universal Studios",
    "Warner Bros. Discovery - HBO, DC Comics",
    "Netflix, Amazon Prime Video",
    "Paramount Global - CBS, MTV",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(5.0), Inches(5.5), items_left, font_size=13, bullet=False)

items_right = [
    ("Study Details:", 16, True, DARK_BLUE),
    "",
    ("Location:", 13, True, MEDIUM_BLUE),
    "Disney Shop, DLF Mall, Noida, India",
    "",
    ("Sample Size:", 13, True, MEDIUM_BLUE),
    "400 respondents (online questionnaire)",
    "",
    ("Sampling:", 13, True, MEDIUM_BLUE),
    "Convenient sampling technique",
    "",
    ("Scale:", 13, True, MEDIUM_BLUE),
    "5-point Likert scale",
    "",
    ("Analysis Tool:", 13, True, MEDIUM_BLUE),
    "IBM AMOS (Version 24)",
    "CB-SEM approach",
]
add_content_box(slide, Inches(5.5), Inches(1.4), Inches(4.3), Inches(5.5), items_right, font_size=12, bullet=False)

# ============================================================
# SLIDE 9: RESEARCH FRAMEWORK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Research Framework: 5 Dimensions of Brand Equity")

slide.shapes.add_picture('figures2/research_framework.png', Inches(0.3), Inches(1.4), Inches(5.5), Inches(5.0))

items = [
    ("Five Hypotheses Tested:", 15, True, DARK_BLUE),
    "",
    "H1: Brand Awareness \u2192 Brand Equity",
    "H2: Brand Image \u2192 Brand Equity",
    "H3: Perceived Quality \u2192 Brand Equity",
    "H4: Brand Association \u2192 Brand Equity",
    "H5: Brand Loyalty \u2192 Brand Equity",
    "",
    ("Theoretical Basis:", 14, True, MEDIUM_BLUE),
    "Aaker (1991) Brand Equity Model",
    "Keller (1993) Customer-Based Brand Equity",
    "",
    ("Source: Gilitwala & Nag, 2022", 11, True, LIGHT_BLUE),
]
add_content_box(slide, Inches(5.8), Inches(1.4), Inches(4.0), Inches(5.5), items, font_size=13, bullet=False)

# ============================================================
# SLIDE 10: CORRELATION MATRIX
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 1: Correlation Matrix", "Examining Relationships Between Variables")

slide.shapes.add_picture('figures2/correlation_heatmap.png', Inches(0.3), Inches(1.4), Inches(5.2), Inches(5.0))

items = [
    ("Key Observations:", 15, True, DARK_BLUE),
    "",
    "Brand Image has the strongest relationship with Brand Equity (highest β in SEM)",
    "",
    "Brand Loyalty also strongly relates to Brand Equity",
    "",
    "Brand Awareness has the weakest (negative) relationship with Brand Equity",
    "",
    ("Why This Matters:", 14, True, MEDIUM_BLUE),
    "Moderate correlations between variables suggest they share variance but are distinct constructs",
    "This supports the applicability of CFA",
    "",
    ("Illustrative correlation pattern based on", 10, False, DARK_GRAY),
    ("reported SEM results in Gilitwala & Nag (2022)", 10, False, DARK_GRAY),
]
add_content_box(slide, Inches(5.6), Inches(1.4), Inches(4.2), Inches(5.5), items, font_size=12, bullet=False)

# ============================================================
# SLIDE 11: RELIABILITY ANALYSIS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 2: Reliability Analysis (Cronbach's Alpha)")

slide.shapes.add_picture('figures2/cronbach_alpha.png', Inches(0.2), Inches(1.4), Inches(5.5), Inches(4.0))

table_data = [
    ['Variable', 'α Value', 'Status'],
    ['Brand Awareness', '0.909', 'Excellent'],
    ['Brand Image', '0.711', 'Acceptable'],
    ['Perceived Quality', '0.889', 'Good'],
    ['Brand Association', '0.943', 'Excellent'],
    ['Brand Loyalty', '0.912', 'Excellent'],
    ['Brand Equity', '0.897', 'Good'],
]

table = slide.shapes.add_table(len(table_data), 3, Inches(5.8), Inches(1.5), Inches(4.0), Inches(3.2)).table
for i, row_data in enumerate(table_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = DARK_BLUE
                paragraph.font.bold = True if j == 1 else False
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE

txBox = slide.shapes.add_textbox(Inches(5.8), Inches(4.9), Inches(4.0), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "All values > 0.70 threshold"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 128, 0)
p2 = tf.add_paragraph()
p2.text = "Questionnaire is reliable and acceptable for the study (Peterson, 1994)"
p2.font.size = Pt(11)
p2.font.color.rgb = DARK_GRAY
p3 = tf.add_paragraph()
p3.text = ""
p3.font.size = Pt(6)
p4 = tf.add_paragraph()
p4.text = "Source: Gilitwala & Nag (2022), Table 1"
p4.font.size = Pt(10)
p4.font.color.rgb = LIGHT_BLUE
p4.font.bold = True

# ============================================================
# SLIDE 12: CFA RESULTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 3: Confirmatory Factor Analysis Results")

items = [
    ("CFA Validity Criteria (Gilitwala & Nag, 2022):", 15, True, DARK_BLUE),
    "",
    ("Convergent Validity:", 14, True, MEDIUM_BLUE),
    "All factor loadings > 0.70 (threshold met)",
    "Average Variance Extracted (AVE) > 0.50 for all constructs",
    "Composite Reliability (CR) > 0.70 for all constructs",
    "",
    ("Discriminant Validity:", 14, True, MEDIUM_BLUE),
    "Square root of AVE > inter-construct correlations",
    "All constructs are distinct from each other",
    "",
    ("What This Means:", 14, True, RGBColor(0, 128, 0)),
    "The 5 brand equity dimensions are reliably measured",
    "Each dimension captures a unique aspect of brand perception",
    "The measurement model is valid for structural analysis",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(4.8), Inches(5.5), items, font_size=13, bullet=False)

cfa_table = [
    ['Construct', 'FL', 'AVE', 'CR', 'α'],
    ['Brand Awareness', '> 0.70', '> 0.50', '> 0.70', '0.909'],
    ['Brand Image', '> 0.70', '> 0.50', '> 0.70', '0.711'],
    ['Perceived Quality', '> 0.70', '> 0.50', '> 0.70', '0.889'],
    ['Brand Association', '> 0.70', '> 0.50', '> 0.70', '0.943'],
    ['Brand Loyalty', '> 0.70', '> 0.50', '> 0.70', '0.912'],
    ['Brand Equity', '> 0.70', '> 0.50', '> 0.70', '0.897'],
]

table = slide.shapes.add_table(len(cfa_table), 5, Inches(5.3), Inches(1.5), Inches(4.5), Inches(3.5)).table
table.columns[0].width = Inches(1.5)
for i, row_data in enumerate(cfa_table):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = DARK_BLUE
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        elif i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 245, 255)

note = slide.shapes.add_textbox(Inches(5.3), Inches(5.2), Inches(4.5), Inches(1.5))
ntf = note.text_frame
ntf.word_wrap = True
p = ntf.paragraphs[0]
p.text = "FL = Factor Loading, AVE = Average Variance Extracted"
p.font.size = Pt(9)
p.font.color.rgb = DARK_GRAY
p2 = ntf.add_paragraph()
p2.text = "CR = Composite Reliability, α = Cronbach's Alpha"
p2.font.size = Pt(9)
p2.font.color.rgb = DARK_GRAY
p3 = ntf.add_paragraph()
p3.text = "Source: Gilitwala & Nag (2022), Table 3"
p3.font.size = Pt(10)
p3.font.color.rgb = LIGHT_BLUE
p3.font.bold = True

# ============================================================
# SLIDE 13: MODEL FIT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 4: Model Fit Assessment")

slide.shapes.add_picture('figures2/model_fit.png', Inches(0.2), Inches(1.4), Inches(6.0), Inches(4.0))

fit_table = [
    ['Index', 'Value', 'Threshold', 'Result'],
    ['CMIN/df', '2.975', '< 3.0', 'Good Fit'],
    ['GFI', '0.880', '> 0.90', 'Marginal'],
    ['CFI', '0.970', '> 0.90', 'Good Fit'],
    ['IFI', '0.971', '> 0.90', 'Good Fit'],
    ['NFI', '0.956', '> 0.90', 'Good Fit'],
    ['TLI', '0.965', '> 0.90', 'Good Fit'],
    ['RMSEA', '0.070', '< 0.08', 'Good Fit'],
]

table = slide.shapes.add_table(len(fit_table), 4, Inches(6.3), Inches(1.5), Inches(3.5), Inches(3.8)).table
for i, row_data in enumerate(fit_table):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            elif j == 3:
                paragraph.font.bold = True
                if cell_text == 'Good Fit':
                    paragraph.font.color.rgb = RGBColor(0, 128, 0)
                else:
                    paragraph.font.color.rgb = RGBColor(200, 100, 0)
            else:
                paragraph.font.color.rgb = DARK_BLUE
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE

note = slide.shapes.add_textbox(Inches(6.3), Inches(5.5), Inches(3.5), Inches(1.5))
ntf = note.text_frame
ntf.word_wrap = True
p = ntf.paragraphs[0]
p.text = "6 out of 7 indices indicate Good Fit"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 128, 0)
p2 = ntf.add_paragraph()
p2.text = "Source: Gilitwala & Nag (2022), Table 5"
p2.font.size = Pt(10)
p2.font.color.rgb = LIGHT_BLUE
p2.font.bold = True

# ============================================================
# SLIDE 14: SEM PATH ANALYSIS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 5: Structural Equation Model (Path Analysis)")

slide.shapes.add_picture('figures2/sem_path.png', Inches(0.2), Inches(1.4), Inches(5.8), Inches(5.0))

items = [
    ("Path Coefficients (β):", 15, True, DARK_BLUE),
    "",
    ("Brand Image \u2192 BE:", 13, True, RGBColor(0, 100, 0)),
    "Highest positive relationship (strongest predictor)",
    "",
    ("Brand Loyalty \u2192 BE:", 13, True, RGBColor(0, 64, 128)),
    "Second strongest, significant relationship",
    "",
    ("Perceived Quality \u2192 BE:", 13, False, DARK_GRAY),
    "Low, insignificant relationship",
    "",
    ("Brand Association \u2192 BE:", 13, False, DARK_GRAY),
    "Low, insignificant relationship",
    "",
    ("Brand Awareness \u2192 BE: β = -0.086", 13, True, RGBColor(200, 0, 0)),
    "Negative effect (saturation effect)",
    "",
    ("Source: Gilitwala & Nag (2022), Table 6", 10, True, LIGHT_BLUE),
    ("Note: β = -0.086 is directly reported;", 9, False, DARK_GRAY),
    ("other β values are illustrative of", 9, False, DARK_GRAY),
    ("reported relative effect sizes", 9, False, DARK_GRAY),
]
add_content_box(slide, Inches(6.0), Inches(1.4), Inches(3.8), Inches(5.5), items, font_size=12, bullet=False)

# ============================================================
# SLIDE 15: BRAND EQUITY IMPACT CHART
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 6: Impact of Dimensions on Disney Brand Equity")

slide.shapes.add_picture('figures2/brand_equity_impact.png', Inches(0.3), Inches(1.4), Inches(6.0), Inches(4.5))

items = [
    ("Key Insights:", 15, True, DARK_BLUE),
    "",
    ("Brand Image is #1 driver", 13, True, RGBColor(0, 100, 0)),
    "How consumers perceive Disney emotionally matters most",
    "",
    ("Brand Loyalty is #2 driver", 13, True, RGBColor(0, 64, 128)),
    "Repeat purchase and devotion are crucial",
    "",
    ("Surprising Finding:", 13, True, RGBColor(200, 0, 0)),
    "Brand Awareness has a negative effect (regression weight = -0.086)!",
    "Possible saturation: everyone knows Disney, but awareness alone does not create equity",
    "",
    ("Source: Gilitwala & Nag (2022)", 10, True, LIGHT_BLUE),
    ("Chart shows illustrative β values", 9, False, DARK_GRAY),
]
add_content_box(slide, Inches(6.3), Inches(1.4), Inches(3.5), Inches(5.5), items, font_size=11, bullet=False)

# ============================================================
# SLIDE 16: SUPPORTING STUDY - MEDIA ADVERTISEMENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Supporting Study: Media Advertisements & Consumer Behaviour", "Sama (2019) - EFA Application")

slide.shapes.add_picture('figures2/media_efa.png', Inches(0.2), Inches(1.4), Inches(5.5), Inches(4.5))

items = [
    ("Study Overview:", 14, True, DARK_BLUE),
    "529 respondents, online survey",
    "5 media: TV, Radio, Newspaper, Magazine, Internet",
    "",
    ("EFA Results:", 14, True, MEDIUM_BLUE),
    "Extraction: PCA with Varimax rotation",
    "2 factors extracted (eigenvalue > 1):",
    "",
    ("Factor 1 - Pre-Purchase:", 13, True, RGBColor(0, 64, 128)),
    "Awareness, Interest, Conviction",
    "",
    ("Factor 2 - Purchase:", 13, True, RGBColor(200, 150, 0)),
    "Purchase, Post-Purchase behaviour",
    "",
    ("Both factors: Cronbach's α > 0.70", 12, True, RGBColor(0, 128, 0)),
    "",
    ("Chart: Illustrative loadings based on", 9, False, DARK_GRAY),
    ("the two-factor structure in Sama (2019)", 9, False, DARK_GRAY),
    ("DOI: 10.1177/0973258618822624", 10, True, LIGHT_BLUE),
]
add_content_box(slide, Inches(5.8), Inches(1.4), Inches(4.0), Inches(5.5), items, font_size=11, bullet=False)

# ============================================================
# SLIDE 17: FACTOR SCORES & INTERPRETATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Interpretation & Business Implications")

items_left = [
    ("Disney Brand Equity Drivers:", 16, True, DARK_BLUE),
    "",
    ("Most Important \u2192 Brand Image:", 14, True, RGBColor(0, 100, 0)),
    "Emotional connection ('magical', 'fun')",
    "Visual identity and character recognition",
    "Nostalgic brand associations",
    "Cross-generational appeal",
    "",
    ("Second Most Important \u2192 Brand Loyalty:", 14, True, RGBColor(0, 64, 128)),
    "Repeat purchases of Disney products",
    "Willingness to pay premium prices",
    "Recommendation to family and friends",
    "Theme park revisits & streaming retention",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(4.8), Inches(5.5), items_left, font_size=12, bullet=False)

items_right = [
    ("Implications for Media Companies:", 16, True, DARK_BLUE),
    "",
    ("For Disney:", 14, True, MEDIUM_BLUE),
    "Invest in brand image (content quality, emotional storytelling)",
    "Build loyalty programs (Disney+, annual passes)",
    "",
    ("For NBCUniversal:", 14, True, MEDIUM_BLUE),
    "Strengthen brand image through franchise experiences",
    "Develop Peacock loyalty features",
    "Universal Studios immersive attractions",
    "",
    ("For All Media Companies:", 14, True, MEDIUM_BLUE),
    "Brand awareness alone is insufficient",
    "Focus on emotional engagement > awareness",
    "Loyalty drives long-term brand equity",
]
add_content_box(slide, Inches(5.3), Inches(1.4), Inches(4.5), Inches(5.5), items_right, font_size=12, bullet=False)

# ============================================================
# SLIDE 18: DISCUSSION & KEY FINDINGS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Discussion & Key Findings")

items = [
    ("Finding 1: Brand Image is the Primary Driver", 15, True, RGBColor(0, 100, 0)),
    "Disney's emotional storytelling, character recognition, and 'magical' brand perception create the strongest impact on brand equity (highest positive relationship)",
    "",
    ("Finding 2: Brand Loyalty is the Second Driver", 15, True, RGBColor(0, 64, 128)),
    "Repeat consumption, willingness to pay premium, and brand devotion significantly predict brand equity (second strongest relationship)",
    "",
    ("Finding 3: Brand Awareness Has Negative Effect", 15, True, RGBColor(200, 0, 0)),
    "Surprising result: regression weight = -0.086. Possible explanation: Disney is so universally known that mere awareness creates a saturation effect without adding equity value",
    "",
    ("Finding 4: Media Advertising Operates in Two Stages", 15, True, MEDIUM_BLUE),
    "Sama (2019) found advertising effects split into Pre-Purchase (awareness, interest, conviction) and Purchase (buying decision, post-purchase) stages",
    "",
    ("Finding 5: CFA/SEM Validates Brand Equity Dimensions", 15, True, DARK_BLUE),
    "The five-factor model (Aaker, 1991) is statistically validated for Disney with strong model fit (CFI = 0.970, RMSEA = 0.070)",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.2), Inches(5.5), items, font_size=12, bullet=False)

# ============================================================
# SLIDE 19: COMMENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Comments")

items = [
    ("Strengths of Gilitwala & Nag (2022):", 16, True, DARK_BLUE),
    "Strong theoretical foundation based on Aaker's (1991) Brand Equity Model",
    "Large sample size (n = 400) exceeds minimum requirements for CFA",
    "Use of CB-SEM with AMOS provides rigorous statistical testing",
    "All reliability and validity criteria met (CR > 0.70, AVE > 0.50)",
    "Open access (CC-BY 4.0) promotes research transparency",
    "",
    ("Strengths of Sama (2019):", 16, True, DARK_BLUE),
    "Comprehensive coverage of 5 media platforms across 5 behaviour stages",
    "Large sample (n = 529) with diverse student population",
    "EFA reveals clear two-factor structure in consumer behaviour",
    "",
    ("Practical Relevance:", 16, True, RGBColor(0, 128, 0)),
    "Both studies provide actionable insights for entertainment marketing",
    "Results can guide budget allocation between brand image vs. awareness campaigns",
    "Factor analysis methodology is well-suited for brand equity measurement",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.2), Inches(5.8), items, font_size=13, bullet=False)

# ============================================================
# SLIDE 20: LIMITATIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Limitations")

items = [
    ("Gilitwala & Nag (2022) Limitations:", 16, True, DARK_BLUE),
    "",
    "Single location: Only Disney Shop at DLF Mall, Noida, India",
    "Convenience sampling: May not represent all Disney consumers globally",
    "GFI = 0.880 is slightly below the recommended 0.90 threshold",
    "Cross-sectional design: Cannot capture changes in brand equity over time",
    "Cultural bias: Indian consumers may perceive Disney differently than Western markets",
    "",
    ("Sama (2019) Limitations:", 16, True, DARK_BLUE),
    "",
    "Student-only sample (529 respondents): Limited generalizability",
    "Did not include emerging media (social media influencers, podcasts)",
    "Exact KMO and Bartlett's test statistics not prominently reported",
    "",
    ("General Limitations:", 16, True, MEDIUM_BLUE),
    "",
    "Self-reported data may contain response bias",
    "Factor analysis assumes linear relationships between variables",
    "Results are context-specific and may not generalize to other entertainment companies like NBCUniversal",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.2), Inches(5.8), items, font_size=13, bullet=False)

# ============================================================
# SLIDE 21: SUGGESTED IMPROVEMENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Suggested Improvements")

items = [
    ("Methodological Improvements:", 16, True, DARK_BLUE),
    "Use probability sampling (stratified random) for better generalizability",
    "Include multi-location data collection (Disney parks, stores across countries)",
    "Apply longitudinal design to track brand equity changes over time",
    "Compare EFA and CFA approaches on the same dataset for robustness",
    "",
    ("Extended Scope:", 16, True, MEDIUM_BLUE),
    "Include competitor comparison: Disney vs. NBCUniversal vs. Warner Bros.",
    "Add digital brand equity dimensions: social media engagement, streaming metrics",
    "Include Net Promoter Score (NPS) and Customer Lifetime Value (CLV)",
    "Study the impact of Disney+ streaming launch on brand equity",
    "",
    ("Advanced Analytics:", 16, True, RGBColor(0, 128, 0)),
    "Apply Multi-Group Analysis (MGA) to compare across demographics",
    "Use higher-order CFA with second-order factor for overall brand equity",
    "Incorporate sentiment analysis from social media as additional data source",
    "Test mediation effects: Image \u2192 Loyalty \u2192 Equity (indirect effects)",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.2), Inches(5.8), items, font_size=13, bullet=False)

# ============================================================
# SLIDE 22: CONCLUSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Conclusion")

items = [
    ("Summary:", 18, True, DARK_BLUE),
    "",
    "Factor Analysis (CFA/SEM) is a powerful tool for understanding brand equity in the entertainment industry",
    "",
    "Disney's brand equity is primarily driven by Brand Image and Brand Loyalty, not just Brand Awareness",
    "",
    "The five-factor model (Awareness, Image, Quality, Association, Loyalty) is validated for Disney with strong model fit indices (CFI = 0.970, RMSEA = 0.070)",
    "",
    "Brand Awareness alone has a negative effect on equity, suggesting entertainment companies should focus on building emotional connections, not just recognition",
    "",
    "Media advertising operates in two distinct stages (Pre-Purchase and Purchase), as confirmed by EFA in Sama (2019)",
    "",
    ("Key Takeaway:", 16, True, RGBColor(0, 128, 0)),
    "For entertainment companies like Disney and NBCUniversal, investing in brand image and customer loyalty yields stronger brand equity than increasing awareness alone",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.2), Inches(5.5), items, font_size=14, bullet=False)

# ============================================================
# SLIDE 23: REFERENCES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "References")

refs = [
    '1. Gilitwala, B. & Nag, A.K. (2022). "Understanding Effective Factors Affecting Brand Equity." Cogent Business & Management, 9(1), 2104431. DOI: 10.1080/23311975.2022.2104431',
    '',
    '2. Sama, R. (2019). "Impact of Media Advertisements on Consumer Behaviour." Journal of Creative Communications, 14(1), 54-68. DOI: 10.1177/0973258618822624',
    '',
    '3. Shrestha, N. (2021). "Factor Analysis as a Tool for Survey Analysis." American J. of Applied Mathematics and Statistics, 9(1), 4-11. DOI: 10.12691/ajams-9-1-2',
    '',
    '4. Aaker, D.A. (1991). Managing Brand Equity: Capitalizing on the Value of a Brand Name. Free Press, New York.',
    '',
    '5. Keller, K.L. (1993). "Conceptualizing, Measuring, and Managing Customer-Based Brand Equity." Journal of Marketing, 57(1), 1-22.',
    '',
    '6. Peterson, R.A. (1994). "A Meta-Analysis of Cronbach\'s Coefficient Alpha." Journal of Consumer Research, 21(2), 381-391.',
    '',
    '7. Tong, X. & Hawley, J.M. (2009). "Measuring Customer-Based Brand Equity." Journal of Fashion Marketing and Management, 13(3), 357-376.',
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.2), Inches(5.8), refs, font_size=11, bullet=False)

# ============================================================
# SLIDE 24: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)

gold_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(10), Inches(0.08))
gold_bar.fill.solid()
gold_bar.fill.fore_color.rgb = ACCENT_GOLD
gold_bar.line.fill.background()

thank_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
tf = thank_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

qa_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(2.5))
tf = qa_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Questions & Discussion"
p.font.size = Pt(28)
p.font.color.rgb = ACCENT_GOLD
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = ""
p2.font.size = Pt(10)
p3 = tf.add_paragraph()
p3.text = "Lana Jalal Gidan"
p3.font.size = Pt(18)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = "SSIE-605 | Binghamton University | Professor Susan Lu"
p4.font.size = Pt(14)
p4.font.bold = True
p4.font.color.rgb = LIGHT_BLUE
p4.alignment = PP_ALIGN.CENTER

# ============================================================
# SAVE
# ============================================================

output_file = "Factor_Analysis_Disney_Presentation.pptx"
prs.save(output_file)

slide_count = len(prs.slides)
print(f"\n{'='*80}")
print(f"PRESENTATION SAVED: {output_file}")
print(f"Total Slides: {slide_count}")
print(f"{'='*80}")
print(f"\nSlide Summary:")
slide_titles = [
    "Title Slide", "Agenda", "Problem Statement", "Article Information",
    "What is Factor Analysis?", "FA Procedure (CFA/SEM)", "EFA vs CFA Comparison",
    "Case Study: Disney Brand Equity", "Research Framework",
    "Correlation Matrix", "Reliability Analysis (Cronbach's Alpha)",
    "CFA Results", "Model Fit Assessment", "SEM Path Analysis",
    "Brand Equity Impact Chart", "Supporting Study: Media Advertisements",
    "Factor Interpretation & Business Implications", "Discussion & Key Findings",
    "Comments", "Limitations", "Suggested Improvements", "Conclusion",
    "References", "Thank You"
]
for i, title in enumerate(slide_titles[:slide_count], 1):
    print(f"  Slide {i}: {title}")
