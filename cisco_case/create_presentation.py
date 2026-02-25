import numpy as np
import matplotlib
matplotlib.use('Agg')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

CISCO_DARK = RGBColor(0, 32, 96)
CISCO_BLUE = RGBColor(0, 188, 235)
CISCO_NAVY = RGBColor(0, 79, 169)
ACCENT_GREEN = RGBColor(39, 174, 96)
ACCENT_RED = RGBColor(231, 76, 60)
ACCENT_ORANGE = RGBColor(243, 156, 18)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_bg(slide, color=CISCO_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title, subtitle=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CISCO_DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = CISCO_BLUE
        p2.alignment = PP_ALIGN.LEFT

def add_accent(slide, top=Inches(1.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CISCO_BLUE
    shape.line.fill.background()

def add_content(slide, left, top, width, height, items, font_size=13, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, size, bold, color = item
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
        elif item == "":
            p.text = ""
            p.font.size = Pt(6)
        else:
            prefix = "  " if bullet else ""
            p.text = f"{prefix}{item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY

def add_table(slide, left, top, width, height, data, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            if i == 0:
                from pptx.oxml.ns import qn
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '002060'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

def add_footer(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SLIDE 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CISCO_DARK)
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.5))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Repositioning Webex with AI"
p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "From Assistant-Led to Agentic Work"
p2.font.size = Pt(22); p2.font.color.rgb = CISCO_BLUE; p2.alignment = PP_ALIGN.CENTER

add_accent(slide, top=Inches(2.6))

txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(3.5))
tf2 = txBox2.text_frame; tf2.word_wrap = True
items = [
    ("Cisco x WiB x SWE Spring 2026 Case Competition", 16, True, WHITE),
    ("", 6, False, WHITE),
    ("What product and platform improvements can Cisco make to Webex,", 13, False, WHITE),
    ("powered by its AI Assistant, to outperform Zoom and Microsoft Teams", 13, False, WHITE),
    ("in the next phase of hybrid work?", 13, False, WHITE),
    ("", 12, False, WHITE),
    ("Presented by: Lana Jalal Gidan", 16, True, CISCO_BLUE),
    ("Watson College of Engineering and Applied Science", 13, False, WHITE),
    ("Binghamton University", 13, False, WHITE),
]
first = True
for text, size, bold, color in items:
    if first: p = tf2.paragraphs[0]; first = False
    else: p = tf2.add_paragraph()
    p.text = text; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

# SLIDE 2: Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Agenda")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Industry Analysis", 16, True, CISCO_DARK),
    "",
    ("Problem Identification", 16, True, CISCO_DARK),
    "",
    ("Recommendations", 16, True, CISCO_DARK),
    "",
    ("Implementation", 16, True, CISCO_DARK),
    "",
    ("Financials & KPIs", 16, True, CISCO_DARK),
    "",
    ("Risks & Mitigations", 16, True, CISCO_DARK),
    "",
    ("Executive Summary", 16, True, CISCO_DARK),
])
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 3: Industry Analysis - Market Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Industry Analysis", "The UC&C Market: $55B and Growing at 8.4% CAGR")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/market_share.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 4: Industry Analysis - Hybrid Work
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "The Hybrid Work Revolution", "48% of Knowledge Workers Now Operate in Hybrid Mode")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/hybrid_work.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 5: Industry Analysis - Financials
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Financial Overview: Cisco vs. Competitors", "Source: SEC 10-K Filings (FY2024)")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/financials.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 6: Industry - AI Feature Comparison
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "AI Feature Gap Analysis", "Webex Trails in Agentic Capabilities")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/ai_comparison.png', Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.8))
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 7: Problem Identification Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Problem Identification", "Three Root-Cause Problems Limiting Webex Competitiveness")
add_accent(slide)

for i, (title, desc, data_pt) in enumerate([
    ("Low User Engagement\n& Adoption Gap",
     "Webex ~150M vs Teams 320M MAU\nMarket share: ~10% vs ~37%",
     "G2 Rating: 4.2 vs 4.5 (Zoom)"),
    ("AI Limited to\nPassive Assistance",
     "Only summarizes & transcribes\nvs. Copilot executes in Office 365",
     "No autonomous task execution"),
    ("Weak Ecosystem\nIntegration",
     "Limited 3rd-party app marketplace\nvs. Teams 2,000+ integrations",
     "Low developer SDK adoption"),
]):
    left = Inches(0.5 + i * 3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(2.8), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 245, 245)
    shape.line.color.rgb = ACCENT_RED
    shape.line.width = Pt(2)
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = f"Problem {i+1}"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT_RED; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = ""; p2.font.size = Pt(4)
    p3 = tf.add_paragraph(); p3.text = title; p3.font.size = Pt(13); p3.font.bold = True; p3.font.color.rgb = CISCO_DARK; p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph(); p4.text = ""; p4.font.size = Pt(6)
    p5 = tf.add_paragraph(); p5.text = desc; p5.font.size = Pt(10); p5.font.color.rgb = DARK_GRAY; p5.alignment = PP_ALIGN.CENTER
    p6 = tf.add_paragraph(); p6.text = ""; p6.font.size = Pt(6)
    p7 = tf.add_paragraph(); p7.text = data_pt; p7.font.size = Pt(9); p7.font.bold = True; p7.font.color.rgb = CISCO_NAVY; p7.alignment = PP_ALIGN.CENTER

add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 8: Problem 1 Detail - G2 Ratings
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Problem 1: Low User Engagement", "Webex Trails Competitors in Satisfaction & Adoption")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/g2_ratings.png', Inches(0.3), Inches(1.4), Inches(5.5), Inches(4.0))
add_content(slide, Inches(5.8), Inches(1.6), Inches(3.8), Inches(4.5), [
    ("Root Cause (5 Whys):", 13, True, CISCO_DARK),
    "",
    "1. Low engagement?",
    "   Less intuitive UX",
    "",
    "2. Less intuitive?",
    "   No differentiated value",
    "",
    "3. No differentiation?",
    "   AI features are basic",
    "",
    "4. Basic AI?",
    "   Passive assistant-only",
    "",
    ("5. ROOT CAUSE:", 12, True, ACCENT_RED),
    ("AI does not perform tasks", 12, True, ACCENT_RED),
], font_size=11)
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 9: Problem 2 & 3 Detail
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Problems 2 & 3: Passive AI + Weak Ecosystem", "Root Causes Traced Through 5 Whys Analysis")
add_accent(slide)
add_content(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(5.5), [
    ("Problem 2: Passive AI", 14, True, ACCENT_RED),
    "",
    "Current Webex AI can:",
    "  - Summarize meetings",
    "  - Transcribe in real-time",
    "  - Translate 100+ languages",
    "",
    "But CANNOT:",
    "  - Execute follow-up tasks",
    "  - Draft documents",
    "  - Automate workflows",
    "  - Coordinate across tools",
    "",
    ("Root Cause: No agentic", 12, True, ACCENT_RED),
    ("architecture exists", 12, True, ACCENT_RED),
], font_size=11)
add_content(slide, Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.5), [
    ("Problem 3: Weak Ecosystem", 14, True, ACCENT_RED),
    "",
    "Integration comparison:",
    "  - Teams: 2,000+ apps",
    "  - Zoom: 2,500+ apps",
    "  - Webex: Significantly fewer",
    "",
    "Self-reinforcing cycle:",
    "  Fewer integrations ->",
    "  Lower stickiness ->",
    "  Fewer developers ->",
    "  Even fewer integrations",
    "",
    ("Root Cause: No open AI agent", 12, True, ACCENT_RED),
    ("platform for developers", 12, True, ACCENT_RED),
], font_size=11)
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 10: Engineering Flowchart - Workflow Orchestrator Agent
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Workflow Orchestrator Agent: Technical Flowchart", "Perceive - Reason - Act Loop | Swimlane Architecture")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/flowchart.png', Inches(0.1), Inches(1.25), Inches(9.8), Inches(6.1))

# SLIDE 11: Why Teams & Zoom Win Today + How Webex Leapfrogs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Competitive Analysis: How Webex Leapfrogs", "Why Teams & Zoom Lead Today + Webex's Agentic Advantage")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/competitive_advantage.png', Inches(0.1), Inches(1.25), Inches(9.8), Inches(6.1))

# SLIDE 12: Recommendation Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Recommendation Overview", "Three Solutions to Make Webex the Agentic Collaboration Leader")
add_accent(slide)

for i, (title, desc, color) in enumerate([
    ("Agentic AI\nTransformation",
     "AI that WORKS, not just assists.\nAutonomous task execution:\nschedule, draft, assign, follow up",
     ACCENT_GREEN),
    ("Webex AI Agent\nPlatform (SDK)",
     "Open developer framework for\nvertical-specific AI agents.\nAgent marketplace ecosystem",
     CISCO_NAVY),
    ("Intelligent Workflow\nEngine",
     "200+ enterprise integrations.\nAI-powered cross-platform\nworkflow automation",
     ACCENT_ORANGE),
]):
    left = Inches(0.5 + i * 3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(2.8), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 255, 245)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = f"Solution {i+1}"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = color; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = ""; p2.font.size = Pt(4)
    p3 = tf.add_paragraph(); p3.text = title; p3.font.size = Pt(14); p3.font.bold = True; p3.font.color.rgb = CISCO_DARK; p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph(); p4.text = ""; p4.font.size = Pt(8)
    p5 = tf.add_paragraph(); p5.text = desc; p5.font.size = Pt(11); p5.font.color.rgb = DARK_GRAY; p5.alignment = PP_ALIGN.CENTER

add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 12: Agentic Evolution
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Solution 1: Agentic AI Transformation", "Evolving from Assistant to Autonomous Agent")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/agentic_evolution.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 13: Solution 1 Technical Details
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Solution 1: Technical Architecture", "How Agentic AI Works in Webex")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), [
    ("Perceive-Reason-Act Architecture:", 14, True, CISCO_DARK),
    "",
    "1. PERCEIVE: Webex AI Codec",
    "   Real-Time Media Models (RMM)",
    "   process audio + detect task intents",
    "",
    "2. REASON: NLP Task Decomposition",
    "   LLM identifies intent, fetches context",
    "   via A2A (Agent-to-Agent) Protocols",
    "   from Jira, Salesforce, SAP, etc.",
    "",
    "3. ACT: Autonomous Execution",
    "   Agent pushes via secure API",
    "   Human-in-the-loop: Approve/Reject",
    "   Confirmation posted to Webex chat",
    "",
    "4. OBSERVE: Splunk Observability",
    "   End-to-end monitoring + learning loop",
    "   Reinforcement from user feedback",
], font_size=11)
add_content(slide, Inches(5.0), Inches(1.5), Inches(4.5), Inches(5.5), [
    ("Feasibility Evidence:", 14, True, ACCENT_GREEN),
    "",
    "R&D Budget: $7.98B (FY2024)",
    "  Sufficient for agentic development",
    "",
    "Existing AI Infrastructure:",
    "  150M+ users already served",
    "  Real-time transcription at scale",
    "",
    "Splunk Acquisition ($28B):",
    "  Enterprise data observability",
    "  Pattern recognition capability",
    "",
    "Security Advantage:",
    "  FedRAMP, HIPAA, SOC 2 Type II",
    "  Zero-trust agent architecture",
    "",
    ("Scalability: Microservices-based", 11, True, CISCO_NAVY),
    ("agents on Cisco global cloud (40+", 11, True, CISCO_NAVY),
    ("data center locations)", 11, True, CISCO_NAVY),
], font_size=11)
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 14: Solution 2 - Agent SDK
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Solution 2: Webex AI Agent Platform", "Open Developer SDK & Agent Marketplace")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), [
    ("Platform Components:", 14, True, CISCO_DARK),
    "",
    "1. Agent SDK",
    "   Python & JavaScript SDKs",
    "   Pre-built agent templates",
    "   (meeting coordinator, PM, etc.)",
    "",
    "2. Agent Runtime",
    "   Sandboxed execution in Webex Cloud",
    "   Enterprise-grade security & isolation",
    "   Container-based (Kubernetes)",
    "",
    "3. Agent Marketplace",
    "   Discover, deploy, monetize agents",
    "   Revenue sharing for developers",
    "",
    "4. Vertical Templates",
    "   Healthcare (HIPAA-compliant)",
    "   Finance (compliance-aware)",
    "   Government (FedRAMP-certified)",
], font_size=11)
add_content(slide, Inches(5.0), Inches(1.5), Inches(4.5), Inches(5.5), [
    ("Adoption Strategy:", 14, True, CISCO_NAVY),
    "",
    "$10M Developer Incentive Fund",
    "  (Modeled on Salesforce Trailhead)",
    "",
    "University Hackathons (50+)",
    "  Partner with CS programs",
    "",
    "Systems Integrator Partners",
    "  Accenture, Deloitte, Wipro",
    "  for initial vertical agents",
    "",
    ("Risk Mitigation:", 14, True, ACCENT_ORANGE),
    "",
    "Primary risk: developer adoption",
    "  Mitigation: incentives + hackathons",
    "",
    "Security risk: 3rd-party agents",
    "  Mitigation: sandboxed execution +",
    "  mandatory security review process",
], font_size=11)
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 15: Solution 3 - Workflow Engine
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Solution 3: Intelligent Workflow Engine", "AI-Powered Cross-Platform Automation")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("How It Works:", 14, True, CISCO_DARK),
    "",
    "User says: 'After every customer call, update the CRM, create a follow-up task, and schedule a review if deal > $50K'",
    "",
    ("Technical Architecture:", 14, True, CISCO_DARK),
    "",
    "1. Intent Recognition: NLP converts requests into workflow DAGs (Directed Acyclic Graphs)",
    "2. Connector Framework: Pre-built connectors for top 50 enterprise SaaS tools",
    "   (Salesforce, HubSpot, Jira, ServiceNow, SAP, Workday, Slack, GitHub, etc.)",
    "3. AI Orchestration Engine: Coordinates multi-agent workflows with error recovery",
    "4. Analytics Dashboard: Real-time ROI tracking powered by Splunk observability",
    "",
    ("Feasibility:", 14, True, ACCENT_GREEN),
    "",
    "Cisco already owns AppDynamics ($3.7B acquisition) + Splunk ($28B) for monitoring",
    "Standard REST/GraphQL APIs available from all major SaaS platforms",
    "Estimated development: 12-18 months for first 50 connectors, scale to 200+",
    "",
    ("SaaS Integration market: $15.6B by 2027 (MarketsandMarkets)", 12, True, CISCO_NAVY),
], font_size=11)
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 16: Agentic AI Market
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "The Agentic AI Opportunity", "Market Growing from $5.1B to $47B by 2030 (CAGR 44.8%)")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/agentic_market.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 17: Implementation Timeline
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Implementation Timeline", "18-Month Phased Roadmap (Q3 2026 - Q4 2027)")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/timeline.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 18: Implementation Details
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Implementation Plan: Phased Approach", "Engineering Perspective from Watson College")
add_accent(slide)
data = [
    ["Phase", "Timeline", "Key Activities", "Resources"],
    ["Phase 1:\nFoundation", "Q3-Q4 2026\n(6 months)", "Core agent framework\nAPI connector library\nSecurity architecture", "$150M R&D\n50 engineers\nAI/ML team"],
    ["Phase 2:\nBeta & SDK", "Q1-Q2 2027\n(6 months)", "Enterprise beta (50 partners)\nSDK developer preview\nHackathon program launch", "$100M investment\n10K developers\n3 SI partners"],
    ["Phase 3:\nScale", "Q3-Q4 2027\n(6 months)", "GA launch agentic AI\nMarketplace launch\n200+ integrations", "$50M marketing\nGlobal rollout\nContinuous iteration"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.5), data,
          col_widths=[1.5, 1.8, 3.5, 2.6])
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 19: Financials
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Financial Projections", "Projected Revenue Impact of Agentic AI Strategy")
add_accent(slide)
slide.shapes.add_picture('cisco_case/figures/revenue_projection.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(4.2))
add_content(slide, Inches(0.5), Inches(5.7), Inches(9), Inches(1.5), [
    "Investment: ~$300M over 18 months (from existing $7.98B R&D budget = 3.8% reallocation)",
    "Projected Return: +$2.1B incremental ARR by FY2029 | ROI: ~600% over 4 years",
    ("Based on IDC UC&C CAGR 8.4% + estimated 3-5% market share gain from agentic differentiation", 10, False, RGBColor(100, 100, 100)),
], font_size=11)
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 20: KPIs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Key Performance Indicators", "Measurable Success Metrics Over 18 Months")
add_accent(slide)

for i, (title, kpis, color) in enumerate([
    ("Agentic AI", ["+35% daily active usage", "25% meeting-to-action rate", "NPS: 4.2 -> 4.5+", "85% task completion accuracy"], ACCENT_GREEN),
    ("Agent Platform", ["10,000+ developers", "500+ marketplace agents", "3 vertical solutions", "200+ enterprise adopters"], CISCO_NAVY),
    ("Workflow Engine", ["200+ integrations", "40% automation rate", "50% less app switching", "+$2.1B ARR by FY2029"], ACCENT_ORANGE),
]):
    left = Inches(0.5 + i * 3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(2.8), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = f"Solution {i+1}"; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = color; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = title; p2.font.size = Pt(14); p2.font.bold = True; p2.font.color.rgb = CISCO_DARK; p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph(); p3.text = ""; p3.font.size = Pt(8)
    for kpi in kpis:
        p4 = tf.add_paragraph(); p4.text = kpi; p4.font.size = Pt(12); p4.font.color.rgb = DARK_GRAY; p4.alignment = PP_ALIGN.CENTER
        p5 = tf.add_paragraph(); p5.text = ""; p5.font.size = Pt(4)

add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 21: Risks and Mitigations
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Risks and Mitigations")
add_accent(slide)
data = [
    ["Risk", "Probability", "Impact", "Mitigation"],
    ["Developer SDK\nadoption slower\nthan expected", "Medium", "High", "$10M incentive fund\n50+ university hackathons\n3 major SI partnerships"],
    ["AI agent errors\nin enterprise\nenvironment", "Medium", "High", "Human-in-the-loop guardrails\nConfigurable autonomy levels\nComprehensive audit logging"],
    ["Competitor response\n(Microsoft / Zoom\naccelerate AI)", "High", "Medium", "First-mover advantage in\nagentic architecture\nSecurity moat (FedRAMP/HIPAA)"],
    ["Integration\nscalability\nchallenges", "Low", "Medium", "Microservices architecture\nCisco's 40+ global data centers\nKubernetes orchestration"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.0), data,
          col_widths=[2.2, 1.2, 1.2, 4.8])
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 22: System Trade-offs (Engineering Focus)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Engineering Trade-offs & Constraints", "Watson College Engineering Analysis")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Autonomy vs. Control", 14, True, CISCO_DARK),
    "More agentic capability increases productivity but raises risk of unintended actions",
    "Solution: Configurable autonomy levels (low / medium / high) per organization",
    "",
    ("Open Platform vs. Security", 14, True, CISCO_DARK),
    "Open Agent SDK increases ecosystem growth but expands attack surface",
    "Solution: Sandboxed execution + mandatory security review for marketplace agents",
    "",
    ("Build vs. Buy for Workflow Engine", 14, True, CISCO_DARK),
    "Building in-house ensures deep integration; acquiring iPaaS is faster",
    "Solution: Hybrid -- build AI orchestration layer, partner for connector framework",
    "",
    ("Latency Constraint", 14, True, CISCO_DARK),
    "Agentic tasks must complete within 2-5 seconds for user satisfaction",
    "Multi-tool workflows may exceed this; use async execution with real-time updates",
], font_size=12)
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 23: Executive Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Executive Summary")
add_accent(slide)

for i, (prob, sol) in enumerate([
    ("Low User Engagement\n& Adoption Gap",
     "Agentic AI Transformation\nAI that executes tasks,\nnot just assists"),
    ("AI Limited to\nPassive Assistance",
     "Webex AI Agent Platform\nOpen SDK + marketplace\nfor custom agents"),
    ("Weak Ecosystem\nIntegration",
     "Intelligent Workflow Engine\n200+ integrations with\nAI orchestration"),
]):
    left = Inches(0.5 + i * 3.1)
    shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(2.8), Inches(2.0))
    shape1.fill.solid(); shape1.fill.fore_color.rgb = RGBColor(255, 240, 240)
    shape1.line.color.rgb = ACCENT_RED; shape1.line.width = Pt(1)
    tf1 = shape1.text_frame; tf1.word_wrap = True
    p1 = tf1.paragraphs[0]; p1.text = prob; p1.font.size = Pt(11); p1.font.color.rgb = ACCENT_RED; p1.alignment = PP_ALIGN.CENTER; p1.font.bold = True

    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.2), Inches(2.8), Inches(2.0))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = RGBColor(240, 255, 240)
    shape2.line.color.rgb = ACCENT_GREEN; shape2.line.width = Pt(1)
    tf2 = shape2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = sol; p2.font.size = Pt(11); p2.font.color.rgb = ACCENT_GREEN; p2.alignment = PP_ALIGN.CENTER; p2.font.bold = True

add_content(slide, Inches(0.5), Inches(6.4), Inches(9), Inches(0.8), [
    ("Projected Impact: +$2.1B incremental ARR by FY2029 | Investment: ~$300M (3.8% of R&D)", 12, True, CISCO_DARK),
])
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 24: Data Sources
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Data Sources & References")
add_accent(slide)
data = [
    ["Source", "URL", "Key Data Used"],
    ["Cisco 10-K FY2024", "investor.cisco.com", "Rev: $53.8B, R&D: $7.98B"],
    ["Microsoft FY2024", "microsoft.com/investor", "Teams: 320M MAU"],
    ["Zoom FY2025 10-K", "investors.zoom.us", "Rev: $4.6B, 3,900+ enterprise"],
    ["IDC UC&C Tracker", "idc.com", "Market: $55B, CAGR 8.4%"],
    ["Gartner 2024", "gartner.com/newsroom", "48% hybrid workers"],
    ["MarketsandMarkets", "marketsandmarkets.com", "Agentic AI: $5.1B-$47B"],
    ["G2 Reviews", "g2.com/categories/video-conferencing", "Webex: 4.2, Zoom: 4.5"],
    ["McKinsey 2024", "mckinsey.com", "58% hybrid opportunity"],
    ["Webex AI Docs", "webex.com/ai", "AI Assistant features"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.0), data,
          col_widths=[2.5, 3.5, 3.4])
add_footer(slide, "Industry Analysis | Problem Identification | Recommendations | Implementation | Financials | Risks | Executive Summary")

# SLIDE 25: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CISCO_DARK)
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(3))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Thank You!"; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = ""; p2.font.size = Pt(10)
p3 = tf.add_paragraph()
p3.text = "Questions?"; p3.font.size = Pt(24); p3.font.color.rgb = CISCO_BLUE; p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph(); p4.text = ""; p4.font.size = Pt(20)
p5 = tf.add_paragraph()
p5.text = "Lana Jalal Gidan"; p5.font.size = Pt(18); p5.font.bold = True; p5.font.color.rgb = WHITE; p5.alignment = PP_ALIGN.CENTER
p6 = tf.add_paragraph()
p6.text = "Watson College of Engineering | Binghamton University"; p6.font.size = Pt(14); p6.font.color.rgb = WHITE; p6.alignment = PP_ALIGN.CENTER

prs.save('cisco_case/Cisco_Webex_AI_Presentation.pptx')
print(f"\nPresentation saved: cisco_case/Cisco_Webex_AI_Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
