"""
Convert Final_Presentation_AutoPlay.pptx to MP4 video.
Renders each slide as a 1280x720 PNG using python-pptx + Pillow,
then stitches with per-slide audio using ffmpeg.
"""

import os, io, subprocess, shutil, tempfile, zipfile, re
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

PPTX        = 'Final_Presentation_AutoPlay.pptx'
OUTPUT      = 'Final_Presentation_Video.mp4'
TMPDIR      = tempfile.mkdtemp()
FRAMES_DIR  = os.path.join(TMPDIR, 'frames')
AUDIO_DIR   = os.path.join(TMPDIR, 'audio')
os.makedirs(FRAMES_DIR)
os.makedirs(AUDIO_DIR)

W, H = 1280, 720   # slide pixel size

# Audio clips per slide (slide index 0-based, duration in seconds)
# Slides 5-15 in 1-based = indices 4-14 in 0-based
AUDIO_MAP = {
     4: ('media5.wav',  85.90),
     5: ('media6.wav',  42.75),
     6: ('media7.wav',  25.59),
     7: ('media8.wav',  56.18),
     8: ('media9.wav',  33.42),
     9: ('media10.wav', 31.94),
    10: ('media11.wav', 68.68),
    11: ('media12.wav', 53.95),
    12: ('media13.wav', 47.26),
    13: ('media14.wav', 101.16),
    14: ('media15.wav', 102.76),
}
DEFAULT_DURATION = 5.0   # seconds for slides without audio

# ── Extract WAV files from PPTX ───────────────────────────────────────────────
print('Extracting audio files...')
with zipfile.ZipFile(PPTX, 'r') as z:
    for slide_idx, (wav_file, dur) in AUDIO_MAP.items():
        z.extract(f'ppt/media/{wav_file}', TMPDIR)
        shutil.move(os.path.join(TMPDIR, 'ppt', 'media', wav_file),
                    os.path.join(AUDIO_DIR, wav_file))

# ── Helpers ────────────────────────────────────────────────────────────────────
def emu_to_px(emu, slide_emu, slide_px):
    return int(emu / slide_emu * slide_px)

def get_bg_color(slide, prs):
    """Return PIL-compatible (R,G,B) background colour."""
    try:
        fill = slide.background.fill
        if fill.type is not None and str(fill.type) == 'SOLID (1)':
            rgb = fill.fore_color.rgb
            return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    try:
        layout_fill = slide.slide_layout.background.fill
        if layout_fill.type is not None and str(layout_fill.type) == 'SOLID (1)':
            rgb = layout_fill.fore_color.rgb
            return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return (15, 36, 74)   # fallback: navy

def get_text_color(shape):
    """Best-effort text colour extraction."""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.type is not None:
                    rgb = run.font.color.rgb
                    return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return (255, 255, 255)

def get_font_size(shape):
    """Best-effort font size (points → pixels at 96 dpi)."""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    return max(12, int(run.font.size.pt * 96 / 72))
            if para.runs == [] and para.font.size:
                return max(12, int(para.font.size.pt * 96 / 72))
    except Exception:
        pass
    return 18

def load_pil_font(size):
    """Load a PIL font, falling back to default if truetype not available."""
    for path in [
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
        '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
        '/nix/store',   # skip nix paths
    ]:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                pass
    return ImageFont.load_default()

def wrap_text(text, font, max_width, draw):
    """Word-wrap text to fit within max_width pixels."""
    words = text.split()
    lines = []
    current = ''
    for word in words:
        test = (current + ' ' + word).strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines

# ── Render each slide ──────────────────────────────────────────────────────────
prs = Presentation(PPTX)
slide_emu_w = prs.slide_width   # in EMUs
slide_emu_h = prs.slide_height

print(f'Rendering {len(prs.slides)} slides at {W}x{H}...')

for slide_idx, slide in enumerate(prs.slides):
    img = Image.new('RGB', (W, H), color=get_bg_color(slide, prs))
    draw = ImageDraw.Draw(img)

    # Sort shapes: pictures first, then text on top
    shapes = list(slide.shapes)
    shapes_sorted = (
        [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE] +
        [s for s in shapes if s.shape_type != MSO_SHAPE_TYPE.PICTURE
         and s.shape_type != MSO_SHAPE_TYPE.MEDIA]
    )

    for shape in shapes_sorted:
        # Position in pixels
        left   = emu_to_px(shape.left   or 0, slide_emu_w, W)
        top    = emu_to_px(shape.top    or 0, slide_emu_h, H)
        width  = emu_to_px(shape.width  or 0, slide_emu_w, W)
        height = emu_to_px(shape.height or 0, slide_emu_h, H)

        # ── Images ──────────────────────────────────────────────────────────
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img_data = shape.image.blob
                pil_img  = Image.open(io.BytesIO(img_data)).convert('RGBA')
                pil_img  = pil_img.resize((max(1,width), max(1,height)),
                                           Image.LANCZOS)
                # Paste with alpha mask if available
                if pil_img.mode == 'RGBA':
                    img.paste(pil_img, (left, top), pil_img)
                else:
                    img.paste(pil_img, (left, top))
            except Exception as e:
                pass   # skip broken images silently

        # ── Text shapes ─────────────────────────────────────────────────────
        elif hasattr(shape, 'text_frame') and shape.has_text_frame:
            text  = shape.text_frame.text.strip()
            if not text:
                continue
            color = get_text_color(shape)
            size  = get_font_size(shape)
            font  = load_pil_font(size)

            # Draw each paragraph
            y_off = top + 4
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    y_off += size + 4
                    continue
                para_size = size
                try:
                    if para.runs and para.runs[0].font.size:
                        para_size = max(10, int(para.runs[0].font.size.pt * 96/72))
                except Exception:
                    pass
                para_font = load_pil_font(para_size)
                lines = wrap_text(para_text, para_font, max(1, width - 8), draw)
                for line in lines:
                    if y_off + para_size > top + height:
                        break
                    # Subtle shadow for readability
                    draw.text((left + 5, y_off + 1), line, font=para_font,
                              fill=(0, 0, 0, 180))
                    draw.text((left + 4, y_off), line, font=para_font, fill=color)
                    y_off += para_size + 3
                y_off += 4   # paragraph spacing

    # Save frame
    frame_path = os.path.join(FRAMES_DIR, f'slide_{slide_idx:02d}.png')
    img.save(frame_path, 'PNG')
    print(f'  Slide {slide_idx+1:2d} rendered → {frame_path}')

# ── Build per-slide video segments then concatenate ──────────────────────────
print('\nBuilding video segments with ffmpeg...')
segment_files = []

for slide_idx in range(len(prs.slides)):
    frame  = os.path.join(FRAMES_DIR, f'slide_{slide_idx:02d}.png')
    seg    = os.path.join(TMPDIR, f'seg_{slide_idx:02d}.mp4')

    if slide_idx in AUDIO_MAP:
        wav_file, dur = AUDIO_MAP[slide_idx]
        audio_path    = os.path.join(AUDIO_DIR, wav_file)
        cmd = [
            'ffmpeg', '-y', '-loop', '1', '-i', frame,
            '-i', audio_path,
            '-c:v', 'libx264', '-tune', 'stillimage',
            '-c:a', 'aac', '-b:a', '128k',
            '-pix_fmt', 'yuv420p',
            '-t', str(dur + 0.5),   # slight buffer after audio ends
            '-shortest',
            seg
        ]
    else:
        cmd = [
            'ffmpeg', '-y', '-loop', '1', '-i', frame,
            '-c:v', 'libx264', '-tune', 'stillimage',
            '-pix_fmt', 'yuv420p',
            '-t', str(DEFAULT_DURATION),
            '-an',   # no audio for silent slides
            seg
        ]

    result = subprocess.run(cmd, capture_output=True)
    if result.returncode != 0:
        print(f'  ERROR slide {slide_idx+1}: {result.stderr.decode()[-200:]}')
    else:
        segment_files.append(seg)
        tag = f'({AUDIO_MAP[slide_idx][1]:.0f}s audio)' if slide_idx in AUDIO_MAP else f'({DEFAULT_DURATION}s silent)'
        print(f'  Slide {slide_idx+1:2d} segment done {tag}')

# ── Concatenate all segments ───────────────────────────────────────────────────
print(f'\nConcatenating {len(segment_files)} segments...')
concat_list = os.path.join(TMPDIR, 'concat.txt')
with open(concat_list, 'w') as f:
    for seg in segment_files:
        f.write(f"file '{seg}'\n")

subprocess.run([
    'ffmpeg', '-y', '-f', 'concat', '-safe', '0',
    '-i', concat_list,
    '-c', 'copy',
    OUTPUT
], check=True, capture_output=True)

shutil.rmtree(TMPDIR)
size_mb = os.path.getsize(OUTPUT) / 1024 / 1024
print(f'\nDone! {OUTPUT}  ({size_mb:.1f} MB)')

total_dur = sum(AUDIO_MAP[i][1] for i in AUDIO_MAP) + \
            (len(prs.slides) - len(AUDIO_MAP)) * DEFAULT_DURATION
print(f'Total video duration: ~{total_dur/60:.1f} minutes')
