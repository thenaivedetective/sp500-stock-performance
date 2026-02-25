from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PURPLE = RGBColor(107, 45, 123)
TEAL = RGBColor(42, 157, 143)
WARM = RGBColor(231, 111, 81)
DARK = RGBColor(38, 70, 83)
GOLD = RGBColor(233, 196, 106)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
LIGHT_BG = RGBColor(250, 247, 252)
DARK_GRAY = RGBColor(50, 50, 50)

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title, subtitle=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = GOLD
        p2.alignment = PP_ALIGN.LEFT

def add_accent(slide, top=Inches(1.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

def add_content(slide, left, top, width, height, items, font_size=13, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, size, bold, color = item
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
        elif item == "":
            p.text = ""
            p.font.size = Pt(6)
        else:
            prefix = "  " if bullet else ""
            p.text = f"{prefix}{item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY

def add_table(slide, left, top, width, height, data, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            if i == 0:
                from pptx.oxml.ns import qn
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '6B2D7B'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

def add_footer(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

FOOTER = "Introduction | Literature Review | Methodology | Results | Discussion | Conclusion | References"

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PURPLE)
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Factor Analysis of Work Engagement"; p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Women & Workforce Engagement"; p2.font.size = Pt(20); p2.font.color.rgb = GOLD; p2.alignment = PP_ALIGN.CENTER

add_accent(slide, top=Inches(2.3))

txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(4.5))
tf2 = txBox2.text_frame; tf2.word_wrap = True
items = [
    ("EFA & CFA of the 9-Item Utrecht Work Engagement Scale (UWES-9)", 14, True, WHITE),
    ("in a Multi-Occupational Female Sample", 14, True, WHITE),
    ("", 8, False, WHITE),
    ("Primary Study:", 13, True, GOLD),
    ("Willmer, M., Westerberg Jacobson, J. & Lindberg, M. (2019)", 12, False, WHITE),
    ("Frontiers in Psychology, 10, Article 2771", 12, False, WHITE),
    ("DOI: 10.3389/fpsyg.2019.02771", 11, False, RGBColor(200, 200, 200)),
    ("", 12, False, WHITE),
    ("Presented by: Lana Jalal Gidan", 16, True, GOLD),
    ("SSIE-605: Applied Multivariate Data Analysis", 13, False, WHITE),
    ("Professor Susan Lu", 13, False, WHITE),
    ("Watson College of Engineering | Binghamton University", 12, False, WHITE),
]
first = True
for text, size, bold, color in items:
    if first: p = tf2.paragraphs[0]; first = False
    else: p = tf2.add_paragraph()
    p.text = text; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 2: AGENDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Agenda")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("1. Introduction & Background", 16, True, DARK),
    "",
    ("2. Literature Review", 16, True, DARK),
    "",
    ("3. Research Methodology", 16, True, DARK),
    "",
    ("4. Results: EFA & CFA", 16, True, DARK),
    "",
    ("5. Discussion", 16, True, DARK),
    "",
    ("6. Conclusion & Implications", 16, True, DARK),
    "",
    ("7. References", 16, True, DARK),
])
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 3: INTRODUCTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Introduction", "Work Engagement & Women in the Workforce")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("What is Work Engagement?", 15, True, PURPLE),
    "",
    "A positive, work-related state of mind characterized by three dimensions:",
    "",
    ("  Vigor: High energy, mental resilience, willingness to invest effort", 12, False, DARK_GRAY),
    ("  Dedication: Sense of significance, enthusiasm, inspiration, pride", 12, False, DARK_GRAY),
    ("  Absorption: Being fully engrossed in work, time passes quickly", 12, False, DARK_GRAY),
    "",
    ("Why Study Women Specifically?", 15, True, PURPLE),
    "",
    "- Most UWES-9 studies use mixed-gender samples with majority male participants",
    "- Gender may influence how work engagement is structured and experienced",
    "- Only one prior Swedish study existed, using 186 IT consultants (63% male)",
    "- Understanding women's engagement patterns is critical for retention & policy",
    "",
    ("(Schaufeli et al., 2002; Hallberg & Schaufeli, 2006)", 10, False, RGBColor(120, 120, 120)),
], font_size=13)
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 4: LITERATURE REVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Literature Review", "The UWES-9 Debate: One Factor vs. Three Factors")
add_accent(slide)
add_content(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(5.5), [
    ("The UWES-9 Instrument:", 14, True, PURPLE),
    "",
    "- 9 items rated on 0-6 scale",
    "  (Never to All the time)",
    "- 3 theoretical subscales:",
    "  Vigor (3 items: 1, 2, 5)",
    "  Dedication (3 items: 3, 4, 7)",
    "  Absorption (3 items: 6, 8, 9)",
    "",
    ("The Core Debate:", 14, True, WARM),
    "",
    "Kulikowski (2017) review of 21 studies:",
    "  - 3 confirmed 1-factor",
    "  - 3 confirmed 3-factor",
    "  - 4 found both equivalent",
    "  - 1 supported neither",
    "No definitive conclusion reached",
], font_size=11)
add_content(slide, Inches(5.0), Inches(1.5), Inches(4.8), Inches(5.5), [
    ("Key Prior Findings:", 14, True, TEAL),
    "",
    "Swedish (Hallberg & Schaufeli, 2006):",
    "  n = 186 IT consultants (37% women)",
    "  Both 1-factor & 3-factor fit equally",
    "  RMSEA = 0.13, CFI = 0.97",
    "",
    "Norwegian (Nerstad et al., 2010):",
    "  n = 1,266 multi-occupational (67% women)",
    "  3-factor supported, but factors",
    "  highly correlated (suggesting 1 factor)",
    "",
    "Finnish (Seppala et al., 2009):",
    "  n = 9,404 multi-occupational",
    "  Both structures reasonable",
    "",
    ("Gap: No all-female, multi-occupational", 11, True, WARM),
    ("Swedish study existed", 11, True, WARM),
], font_size=11)
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 5: DEMOGRAPHICS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Sample Demographics", "n = 702 Women, Aged 26-37, Multi-Occupational Swedish Sample")
add_accent(slide)
slide.shapes.add_picture('figures4_new/demographics.png', Inches(0.2), Inches(1.4), Inches(9.6), Inches(5.5))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 6: METHODOLOGY - VISUAL FLOWCHART
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Research Methodology", "Split-Sample Design: EFA + CFA")
add_accent(slide)
slide.shapes.add_picture('figures4_new/split_sample_flowchart.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 7: KMO & BARTLETT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Sampling Adequacy & Reliability", "KMO = 0.922 (Marvellous) | Cronbach's Alpha = 0.947")
add_accent(slide)
slide.shapes.add_picture('figures4_new/kmo_bartlett.png', Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 8: SUBSCALE MEANS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "UWES-9 Subscale Scores", "Overall Mean = 4.06 | Dedication Highest at 4.24")
add_accent(slide)
slide.shapes.add_picture('figures4_new/subscale_means.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 9: ITEM MEANS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Individual Item Mean Scores", "Item 7 (Pride) Highest | Item 9 (Immersion) Lowest")
add_accent(slide)
slide.shapes.add_picture('figures4_new/item_means.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 10: SKEWNESS & KURTOSIS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Item Distribution Properties", "Skewness & Kurtosis Within Acceptable Range for ML Estimation")
add_accent(slide)
slide.shapes.add_picture('figures4_new/skewness_kurtosis.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 11: EFA RESULTS - FACTOR LOADINGS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "EFA Results: One-Factor Solution", "All Loadings > 0.65 | One Factor Explains > 70% Variance")
add_accent(slide)
slide.shapes.add_picture('figures4_new/factor_loadings.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 12: EFA SUMMARY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "EFA Summary", "Multiple Tests Converge on One-Factor Solution")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("EFA Key Findings (n = 341):", 15, True, PURPLE),
    "",
    "1. Unrotated EFA: One factor explains > 70% of variance",
    "   - Chi-squared = 332.43 (df = 27), p < 0.001",
    "",
    "2. Scree Plot: Strongly favors one-factor structure",
    "",
    "3. Velicer's MAP Test: Both original and revised versions point to one factor",
    "",
    "4. Parallel Analysis: Raw data eigenvalue exceeded 95th percentile for 4 factors",
    "   - This was the only test that disagreed with the one-factor solution",
    "",
    "5. Forced 3-Factor Promax Rotation:",
    "   - Chi-squared = 45.72 (df = 12), p < 0.001",
    "   - Items did NOT load on their expected theoretical factors",
    "   - Dedication had 4 items (3, 4, 5, 6) instead of 3",
    "   - Vigor had only 2 items (1, 2) instead of 3",
    "",
    ("Conclusion: EFA strongly supports a single 'Work Engagement' factor", 13, True, TEAL),
], font_size=12)

shape_tc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(9.0), Inches(1.2))
shape_tc.fill.solid()
shape_tc.fill.fore_color.rgb = RGBColor(255, 243, 224)
shape_tc.line.color.rgb = WARM
shape_tc.line.width = Pt(2)
tf_tc = shape_tc.text_frame
tf_tc.word_wrap = True
tf_tc.margin_left = Inches(0.15)
tf_tc.margin_right = Inches(0.15)
p_tc = tf_tc.paragraphs[0]
p_tc.text = "Technical Commentary:"
p_tc.font.size = Pt(11)
p_tc.font.bold = True
p_tc.font.color.rgb = WARM
p_tc2 = tf_tc.add_paragraph()
p_tc2.text = ("While EFA and MAP tests supported a 1-factor solution, Parallel Analysis suggested 4 factors. "
              "This often occurs in multivariate data with high inter-item correlations (0.52-0.85), "
              "leading to 'over-factoring' -- a known limitation of Parallel Analysis in highly correlated datasets.")
p_tc2.font.size = Pt(10)
p_tc2.font.color.rgb = DARK_GRAY

add_footer(slide, FOOTER)

# ============================================================
# SLIDE 13: ONE-FACTOR CFA MODEL (Figure 1)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "CFA Model 1: One-Factor Structure", "Single Latent Factor: Work Engagement (Figure 1)")
add_accent(slide)
slide.shapes.add_picture('figures4_new/one_factor_model.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 14: THREE-FACTOR CFA MODEL (Figure 3)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "CFA Model 3: Three-Factor Structure", "Vigor, Dedication, Absorption (Figure 3)")
add_accent(slide)
slide.shapes.add_picture('figures4_new/three_factor_model.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 15: CFA COEFFICIENTS COMPARISON
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "CFA Standardized Coefficients", "All Three Models Compared (Table 4)")
add_accent(slide)
slide.shapes.add_picture('figures4_new/cfa_coefficients.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 16: CFA COEFFICIENTS TABLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "CFA Coefficients: Detailed Results", "All p-values < 0.0001 (Table 4)")
add_accent(slide)
data = [
    ["Item", "Subscale", "1-Factor Coef.", "2-Factor Coef.", "3-Factor Coef.", "z-value (1F)"],
    ["Item 1", "Vigor", "0.79", "0.80", "0.89", "50.42"],
    ["Item 2", "Vigor", "0.82", "0.83", "0.92", "59.95"],
    ["Item 3", "Dedication", "0.92", "0.92", "0.94", "132.99"],
    ["Item 4", "Dedication", "0.90", "0.90", "0.93", "109.15"],
    ["Item 5", "Vigor", "0.81", "0.76", "0.74", "55.85"],
    ["Item 6", "Absorption", "0.87", "0.89", "0.99", "83.55"],
    ["Item 7", "Dedication", "0.76", "0.77", "0.75", "44.83"],
    ["Item 8", "Absorption", "0.81", "0.83", "0.84", "57.54"],
    ["Item 9", "Absorption", "0.69", "0.76", "0.73", "33.19"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.0), data,
          col_widths=[1.0, 1.3, 1.5, 1.5, 1.5, 1.5])
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 17: GOODNESS-OF-FIT CHARTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "CFA Goodness-of-Fit Statistics", "None of the Three Models Achieved Good Fit")
add_accent(slide)
slide.shapes.add_picture('figures4_new/goodness_of_fit.png', Inches(0.2), Inches(1.3), Inches(9.6), Inches(5.8))
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 18: GOODNESS-OF-FIT TABLE WITH THRESHOLDS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Goodness-of-Fit: Full Table with Thresholds", "Table 5 | Red = Failed to Meet Acceptable Threshold")
add_accent(slide)
gof_data = [
    ["Fit Statistic", "One-Factor", "Two-Factor", "Three-Factor", "Good Fit Threshold"],
    ["Chi2 (df)", "633.90 (27)", "354.49 (26)", "247.76 (24)", "Non-significant"],
    ["RMSEA", "0.181", "0.192", "0.167", "< 0.05 (pref.) / < 0.08"],
    ["RMSEA 90% CI", "0.169-0.194", "0.175-0.192", "0.154-0.180", "--"],
    ["AIC", "16221.47", "8246.29", "8143.56", "Lower = better"],
    ["BIC", "16343.70", "8353.66", "8258.60", "Lower = better"],
    ["CFI", "0.895", "0.882", "0.920", "> 0.95"],
    ["TLI", "0.860", "0.837", "0.880", "> 0.95"],
    ["SRMR", "0.046", "0.049", "0.065", "< 0.08"],
]
rows = len(gof_data)
cols = len(gof_data[0])
table_shape = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.8))
table = table_shape.table
col_ws = [1.8, 1.6, 1.6, 1.6, 2.8]
for ci, w in enumerate(col_ws):
    table.columns[ci].width = Inches(w)
RED = RGBColor(192, 57, 43)
for i, row in enumerate(gof_data):
    for j, val in enumerate(row):
        cell = table.cell(i, j)
        cell.text = str(val)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(10)
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            elif i == 2 and j in (1, 2, 3):
                paragraph.font.bold = True
                paragraph.font.color.rgb = RED
            elif i in (6, 7) and j in (1, 2, 3):
                paragraph.font.bold = True
                paragraph.font.color.rgb = RED
        if i == 0:
            from pptx.oxml.ns import qn
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn('a:solidFill'), {})
            srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '6B2D7B'})
            solidFill.append(srgbClr)
            tcPr.append(solidFill)
        elif j == 4:
            from pptx.oxml.ns import qn
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn('a:solidFill'), {})
            srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFF3E0'})
            solidFill.append(srgbClr)
            tcPr.append(solidFill)
        cell.text_frame.paragraphs[0].font.size = Pt(10)

shape_note = slide.shapes.add_textbox(Inches(0.3), Inches(6.4), Inches(9.4), Inches(0.5))
tf_note = shape_note.text_frame
tf_note.word_wrap = True
p_note = tf_note.paragraphs[0]
p_note.text = "Red values = Failed to meet acceptable fit thresholds. Only SRMR met the < 0.08 criterion for all models."
p_note.font.size = Pt(10)
p_note.font.bold = True
p_note.font.color.rgb = RED
p_note.alignment = PP_ALIGN.CENTER
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 19: KEY FINDINGS SUMMARY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Key Findings Summary", "EFA vs. CFA Results")
add_accent(slide)

for i, (title, findings, color) in enumerate([
    ("EFA Results", [
        "One factor explains > 70% variance",
        "KMO = 0.922 (Marvellous)",
        "All loadings: 0.65 - 0.93",
        "MAP test: 1-factor",
        "Forced 3-factor: items misload",
    ], TEAL),
    ("CFA Results", [
        "RMSEA: 0.167 - 0.192 (poor)",
        "CFI: 0.882 - 0.920 (below 0.95)",
        "TLI: 0.837 - 0.880 (below 0.95)",
        "SRMR: 0.046 - 0.065 (acceptable)",
        "No model shows good overall fit",
    ], WARM),
    ("Interpretation", [
        "EFA supports 1-factor structure",
        "CFA cannot confirm any model",
        "3 subscales highly correlated",
        "Possible gender/culture effects",
        "Further research needed",
    ], PURPLE),
]):
    left = Inches(0.3 + i * 3.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.0), Inches(5.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2.5)
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = color; p.alignment = PP_ALIGN.CENTER
    for finding in findings:
        p2 = tf.add_paragraph(); p2.text = ""; p2.font.size = Pt(6)
        p3 = tf.add_paragraph(); p3.text = finding; p3.font.size = Pt(12); p3.font.color.rgb = DARK_GRAY; p3.alignment = PP_ALIGN.CENTER

add_footer(slide, FOOTER)

# ============================================================
# SLIDE 20: DISCUSSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Discussion", "Why Did CFA Fail to Confirm the EFA Results?")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Possible Explanations:", 15, True, PURPLE),
    "",
    "1. Gender Effects: All-female sample may experience work engagement differently",
    "   than mixed-gender samples used in original UWES validation",
    "",
    "2. Cultural Context: Swedish work culture emphasizes egalitarianism and",
    "   work-life balance, potentially affecting how engagement manifests",
    "",
    "3. Subscale Overlap: Inter-factor correlations of 0.79-0.84 suggest the three",
    "   dimensions are not clearly distinct -- they may be measuring the same thing",
    "",
    "4. Instrument Limitation: Shirom (2003) argued the three UWES dimensions",
    "   were not theoretically deduced and overlap conceptually",
    "",
    ("Comparison with Similar Studies:", 15, True, TEAL),
    "",
    "- Wefald et al. (2012): Similar sample size (n = 382), similar education level,",
    "  RMSEA = 0.18 and 0.16 -- almost identical to our 0.181 and 0.167",
    "- This suggests the issue may be with the UWES-9 instrument itself,",
    "  not specific to this sample",
], font_size=12)
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 21: DISCUSSION - WOMEN IN WORKFORCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Women & Work Engagement", "Implications for Workforce Policy & Research")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("What This Study Tells Us About Women's Engagement:", 14, True, PURPLE),
    "",
    "- Women in this sample reported moderately high engagement (Mean = 4.06/6.0)",
    "- Dedication (meaning, pride, enthusiasm) scored highest at 4.24",
    "- Item 7 'I am proud of the work that I do' had the highest individual score (4.62)",
    "- This suggests pride and purpose are strongest drivers for women's engagement",
    "",
    ("Practical Implications:", 14, True, TEAL),
    "",
    "- Organizations should focus on meaningfulness and purpose to engage women",
    "- Standard 3-subscale UWES-9 scoring may not be appropriate for all-female populations",
    "- Total score (1-factor) may be more reliable than subscale scores for women",
    "- Work engagement instruments may need gender-specific validation",
    "",
    ("Limitations:", 14, True, WARM),
    "",
    "- No data on specific occupations (only education level)",
    "- Swedish-speaking women only (though 21.6% had immigrant background)",
    "- Age range limited to 26-37 years",
    "- University-educated majority (61%) may not represent all working women",
], font_size=11)
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 22: METHODOLOGICAL IMPROVEMENTS (Lu Requirement)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Technical Improvements to the UWES-9 Framework", "Proposed Methodological Enhancements")
add_accent(slide)

for i, (title, items_list, color) in enumerate([
    ("Bifactor Modeling", [
        "Since the three subscales (Vigor, Dedication, Absorption) are highly correlated (0.79-0.84), a Bifactor Model would allow for:",
        "",
        "- A general 'Engagement' factor accounting for shared variance across all 9 items",
        "- Specific sub-factors capturing unique variance of each subscale",
        "- Better separation of general vs. specific engagement components",
        "- This addresses the core problem: are the subscales truly distinct?",
    ], PURPLE),
    ("Multigroup CFA & Invariance Testing", [
        "Future research should use Multigroup CFA to determine if the factor structure is invariant across:",
        "",
        "- Gender: Compare all-female vs. mixed-gender samples",
        "- Nationality: Test across Swedish, Norwegian, Finnish populations",
        "- Occupation: Compare engagement structure across professions",
        "- This would reveal whether poor fit is due to instrument or sample",
    ], TEAL),
]):
    top = Inches(1.5) + i * Inches(2.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), top, Inches(9.4), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = color
    for item in items_list:
        p2 = tf.add_paragraph()
        p2.text = item
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY
        if item == "":
            p2.font.size = Pt(4)

add_footer(slide, FOOTER)

# ============================================================
# SLIDE 23: CONCLUSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Conclusion")
add_accent(slide)
add_content(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Study Summary:", 15, True, PURPLE),
    "",
    "- EFA on n = 341 women: One-factor solution best fits the data",
    "  (KMO = 0.922, > 70% variance explained, all loadings 0.65-0.93)",
    "",
    "- CFA on n = 342 women: Could NOT confirm any model (1-, 2-, or 3-factor)",
    "  (RMSEA never below 0.167, CFI/TLI below 0.95 for all models)",
    "",
    ("Key Takeaway:", 15, True, TEAL),
    "",
    "The optimal factor structure of the UWES-9 remains unresolved.",
    "Further research is needed to understand how gender, nationality,",
    "and occupation influence the measurement of work engagement.",
    "",
    ("Future Directions:", 15, True, WARM),
    "",
    "- Test UWES-9 in larger all-female samples across different cultures",
    "- Develop gender-specific engagement instruments",
    "- Investigate whether a shorter version (UWES-3) performs better",
    "- Examine the role of occupation type in factor structure",
], font_size=12)
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 23: METHODOLOGICAL NOTES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "Factor Analysis Methods Applied", "EFA & CFA Techniques Used in This Study")
add_accent(slide)
data = [
    ["Method", "Purpose", "Key Finding"],
    ["KMO Test", "Sampling adequacy", "0.922 (Marvellous)"],
    ["Bartlett's Test", "Sphericity", "p < 0.001 (Significant)"],
    ["EFA (ML extraction)", "Identify factor structure", "1 factor, > 70% variance"],
    ["Scree Plot", "Visual eigenvalue analysis", "Supports 1-factor"],
    ["Velicer's MAP", "Minimum average partial", "Supports 1-factor"],
    ["Parallel Analysis", "Compare with random data", "Suggested 4 factors"],
    ["CFA: 1-Factor", "Test unidimensional model", "RMSEA = 0.181 (poor)"],
    ["CFA: 2-Factor", "Test V+D / A split", "RMSEA = 0.192 (poor)"],
    ["CFA: 3-Factor", "Test V / D / A structure", "RMSEA = 0.167 (poor)"],
    ["Cronbach's Alpha", "Internal consistency", "0.947 (Excellent)"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.0), data,
          col_widths=[2.5, 3.0, 3.9])
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 24: REFERENCES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_title_bar(slide, "References")
add_accent(slide)
add_content(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5), [
    ("Primary Study:", 12, True, PURPLE),
    "Willmer, M., Westerberg Jacobson, J. & Lindberg, M. (2019). Exploratory and Confirmatory Factor",
    "Analysis of the 9-Item UWES in a Multi-Occupational Female Sample. Frontiers in Psychology, 10, 2771.",
    "DOI: 10.3389/fpsyg.2019.02771",
    "",
    ("Supporting References:", 12, True, TEAL),
    "Schaufeli, W.B., Salanova, M., Gonzalez-Roma, V. & Bakker, A.B. (2002). The measurement of",
    "engagement and burnout. Journal of Happiness Studies, 3, 71-92.",
    "",
    "Kulikowski, K. (2017). Do we all agree on how to measure work engagement? Factorial validity of",
    "UWES-9 and UWES-17. Polish Psychological Bulletin, 48(3), 350-360.",
    "",
    "Hair, J.F., Black, W.C., Babin, B.J. & Anderson, R.E. (2019). Multivariate Data Analysis (8th Ed.).",
    "Cengage Learning.",
], font_size=11)
add_footer(slide, FOOTER)

# ============================================================
# SLIDE 25: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PURPLE)
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(3))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Thank You!"; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = ""; p2.font.size = Pt(10)
p3 = tf.add_paragraph()
p3.text = "Questions?"; p3.font.size = Pt(24); p3.font.color.rgb = GOLD; p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph(); p4.text = ""; p4.font.size = Pt(20)
p5 = tf.add_paragraph()
p5.text = "Lana Jalal Gidan"; p5.font.size = Pt(18); p5.font.bold = True; p5.font.color.rgb = WHITE; p5.alignment = PP_ALIGN.CENTER
p6 = tf.add_paragraph()
p6.text = "SSIE-605 | Professor Susan Lu"; p6.font.size = Pt(14); p6.font.color.rgb = WHITE; p6.alignment = PP_ALIGN.CENTER
p7 = tf.add_paragraph()
p7.text = "Watson College of Engineering | Binghamton University"; p7.font.size = Pt(13); p7.font.color.rgb = WHITE; p7.alignment = PP_ALIGN.CENTER

prs.save('Factor_Analysis_Women_WorkEngagement_Presentation.pptx')
print(f"\nPresentation saved: Factor_Analysis_Women_WorkEngagement_Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
