import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BLUE = RGBColor(0, 32, 96)
MEDIUM_BLUE = RGBColor(0, 70, 140)
LIGHT_BLUE = RGBColor(0, 112, 192)
ACCENT_GOLD = RGBColor(255, 192, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(240, 240, 240)

os.makedirs("figures", exist_ok=True)

def add_slide_background(slide, color=DARK_BLUE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=None):
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = ACCENT_GOLD
        p2.alignment = PP_ALIGN.LEFT

def add_content_box(slide, left, top, width, height, text_items, font_size=16, bullet=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(text_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            p.text = item[0]
            p.font.size = Pt(item[1])
            if len(item) > 2:
                p.font.bold = item[2]
            if len(item) > 3:
                p.font.color.rgb = item[3]
            else:
                p.font.color.rgb = DARK_GRAY
        else:
            prefix = "\u2022 " if bullet else ""
            p.text = f"{prefix}{item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    return txBox

def create_scree_plot():
    eigenvalues = [2.731, 2.218, 0.442, 0.341, 0.183, 0.085]
    components = list(range(1, 7))

    fig, ax = plt.subplots(figsize=(7, 4.5))
    ax.plot(components, eigenvalues, 'bo-', markersize=10, linewidth=2.5, color='#004080')
    ax.axhline(y=1.0, color='red', linestyle='--', linewidth=1.5, label='Eigenvalue = 1 (Kaiser Criterion)')
    ax.fill_between(components, eigenvalues, alpha=0.15, color='#004080')
    ax.set_xlabel('Component Number', fontsize=13, fontweight='bold')
    ax.set_ylabel('Eigenvalue', fontsize=13, fontweight='bold')
    ax.set_title('Scree Plot', fontsize=15, fontweight='bold', color='#004080')
    ax.set_xticks(components)
    ax.legend(fontsize=11)
    ax.grid(True, alpha=0.3)
    plt.tight_layout()
    plt.savefig('figures/scree_plot.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_variance_explained_chart():
    components = ['Comp 1', 'Comp 2', 'Comp 3', 'Comp 4', 'Comp 5', 'Comp 6']
    variance_pct = [45.52, 36.97, 7.36, 5.69, 3.04, 1.42]
    cumulative = [45.52, 82.49, 89.85, 95.54, 98.58, 100.00]

    fig, ax1 = plt.subplots(figsize=(7, 4.5))
    bars = ax1.bar(components, variance_pct, color='#004080', alpha=0.8, label='% of Variance')
    ax2 = ax1.twinx()
    ax2.plot(components, cumulative, 'ro-', markersize=8, linewidth=2.5, label='Cumulative %')
    ax2.axhline(y=60, color='green', linestyle='--', alpha=0.7, label='60% Threshold')

    ax1.set_ylabel('% of Variance', fontsize=12, fontweight='bold')
    ax2.set_ylabel('Cumulative %', fontsize=12, fontweight='bold')
    ax1.set_title('Total Variance Explained', fontsize=15, fontweight='bold', color='#004080')

    for bar, val in zip(bars, variance_pct):
        ax1.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 0.5,
                f'{val:.1f}%', ha='center', va='bottom', fontweight='bold', fontsize=10)

    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc='center right', fontsize=10)

    ax1.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig('figures/variance_explained.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_rotated_component_heatmap():
    variables = ['V1\n(Cavities)', 'V2\n(Shiny Teeth)', 'V3\n(Gums)', 'V4\n(Breath)', 'V5\n(Decay)', 'V6\n(Attractive)']
    data = np.array([
        [0.962, -0.027],
        [-0.057, 0.848],
        [0.934, -0.146],
        [-0.098, 0.854],
        [-0.933, -0.084],
        [0.083, 0.885]
    ])

    fig, ax = plt.subplots(figsize=(5, 5))
    sns.heatmap(data, annot=True, fmt='.3f', cmap='RdBu_r', center=0,
                xticklabels=['Factor 1\n(Health)', 'Factor 2\n(Cosmetic)'],
                yticklabels=variables,
                vmin=-1, vmax=1, linewidths=0.5, ax=ax,
                annot_kws={'size': 12, 'fontweight': 'bold'})
    ax.set_title('Rotated Component Matrix\n(Varimax Rotation)', fontsize=14, fontweight='bold', color='#004080')
    plt.tight_layout()
    plt.savefig('figures/rotated_matrix.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_factor_loading_plot():
    labels = ['V1', 'V2', 'V3', 'V4', 'V5', 'V6']
    f1 = [0.962, -0.057, 0.934, -0.098, -0.933, 0.083]
    f2 = [-0.027, 0.848, -0.146, 0.854, -0.084, 0.885]

    fig, ax = plt.subplots(figsize=(6, 5))

    colors = ['#004080', '#FF6600', '#004080', '#FF6600', '#004080', '#FF6600']
    for i, label in enumerate(labels):
        ax.scatter(f1[i], f2[i], s=120, c=colors[i], zorder=5, edgecolors='black')
        ax.annotate(label, (f1[i], f2[i]), textcoords="offset points",
                   xytext=(10, 5), fontsize=12, fontweight='bold')

    ax.axhline(y=0, color='gray', linestyle='-', linewidth=0.5)
    ax.axvline(x=0, color='gray', linestyle='-', linewidth=0.5)
    ax.set_xlim(-1.1, 1.1)
    ax.set_ylim(-1.1, 1.1)
    ax.set_xlabel('Factor 1 (Health Benefits)', fontsize=12, fontweight='bold')
    ax.set_ylabel('Factor 2 (Cosmetic Benefits)', fontsize=12, fontweight='bold')
    ax.set_title('Factor Loading Plot', fontsize=14, fontweight='bold', color='#004080')
    ax.grid(True, alpha=0.3)

    from matplotlib.patches import Patch
    legend_elements = [Patch(facecolor='#004080', label='Health Factor'),
                      Patch(facecolor='#FF6600', label='Cosmetic Factor')]
    ax.legend(handles=legend_elements, loc='lower left', fontsize=10)

    plt.tight_layout()
    plt.savefig('figures/factor_loading_plot.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_communalities_chart():
    variables = ['V1\nCavities', 'V2\nShiny', 'V3\nGums', 'V4\nBreath', 'V5\nDecay', 'V6\nAttractive']
    communalities = [0.926, 0.723, 0.894, 0.739, 0.878, 0.790]

    fig, ax = plt.subplots(figsize=(7, 4))
    bars = ax.barh(variables, communalities, color='#004080', alpha=0.85, edgecolor='white')
    ax.axvline(x=0.5, color='red', linestyle='--', linewidth=1.5, label='Minimum Threshold (0.5)')

    for bar, val in zip(bars, communalities):
        ax.text(val + 0.01, bar.get_y() + bar.get_height()/2,
               f'{val:.3f}', va='center', fontweight='bold', fontsize=11)

    ax.set_xlabel('Communality Value', fontsize=12, fontweight='bold')
    ax.set_title('Communalities (Extraction)', fontsize=14, fontweight='bold', color='#004080')
    ax.set_xlim(0, 1.1)
    ax.legend(fontsize=10)
    ax.grid(axis='x', alpha=0.3)
    plt.tight_layout()
    plt.savefig('figures/communalities.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_kmo_visual():
    categories = ['Unacceptable\n< 0.5', 'Miserable\n0.5-0.6', 'Mediocre\n0.6-0.7', 'Middling\n0.7-0.8', 'Meritorious\n0.8-0.9', 'Marvelous\n> 0.9']
    values = [0.4, 0.55, 0.65, 0.75, 0.85, 0.95]
    colors_list = ['#FF0000', '#FF6600', '#FFC000', '#92D050', '#00B050', '#006400']

    fig, ax = plt.subplots(figsize=(8, 3.5))
    bars = ax.barh(categories, values, color=colors_list, edgecolor='white', height=0.6)

    ax.axvline(x=0.813, color='blue', linestyle='--', linewidth=2.5, label='Our Study KMO = 0.813')
    ax.set_xlabel('KMO Value', fontsize=12, fontweight='bold')
    ax.set_title('KMO Sampling Adequacy Scale', fontsize=14, fontweight='bold', color='#004080')
    ax.legend(fontsize=12, loc='lower right')
    ax.set_xlim(0, 1.05)
    plt.tight_layout()
    plt.savefig('figures/kmo_scale.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_fa_procedure_flowchart():
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis('off')

    steps = [
        (5, 9.2, 'Step 1: Formulate the Problem'),
        (5, 8.0, 'Step 2: Construct Correlation Matrix'),
        (5, 6.8, 'Step 3: Determine FA Method (PCA/PAF)'),
        (5, 5.6, 'Step 4: Determine Number of Factors'),
        (5, 4.4, 'Step 5: Rotate Factors (Varimax)'),
        (5, 3.2, 'Step 6: Interpret Factors'),
        (5, 2.0, 'Step 7: Calculate Factor Scores'),
        (5, 0.8, 'Step 8: Evaluate Model Fit'),
    ]

    colors_flow = ['#001F60', '#002B80', '#003BA0', '#004BC0', '#005BE0', '#0070C0', '#0085E0', '#009BFF']

    for i, (x, y, text) in enumerate(steps):
        box = plt.Rectangle((1.2, y - 0.4), 7.6, 0.7, facecolor=colors_flow[i],
                           edgecolor='white', linewidth=1.5, alpha=0.9)
        ax.add_patch(box)
        ax.text(x, y, text, ha='center', va='center', fontsize=12,
               fontweight='bold', color='white')

        if i < len(steps) - 1:
            ax.annotate('', xy=(5, steps[i+1][1] + 0.35), xytext=(5, y - 0.45),
                       arrowprops=dict(arrowstyle='->', color='#FFC000', lw=2))

    ax.set_title('Factor Analysis Procedure', fontsize=16, fontweight='bold', color='#004080', pad=10)
    plt.tight_layout()
    plt.savefig('figures/fa_procedure.png', dpi=200, bbox_inches='tight')
    plt.close()

def create_correlation_heatmap():
    labels = ['V1', 'V2', 'V3', 'V4', 'V5', 'V6']
    corr = np.array([
        [1.000, -0.053, 0.873, -0.086, -0.858, 0.004],
        [-0.053, 1.000, -0.155, 0.572, 0.020, 0.640],
        [0.873, -0.155, 1.000, -0.248, -0.778, -0.018],
        [-0.086, 0.572, -0.248, 1.000, -0.007, 0.640],
        [-0.858, 0.020, -0.778, -0.007, 1.000, -0.136],
        [0.004, 0.640, -0.018, 0.640, -0.136, 1.000]
    ])

    fig, ax = plt.subplots(figsize=(6, 5))
    mask = np.triu(np.ones_like(corr, dtype=bool), k=1)
    sns.heatmap(corr, mask=mask, annot=True, fmt='.3f', cmap='RdBu_r', center=0,
                xticklabels=labels, yticklabels=labels,
                vmin=-1, vmax=1, linewidths=0.5, ax=ax,
                annot_kws={'size': 11})
    ax.set_title('Correlation Matrix\n(Toothpaste Survey Variables)', fontsize=14, fontweight='bold', color='#004080')
    plt.tight_layout()
    plt.savefig('figures/correlation_matrix.png', dpi=200, bbox_inches='tight')
    plt.close()

print("Generating all visualizations...")
create_scree_plot()
create_variance_explained_chart()
create_rotated_component_heatmap()
create_factor_loading_plot()
create_communalities_chart()
create_kmo_visual()
create_fa_procedure_flowchart()
create_correlation_heatmap()
print("All visualizations created!")

print("\nBuilding PowerPoint presentation...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(9), Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_GOLD
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Factor Analysis in Marketing"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Applications, Methodology & Real-World Case Study"
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT_GOLD
p2.alignment = PP_ALIGN.CENTER

p3 = tf.add_paragraph()
p3.text = ""
p3.font.size = Pt(12)

p4 = tf.add_paragraph()
p4.text = "SSIE 605 \u2013 Applied Multivariate Data Analysis"
p4.font.size = Pt(18)
p4.font.color.rgb = WHITE
p4.alignment = PP_ALIGN.CENTER

p5 = tf.add_paragraph()
p5.text = "Professor Susan Lu"
p5.font.size = Pt(16)
p5.font.color.rgb = RGBColor(180, 210, 255)
p5.alignment = PP_ALIGN.CENTER

p6 = tf.add_paragraph()
p6.text = ""
p6.font.size = Pt(10)

p7 = tf.add_paragraph()
p7.text = "Lana Jalal Gidan"
p7.font.size = Pt(22)
p7.font.bold = True
p7.font.color.rgb = WHITE
p7.alignment = PP_ALIGN.CENTER

p8 = tf.add_paragraph()
p8.text = "Binghamton University"
p8.font.size = Pt(14)
p8.font.color.rgb = RGBColor(180, 210, 255)
p8.alignment = PP_ALIGN.CENTER

shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.1), Inches(9), Inches(0.08))
shape2.fill.solid()
shape2.fill.fore_color.rgb = ACCENT_GOLD
shape2.line.fill.background()

# ============================================================
# SLIDE 2: AGENDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Agenda")

agenda_items = [
    "1.  Problem Statement",
    "2.  Article Information",
    "3.  What is Factor Analysis?",
    "4.  Factor Analysis Methodology",
    "5.  Case Study: Consumer Preferences (Toothpaste Market)",
    "6.  Statistical Results & Numerical Outputs",
    "7.  Factor Interpretation & Visualization",
    "8.  Discussion & Key Findings",
    "9.  Comments, Limitations & Improvements",
    "10. Conclusion"
]

for i, item in enumerate(agenda_items):
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8 + i * 0.5), Inches(8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True if i == 0 else False

# ============================================================
# SLIDE 3: PROBLEM STATEMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Problem Statement")

items = [
    "Marketing researchers collect large volumes of consumer data with many variables",
    "Too many correlated variables lead to multicollinearity and overfitting",
    "Understanding what really drives consumer decisions requires identifying hidden patterns",
    "Factor Analysis helps reduce complexity by uncovering latent (hidden) factors behind observed consumer behaviors",
    "Key Question: Can we identify a small number of underlying factors that explain consumer preferences in marketing surveys?"
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=17)

# ============================================================
# SLIDE 4: ARTICLE INFORMATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Article Information")

article_info = [
    ("Primary Reference:", 14, True, DARK_BLUE),
    "",
    ("Title: \"Factor Analysis as a Tool for Survey Analysis\"", 15, False, DARK_GRAY),
    ("Author: Noora Shrestha", 15, False, DARK_GRAY),
    ("Journal: American Journal of Applied Mathematics and Statistics, Vol. 9, No. 1, pp. 4-11", 14, False, DARK_GRAY),
    ("Year: 2021", 15, False, DARK_GRAY),
    ("DOI: 10.12691/ajams-9-1-2", 15, True, LIGHT_BLUE),
    "",
    ("Supplementary Reference:", 14, True, DARK_BLUE),
    "",
    ("Title: \"An Easy Approach to Exploratory Factor Analysis: Marketing Perspective\"", 15, False, DARK_GRAY),
    ("Authors: Noor Ul Hadi, Naziruddin Abdullah, Ilham Sentosa", 14, False, DARK_GRAY),
    ("Journal: Journal of Educational and Social Research, Vol. 6, No. 1", 14, False, DARK_GRAY),
    ("Year: 2016       DOI: 10.5901/jesr.2016.v6n1p215", 14, True, LIGHT_BLUE),
    "",
    ("Case Study Data: Toothpaste Consumer Survey (Malhotra, 2010)", 14, True, MEDIUM_BLUE),
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), article_info, bullet=False)

# ============================================================
# SLIDE 5: WHAT IS FACTOR ANALYSIS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "What is Factor Analysis?")

items = [
    "A statistical method used for data reduction and summarization",
    "Identifies underlying dimensions (latent factors) that explain correlations among observed variables",
    "An interdependence technique \u2013 no distinction between dependent and independent variables",
    "Reduces a large set of correlated variables into a smaller set of uncorrelated factors",
    "Widely used in marketing for: customer segmentation, brand perception, satisfaction surveys, and product positioning"
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(2.8), items, font_size=16)

txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Mathematical Model:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p2 = tf2.add_paragraph()
p2.text = "Xi = Ai1F1 + Ai2F2 + ... + AimFm + ViUi"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = MEDIUM_BLUE
p2.alignment = PP_ALIGN.CENTER
p3 = tf2.add_paragraph()
p3.text = "Where Xi = observed variable, Aij = factor loading, Fj = common factor, Ui = unique factor"
p3.font.size = Pt(12)
p3.font.color.rgb = DARK_GRAY

# ============================================================
# SLIDE 6: FA PROCEDURE (FLOWCHART)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Analysis Procedure")

slide.shapes.add_picture('figures/fa_procedure.png', Inches(1.0), Inches(1.5), Inches(8), Inches(5.5))

# ============================================================
# SLIDE 7: FA vs PCA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Analysis vs. PCA")

left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.5), Inches(4.5), Inches(5.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
left_box.line.color.rgb = DARK_BLUE

tf = left_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "Factor Analysis (FA)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
for item in [
    "Produces latent factors",
    "Analyzes covariance (communality)",
    "Close approximation to R matrix",
    "Explains covariance with minimum factors",
    "Multiple solutions possible",
    "Goals: data reduction, describe & explore relationships"
]:
    p2 = tf.add_paragraph()
    p2.text = f"\u2022 {item}"
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_GRAY
    p2.space_after = Pt(4)

right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.5))
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(255, 240, 230)
right_box.line.color.rgb = RGBColor(200, 100, 0)

tf2 = right_box.text_frame
tf2.word_wrap = True
tf2.margin_left = Inches(0.2)
tf2.margin_top = Inches(0.2)
p = tf2.paragraphs[0]
p.text = "Principal Component Analysis (PCA)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(200, 100, 0)
for item in [
    "Produces components",
    "Analyzes total variance",
    "Reproduces R matrix completely",
    "Extracts max variance with fewest factors",
    "Gives a unique solution",
    "Goal: variance reduction, fewer variables"
]:
    p2 = tf2.add_paragraph()
    p2.text = f"\u2022 {item}"
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_GRAY
    p2.space_after = Pt(4)

# ============================================================
# SLIDE 8: CASE STUDY INTRO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Case Study: Toothpaste Consumer Preferences")

items = [
    ("Study Objective:", 16, True, DARK_BLUE),
    "Determine the underlying benefits consumers seek when purchasing toothpaste",
    "",
    ("Survey Design:", 16, True, DARK_BLUE),
    "30 respondents rated 6 statements on a 7-point Likert scale (1 = Strongly Disagree to 7 = Strongly Agree)",
    "",
    ("Variables Measured:", 16, True, DARK_BLUE),
    "V1: Important to buy toothpaste that prevents cavities",
    "V2: I like toothpaste that gives shiny teeth",
    "V3: Toothpaste should strengthen gums",
    "V4: I prefer toothpaste that freshens breath",
    "V5: Prevention of tooth decay is NOT an important benefit",
    "V6: Most important consideration is attractive teeth",
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5), items, font_size=15, bullet=False)

# ============================================================
# SLIDE 9: CORRELATION MATRIX
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 1: Correlation Matrix")

slide.shapes.add_picture('figures/correlation_matrix.png', Inches(0.3), Inches(1.4), Inches(5.0), Inches(4.5))

items = [
    ("Key Observations:", 15, True, DARK_BLUE),
    "V1 & V3 highly correlated (r = 0.873) \u2192 both relate to dental health",
    "V2, V4, V6 correlated with each other \u2192 cosmetic/aesthetic concerns",
    "V1 & V5 strongly negative (r = -0.858) \u2192 opposite statements about decay prevention",
    "Pattern suggests TWO distinct groups of variables",
    "",
    ("Conclusion:", 14, True, MEDIUM_BLUE),
    "Strong correlations indicate factor analysis is appropriate"
]
add_content_box(slide, Inches(5.3), Inches(1.4), Inches(4.5), Inches(5.5), items, font_size=13, bullet=False)

# ============================================================
# SLIDE 10: KMO AND BARTLETT'S TEST
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 2: KMO & Bartlett's Test of Sphericity")

slide.shapes.add_picture('figures/kmo_scale.png', Inches(0.3), Inches(1.4), Inches(6.5), Inches(3.0))

result_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.6), Inches(9), Inches(2.5))
result_box.fill.solid()
result_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
result_box.line.color.rgb = DARK_BLUE

tf = result_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "Results from Shrestha (2021):"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

results = [
    ("KMO Measure of Sampling Adequacy: 0.813 (Meritorious)", Pt(15)),
    ("Bartlett's Test of Sphericity:", Pt(15)),
    ("    Chi-Square = 637.65,    Significance: p < 0.0001", Pt(15)),
    ("Conclusion: Data is SUITABLE for Factor Analysis", Pt(15)),
]
for text, size in results:
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.size = size
    p2.font.color.rgb = DARK_GRAY
    if "Conclusion" in text:
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0, 128, 0)

# ============================================================
# SLIDE 11: COMMUNALITIES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 3: Communalities")

slide.shapes.add_picture('figures/communalities.png', Inches(0.3), Inches(1.4), Inches(5.5), Inches(3.8))

items = [
    ("What are Communalities?", 15, True, DARK_BLUE),
    "Proportion of variance in a variable explained by the extracted factors",
    "",
    ("Interpretation:", 15, True, DARK_BLUE),
    "All values exceed 0.5 threshold",
    "V1 (Cavities): 0.926 \u2192 92.6% of variance explained",
    "V3 (Gums): 0.894 \u2192 89.4% explained",
    "V5 (Decay): 0.878 \u2192 87.8% explained",
    "",
    ("Conclusion:", 14, True, RGBColor(0, 128, 0)),
    "All variables are well-represented by the factor solution"
]
add_content_box(slide, Inches(5.8), Inches(1.4), Inches(4.0), Inches(5.5), items, font_size=12, bullet=False)

# ============================================================
# SLIDE 12: EIGENVALUES & VARIANCE EXPLAINED
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 4: Eigenvalues & Total Variance Explained")

slide.shapes.add_picture('figures/variance_explained.png', Inches(0.2), Inches(1.4), Inches(5.5), Inches(3.5))

table_data = [
    ['Comp', 'Eigenvalue', '% Variance', 'Cumulative %'],
    ['1', '2.731', '45.52%', '45.52%'],
    ['2', '2.218', '36.97%', '82.49%'],
    ['3', '0.442', '7.36%', '89.85%'],
    ['4', '0.341', '5.69%', '95.54%'],
    ['5', '0.183', '3.04%', '98.58%'],
    ['6', '0.085', '1.42%', '100.00%'],
]

table = slide.shapes.add_table(len(table_data), 4, Inches(5.8), Inches(1.5), Inches(4.0), Inches(3.0)).table

for i, row_data in enumerate(table_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            elif i <= 2:
                paragraph.font.bold = True
                paragraph.font.color.rgb = DARK_BLUE
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        elif i <= 2:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(220, 235, 255)

txBox = slide.shapes.add_textbox(Inches(5.8), Inches(4.7), Inches(4.0), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Kaiser's Rule: Retain factors with eigenvalue > 1.0"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p2 = tf.add_paragraph()
p2.text = "\u2192 2 factors retained"
p2.font.size = Pt(14)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0, 128, 0)
p3 = tf.add_paragraph()
p3.text = "\u2192 Explain 82.49% of total variance"
p3.font.size = Pt(14)
p3.font.bold = True
p3.font.color.rgb = RGBColor(0, 128, 0)
p4 = tf.add_paragraph()
p4.text = "(Exceeds 60% minimum threshold)"
p4.font.size = Pt(12)
p4.font.color.rgb = DARK_GRAY

# ============================================================
# SLIDE 13: SCREE PLOT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 4 (cont.): Scree Plot")

slide.shapes.add_picture('figures/scree_plot.png', Inches(1.0), Inches(1.4), Inches(7.5), Inches(5.0))

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "The scree plot confirms 2 factors: sharp drop after Component 2, then leveling off (the \"elbow\")"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 14: ROTATED COMPONENT MATRIX
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 5: Rotated Component Matrix (Varimax)")

slide.shapes.add_picture('figures/rotated_matrix.png', Inches(0.2), Inches(1.4), Inches(4.5), Inches(5.0))

items = [
    ("Factor 1 \u2013 Health Benefits:", 16, True, RGBColor(0, 64, 128)),
    "V1 (Cavities): 0.962",
    "V3 (Gums): 0.934",
    "V5 (Decay): -0.933 (reverse coded)",
    "",
    ("Factor 2 \u2013 Cosmetic Benefits:", 16, True, RGBColor(200, 100, 0)),
    "V6 (Attractive teeth): 0.885",
    "V4 (Fresh breath): 0.854",
    "V2 (Shiny teeth): 0.848",
    "",
    ("Rotation Method:", 14, True, DARK_BLUE),
    "Varimax with Kaiser Normalization",
    "Converged in 3 iterations",
]
add_content_box(slide, Inches(4.8), Inches(1.4), Inches(5.0), Inches(5.5), items, font_size=14, bullet=False)

# ============================================================
# SLIDE 15: FACTOR LOADING PLOT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 6: Factor Loading Plot & Interpretation")

slide.shapes.add_picture('figures/factor_loading_plot.png', Inches(0.5), Inches(1.4), Inches(5.0), Inches(5.0))

items = [
    ("Visual Interpretation:", 15, True, DARK_BLUE),
    "",
    "V1, V3, V5 cluster along Factor 1 axis \u2192 Health dimension",
    "",
    "V2, V4, V6 cluster along Factor 2 axis \u2192 Cosmetic dimension",
    "",
    "Clear separation confirms two distinct consumer benefit dimensions",
    "",
    ("Marketing Insight:", 15, True, RGBColor(0, 128, 0)),
    "",
    "Consumers evaluate toothpaste on TWO main dimensions: health protection and cosmetic appeal"
]
add_content_box(slide, Inches(5.5), Inches(1.4), Inches(4.3), Inches(5.5), items, font_size=13, bullet=False)

# ============================================================
# SLIDE 16: FACTOR SCORES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Step 7: Factor Scores")

items = [
    ("Factor Score Equation:", 16, True, DARK_BLUE),
    "",
    ("Fi = Wi1X1 + Wi2X2 + Wi3X3 + ... + WikXk", 18, True, MEDIUM_BLUE),
    "",
    ("Component Score Coefficient Matrix:", 16, True, DARK_BLUE),
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.4), Inches(2.0), items, font_size=15, bullet=False)

coeff_data = [
    ['Variable', 'Factor 1', 'Factor 2'],
    ['V1', '0.358', '0.011'],
    ['V2', '-0.001', '0.375'],
    ['V3', '0.345', '-0.043'],
    ['V4', '-0.017', '0.377'],
    ['V5', '-0.350', '-0.059'],
    ['V6', '0.052', '0.395'],
]

table = slide.shapes.add_table(len(coeff_data), 3, Inches(1.5), Inches(3.5), Inches(4.0), Inches(2.8)).table
for i, row_data in enumerate(coeff_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE

items2 = [
    ("Application:", 14, True, DARK_BLUE),
    "Factor scores replace original 6 variables with 2 factor scores per respondent",
    "These can be used in subsequent regression, clustering, or segmentation analysis"
]
add_content_box(slide, Inches(6.0), Inches(3.5), Inches(3.8), Inches(3.0), items2, font_size=13, bullet=False)

# ============================================================
# SLIDE 17: DISCUSSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Discussion & Key Findings")

items = [
    ("Key Finding 1: Two-Factor Solution", 16, True, DARK_BLUE),
    "Factor Analysis successfully reduced 6 correlated variables to 2 independent factors, explaining 82.49% of total variance",
    "",
    ("Key Finding 2: Clear Factor Structure", 16, True, DARK_BLUE),
    "Factor 1 (Health): Cavity prevention, gum strength, decay prevention \u2192 functional benefits",
    "Factor 2 (Cosmetic): Shiny teeth, fresh breath, attractive teeth \u2192 aesthetic benefits",
    "",
    ("Key Finding 3: High Data Suitability", 16, True, DARK_BLUE),
    "KMO = 0.813 (Meritorious) and Bartlett's test significant (p < 0.0001) confirm data appropriateness",
    "",
    ("Key Finding 4: Marketing Implications", 16, True, DARK_BLUE),
    "Toothpaste brands can position products along Health vs. Cosmetic dimensions",
    "Different consumer segments prioritize different benefit dimensions"
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5), items, font_size=14, bullet=False)

# ============================================================
# SLIDE 18: COMMENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Comments")

items = [
    "The study demonstrates that Factor Analysis is a powerful tool for uncovering hidden patterns in marketing survey data",
    "",
    "The toothpaste example is a classic marketing case that clearly shows how correlated variables group into meaningful factors",
    "",
    "The two-factor solution (Health vs. Cosmetic) aligns well with marketing theory \u2013 consumers evaluate products on functional vs. emotional/aesthetic dimensions",
    "",
    "The high communalities (all > 0.72) indicate that the factor solution captures most of the information in the original variables",
    "",
    "Varimax rotation produced a clean, interpretable solution where each variable loads strongly on only one factor \u2013 this is the ideal \"simple structure\" in Factor Analysis",
    "",
    "The methodology used (PCA extraction + Varimax rotation) follows standard practice recommended by Hair et al. (2006) and Tabachnick & Fidell (2007)"
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5), items, font_size=15)

# ============================================================
# SLIDE 19: LIMITATIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Limitations")

items = [
    "Small sample size (n = 30): Below the recommended minimum of 50 observations; a ratio of at least 5:1 (observations to variables) is recommended",
    "",
    "Limited to 6 variables: A broader set of toothpaste attributes (price, brand loyalty, packaging, ingredients) could reveal additional factors",
    "",
    "Only Varimax (orthogonal) rotation was used: Oblique rotation (e.g., Promax) should be explored if factors are expected to be correlated \u2013 health and cosmetic concerns may not be fully independent",
    "",
    "Cross-sectional design: The study captures preferences at one point in time; consumer preferences change over time",
    "",
    "Subjective factor naming: The labeling of factors as \"Health\" and \"Cosmetic\" involves researcher interpretation, which introduces potential bias",
    "",
    "No Confirmatory Factor Analysis (CFA): EFA results were not validated using CFA on an independent sample"
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5), items, font_size=14)

# ============================================================
# SLIDE 20: IMPROVEMENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Suggested Improvements")

items = [
    "Increase sample size to at least 200-300 respondents for more reliable factor extraction and stable loadings",
    "",
    "Include more variables (10-15) covering additional attributes like price sensitivity, brand preference, ingredient awareness, and packaging appeal",
    "",
    "Apply both orthogonal (Varimax) and oblique (Promax) rotation methods and compare the solutions for robustness",
    "",
    "Use Parallel Analysis (Horn, 1965) in addition to Kaiser's eigenvalue > 1 rule for more accurate determination of the number of factors",
    "",
    "Follow up with Confirmatory Factor Analysis (CFA) using Structural Equation Modeling (SEM) to validate the factor structure",
    "",
    "Conduct longitudinal studies to track how consumer preference factors evolve over time",
    "",
    "Apply the factor solution to real marketing decisions \u2013 use factor scores for customer segmentation and targeted advertising"
]
add_content_box(slide, Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5), items, font_size=14)

# ============================================================
# SLIDE 21: CONCLUSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Conclusion")

items = [
    "Factor Analysis successfully reduced 6 correlated consumer preference variables into 2 meaningful, independent factors",
    "",
    "The two factors \u2013 Health Benefits and Cosmetic Benefits \u2013 explain 82.49% of the total variance in the data",
    "",
    "KMO = 0.813 and significant Bartlett's test confirmed data suitability for factor analysis",
    "",
    "The rotated component matrix showed clean factor loadings, making interpretation straightforward",
    "",
    "Factor Analysis proves to be a valuable tool in marketing research for uncovering hidden dimensions of consumer decision-making",
    "",
    "The methodology can be applied to any marketing survey to reduce complexity and identify key drivers of consumer behavior"
]
add_content_box(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5), items, font_size=16)

# ============================================================
# SLIDE 22: REFERENCES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "References")

refs = [
    ("1. Shrestha, N. (2021). \"Factor Analysis as a Tool for Survey Analysis.\" American Journal of Applied Mathematics and Statistics, 9(1), 4-11. DOI: 10.12691/ajams-9-1-2", 12),
    ("", 8),
    ("2. Hadi, N.U., Abdullah, N., & Sentosa, I. (2016). \"An Easy Approach to Exploratory Factor Analysis: Marketing Perspective.\" Journal of Educational and Social Research, 6(1), 215. DOI: 10.5901/jesr.2016.v6n1p215", 12),
    ("", 8),
    ("3. Malhotra, N.K. (2010). Marketing Research: An Applied Orientation. 6th Edition, Pearson.", 12),
    ("", 8),
    ("4. Hair, J.F., Black, W.C., Babin, B.J., & Anderson, R.E. (2006). Multivariate Data Analysis. 7th Edition, Pearson.", 12),
    ("", 8),
    ("5. Tabachnick, B.G. & Fidell, L.S. (2007). Using Multivariate Statistics. 5th Edition, Pearson Education.", 12),
    ("", 8),
    ("6. Kaiser, H.F. (1974). \"An Index of Factorial Simplicity.\" Psychometrika, 39(1), 31-36. DOI: 10.1007/BF02291575", 12),
    ("", 8),
    ("7. Bartlett, M.S. (1951). \"The Effect of Standardization on a Chi-Square Approximation in Factor Analysis.\" Biometrika, 38(3/4), 337-344.", 12),
]

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, (text, size) in enumerate(refs):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = DARK_GRAY

# ============================================================
# SLIDE 23: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(9), Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_GOLD
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = ""
p2.font.size = Pt(16)

p3 = tf.add_paragraph()
p3.text = "Questions?"
p3.font.size = Pt(28)
p3.font.color.rgb = ACCENT_GOLD
p3.alignment = PP_ALIGN.CENTER

p4 = tf.add_paragraph()
p4.text = ""
p4.font.size = Pt(20)

p5 = tf.add_paragraph()
p5.text = "Lana Jalal Gidan"
p5.font.size = Pt(20)
p5.font.color.rgb = WHITE
p5.alignment = PP_ALIGN.CENTER

p6 = tf.add_paragraph()
p6.text = "SSIE 605 \u2013 Binghamton University"
p6.font.size = Pt(16)
p6.font.color.rgb = RGBColor(180, 210, 255)
p6.alignment = PP_ALIGN.CENTER

shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.1), Inches(9), Inches(0.08))
shape2.fill.solid()
shape2.fill.fore_color.rgb = ACCENT_GOLD
shape2.line.fill.background()

output_path = "Factor_Analysis_Presentation.pptx"
prs.save(output_path)
print(f"\n{'='*80}")
print(f"PRESENTATION SAVED: {output_path}")
print(f"Total Slides: {len(prs.slides)}")
print(f"{'='*80}")
print("\nSlide Summary:")
slide_names = [
    "Title Slide", "Agenda", "Problem Statement", "Article Information",
    "What is Factor Analysis?", "FA Procedure (Flowchart)", "FA vs PCA",
    "Case Study Intro", "Correlation Matrix", "KMO & Bartlett's Test",
    "Communalities", "Eigenvalues & Variance Explained", "Scree Plot",
    "Rotated Component Matrix", "Factor Loading Plot", "Factor Scores",
    "Discussion & Key Findings", "Comments", "Limitations",
    "Suggested Improvements", "Conclusion", "References", "Thank You"
]
for i, name in enumerate(slide_names, 1):
    print(f"  Slide {i}: {name}")
