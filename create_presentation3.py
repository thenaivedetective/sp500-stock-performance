import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BLUE = RGBColor(0, 32, 96)
MEDIUM_BLUE = RGBColor(0, 70, 140)
LIGHT_BLUE = RGBColor(0, 112, 192)
ACCENT_GOLD = RGBColor(255, 192, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(240, 240, 240)

os.makedirs("figures3", exist_ok=True)

def add_slide_background(slide, color=DARK_BLUE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=None):
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = ACCENT_GOLD
        p2.alignment = PP_ALIGN.LEFT

def add_content_box(slide, left, top, width, height, items, font_size=14, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, size, bold, color = item
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
        elif item == "":
            p.text = ""
            p.font.size = Pt(6)
        else:
            prefix = "\u2022 " if bullet else ""
            p.text = f"{prefix}{item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY

def add_table(slide, left, top, width, height, data, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            if i == 0:
                from pptx.oxml.ns import qn
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '002060'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

def add_accent_bar(slide, top=Inches(1.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()

# ============================================================
# FIGURE 1: Scree Plot (Only published eigenvalues from Table 3)
# ============================================================
eigenvalues_published = [3.262, 2.057, 1.300]
components_published = [1, 2, 3]

fig, ax = plt.subplots(figsize=(10, 6))
ax.bar(components_published, eigenvalues_published, color=['#2E86C1', '#E67E22', '#27AE60'],
       alpha=0.85, width=0.5, edgecolor='black', linewidth=0.8)
ax.axhline(y=1, color='r', linestyle='--', linewidth=1.5, label='Kaiser Criterion (Eigenvalue = 1)')
for i, ev in enumerate(eigenvalues_published):
    ax.annotate(f'{ev:.3f}', (components_published[i], ev), textcoords="offset points",
                xytext=(0, 12), ha='center', fontsize=12, fontweight='bold')
ax.set_xlabel('Component Number', fontsize=13, fontweight='bold')
ax.set_ylabel('Eigenvalue', fontsize=13, fontweight='bold')
ax.set_title('Eigenvalues of Extracted Factors (Kaiser Criterion > 1)\n(Wang & Song, 2025 — Table 3, n=100 companies)',
             fontsize=14, fontweight='bold')
ax.set_xticks(components_published)
ax.set_xticklabels(['C1: Solvency', 'C2: Profitability', 'C3: Oper. Dev.'])
ax.legend(fontsize=11)
ax.grid(True, alpha=0.3, axis='y')
ax.text(0.5, -0.12, 'Note: Only eigenvalues for the 3 retained factors are reported in Table 3',
        transform=ax.transAxes, ha='center', fontsize=10, fontstyle='italic', color='gray')
plt.tight_layout()
plt.savefig('figures3/scree_plot.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 2: Variance Explained (Real data from Table 3)
# ============================================================
factors = ['F1\nSolvency', 'F2\nProfitability', 'F3\nOper. Dev.\nCapability']
variance_before = [36.248, 22.852, 14.449]
variance_after = [30.369, 28.025, 15.155]
cumulative_after = [30.369, 58.395, 73.550]

fig, ax = plt.subplots(figsize=(10, 6))
x = np.arange(len(factors))
width = 0.3

bars1 = ax.bar(x - width/2, variance_before, width, label='Before Rotation', color='#2E86C1', alpha=0.85)
bars2 = ax.bar(x + width/2, variance_after, width, label='After Rotation', color='#E67E22', alpha=0.85)

ax2 = ax.twinx()
ax2.plot(x, cumulative_after, 'rs-', linewidth=2.5, markersize=10, label='Cumulative % (After Rotation)')
for i, cum in enumerate(cumulative_after):
    ax2.annotate(f'{cum:.1f}%', (x[i], cum), textcoords="offset points",
                 xytext=(0, 12), ha='center', fontsize=11, fontweight='bold', color='red')

for bars in [bars1, bars2]:
    for bar in bars:
        height = bar.get_height()
        ax.annotate(f'{height:.1f}%', xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 5), textcoords="offset points", ha='center', fontsize=9)

ax.set_xlabel('Factor', fontsize=13, fontweight='bold')
ax.set_ylabel('Individual Variance Explained (%)', fontsize=12, fontweight='bold')
ax2.set_ylabel('Cumulative Variance (%)', fontsize=12, fontweight='bold', color='red')
ax.set_title('Variance Explained Before & After Rotation\n(Wang & Song, 2025 — Table 3, Total = 73.55%)',
             fontsize=14, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(factors)
ax.legend(loc='upper left', fontsize=10)
ax2.legend(loc='upper right', fontsize=10)
ax.grid(True, alpha=0.3, axis='y')
plt.tight_layout()
plt.savefig('figures3/variance_explained.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 3: Rotated Factor Loading Matrix Heatmap (Table 4)
# ============================================================
variables = ['X1: Inventory\nTurnover', 'X2: A/R\nTurnover', 'X3: Current\nRatio',
             'X4: Quick\nRatio', 'X5: Debt-to-Asset\n(Reciprocal)', 'X6: Asset\nReturn Ratio',
             'X7: Return on\nAssets (ROA)', 'X8: EBIT\nMargin', 'X9: Revenue\nGrowth Rate']
factor_names = ['F1: Solvency', 'F2: Profitability', 'F3: Oper. Dev. Capability']

loadings = np.array([
    [0.001, 0.264, -0.479],
    [-0.067, 0.070, 0.732],
    [0.956, 0.086, -0.020],
    [0.929, 0.128, -0.154],
    [0.945, 0.116, -0.034],
    [0.010, 0.955, -0.065],
    [0.144, 0.940, -0.071],
    [0.181, 0.778, 0.103],
    [-0.082, 0.097, 0.744]
])

fig, ax = plt.subplots(figsize=(9, 8))
mask = np.abs(loadings) < 0.30
cmap = sns.diverging_palette(240, 10, as_cmap=True)
sns.heatmap(loadings, annot=True, fmt='.3f', cmap=cmap, center=0,
            xticklabels=factor_names, yticklabels=variables,
            vmin=-1, vmax=1, linewidths=0.5, ax=ax,
            annot_kws={'size': 11, 'fontweight': 'bold'})
ax.set_title('Rotated Factor Loading Matrix (Varimax)\n(Wang & Song, 2025 — Table 4, n=100 companies)',
             fontsize=13, fontweight='bold', pad=15)
ax.tick_params(axis='y', labelsize=10)
ax.tick_params(axis='x', labelsize=11)
plt.tight_layout()
plt.savefig('figures3/factor_loading_matrix.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 4: Factor Loading Bar Chart (Table 4 - Primary Loadings)
# ============================================================
fig, axes = plt.subplots(1, 3, figsize=(14, 5), sharey=True)
factor_colors = ['#2E86C1', '#E67E22', '#27AE60']

f1_vars = ['X3: Current Ratio', 'X5: Debt-to-Asset (R)', 'X4: Quick Ratio']
f1_loads = [0.956, 0.945, 0.929]
axes[0].barh(f1_vars, f1_loads, color=factor_colors[0], alpha=0.85, height=0.5)
for i, v in enumerate(f1_loads):
    axes[0].text(v + 0.02, i, f'{v:.3f}', va='center', fontsize=11, fontweight='bold')
axes[0].set_title('F1: Solvency Factor', fontsize=12, fontweight='bold', color=factor_colors[0])
axes[0].set_xlim(0, 1.1)

f2_vars = ['X6: Asset Return Ratio', 'X7: ROA', 'X8: EBIT Margin']
f2_loads = [0.955, 0.940, 0.778]
axes[1].barh(f2_vars, f2_loads, color=factor_colors[1], alpha=0.85, height=0.5)
for i, v in enumerate(f2_loads):
    axes[1].text(v + 0.02, i, f'{v:.3f}', va='center', fontsize=11, fontweight='bold')
axes[1].set_title('F2: Profitability Factor', fontsize=12, fontweight='bold', color=factor_colors[1])
axes[1].set_xlim(0, 1.1)

f3_vars = ['X9: Revenue Growth', 'X2: A/R Turnover']
f3_loads = [0.744, 0.732]
axes[2].barh(f3_vars, f3_loads, color=factor_colors[2], alpha=0.85, height=0.5)
for i, v in enumerate(f3_loads):
    axes[2].text(v + 0.02, i, f'{v:.3f}', va='center', fontsize=11, fontweight='bold')
axes[2].set_title('F3: Oper. Dev. Capability', fontsize=12, fontweight='bold', color=factor_colors[2])
axes[2].set_xlim(0, 1.1)

plt.suptitle('Primary Factor Loadings by Factor (Wang & Song, 2025 — Table 4)',
             fontsize=13, fontweight='bold', y=1.02)
plt.tight_layout()
plt.savefig('figures3/factor_loadings_bar.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 5: Component Score Coefficient Matrix (Table 5)
# ============================================================
score_coefficients = np.array([
    [-0.062, 0.108, -0.358],
    [0.026, 0.038, 0.544],
    [0.367, -0.049, 0.062],
    [0.342, -0.029, -0.040],
    [0.359, -0.036, 0.051],
    [-0.086, 0.397, -0.044],
    [-0.033, 0.379, -0.038],
    [0.010, 0.309, 0.096],
    [0.019, 0.051, 0.553]
])

var_labels_short = ['X1: Inv. Turn.', 'X2: A/R Turn.', 'X3: Curr. Ratio',
                    'X4: Quick Ratio', 'X5: D/A Recip.', 'X6: Asset Ret.',
                    'X7: ROA', 'X8: EBIT Marg.', 'X9: Rev. Growth']

fig, ax = plt.subplots(figsize=(8, 7))
sns.heatmap(score_coefficients, annot=True, fmt='.3f', cmap='RdYlGn', center=0,
            xticklabels=factor_names, yticklabels=var_labels_short,
            linewidths=0.5, ax=ax, annot_kws={'size': 10, 'fontweight': 'bold'})
ax.set_title('Component Score Coefficient Matrix\n(Wang & Song, 2025 — Table 5)',
             fontsize=13, fontweight='bold', pad=15)
ax.tick_params(axis='y', labelsize=10)
ax.tick_params(axis='x', labelsize=11)
plt.tight_layout()
plt.savefig('figures3/score_coefficients.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 6: KMO and Bartlett's Test Visual (Table 2)
# ============================================================
fig, axes = plt.subplots(1, 2, figsize=(12, 5))

categories = ['Poor\n(<0.5)', 'Mediocre\n(0.5-0.6)', 'Middling\n(0.6-0.7)',
              'Meritorious\n(0.7-0.8)', 'Great\n(0.8-0.9)', 'Superb\n(>0.9)']
values = [0.5, 0.6, 0.7, 0.8, 0.9, 1.0]
colors_kmo = ['#E74C3C', '#E67E22', '#F1C40F', '#2ECC71', '#27AE60', '#1ABC9C']
bars = axes[0].bar(categories, values, color=colors_kmo, alpha=0.4, edgecolor='gray')
bars[3].set_alpha(0.9)
axes[0].axhline(y=0.706, color='red', linestyle='--', linewidth=2.5, label=f'This Study: KMO = 0.706')
axes[0].set_ylabel('KMO Value', fontsize=12, fontweight='bold')
axes[0].set_title('KMO Measure of Sampling Adequacy', fontsize=13, fontweight='bold')
axes[0].legend(fontsize=11, loc='upper left')
axes[0].set_ylim(0, 1.1)
axes[0].grid(True, alpha=0.3, axis='y')

test_info = ['KMO = 0.706', "Bartlett's \u03C7\u00B2 = 573.599", 'df = 36', 'p = 0.000']
y_pos = [0.8, 0.6, 0.4, 0.2]
bar_colors = ['#2ECC71', '#3498DB', '#9B59B6', '#E74C3C']
bars2 = axes[1].barh(test_info[::-1], [0.706, 573.599/700, 36/50, 1.0][::-1],
                      color=bar_colors[::-1], alpha=0.8, height=0.5)
axes[1].set_xlim(0, 1.2)
axes[1].set_title("Bartlett's Sphericity Test Results", fontsize=13, fontweight='bold')
axes[1].tick_params(axis='y', labelsize=12)
axes[1].set_xlabel('Normalized Scale', fontsize=11)

result_text = "Result: Data is suitable\nfor factor analysis"
axes[1].text(0.6, 0.5, result_text, transform=axes[1].transAxes,
             fontsize=12, fontweight='bold', color='green',
             bbox=dict(boxstyle='round', facecolor='lightgreen', alpha=0.3),
             ha='center', va='center')

plt.suptitle('Suitability Tests (Wang & Song, 2025 — Table 2)', fontsize=14, fontweight='bold', y=1.02)
plt.tight_layout()
plt.savefig('figures3/kmo_bartlett.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 7: Factor Score Ranges (Table 6 partial)
# ============================================================
fig, ax = plt.subplots(figsize=(10, 6))
factor_labels = ['F1: Solvency\n(Debt Repayment)', 'F2: Profitability', 'F3: Operational\nDev. Capability']
mins = [-24.85, -0.32, -156.91]
maxs = [46.09, 66.24, 947.95]
colors_range = ['#2E86C1', '#E67E22', '#27AE60']

for i, (label, mn, mx) in enumerate(zip(factor_labels, mins, maxs)):
    ax.barh(label, mx - mn, left=mn, color=colors_range[i], alpha=0.7, height=0.5,
            edgecolor='black', linewidth=0.8)
    ax.plot(mn, i, 'v', color='red', markersize=12)
    ax.plot(mx, i, '^', color='green', markersize=12)
    ax.text(mn - 15, i, f'{mn}', ha='right', va='center', fontsize=10, fontweight='bold', color='red')
    ax.text(mx + 15, i, f'{mx}', ha='left', va='center', fontsize=10, fontweight='bold', color='green')

ax.axvline(x=0, color='black', linestyle='-', linewidth=1)
ax.set_title('Factor Score Ranges Across 100 Companies\n(Wang & Song, 2025 — Table 6)',
             fontsize=14, fontweight='bold')
ax.set_xlabel('Factor Score', fontsize=12, fontweight='bold')
ax.legend(['Min Score', 'Max Score'], loc='lower right', fontsize=11)
ax.grid(True, alpha=0.3, axis='x')
plt.tight_layout()
plt.savefig('figures3/factor_score_ranges.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 8: Variable Definitions Overview
# ============================================================
fig, ax = plt.subplots(figsize=(10, 7))
ax.axis('off')

dim_colors = {'Operational Capability': '#3498DB', 'Debt-paying Ability': '#E74C3C',
              'Profitability': '#2ECC71', 'Development Capability': '#9B59B6'}
dim_data = [
    ('Operational Capability', 'X1', 'Inventory Turnover Ratio'),
    ('Operational Capability', 'X2', 'Accounts Receivable Turnover Ratio'),
    ('Debt-paying Ability', 'X3', 'Current Ratio'),
    ('Debt-paying Ability', 'X4', 'Quick Ratio'),
    ('Debt-paying Ability', 'X5', 'Reciprocal of Debt-to-Asset Ratio'),
    ('Profitability', 'X6', 'Asset Return Ratio'),
    ('Profitability', 'X7', 'Return on Assets (ROA)'),
    ('Profitability', 'X8', 'EBIT Margin'),
    ('Development Capability', 'X9', 'Revenue Growth Rate')
]

y_start = 0.92
for dim_name, color in dim_colors.items():
    items = [(d[1], d[2]) for d in dim_data if d[0] == dim_name]
    ax.text(0.02, y_start, dim_name, fontsize=13, fontweight='bold', color=color,
            transform=ax.transAxes, verticalalignment='top')
    for j, (var_id, var_name) in enumerate(items):
        y_pos = y_start - 0.045 * (j + 1)
        ax.text(0.06, y_pos, f'{var_id}:', fontsize=11, fontweight='bold',
                transform=ax.transAxes, verticalalignment='top', color='#333333')
        ax.text(0.14, y_pos, var_name, fontsize=11,
                transform=ax.transAxes, verticalalignment='top', color='#555555')
    y_start = y_start - 0.045 * (len(items) + 1) - 0.02

ax.text(0.5, 0.01, 'Source: Wang & Song (2025) — Table 1',
        fontsize=10, fontstyle='italic', ha='center', transform=ax.transAxes, color='gray')
ax.set_title('Financial Performance Evaluation Indicators\n(9 Variables Across 4 Dimensions)',
             fontsize=14, fontweight='bold', pad=20)
plt.tight_layout()
plt.savefig('figures3/variable_definitions.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 9: EFA Procedure Flowchart
# ============================================================
fig, ax = plt.subplots(figsize=(10, 7))
ax.axis('off')
ax.set_xlim(0, 10)
ax.set_ylim(0, 10)

steps = [
    ('Step 1', 'Data Collection', '100 listed real estate companies\n2023 financial data, 9 indicators'),
    ('Step 2', 'Suitability Testing', 'KMO = 0.706 (Meritorious)\nBartlett\'s \u03C7\u00B2 = 573.599, p = 0.000'),
    ('Step 3', 'Factor Extraction', 'Principal Component Analysis\nKaiser Criterion: Eigenvalue > 1'),
    ('Step 4', 'Factor Rotation', 'Varimax Rotation\n3 factors extracted'),
    ('Step 5', 'Factor Interpretation', 'F1: Solvency, F2: Profitability\nF3: Operational Dev. Capability'),
    ('Step 6', 'Scoring & Evaluation', 'Comprehensive scores for 100\ncompanies using factor weights')
]
colors_steps = ['#2E86C1', '#3498DB', '#1ABC9C', '#27AE60', '#E67E22', '#E74C3C']
y_positions = [8.5, 7.0, 5.5, 4.0, 2.5, 1.0]

for i, (step_num, title, desc) in enumerate(steps):
    rect = plt.Rectangle((1, y_positions[i] - 0.4), 8, 1.1, linewidth=2,
                          edgecolor=colors_steps[i], facecolor=colors_steps[i], alpha=0.15,
                          transform=ax.transData, zorder=1)
    ax.add_patch(rect)
    ax.text(1.3, y_positions[i] + 0.4, f'{step_num}: {title}',
            fontsize=13, fontweight='bold', color=colors_steps[i], verticalalignment='top')
    ax.text(1.3, y_positions[i] - 0.05, desc,
            fontsize=10, color='#333333', verticalalignment='top')
    if i < len(steps) - 1:
        ax.annotate('', xy=(5, y_positions[i] - 0.5), xytext=(5, y_positions[i+1] + 0.8),
                    arrowprops=dict(arrowstyle='->', color='gray', lw=2))

ax.set_title('Factor Analysis Procedure\n(Wang & Song, 2025)', fontsize=14, fontweight='bold', pad=15)
plt.tight_layout()
plt.savefig('figures3/efa_procedure.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 10: Eigenvalue Comparison (Before vs After Rotation)
# ============================================================
fig, ax = plt.subplots(figsize=(10, 6))
x = np.arange(3)
width = 0.35
eigenvalues_before = [3.262, 2.057, 1.300]
eigenvalues_after = [2.733, 2.522, 1.364]

bars1 = ax.bar(x - width/2, eigenvalues_before, width, label='Before Rotation',
               color='#2E86C1', alpha=0.85, edgecolor='black', linewidth=0.5)
bars2 = ax.bar(x + width/2, eigenvalues_after, width, label='After Rotation',
               color='#E67E22', alpha=0.85, edgecolor='black', linewidth=0.5)

for bar in bars1:
    h = bar.get_height()
    ax.text(bar.get_x() + bar.get_width()/2, h + 0.05, f'{h:.3f}', ha='center', fontsize=11, fontweight='bold')
for bar in bars2:
    h = bar.get_height()
    ax.text(bar.get_x() + bar.get_width()/2, h + 0.05, f'{h:.3f}', ha='center', fontsize=11, fontweight='bold')

ax.axhline(y=1, color='r', linestyle='--', linewidth=1.5, alpha=0.7, label='Kaiser Criterion')
ax.set_xlabel('Factor', fontsize=13, fontweight='bold')
ax.set_ylabel('Eigenvalue', fontsize=13, fontweight='bold')
ax.set_title('Eigenvalues Before & After Varimax Rotation\n(Wang & Song, 2025 — Table 3)',
             fontsize=14, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(['F1: Solvency', 'F2: Profitability', 'F3: Oper. Dev.\nCapability'])
ax.legend(fontsize=11)
ax.grid(True, alpha=0.3, axis='y')
plt.tight_layout()
plt.savefig('figures3/eigenvalue_comparison.png', dpi=200, bbox_inches='tight')
plt.close()

print("All 10 figures generated successfully in figures3/")

# ============================================================
# BUILD POWERPOINT PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SLIDE 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(3))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 25, 75)
shape.line.fill.background()

add_content_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5), [
    ("SSIE-605: Applied Multivariate Data Analysis", 14, False, ACCENT_GOLD)
])
add_content_box(slide, Inches(0.5), Inches(2.7), Inches(9), Inches(2.5), [
    ("Factor Analysis in Financial\nPerformance Evaluation", 36, True, WHITE),
    ("", 6, False, WHITE),
    ("Evaluating Listed Real Estate Companies Using\nExploratory Factor Analysis", 18, False, ACCENT_GOLD),
])
add_content_box(slide, Inches(0.5), Inches(5.7), Inches(9), Inches(1.5), [
    ("Lana Jalal Gidan", 16, True, WHITE),
    ("Binghamton University | Professor Susan Lu", 13, False, RGBColor(180, 198, 231)),
])
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.5), Inches(9), Inches(0.04))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GOLD
bar.line.fill.background()

# SLIDE 2: Presentation Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Presentation Overview")
add_accent_bar(slide)

sections = [
    ("1. Introduction", "Why factor analysis matters in financial performance evaluation"),
    ("2. Study Background", "Real estate industry challenges and research context"),
    ("3. Research Design", "Sample selection, variables, and methodology"),
    ("4. Factor Analysis Methodology", "Step-by-step EFA procedure"),
    ("5. Suitability Testing", "KMO and Bartlett's test results"),
    ("6. Factor Extraction & Rotation", "Eigenvalues, variance explained, and factor loadings"),
    ("7. Factor Interpretation", "Naming and understanding the three extracted factors"),
    ("8. Company Scoring & Evaluation", "Comprehensive financial performance scores"),
    ("9. Key Findings & Implications", "Practical takeaways for industry and investors"),
    ("10. Conclusion", "Summary and significance of factor analysis in finance"),
]
for i, (title, desc) in enumerate(sections):
    y = Inches(1.5) + Inches(i * 0.55)
    add_content_box(slide, Inches(0.5), y, Inches(9), Inches(0.55), [
        (f"{title}  |  {desc}", 13, False, DARK_GRAY)
    ])

# SLIDE 3: What is Factor Analysis?
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "What is Factor Analysis?", "Reducing Complexity in Financial Data")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Definition", 18, True, DARK_BLUE),
    "Factor Analysis is a statistical method that reduces many observed variables into",
    "a smaller number of underlying 'factors' that explain shared patterns in the data.",
    "",
    ("Why Use It in Finance?", 18, True, DARK_BLUE),
    "Financial performance involves many correlated indicators (profitability, debt ratios, growth, etc.)",
    "Factor analysis identifies the key underlying dimensions that drive company performance",
    "Reduces 9+ financial indicators into 2-3 core factors for clearer evaluation",
    "",
    ("Two Main Types", 18, True, DARK_BLUE),
    "Exploratory Factor Analysis (EFA): Discovers the underlying factor structure",
    "Confirmatory Factor Analysis (CFA): Tests a pre-specified factor structure",
    "",
    ("This Study", 16, True, MEDIUM_BLUE),
    "Uses EFA with Principal Component Analysis to evaluate 100 real estate companies (Wang & Song, 2025)",
], font_size=14, bullet=False)

# SLIDE 4: Study Background
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Study Background", "China's Real Estate Industry Context")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("The Problem", 18, True, DARK_BLUE),
    "Real estate is a critical economic driver in China, linked to employment and related industries",
    "The Evergrande Group crisis (2021) exposed severe financial risks across the sector",
    "The pandemic (2022) caused dramatic market contraction and broken capital chains",
    "",
    ("Why This Research Matters", 18, True, DARK_BLUE),
    "Traditional single-indicator analysis cannot capture the complexity of financial health",
    "Real estate companies face intertwined pressures: market fluctuations, capital turnover,",
    "   cost control, regulatory policies, and misappropriation risks",
    "Scientific multi-dimensional evaluation is needed for corporate managers and investors",
    "",
    ("Research Goal", 18, True, MEDIUM_BLUE),
    "Use factor analysis to identify the core dimensions driving financial performance",
    "of listed real estate companies and evaluate their relative standings",
], font_size=14, bullet=False)

# SLIDE 5: Research Design
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Research Design", "Sample, Data Sources, and Approach")
add_accent_bar(slide)

data = [
    ["Aspect", "Details"],
    ["Sample Size", "100 listed real estate companies"],
    ["Data Period", "2023 financial year data"],
    ["Data Sources", "CSMAR database, Wind database, company reports"],
    ["Exclusions", "ST-labeled companies; incomplete data removed"],
    ["Coverage", "Multiple regions, sizes, and development stages"],
    ["Variables", "9 financial indicators across 4 dimensions"],
    ["Method", "Exploratory Factor Analysis (PCA extraction)"],
    ["Rotation", "Varimax (orthogonal rotation)"],
    ["Software", "SPSS and Stata"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5), data,
          col_widths=[2.5, 6.5])

# SLIDE 6: Variable Definitions
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Financial Performance Indicators", "9 Variables Across 4 Dimensions (Table 1)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/variable_definitions.png',
                         Inches(0.5), Inches(1.4), Inches(9), Inches(5.8))

# SLIDE 7: Variable Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Variable Definitions — Detail", "Table 1: Wang & Song (2025)")
add_accent_bar(slide)

var_data = [
    ["Dimension", "Variable", "Name", "Formula"],
    ["Operational Cap.", "X1", "Inventory Turnover", "COGS / Ending Inventory"],
    ["Operational Cap.", "X2", "A/R Turnover", "Revenue / Ending A/R"],
    ["Debt-paying", "X3", "Current Ratio", "Current Assets / Current Liabilities"],
    ["Debt-paying", "X4", "Quick Ratio", "Quick Assets / Current Liabilities"],
    ["Debt-paying", "X5", "Debt-to-Asset (R)", "Total Assets / Total Liabilities"],
    ["Profitability", "X6", "Asset Return Ratio", "(Profit + Fin. Exp.) / Total Assets"],
    ["Profitability", "X7", "ROA", "Net Profit / Total Assets"],
    ["Profitability", "X8", "EBIT Margin", "(Net Profit + Tax + Fin. Exp.) / Revenue"],
    ["Development", "X9", "Revenue Growth Rate", "(Current Rev. - Prior Rev.) / Prior Rev."],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.2), var_data,
          col_widths=[1.8, 1.0, 2.0, 4.6])

# SLIDE 8: EFA Procedure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Analysis Procedure", "Step-by-Step Methodology")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/efa_procedure.png',
                         Inches(0.5), Inches(1.4), Inches(9), Inches(5.8))

# SLIDE 9: KMO and Bartlett's Test
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Suitability Testing", "KMO and Bartlett's Sphericity Test (Table 2)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/kmo_bartlett.png',
                         Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.0))
add_content_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.8), [
    ("KMO = 0.706 (Meritorious) confirms sampling adequacy; Bartlett's p = 0.000 confirms variable correlations", 12, False, DARK_GRAY)
])

# SLIDE 10: KMO/Bartlett Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Suitability Test Results — Table", "Table 2: Wang & Song (2025)")
add_accent_bar(slide)

kmo_data = [
    ["Test", "Measure", "Value"],
    ["KMO", "Sampling Adequacy", "0.706"],
    ["Bartlett's Test", "Approx. Chi-Square", "573.599"],
    ["Bartlett's Test", "Degrees of Freedom", "36"],
    ["Bartlett's Test", "Significance", "0.000"],
]
add_table(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(3), kmo_data,
          col_widths=[2.5, 2.5, 2.0])

add_content_box(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(2), [
    ("Interpretation:", 16, True, DARK_BLUE),
    "KMO > 0.7 indicates the data is suitable for factor analysis (Kaiser, 1974)",
    "Bartlett's test significance (p < 0.05) confirms variables are correlated",
    "Both tests confirm factor analysis is appropriate for these 9 financial indicators",
], font_size=14, bullet=True)

# SLIDE 11: Scree Plot
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Scree Plot", "Eigenvalue Analysis for Factor Extraction (Table 3)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/scree_plot.png',
                         Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_content_box(slide, Inches(0.5), Inches(6.8), Inches(9), Inches(0.5), [
    ("3 factors retained with eigenvalues > 1: 3.262, 2.057, 1.300 (only retained factors reported in paper)", 12, False, DARK_GRAY)
])

# SLIDE 12: Eigenvalue Comparison
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Eigenvalues: Before vs. After Rotation", "Effect of Varimax Rotation (Table 3)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/eigenvalue_comparison.png',
                         Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_content_box(slide, Inches(0.5), Inches(6.8), Inches(9), Inches(0.5), [
    ("Rotation redistributes variance more evenly: F1 decreases (3.262 to 2.733), F2 increases (2.057 to 2.522)", 12, False, DARK_GRAY)
])

# SLIDE 13: Variance Explained
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Total Variance Explained", "Before and After Rotation (Table 3)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/variance_explained.png',
                         Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_content_box(slide, Inches(0.5), Inches(6.8), Inches(9), Inches(0.5), [
    ("3 factors explain 73.55% of total variance in financial performance indicators", 12, False, DARK_GRAY)
])

# SLIDE 14: Variance Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Variance Explained — Full Table", "Table 3: Wang & Song (2025)")
add_accent_bar(slide)

var_table = [
    ["Component", "Eigenvalue\n(Before)", "Variance %\n(Before)", "Cumul. %\n(Before)",
     "Eigenvalue\n(After)", "Variance %\n(After)", "Cumul. %\n(After)"],
    ["1", "3.262", "36.248%", "36.248%", "2.733", "30.369%", "30.369%"],
    ["2", "2.057", "22.852%", "59.101%", "2.522", "28.025%", "58.395%"],
    ["3", "1.300", "14.449%", "73.550%", "1.364", "15.155%", "73.550%"],
]
add_table(slide, Inches(0.3), Inches(1.8), Inches(9.4), Inches(2.5), var_table,
          col_widths=[1.2, 1.3, 1.3, 1.3, 1.3, 1.3, 1.3])

add_content_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(2.5), [
    ("Key Observations:", 16, True, DARK_BLUE),
    "Before rotation: F1 dominates with 36.25% variance; uneven distribution",
    "After Varimax rotation: Variance is redistributed more evenly across factors",
    "F1 drops from 36.25% to 30.37%; F2 rises from 22.85% to 28.03%",
    "Total cumulative variance remains unchanged at 73.55%",
    "This redistribution makes factor interpretation clearer and more meaningful",
], font_size=13, bullet=True)

# SLIDE 15: Factor Loading Matrix Heatmap
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Rotated Factor Loading Matrix", "Varimax Rotation Results (Table 4)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/factor_loading_matrix.png',
                         Inches(0.5), Inches(1.4), Inches(9), Inches(5.8))

# SLIDE 16: Factor Loadings Bar Chart
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Primary Factor Loadings", "Which Variables Define Each Factor (Table 4)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/factor_loadings_bar.png',
                         Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.0))
add_content_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.8), [
    ("F1: Solvency (0.929-0.956) | F2: Profitability (0.778-0.955) | F3: Oper. Dev. (0.732-0.744)", 12, True, DARK_GRAY)
])

# SLIDE 17: Factor Loading Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Rotated Component Matrix — Table", "Table 4: Wang & Song (2025)")
add_accent_bar(slide)

loading_data = [
    ["Variable", "F1: Solvency", "F2: Profitability", "F3: Oper. Dev."],
    ["X1: Inventory Turnover", "0.001", "0.264", "-0.479"],
    ["X2: A/R Turnover", "-0.067", "0.070", "0.732"],
    ["X3: Current Ratio", "0.956", "0.086", "-0.020"],
    ["X4: Quick Ratio", "0.929", "0.128", "-0.154"],
    ["X5: Debt-to-Asset (R)", "0.945", "0.116", "-0.034"],
    ["X6: Asset Return Ratio", "0.010", "0.955", "-0.065"],
    ["X7: ROA", "0.144", "0.940", "-0.071"],
    ["X8: EBIT Margin", "0.181", "0.778", "0.103"],
    ["X9: Revenue Growth", "-0.082", "0.097", "0.744"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), loading_data,
          col_widths=[3.0, 2.0, 2.0, 2.0])

# SLIDE 18: Factor Interpretation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Interpretation", "Understanding the Three Extracted Factors")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Factor 1: Solvency Factor (30.37% variance)", 16, True, RGBColor(46, 134, 193)),
    "Loaded by: Current Ratio (0.956), Debt-to-Asset Reciprocal (0.945), Quick Ratio (0.929)",
    "Measures the company's ability to repay debts and meet financial obligations",
    "Reflects short-term liquidity and overall financial stability",
    "",
    ("Factor 2: Profitability Factor (28.03% variance)", 16, True, RGBColor(230, 126, 34)),
    "Loaded by: Asset Return Ratio (0.955), ROA (0.940), EBIT Margin (0.778)",
    "Measures the company's capacity to generate profits from its assets and operations",
    "Reflects long-term earning power and market competitiveness",
    "",
    ("Factor 3: Operational Development Capability (15.16% variance)", 16, True, RGBColor(39, 174, 96)),
    "Loaded by: Revenue Growth Rate (0.744), A/R Turnover (0.732)",
    "Measures the company's growth trajectory and operational efficiency",
    "Reflects expansion momentum and asset utilization effectiveness",
], font_size=13, bullet=False)

# SLIDE 19: Score Coefficients
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Component Score Coefficients", "Building the Scoring Equations (Table 5)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/score_coefficients.png',
                         Inches(1), Inches(1.4), Inches(8), Inches(5.5))

# SLIDE 20: Scoring Equations
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Scoring Equations", "From Table 5 Coefficients")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Factor Score Equations (from Table 5):", 16, True, DARK_BLUE),
    "",
    ("F1 = -0.062X1 + 0.026X2 + 0.367X3 + 0.342X4 + 0.359X5", 13, False, RGBColor(46, 134, 193)),
    ("      - 0.086X6 - 0.033X7 + 0.010X8 + 0.019X9", 13, False, RGBColor(46, 134, 193)),
    "",
    ("F2 = 0.108X1 + 0.038X2 - 0.049X3 - 0.029X4 - 0.036X5", 13, False, RGBColor(230, 126, 34)),
    ("      + 0.397X6 + 0.379X7 + 0.309X8 + 0.051X9", 13, False, RGBColor(230, 126, 34)),
    "",
    ("F3 = -0.358X1 + 0.544X2 + 0.062X3 - 0.040X4 + 0.051X5", 13, False, RGBColor(39, 174, 96)),
    ("      - 0.044X6 - 0.038X7 + 0.096X8 + 0.553X9", 13, False, RGBColor(39, 174, 96)),
    "",
    ("Comprehensive Score:", 16, True, DARK_BLUE),
    ("F = (30.369/73.550) x F1 + (28.025/73.550) x F2 + (15.155/73.550) x F3", 13, True, DARK_GRAY),
    ("F = 0.413 x F1 + 0.381 x F2 + 0.206 x F3", 14, True, MEDIUM_BLUE),
], font_size=13, bullet=False)

# SLIDE 21: Factor Score Ranges
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Score Ranges", "Performance Distribution Across 100 Companies (Table 6)")
add_accent_bar(slide)
slide.shapes.add_picture('figures3/factor_score_ranges.png',
                         Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
add_content_box(slide, Inches(0.5), Inches(6.8), Inches(9), Inches(0.5), [
    ("Wide score ranges indicate significant differentiation in financial performance among companies", 12, False, DARK_GRAY)
])

# SLIDE 22: Key Findings
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Key Findings", "What the Analysis Revealed")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("1. Three Core Dimensions Drive Financial Performance", 15, True, DARK_BLUE),
    "Solvency, Profitability, and Operational Development Capability explain 73.55% of variance",
    "",
    ("2. Significant Differentiation Among Companies", 15, True, DARK_BLUE),
    "Solvency scores range from -24.85 to 46.09 — some companies face serious debt risk",
    "Profitability scores range from -0.32 to 66.24 — large gap between top and bottom performers",
    "Operational development scores range from -156.91 to 947.95 — extreme variation in growth",
    "",
    ("3. Debt Risk is a Major Concern", 15, True, DARK_BLUE),
    "Some companies have very low solvency scores, indicating capital chain vulnerabilities",
    "This aligns with the post-Evergrande industry reality",
    "",
    ("4. Factor Analysis Provides Superior Evaluation", 15, True, DARK_BLUE),
    "Reduces 9 indicators into 3 interpretable factors for clearer assessment",
    "Weighted comprehensive scoring enables fair cross-company comparison",
], font_size=13, bullet=False)

# SLIDE 23: Practical Implications
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Practical Implications", "Recommendations from the Study")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("For Company Managers:", 16, True, DARK_BLUE),
    "Strengthen profitability through better cost control and revenue diversification",
    "Improve operational development capability via faster asset turnover",
    "Pay close attention to debt risks and maintain healthy capital chains",
    "Use multi-dimensional evaluation (not just single metrics) for strategic decisions",
    "",
    ("For Investors:", 16, True, DARK_BLUE),
    "Look beyond simple profit metrics — solvency and growth are equally important",
    "Companies with balanced scores across all 3 factors are more financially stable",
    "Low solvency scores signal potential bankruptcy or capital chain risks",
    "",
    ("For Policy Makers:", 16, True, DARK_BLUE),
    "The significant differentiation highlights need for industry-level risk monitoring",
    "Factor analysis can serve as an early warning system for sector-wide financial stress",
], font_size=13, bullet=False)

# SLIDE 24: Why Factor Analysis Works for Finance
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Why Factor Analysis Works for Finance", "Advantages of the Method")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Advantages:", 18, True, DARK_BLUE),
    "Handles multicollinearity — financial ratios are often highly correlated",
    "Reduces dimensionality — 9 variables condensed into 3 interpretable factors",
    "Objective weighting — factor weights come from data, not subjective judgment",
    "Comprehensive scoring — enables fair ranking of companies on multiple dimensions",
    "",
    ("Limitations:", 18, True, RGBColor(192, 57, 43)),
    "Requires sufficient sample size (this study uses 100 companies)",
    "Results depend on variable selection (only 9 indicators used here)",
    "Historical data — reflects past performance, not predictive of future",
    "Cross-sectional analysis — captures one year (2023) only",
    "",
    ("Broader Applications in Finance:", 18, True, MEDIUM_BLUE),
    "Portfolio risk assessment and diversification analysis",
    "Credit scoring and loan default prediction",
    "Stock return factor models (e.g., Fama-French factors)",
    "Macroeconomic forecasting using latent factor models",
], font_size=13, bullet=True)

# SLIDE 25: Summary of Real Data Used
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Summary: All Real Published Data", "Data Source Verification")
add_accent_bar(slide)

summary_data = [
    ["Data Element", "Value", "Source"],
    ["Sample Size", "100 companies", "Wang & Song (2025)"],
    ["KMO", "0.706", "Table 2"],
    ["Bartlett's Chi-Square", "573.599 (df=36, p=0.000)", "Table 2"],
    ["Eigenvalues (before rot.)", "3.262, 2.057, 1.300", "Table 3"],
    ["Eigenvalues (after rot.)", "2.733, 2.522, 1.364", "Table 3"],
    ["Total Variance Explained", "73.550%", "Table 3"],
    ["Factor Loadings", "9 variables x 3 factors", "Table 4"],
    ["Score Coefficients", "9 variables x 3 factors", "Table 5"],
    ["Score Ranges", "Min/Max per factor", "Table 6"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5), summary_data,
          col_widths=[3.2, 3.5, 2.7])

# SLIDE 26: References
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "References")
add_accent_bar(slide)
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), [
    ("Primary Study:", 16, True, DARK_BLUE),
    "Wang, Y. & Song, X. (2025). \"Research on the Financial Performance Evaluation of",
    "   Listed Companies in the Real Estate Industry Based on Factor Analysis Method.\"",
    "   Academic Journal of Business & Management, 7(5), 78-84.",
    "   DOI: 10.25236/AJBM.2025.070510",
    "",
    ("Supporting References:", 16, True, DARK_BLUE),
    "Kaiser, H.F. (1974). \"An Index of Factorial Simplicity.\" Psychometrika, 39, 31-36.",
    "",
    "Hair, J.F., Black, W.C., Babin, B.J. & Anderson, R.E. (2019). Multivariate Data",
    "   Analysis, 8th Ed. Cengage Learning.",
    "",
    "Shrestha, N. (2021). \"Factor Analysis as a Tool for Survey Analysis.\"",
    "   American Journal of Applied Mathematics and Statistics, 9(1), 4-11.",
    "   DOI: 10.12691/ajams-9-1-2",
], font_size=12, bullet=False)

# SLIDE 27: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(3))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 25, 75)
shape.line.fill.background()

add_content_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(2.5), [
    ("Thank You!", 40, True, WHITE),
    ("", 10, False, WHITE),
    ("Factor Analysis in Financial Performance Evaluation", 18, False, ACCENT_GOLD),
    ("", 6, False, WHITE),
    ("Lana Jalal Gidan", 16, True, WHITE),
    ("SSIE-605 | Binghamton University | Professor Susan Lu", 13, False, RGBColor(180, 198, 231)),
])
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.7), Inches(9), Inches(0.04))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GOLD
bar.line.fill.background()

output_file = "Factor_Analysis_Financial_Performance_Presentation.pptx"
prs.save(output_file)
print(f"\nPresentation saved as: {output_file}")
print(f"Total slides: {len(prs.slides)}")
print("All data sourced from: Wang & Song (2025), DOI: 10.25236/AJBM.2025.070510")
