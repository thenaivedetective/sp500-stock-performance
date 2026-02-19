import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BLUE = RGBColor(0, 32, 96)
MEDIUM_BLUE = RGBColor(0, 70, 140)
LIGHT_BLUE = RGBColor(0, 112, 192)
ACCENT_GOLD = RGBColor(255, 192, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(240, 240, 240)

os.makedirs("figures5", exist_ok=True)

def add_slide_background(slide, color=DARK_BLUE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=None):
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = ACCENT_GOLD
        p2.alignment = PP_ALIGN.LEFT

def add_content_box(slide, left, top, width, height, items, font_size=14, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, size, bold, color = item
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
        elif item == "":
            p.text = ""
            p.font.size = Pt(6)
        else:
            prefix = "\u2022 " if bullet else ""
            p.text = f"{prefix}{item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY

def add_table(slide, left, top, width, height, data, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            if i == 0:
                from pptx.oxml.ns import qn
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '002060'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

def add_accent_bar(slide, top=Inches(1.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()

# ============================================================
# FIGURE 1: Demographic Distribution (Table 2)
# ============================================================
fig, axes = plt.subplots(2, 2, figsize=(14, 10))

gender_labels = ['Female', 'Male']
gender_pcts = [43.70, 56.30]
colors_g = ['#E74C3C', '#3498DB']
axes[0,0].pie(gender_pcts, labels=gender_labels, autopct='%1.1f%%', colors=colors_g,
              startangle=90, textprops={'fontsize': 12, 'fontweight': 'bold'})
axes[0,0].set_title('Gender Distribution', fontsize=13, fontweight='bold')

desig_labels = ['Junior\nManagers', 'Middle-level\nManagers', 'Top\nManagers']
desig_pcts = [42.67, 34.70, 25.19]
colors_d = ['#3498DB', '#2E86C1', '#1F618D']
bars = axes[0,1].barh(desig_labels, desig_pcts, color=colors_d, edgecolor='black', linewidth=0.5)
for i, v in enumerate(desig_pcts):
    axes[0,1].text(v + 0.8, i, f'{v}%', va='center', fontsize=11, fontweight='bold')
axes[0,1].set_xlabel('Percentage (%)', fontsize=11, fontweight='bold')
axes[0,1].set_title('Designation', fontsize=13, fontweight='bold')
axes[0,1].set_xlim(0, 55)

age_labels = ['25-35', '35-50', 'Above 50']
age_pcts = [35.73, 26.73, 37.53]
colors_a = ['#2ECC71', '#27AE60', '#1E8449']
bars = axes[1,0].bar(age_labels, age_pcts, color=colors_a, edgecolor='black', linewidth=0.5, width=0.5)
for bar, v in zip(bars, age_pcts):
    axes[1,0].text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.8,
                   f'{v}%', ha='center', fontsize=11, fontweight='bold')
axes[1,0].set_ylabel('Percentage (%)', fontsize=11, fontweight='bold')
axes[1,0].set_title('Age Group', fontsize=13, fontweight='bold')
axes[1,0].set_ylim(0, 48)

exp_labels = ['< 10 years', '10-20 years', '> 20 years']
exp_pcts = [46.78, 28.02, 25.19]
colors_e = ['#E67E22', '#D35400', '#A04000']
bars = axes[1,1].barh(exp_labels, exp_pcts, color=colors_e, edgecolor='black', linewidth=0.5)
for i, v in enumerate(exp_pcts):
    axes[1,1].text(v + 0.8, i, f'{v}%', va='center', fontsize=11, fontweight='bold')
axes[1,1].set_xlabel('Percentage (%)', fontsize=11, fontweight='bold')
axes[1,1].set_title('Work Experience', fontsize=13, fontweight='bold')
axes[1,1].set_xlim(0, 58)

plt.suptitle('Demographic Profile of Respondents (n=389)\n(Gupta et al., 2025 \u2014 Table 2)',
             fontsize=14, fontweight='bold', y=1.02)
plt.tight_layout()
plt.savefig('figures5/demographics.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 2: Scree Plot / Eigenvalues (Table 6 - PCA)
# ============================================================
components = list(range(1, 16))
eigenvalues = [11.389, 3.348, 3.153, 2.815, 2.472, 2.181, 1.965, 1.067,
               0.770, 0.765, 0.671, 0.638, 0.620, 0.573, 0.564]

fig, ax = plt.subplots(figsize=(12, 6))
colors_eig = ['#002060' if e >= 1.0 else '#B0B0B0' for e in eigenvalues]
bars = ax.bar(components, eigenvalues, color=colors_eig, edgecolor='black', linewidth=0.8, width=0.6)
ax.plot(components, eigenvalues, 'ro-', linewidth=2, markersize=8, color='#E74C3C')
ax.axhline(y=1.0, color='red', linestyle='--', linewidth=2, label='Kaiser Criterion (Eigenvalue = 1.0)')
for i, (comp, eig) in enumerate(zip(components[:8], eigenvalues[:8])):
    ax.text(comp, eig + 0.25, f'{eig:.3f}', ha='center', fontsize=9, fontweight='bold')
ax.set_xlabel('Component Number', fontsize=13, fontweight='bold')
ax.set_ylabel('Eigenvalue', fontsize=13, fontweight='bold')
ax.set_title('Scree Plot: Eigenvalues by Component (PCA)\n(Gupta et al., 2025 \u2014 Table 6, 8 components with eigenvalue > 1)',
             fontsize=14, fontweight='bold')
ax.set_xticks(components)
ax.legend(fontsize=11, loc='upper right')
ax.grid(True, alpha=0.3, axis='y')
ax.annotate('8 components\nretained (66.02%)', xy=(8, 1.067), xytext=(11, 4),
            fontsize=12, fontweight='bold', color='#002060',
            arrowprops=dict(arrowstyle='->', color='#002060', lw=2))
plt.tight_layout()
plt.savefig('figures5/scree_plot.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 3: Variance Explained (Table 6)
# ============================================================
factor_names_var = ['Comp 1', 'Comp 2', 'Comp 3', 'Comp 4', 'Comp 5', 'Comp 6', 'Comp 7', 'Comp 8']
var_pcts = [26.485, 7.785, 7.332, 6.546, 5.750, 5.071, 4.570, 2.481]
cum_pcts = [26.485, 34.271, 41.603, 48.149, 53.899, 58.970, 63.539, 66.020]

fig, ax1 = plt.subplots(figsize=(12, 6))
colors_var = ['#002060', '#0047AB', '#0070C0', '#2E86C1', '#3498DB', '#5DADE2', '#85C1E9', '#AED6F1']
bars = ax1.bar(factor_names_var, var_pcts, color=colors_var, edgecolor='black', linewidth=0.8, width=0.6)
for bar, v in zip(bars, var_pcts):
    ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3,
             f'{v:.1f}%', ha='center', fontsize=10, fontweight='bold')
ax1.set_ylabel('Individual Variance (%)', fontsize=12, fontweight='bold', color='#002060')
ax1.set_ylim(0, 35)

ax2 = ax1.twinx()
ax2.plot(factor_names_var, cum_pcts, 'ro-', linewidth=2.5, markersize=10, color='#E74C3C')
for i, (name, cum) in enumerate(zip(factor_names_var, cum_pcts)):
    ax2.text(i, cum + 1.5, f'{cum:.1f}%', ha='center', fontsize=9, fontweight='bold', color='#E74C3C')
ax2.set_ylabel('Cumulative Variance (%)', fontsize=12, fontweight='bold', color='#E74C3C')
ax2.set_ylim(0, 80)

ax1.set_title('Total Variance Explained by PCA Components\n(Gupta et al., 2025 \u2014 Table 6, Total = 66.02%)',
              fontsize=14, fontweight='bold')
ax1.grid(True, alpha=0.3, axis='y')
plt.tight_layout()
plt.savefig('figures5/variance_explained.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 4: Cronbach's Alpha by Construct (Table 3)
# ============================================================
constructs = ['Green\nProcurement', 'Green\nLogistics', 'Green Product\n& Design',
              'Regulatory\nFramework', 'Top Mgmt.\nPerformance', 'Low Carbon\nPerformance',
              'Sustainable\nSociety', 'Sustainable\nManufacturing']
alphas = [0.873, 0.836, 0.896, 0.871, 0.889, 0.882, 0.894, 0.729]
colors_alpha = ['#002060', '#0047AB', '#0070C0', '#2E86C1', '#E67E22', '#D35400', '#27AE60', '#2ECC71']

fig, ax = plt.subplots(figsize=(12, 6))
bars = ax.bar(constructs, alphas, color=colors_alpha, edgecolor='black', linewidth=0.8, width=0.6)
ax.axhline(y=0.70, color='red', linestyle='--', linewidth=1.5, label='Acceptable Threshold (0.70)')
ax.axhline(y=0.80, color='green', linestyle='--', linewidth=1.5, label='Good Threshold (0.80)')
for bar, alpha in zip(bars, alphas):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.008,
            f'{alpha:.3f}', ha='center', va='bottom', fontsize=10, fontweight='bold')
ax.set_ylabel("Cronbach's Alpha", fontsize=13, fontweight='bold')
ax.set_title("Reliability Analysis: Cronbach's Alpha by Construct\n(Gupta et al., 2025 \u2014 Table 3, n=389)",
             fontsize=14, fontweight='bold')
ax.set_ylim(0.5, 1.0)
ax.legend(fontsize=10, loc='lower right')
ax.grid(True, alpha=0.3, axis='y')
plt.tight_layout()
plt.savefig('figures5/cronbach_alpha.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 5: Factor Loading Ranges (Table 3)
# ============================================================
construct_names = ['Green\nProcurement', 'Green\nLogistics', 'Green Product\n& Design',
                   'Regulatory\nFramework', 'Top Mgmt.\nPerformance', 'Low Carbon\nPerformance',
                   'Sustainable\nSociety', 'Sustainable\nMfg.']
loading_min = [0.681, 0.604, 0.780, 0.673, 0.736, 0.622, 0.738, 0.682]
loading_max = [0.767, 0.781, 0.819, 0.771, 0.844, 0.786, 0.792, 0.841]
loading_mid = [(mn + mx) / 2 for mn, mx in zip(loading_min, loading_max)]
loading_err = [(mx - mn) / 2 for mn, mx in zip(loading_min, loading_max)]
n_items = [6, 5, 5, 6, 5, 7, 5, 2]

fig, ax = plt.subplots(figsize=(12, 7))
colors_load = ['#002060', '#0047AB', '#0070C0', '#2E86C1', '#E67E22', '#D35400', '#27AE60', '#2ECC71']
y_pos = range(len(construct_names))
bars = ax.barh(y_pos, loading_mid, xerr=loading_err, color=colors_load,
               edgecolor='black', linewidth=0.8, height=0.5, capsize=8, ecolor='black')
for i in range(len(construct_names)):
    ax.text(loading_max[i] + 0.01, i,
            f'{loading_min[i]:.3f}\u2013{loading_max[i]:.3f} ({n_items[i]} items)',
            va='center', fontsize=10, fontweight='bold')
ax.axvline(x=0.50, color='red', linestyle='--', linewidth=1.5, label='Minimum Threshold (0.50)')
ax.set_yticks(y_pos)
ax.set_yticklabels(construct_names, fontsize=10)
ax.set_xlabel('Factor Loading', fontsize=13, fontweight='bold')
ax.set_title('Standardized Factor Loading Ranges by Construct\n(Gupta et al., 2025 \u2014 Table 3, All loadings > 0.50)',
             fontsize=14, fontweight='bold')
ax.set_xlim(0.4, 0.95)
ax.legend(fontsize=10, loc='lower right')
ax.grid(True, alpha=0.3, axis='x')
plt.tight_layout()
plt.savefig('figures5/factor_loadings.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 6: Composite Reliability & AVE (Table 3)
# ============================================================
construct_short = ['GP', 'GL', 'GD', 'RF', 'TMC', 'CP', 'SS', 'SM']
cr_values = [0.874, 0.832, 0.898, 0.871, 0.887, 0.874, 0.892, 0.730]
ave_values = [0.531, 0.502, 0.633, 0.531, 0.613, 0.505, 0.582, 0.576]

fig, ax = plt.subplots(figsize=(12, 6))
x = np.arange(len(construct_short))
width = 0.35
bars1 = ax.bar(x - width/2, cr_values, width, label='Composite Reliability (CR)',
               color='#002060', edgecolor='black', linewidth=0.8)
bars2 = ax.bar(x + width/2, ave_values, width, label='Average Variance Extracted (AVE)',
               color='#E67E22', edgecolor='black', linewidth=0.8)
ax.axhline(y=0.70, color='blue', linestyle='--', linewidth=1, alpha=0.5, label='CR Threshold (0.70)')
ax.axhline(y=0.50, color='red', linestyle='--', linewidth=1, alpha=0.5, label='AVE Threshold (0.50)')
for bar, val in zip(bars1, cr_values):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.008,
            f'{val:.3f}', ha='center', fontsize=8, fontweight='bold')
for bar, val in zip(bars2, ave_values):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.008,
            f'{val:.3f}', ha='center', fontsize=8, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(construct_short, fontsize=11, fontweight='bold')
ax.set_ylabel('Value', fontsize=13, fontweight='bold')
ax.set_title('Convergent Validity: Composite Reliability & AVE\n(Gupta et al., 2025 \u2014 Table 3, All CR > 0.70, All AVE > 0.50)',
             fontsize=14, fontweight='bold')
ax.set_ylim(0, 1.05)
ax.legend(fontsize=9, loc='upper right')
ax.grid(True, alpha=0.3, axis='y')
plt.tight_layout()
plt.savefig('figures5/cr_ave.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 7: SEM Path Coefficients (Table 7)
# ============================================================
hypotheses = ['H1: GP \u2192 TMP', 'H2: GL \u2192 TMP', 'H3: GD \u2192 TMP',
              'H4: RF \u2192 TMP', 'H5: TMP \u2192 CP', 'H6: CP \u2192 SS', 'H7: CP \u2192 SM']
beta_values = [0.130, 0.200, 0.274, 0.100, 0.388, 0.359, 0.295]
p_values_str = ['p=0.000', 'p=0.000', 'p=0.000', 'p=0.007', 'p=0.000', 'p=0.000', 'p=0.000']
colors_sem = ['#002060', '#0047AB', '#0070C0', '#2E86C1', '#E74C3C', '#27AE60', '#2ECC71']

fig, ax = plt.subplots(figsize=(12, 7))
bars = ax.barh(hypotheses, beta_values, color=colors_sem, edgecolor='black', linewidth=0.8, height=0.5)
for bar, beta, pv in zip(bars, beta_values, p_values_str):
    ax.text(bar.get_width() + 0.008, bar.get_y() + bar.get_height()/2,
            f'\u03B2 = {beta:.3f} ({pv})', va='center', fontsize=11, fontweight='bold')
ax.set_xlabel('Standardized Path Coefficient (\u03B2)', fontsize=13, fontweight='bold')
ax.set_title('SEM Results: Hypothesis Testing\n(Gupta et al., 2025 \u2014 Table 7, All hypotheses supported)',
             fontsize=14, fontweight='bold')
ax.set_xlim(0, 0.55)
ax.grid(True, alpha=0.3, axis='x')
ax.text(0.5, -0.1, 'All path coefficients significant at p < 0.05',
        transform=ax.transAxes, ha='center', fontsize=11, fontstyle='italic', color='gray')
plt.tight_layout()
plt.savefig('figures5/sem_paths.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 8: R-Squared Values (Table 7)
# ============================================================
endogenous = ['Top Management\nPerformance', 'Low Carbon\nPerformance',
              'Sustainable\nSociety', 'Sustainable\nManufacturing']
r_squared = [20.7, 19.5, 16.1, 14.7]
colors_r2 = ['#E67E22', '#E74C3C', '#27AE60', '#2ECC71']

fig, ax = plt.subplots(figsize=(10, 6))
bars = ax.bar(endogenous, r_squared, color=colors_r2, edgecolor='black', linewidth=0.8, width=0.5)
for bar, val in zip(bars, r_squared):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
            f'{val}%', ha='center', fontsize=14, fontweight='bold')
ax.set_ylabel('R\u00B2 (%)', fontsize=13, fontweight='bold')
ax.set_title('Variance Explained (R\u00B2) in Endogenous Constructs\n(Gupta et al., 2025 \u2014 Table 7)',
             fontsize=14, fontweight='bold')
ax.set_ylim(0, 28)
ax.grid(True, alpha=0.3, axis='y')
plt.tight_layout()
plt.savefig('figures5/r_squared.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 9: Effect Size Comparison
# ============================================================
effect_paths = ['TMP \u2192 CP', 'CP \u2192 SM', 'CP \u2192 SS']
effect_sizes = [0.215, 0.139, 0.176]
effect_labels_cat = ['Medium', 'Small', 'Medium']
colors_eff = ['#E67E22', '#3498DB', '#27AE60']

fig, ax = plt.subplots(figsize=(10, 5))
bars = ax.bar(effect_paths, effect_sizes, color=colors_eff, edgecolor='black', linewidth=0.8, width=0.4)
ax.axhline(y=0.02, color='gray', linestyle=':', linewidth=1.5, label='Small (f\u00B2 = 0.02)')
ax.axhline(y=0.15, color='orange', linestyle=':', linewidth=1.5, label='Medium (f\u00B2 = 0.15)')
ax.axhline(y=0.35, color='red', linestyle=':', linewidth=1.5, label='Large (f\u00B2 = 0.35)')
for bar, val, cat in zip(bars, effect_sizes, effect_labels_cat):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.008,
            f'f\u00B2 = {val:.3f}\n({cat})', ha='center', fontsize=12, fontweight='bold')
ax.set_ylabel('Effect Size (f\u00B2)', fontsize=13, fontweight='bold')
ax.set_title("Cohen's f\u00B2 Effect Sizes for Structural Paths\n(Gupta et al., 2025)",
             fontsize=14, fontweight='bold')
ax.set_ylim(0, 0.42)
ax.legend(fontsize=10, loc='upper right')
ax.grid(True, alpha=0.3, axis='y')
plt.tight_layout()
plt.savefig('figures5/effect_sizes.png', dpi=200, bbox_inches='tight')
plt.close()

# ============================================================
# FIGURE 10: Radar Summary of Construct Properties
# ============================================================
radar_constructs = ['GP', 'GL', 'GD', 'RF', 'TMC', 'CP', 'SS', 'SM']
radar_alpha = [0.873, 0.836, 0.896, 0.871, 0.889, 0.882, 0.894, 0.729]
radar_cr = [0.874, 0.832, 0.898, 0.871, 0.887, 0.874, 0.892, 0.730]
radar_ave = [0.531, 0.502, 0.633, 0.531, 0.613, 0.505, 0.582, 0.576]

angles = np.linspace(0, 2 * np.pi, len(radar_constructs), endpoint=False).tolist()
angles += angles[:1]

fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))
values_a = radar_alpha + radar_alpha[:1]
values_cr = radar_cr + radar_cr[:1]
values_ave = radar_ave + radar_ave[:1]

ax.plot(angles, values_a, 'o-', linewidth=2, color='#002060', label="Cronbach's \u03B1")
ax.fill(angles, values_a, alpha=0.1, color='#002060')
ax.plot(angles, values_cr, 's-', linewidth=2, color='#E67E22', label='Composite Reliability')
ax.fill(angles, values_cr, alpha=0.1, color='#E67E22')
ax.plot(angles, values_ave, '^-', linewidth=2, color='#27AE60', label='AVE')
ax.fill(angles, values_ave, alpha=0.1, color='#27AE60')

ax.set_xticks(angles[:-1])
ax.set_xticklabels(radar_constructs, fontsize=11, fontweight='bold')
ax.set_ylim(0, 1.0)
ax.set_title('Radar Summary: Psychometric Properties\n(Gupta et al., 2025 \u2014 Table 3)',
             fontsize=14, fontweight='bold', pad=20)
ax.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1), fontsize=10)
plt.tight_layout()
plt.savefig('figures5/radar_summary.png', dpi=200, bbox_inches='tight')
plt.close()

print("All 10 figures created successfully in figures5/")

# ============================================================
# BUILD POWERPOINT PRESENTATION (27 slides)
# ============================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SLIDE 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, DARK_BLUE)
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Factor Analysis in Supply Chain Management"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Green Supply Chain Management Practices & Sustainability"
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT_GOLD
p2.alignment = PP_ALIGN.CENTER

add_accent_bar(slide, top=Inches(2.8))

txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(3.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
items = [
    ("Based on: Gupta, P., Sharma, Y., Chauhan, A., Parewa, B., Rai, P. & Naik, N. (2025)", 14, False, WHITE),
    ("Scientific Reports, 15, 14909", 14, True, ACCENT_GOLD),
    ("DOI: 10.1038/s41598-025-95940-9", 12, False, WHITE),
    ("", 6, False, WHITE),
    ("Presented by: Lana Jalal Gidan", 16, True, WHITE),
    ("SSIE-605: Applied Multivariate Data Analysis", 14, False, WHITE),
    ("Binghamton University \u2014 Professor Susan Lu", 14, False, WHITE),
]
first = True
for text, size, bold, color in items:
    if first:
        p = tf2.paragraphs[0]
        first = False
    else:
        p = tf2.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

# SLIDE 2: Presentation Outline
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Presentation Outline")
add_accent_bar(slide)
items = [
    "1. Introduction to Green Supply Chain Management (GSCM)",
    "2. Research Context & Study Overview",
    "3. Methodology: Sample & Survey Design",
    "4. Demographic Profile of Respondents",
    "5. PCA Eigenvalue Decomposition & Scree Plot",
    "6. Total Variance Explained by Components",
    "7. Measurement Model: Factor Loadings",
    "8. Reliability Analysis: Cronbach's Alpha",
    "9. Convergent Validity: CR & AVE",
    "10. SEM: Hypothesis Testing & Path Coefficients",
    "11. Variance Explained (R\u00B2) & Effect Sizes",
    "12. Key Findings & Practical Implications",
    "13. Conclusions & Future Research Directions",
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=14, bullet=True)

# SLIDE 3: Introduction to GSCM
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Introduction to Green Supply Chain Management",
              "Integrating Environmental Sustainability into Supply Chains")
add_accent_bar(slide)
items = [
    ("What is Green Supply Chain Management (GSCM)?", 16, True, DARK_BLUE),
    "",
    "GSCM integrates environmental thinking into supply chain management",
    "Encompasses procurement, manufacturing, distribution, and end-of-life management",
    "Aims to minimize ecological footprint while maintaining competitive advantage",
    "",
    ("Why Factor Analysis in GSCM Research?", 16, True, DARK_BLUE),
    "",
    "Identifies underlying dimensions that drive green practices adoption",
    "Reduces complex multi-item questionnaires to interpretable factors",
    "Validates measurement instruments for organizational sustainability",
    "Enables Structural Equation Modeling (SEM) for hypothesis testing",
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=13, bullet=False)

# SLIDE 4: Research Context
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Research Context & Study Overview",
              "Gupta et al. (2025) \u2014 Scientific Reports")
add_accent_bar(slide)
items = [
    ("Study Focus:", 15, True, DARK_BLUE),
    "Examines the impact of GSCM practices on top management performance,",
    "low-carbon performance, sustainable manufacturing, and sustainable society",
    "",
    ("Research Setting:", 15, True, DARK_BLUE),
    "Indian bag-manufacturing firms implementing GSCM practices",
    "Data collection: March\u2013September 2020 via structured questionnaire",
    "480 responses received; 389 valid after screening (81% response rate)",
    "",
    ("Analytical Approach:", 15, True, DARK_BLUE),
    "Principal Component Analysis (PCA) for dimensionality assessment",
    "Confirmatory Factor Analysis (CFA) via SmartPLS for measurement model",
    "Structural Equation Modeling (SEM) for hypothesis testing",
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=13, bullet=False)

# SLIDE 5: Methodology
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Methodology: Sample & Survey Design",
              "Rigorous Data Collection from Manufacturing Sector")
add_accent_bar(slide)
data = [
    ["Parameter", "Details"],
    ["Sample Size", "n = 389 valid responses (from 480 collected)"],
    ["Population", "Managers in Indian bag-manufacturing firms"],
    ["Respondent Levels", "Junior (42.67%), Middle (34.70%), Top (25.19%)"],
    ["Scale", "5-point Likert (Strongly Disagree to Strongly Agree)"],
    ["Total Items", "41 items across 8 constructs"],
    ["Sampling Method", "Convenience + Snowball sampling"],
    ["Data Platform", "Google Forms (mandatory fields, no missing data)"],
    ["Analysis Software", "SmartPLS (SEM analysis)"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5), data,
          col_widths=[2.5, 6.5])

# SLIDE 6: Eight Constructs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Eight GSCM Constructs Measured",
              "41 Items Across 4 GSCM Practices + 4 Outcome Variables")
add_accent_bar(slide)
data = [
    ["Construct", "Abbr.", "Items", "Role"],
    ["Green Procurement", "GP", "6", "GSCM Practice"],
    ["Green Logistics", "GL", "5", "GSCM Practice"],
    ["Green Product & Product Design", "GD", "5", "GSCM Practice"],
    ["Regulatory Framework", "RF", "6", "GSCM Practice"],
    ["Top Management Performance", "TMC", "5", "Mediator"],
    ["Low Carbon Performance", "CP", "7", "Mediator"],
    ["Sustainable Society", "SS", "5", "Outcome"],
    ["Sustainable Manufacturing", "SM", "2", "Outcome"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5), data,
          col_widths=[3.5, 1.0, 1.0, 3.5])

# SLIDE 7: Demographics Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Demographic Profile of Respondents",
              "Table 2: n=389 Manufacturing Managers from India")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/demographics.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 8: Scree Plot
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "PCA Eigenvalue Decomposition",
              "Table 6: 8 Components with Eigenvalue > 1.0")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/scree_plot.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 9: Eigenvalue Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Eigenvalue Details: Kaiser Criterion",
              "Table 6: Initial Eigenvalues from PCA")
add_accent_bar(slide)
data = [
    ["Component", "Eigenvalue", "% Variance", "Cumulative %"],
    ["1", "11.389", "26.485%", "26.485%"],
    ["2", "3.348", "7.785%", "34.271%"],
    ["3", "3.153", "7.332%", "41.603%"],
    ["4", "2.815", "6.546%", "48.149%"],
    ["5", "2.472", "5.750%", "53.899%"],
    ["6", "2.181", "5.071%", "58.970%"],
    ["7", "1.965", "4.570%", "63.539%"],
    ["8", "1.067", "2.481%", "66.020%"],
]
add_table(slide, Inches(1.5), Inches(1.5), Inches(7), Inches(4.5), data,
          col_widths=[1.5, 1.5, 2.0, 2.0])
add_content_box(slide, Inches(1.5), Inches(6.3), Inches(7), Inches(0.8), [
    ("8 components retained explaining 66.02% of total variance", 13, True, DARK_BLUE)
])

# SLIDE 10: Variance Explained Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Total Variance Explained",
              "Table 6: Individual & Cumulative Variance by Component")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/variance_explained.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 11: Factor Loadings Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Measurement Model: Factor Loadings",
              "Table 3: Standardized Factor Loadings (CFA)")
add_accent_bar(slide)
data = [
    ["Construct", "Items", "Loading Range", "Highest", "Lowest"],
    ["Green Procurement (GP)", "6", "0.681\u20130.767", "GP2: 0.767", "GP1: 0.681"],
    ["Green Logistics (GL)", "5", "0.604\u20130.781", "GL3: 0.781", "GL1: 0.604"],
    ["Green Product & Design (GD)", "5", "0.780\u20130.819", "GD1: 0.819", "GD3: 0.780"],
    ["Regulatory Framework (RF)", "6", "0.673\u20130.771", "RF2: 0.771", "RF6: 0.673"],
    ["Top Mgmt. Performance (TMC)", "5", "0.736\u20130.844", "TMC4: 0.844", "TMC3: 0.736"],
    ["Low Carbon Perf. (CP)", "7", "0.622\u20130.786", "CP5: 0.786", "CP7: 0.622"],
    ["Sustainable Society (SS)", "5", "0.738\u20130.792", "SS4: 0.792", "SS1: 0.738"],
    ["Sustainable Mfg. (SM)", "2", "0.682\u20130.841", "SM2: 0.841", "SM1: 0.682"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.5), data,
          col_widths=[2.5, 0.7, 1.6, 2.3, 2.3])

# SLIDE 12: Factor Loadings Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Factor Loading Ranges Visualization",
              "Table 3: All Loadings Above 0.50 Threshold")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/factor_loadings.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 13: Cronbach's Alpha Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Reliability Analysis: Cronbach's Alpha",
              "Table 3: Internal Consistency by Construct")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/cronbach_alpha.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 14: Reliability Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Reliability & Validity Summary",
              "Table 3: Cronbach's \u03B1, Composite Reliability, AVE")
add_accent_bar(slide)
data = [
    ["Construct", "Cronbach's \u03B1", "C.R.", "AVE", "Assessment"],
    ["Green Procurement", "0.873", "0.874", "0.531", "Good"],
    ["Green Logistics", "0.836", "0.832", "0.502", "Good"],
    ["Green Product & Design", "0.896", "0.898", "0.633", "Good"],
    ["Regulatory Framework", "0.871", "0.871", "0.531", "Good"],
    ["Top Mgmt. Performance", "0.889", "0.887", "0.613", "Good"],
    ["Low Carbon Performance", "0.882", "0.874", "0.505", "Good"],
    ["Sustainable Society", "0.894", "0.892", "0.582", "Good"],
    ["Sustainable Manufacturing", "0.729", "0.730", "0.576", "Acceptable"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.5), data,
          col_widths=[2.8, 1.3, 1.3, 1.3, 2.7])

# SLIDE 15: CR & AVE Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Convergent Validity: CR & AVE",
              "Table 3: All Constructs Exceed Minimum Thresholds")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/cr_ave.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 16: Common Method Bias
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Common Method Bias Assessment",
              "Harman's Single-Factor Test & PCA Validation")
add_accent_bar(slide)
items = [
    ("Harman's Single-Factor Test (Table 6):", 16, True, DARK_BLUE),
    "",
    "All 41 items subjected to unrotated Principal Component Analysis",
    "First component explains only 26.48% of total variance",
    "Below the 50% threshold \u2192 No common method bias detected",
    "",
    ("Additional Safeguards:", 16, True, DARK_BLUE),
    "",
    "Anonymous survey responses to reduce social desirability bias",
    "Mandatory fields in Google Forms eliminated missing data issues",
    "Mixed positive/negative items to reduce acquiescence bias",
    "",
    ("Conclusion: Data is free from significant common method bias", 14, True, RGBColor(0, 128, 0)),
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=13, bullet=False)

# SLIDE 17: SEM Model Description
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Structural Equation Model: Research Hypotheses",
              "7 Hypothesized Relationships Tested via SmartPLS SEM")
add_accent_bar(slide)
data = [
    ["Hypothesis", "Path", "Direction"],
    ["H1", "Green Procurement \u2192 Top Mgmt. Performance", "Positive"],
    ["H2", "Green Logistics \u2192 Top Mgmt. Performance", "Positive"],
    ["H3", "Green Product & Design \u2192 Top Mgmt. Performance", "Positive"],
    ["H4", "Regulatory Framework \u2192 Top Mgmt. Performance", "Positive"],
    ["H5", "Top Mgmt. Performance \u2192 Low Carbon Performance", "Positive"],
    ["H6", "Low Carbon Performance \u2192 Sustainable Society", "Positive"],
    ["H7", "Low Carbon Performance \u2192 Sustainable Mfg.", "Positive"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.5), data,
          col_widths=[1.2, 5.5, 2.7])

# SLIDE 18: SEM Path Coefficients Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "SEM Results: Path Coefficients",
              "Table 7: All 7 Hypotheses Supported (p < 0.05)")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/sem_paths.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 19: SEM Results Table
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "SEM Hypothesis Testing: Complete Results",
              "Table 7: Path Coefficients, Standard Errors, and R\u00B2")
add_accent_bar(slide)
data = [
    ["Path", "\u03B2", "SE", "p-value", "R\u00B2", "Result"],
    ["GP \u2192 TMP", "0.130", "0.038", "0.000", "20.7%", "Supported"],
    ["GL \u2192 TMP", "0.200", "0.046", "0.000", "\u2014", "Supported"],
    ["GD \u2192 TMP", "0.274", "0.044", "0.000", "\u2014", "Supported"],
    ["RF \u2192 TMP", "0.100", "0.037", "0.007", "\u2014", "Supported"],
    ["TMP \u2192 CP", "0.388", "0.036", "0.000", "19.5%", "Supported"],
    ["CP \u2192 SS", "0.359", "0.032", "0.000", "16.1%", "Supported"],
    ["CP \u2192 SM", "0.295", "0.037", "0.000", "14.7%", "Supported"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.2), data,
          col_widths=[2.0, 1.0, 1.0, 1.0, 1.0, 3.4])

# SLIDE 20: R-Squared Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Variance Explained (R\u00B2) in Outcome Variables",
              "Table 7: Predictive Power of the SEM Model")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/r_squared.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 21: Effect Size Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Effect Size Analysis: Cohen's f\u00B2",
              "Practical Significance of Structural Relationships")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/effect_sizes.png', Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# SLIDE 22: Radar Summary Figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Psychometric Properties: Radar Summary",
              "Table 3: \u03B1, CR, and AVE Across All Constructs")
add_accent_bar(slide)
slide.shapes.add_picture('figures5/radar_summary.png', Inches(1.0), Inches(1.4), Inches(8.0), Inches(5.8))

# SLIDE 23: Key Findings - GSCM Practices
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Key Findings: GSCM Practice Effects",
              "Which Green Practices Matter Most?")
add_accent_bar(slide)
items = [
    ("Ranking of GSCM Practices on Top Management Performance:", 15, True, DARK_BLUE),
    "",
    "1. Green Product & Design (\u03B2 = 0.274) \u2014 Strongest influence",
    "   Product innovation and eco-design drive leadership commitment",
    "",
    "2. Green Logistics (\u03B2 = 0.200) \u2014 Second strongest",
    "   Efficient transportation and distribution practices matter",
    "",
    "3. Green Procurement (\u03B2 = 0.130) \u2014 Moderate influence",
    "   Sustainable supplier selection contributes to performance",
    "",
    "4. Regulatory Framework (\u03B2 = 0.100) \u2014 Significant but smallest",
    "   Government regulations play a supporting role (p = 0.007)",
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=13, bullet=False)

# SLIDE 24: Key Findings - Downstream Effects
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Key Findings: Sustainability Outcomes",
              "Cascading Effects from Management to Society")
add_accent_bar(slide)
items = [
    ("Cascading Impact Chain:", 15, True, DARK_BLUE),
    "",
    "GSCM Practices \u2192 Top Management Performance (R\u00B2 = 20.7%)",
    "Top Management \u2192 Low Carbon Performance (\u03B2 = 0.388, medium effect)",
    "Low Carbon \u2192 Sustainable Society (\u03B2 = 0.359, medium effect)",
    "Low Carbon \u2192 Sustainable Manufacturing (\u03B2 = 0.295, small effect)",
    "",
    ("Critical Insight:", 15, True, DARK_BLUE),
    "",
    "Top management commitment is the key mediator between",
    "green supply chain practices and environmental outcomes",
    "",
    "Organizations with strong senior leadership support are more",
    "likely to achieve carbon reduction and sustainability goals",
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=13, bullet=False)

# SLIDE 25: Practical Implications
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Practical Implications",
              "Actionable Insights for Manufacturing Firms")
add_accent_bar(slide)
items = [
    ("For Manufacturing Organizations:", 15, True, DARK_BLUE),
    "",
    "Prioritize green product design and eco-innovation initiatives",
    "Invest in green logistics to reduce environmental footprint",
    "Develop sustainable procurement policies with green suppliers",
    "",
    ("For Policymakers:", 15, True, DARK_BLUE),
    "",
    "Regulatory frameworks support GSCM adoption (\u03B2 = 0.100, p = 0.007)",
    "Subsidies and green labels facilitate smoother green product rollout",
    "",
    ("For Researchers:", 15, True, DARK_BLUE),
    "",
    "41-item instrument with 8 constructs shows acceptable-to-good reliability",
    "SmartPLS SEM effectively models complex GSCM relationships",
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=13, bullet=False)

# SLIDE 26: Limitations & Future Research
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "Limitations & Future Research",
              "Opportunities for Further Investigation")
add_accent_bar(slide)
items = [
    ("Study Limitations:", 15, True, DARK_BLUE),
    "",
    "Single industry focus (bag manufacturing) limits generalizability",
    "Cross-sectional design; longitudinal studies would capture trends",
    "Self-reported data from managers may introduce response bias",
    "Indian context may not generalize to other economies",
    "",
    ("Future Research Directions:", 15, True, DARK_BLUE),
    "",
    "Extend to other manufacturing sectors (automotive, electronics)",
    "Include financial performance as an additional outcome variable",
    "Conduct multi-country comparative studies on GSCM adoption",
    "Explore moderating effects of firm size and industry type",
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=13, bullet=False)

# SLIDE 27: References
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide, WHITE)
add_title_bar(slide, "References")
add_accent_bar(slide)
items = [
    ("Primary Study:", 14, True, DARK_BLUE),
    "",
    ("1. Gupta, P., Sharma, Y., Chauhan, A., Parewa, B., Rai, P. & Naik, N. (2025).", 11, False, DARK_GRAY),
    ("   \"Investigation of green supply chain management practices and sustainability", 11, False, DARK_GRAY),
    ("   in Indian manufacturing enterprises using a structural equation modelling approach.\"", 11, False, DARK_GRAY),
    ("   Scientific Reports, 15, 14909. DOI: 10.1038/s41598-025-95940-9", 11, True, MEDIUM_BLUE),
    "",
    ("Supporting References:", 14, True, DARK_BLUE),
    "",
    ("2. Hair, J.F., Black, W.C., Babin, B.J. & Anderson, R.E. (2019).", 11, False, DARK_GRAY),
    ("   Multivariate Data Analysis (8th Ed.). Cengage Learning.", 11, False, DARK_GRAY),
    "",
    ("3. Kaiser, H.F. (1974). \"An Index of Factorial Simplicity.\"", 11, False, DARK_GRAY),
    ("   Psychometrika, 39, 31\u201336.", 11, False, DARK_GRAY),
]
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), items, font_size=11, bullet=False)

prs.save('Factor_Analysis_GSCM_Presentation.pptx')
print("\nPresentation saved: Factor_Analysis_GSCM_Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
