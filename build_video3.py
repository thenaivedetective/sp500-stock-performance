"""
One-pass video builder: render slide → encode segment → repeat → concat.
All intermediate files stay in workspace so /tmp clears don't matter.
"""
import os, io, subprocess, zipfile, shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

PPTX   = 'Final_Presentation_AutoPlay.pptx'
OUTPUT = 'Final_Presentation_Video.mp4'
WORK   = 'video_build'          # workspace subfolder, survives restarts
os.makedirs(WORK, exist_ok=True)

W, H = 1280, 720

AUDIO_MAP = {
     4: ('media5.wav',   85.90),
     5: ('media6.wav',   42.75),
     6: ('media7.wav',   25.59),
     7: ('media8.wav',   56.18),
     8: ('media9.wav',   33.42),
     9: ('media10.wav',  31.94),
    10: ('media11.wav',  68.68),
    11: ('media12.wav',  53.95),
    12: ('media13.wav',  47.26),
    13: ('media14.wav', 101.16),
    14: ('media15.wav', 102.76),
}

# ── Extract WAV files (skip if already done) ──────────────────────────────────
for _, (wav, _) in AUDIO_MAP.items():
    dst = f'{WORK}/{wav}'
    if not os.path.exists(dst):
        with zipfile.ZipFile(PPTX, 'r') as z:
            open(dst, 'wb').write(z.read(f'ppt/media/{wav}'))
        print(f'Extracted {wav}')

# ── Font helper ───────────────────────────────────────────────────────────────
_fc = {}
def font(sz):
    sz = max(10, sz)
    if sz in _fc: return _fc[sz]
    for p in ['/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
              '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf']:
        if os.path.exists(p):
            try: f = ImageFont.truetype(p, sz); _fc[sz] = f; return f
            except: pass
    f = ImageFont.load_default(); _fc[sz] = f; return f

def ep(e, ed, pd): return int(int(e) / int(ed) * pd)

def slide_bg(slide):
    for src in [slide.background, slide.slide_layout.background,
                slide.slide_layout.slide_master.background]:
        try:
            fill = src.fill
            if fill.type and str(fill.type) == 'SOLID (1)':
                rgb = fill.fore_color.rgb
                return (rgb.red, rgb.green, rgb.blue)
        except: pass
    return (15, 36, 74)

def render_slide(slide, ew, eh):
    bg  = slide_bg(slide)
    img = Image.new('RGB', (W, H), color=bg)
    draw = ImageDraw.Draw(img)
    pics  = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    texts = [s for s in slide.shapes if s.shape_type not in
             (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA)]
    for shape in pics:
        try:
            l=ep(shape.left or 0,ew,W); t=ep(shape.top or 0,eh,H)
            w=ep(shape.width or 0,ew,W); h=ep(shape.height or 0,eh,H)
            if w<1 or h<1: continue
            pi = Image.open(io.BytesIO(shape.image.blob))
            pi = pi.resize((w,h), Image.LANCZOS)
            img.paste(pi,(l,t), pi if pi.mode=='RGBA' else None)
        except: pass
    for shape in texts:
        if not getattr(shape,'has_text_frame',False) or not shape.has_text_frame:
            continue
        try:
            l=ep(shape.left or 0,ew,W); t=ep(shape.top or 0,eh,H)
            w=ep(shape.width or 0,ew,W); h=ep(shape.height or 0,eh,H)
            y = t+2
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if not txt: y+=9; continue
                try:
                    fsz=(max(10,int(para.runs[0].font.size.pt*96/72))
                         if para.runs and para.runs[0].font.size else 18)
                    r=para.runs[0]; fc=r.font.color
                    col=((fc.rgb.red,fc.rgb.green,fc.rgb.blue)
                         if fc and fc.type else (255,255,255))
                except: fsz=18; col=(255,255,255)
                f=font(fsz); mw=max(1,w-8)
                words=txt.split(); lines=[]; cur=''
                for word in words:
                    test=(cur+' '+word).strip()
                    if draw.textbbox((0,0),test,font=f)[2]<=mw: cur=test
                    else:
                        if cur: lines.append(cur)
                        cur=word
                if cur: lines.append(cur)
                for line in lines:
                    if y+fsz > t+h+10: break
                    draw.text((l+5,y+1),line,font=f,fill=(0,0,0))
                    draw.text((l+4,y),  line,font=f,fill=col)
                    y+=fsz+2
                y+=3
        except: pass
    return img

# ── Main loop: render + encode each slide ────────────────────────────────────
prs = Presentation(PPTX)
ew, eh = int(prs.slide_width), int(prs.slide_height)
print(f'Processing {len(prs.slides)} slides...')

seg_list = []
for si, slide in enumerate(prs.slides):
    seg = f'{WORK}/seg_{si:02d}.mp4'
    seg_list.append(seg)
    if os.path.exists(seg):
        print(f'  Slide {si+1:2d} segment exists, skip')
        continue

    frame_path = f'{WORK}/slide_{si:02d}.png'
    img = render_slide(slide, ew, eh)
    img.save(frame_path, 'PNG')

    if si in AUDIO_MAP:
        wav, dur = AUDIO_MAP[si]
        cmd = ['ffmpeg','-y','-loop','1','-i',frame_path,
               '-i',f'{WORK}/{wav}',
               '-c:v','libx264','-preset','ultrafast','-tune','stillimage',
               '-c:a','aac','-b:a','128k','-pix_fmt','yuv420p',
               '-t',str(dur+0.5),'-shortest',seg]
        tag = f'{dur:.0f}s audio'
    else:
        cmd = ['ffmpeg','-y','-loop','1','-i',frame_path,
               '-c:v','libx264','-preset','ultrafast','-tune','stillimage',
               '-pix_fmt','yuv420p','-t','5','-an',seg]
        tag = '5s silent'

    r = subprocess.run(cmd, capture_output=True)
    os.remove(frame_path)   # clean up frame immediately
    if r.returncode == 0:
        sz = os.path.getsize(seg)/1024
        print(f'  Slide {si+1:2d} done  ({tag}, {sz:.0f} KB)')
    else:
        print(f'  ERROR slide {si+1}: {r.stderr.decode()[-120:]}')

# ── Concatenate ────────────────────────────────────────────────────────────────
ready = [s for s in seg_list if os.path.exists(s)]
print(f'\nConcatenating {len(ready)}/{len(seg_list)} segments...')
if len(ready) == len(seg_list):
    cl = f'{WORK}/concat.txt'
    open(cl,'w').write('\n'.join(f"file '{s}'" for s in ready)+'\n')
    r = subprocess.run(['ffmpeg','-y','-f','concat','-safe','0',
                        '-i',cl,'-c','copy',OUTPUT], capture_output=True)
    if r.returncode == 0:
        mb = os.path.getsize(OUTPUT)/1024/1024
        print(f'Done!  {OUTPUT}  ({mb:.1f} MB)')
    else:
        print('Concat error:', r.stderr.decode()[-300:])
else:
    missing = [s for s in seg_list if not os.path.exists(s)]
    print(f'Missing segments: {missing}')
    print('Re-run the script to continue from where it left off.')
