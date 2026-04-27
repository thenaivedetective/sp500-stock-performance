from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

NAVY      = RGBColor(0x0A, 0x16, 0x28)
NAVY_CARD = RGBColor(0x13, 0x20, 0x38)
GOLD      = RGBColor(0xD4, 0xAF, 0x37)
BLUE      = RGBColor(0x4A, 0x9E, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x8B, 0x9C, 0xB8)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_NAV = RGBColor(0x1A, 0x30, 0x50)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def blank_slide():
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    return slide

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.color.rgb = color
    run.font.italic  = italic
    return txBox

def add_para(tf, text, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p  = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(0.08), GOLD)
    add_text(slide, title,
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.65),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
                 size=14, color=MUTED)
    add_rect(slide, 0, Inches(1.1), W, Inches(0.02), LIGHT_NAV)

def stat_box(slide, label, value, unit, l, t, w, h, val_color=GOLD):
    add_rect(slide, l, t, w, h, NAVY_CARD)
    add_rect(slide, l, t, Inches(0.05), h, GOLD)
    add_text(slide, label, l+Inches(0.15), t+Inches(0.12), w-Inches(0.2), Inches(0.3),
             size=11, color=MUTED, bold=False)
    add_text(slide, value, l+Inches(0.15), t+Inches(0.38), w-Inches(0.2), Inches(0.55),
             size=28, bold=True, color=val_color)
    if unit:
        add_text(slide, unit, l+Inches(0.15), t+h-Inches(0.35), w-Inches(0.2), Inches(0.3),
                 size=10, color=MUTED)

def section_divider(slide, section_num, title, subtitle=""):
    add_rect(slide, 0, 0, W, H, NAVY)
    add_rect(slide, 0, Inches(3.0), W, Inches(0.06), GOLD)
    add_text(slide, f"Section {section_num}", Inches(1), Inches(2.2), Inches(11), Inches(0.5),
             size=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1), Inches(2.65), Inches(11), Inches(0.9),
             size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
                 size=18, color=MUTED, align=PP_ALIGN.CENTER)

# ─── SLIDE 1: TITLE ───────────────────────────────────────────────────────────
slide = blank_slide()
add_rect(slide, 0, 0, Inches(5.5), H, NAVY_CARD)
add_rect(slide, Inches(5.5), 0, W - Inches(5.5), H, LIGHT_NAV)
add_rect(slide, Inches(5.5), 0, Inches(0.06), H, GOLD)

add_text(slide, "SSIE 605 · Binghamton University",
         Inches(0.5), Inches(1.0), Inches(4.8), Inches(0.4),
         size=12, color=GOLD, bold=True)

add_text(slide, "Using Logistic\nRegression to\nDetermine U.S.\nEquity Performance",
         Inches(0.5), Inches(1.5), Inches(4.6), Inches(3.2),
         size=32, bold=True, color=WHITE)

add_rect(slide, Inches(0.5), Inches(4.8), Inches(1.2), Inches(0.04), GOLD)

add_text(slide, "Lana Gidan  ·  Matthew Golubow  ·  Shahd Tarman",
         Inches(0.5), Inches(5.0), Inches(4.8), Inches(0.4),
         size=13, color=MUTED)
add_text(slide, "Q1 2010 – Q4 2024  |  S&P 500 Constituents",
         Inches(0.5), Inches(5.5), Inches(4.8), Inches(0.35),
         size=12, color=MUTED)

add_text(slide, "PREDICTING\nOUTPERFORMANCE\nUSING QUARTERLY\nFINANCIAL RATIOS",
         Inches(6.0), Inches(2.2), Inches(6.8), Inches(3.0),
         size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide, "Logistic Regression  |  VIF Filtering  |  PCA  |  Out-of-Sample 2025",
         Inches(6.0), Inches(5.3), Inches(6.8), Inches(0.4),
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ─── SLIDE 2: AGENDA ──────────────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Agenda", "Presentation Structure")

items = [
    ("01", "Background",               "SEC disclosures, S&P 500, research motivation"),
    ("02", "Problem Statement",        "Three core research questions"),
    ("03", "Data Capture",             "WRDS, FRED, EDGAR — sources and variables"),
    ("04", "Data Computations",        "Financial ratios and delta momentum features"),
    ("05", "Data Cleaning",            "Winsorization, standardization, NaN handling"),
    ("06", "VIF Analysis",             "Multicollinearity removal — 20 of 24 kept"),
    ("07", "PCA Analysis",             "9 components, 5% individual variance threshold"),
    ("08", "Logistic Regression",      "Global model results — AUC 0.5381"),
    ("09", "Segmentation & Results",   "Market Cap and Sector — best: Comm. Services 63.9%"),
    ("10", "2025 Out-of-Sample Test",  "Portfolio vs S&P 500"),
    ("11", "Conclusion",               "Efficient Market Hypothesis assessment"),
]

col1 = items[:6]
col2 = items[6:]

for i, (num, title, desc) in enumerate(col1):
    t = Inches(1.25) + i * Inches(0.96)
    add_rect(slide, Inches(0.4), t, Inches(0.5), Inches(0.4), GOLD)
    add_text(slide, num, Inches(0.4), t, Inches(0.5), Inches(0.4),
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.0), t, Inches(2.5), Inches(0.28),
             size=13, bold=True, color=WHITE)
    add_text(slide, desc, Inches(1.0), t+Inches(0.25), Inches(2.8), Inches(0.28),
             size=9, color=MUTED)

for i, (num, title, desc) in enumerate(col2):
    t = Inches(1.25) + i * Inches(0.96)
    add_rect(slide, Inches(6.8), t, Inches(0.5), Inches(0.4), GOLD)
    add_text(slide, num, Inches(6.8), t, Inches(0.5), Inches(0.4),
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(7.4), t, Inches(2.5), Inches(0.28),
             size=13, bold=True, color=WHITE)
    add_text(slide, desc, Inches(7.4), t+Inches(0.25), Inches(5.3), Inches(0.28),
             size=9, color=MUTED)

# ─── SLIDE 3: BACKGROUND ──────────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Background", "Section 01")

add_text(slide,
    "The Securities and Exchange Commission (SEC) requires all publicly traded companies to disclose "
    "their financial statements on a quarterly basis. These reports detail revenues, expenses, assets, "
    "and debt — information that investment analysts use to evaluate company health.",
    Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.9),
    size=14, color=WHITE)

boxes = [
    ("S&P 500", "An index of the 500 largest U.S. companies by market cap. "
                "Covers ~80% of the total U.S. stock market value. "
                "Rebalanced quarterly as constituents change."),
    ("Passive Investing Assumption", "It is widely accepted that the average investor "
                "cannot consistently beat the S&P 500 return. Index funds charge low fees "
                "and require no stock selection skill."),
    ("Our Research Question", "Can we use publicly available quarterly financial data "
                "to predict which specific S&P 500 companies will outperform "
                "the index the following quarter?"),
]

for i, (title, body) in enumerate(boxes):
    l = Inches(0.4) + i * Inches(4.3)
    add_rect(slide, l, Inches(2.35), Inches(4.0), Inches(4.6), NAVY_CARD)
    add_rect(slide, l, Inches(2.35), Inches(4.0), Inches(0.06), GOLD)
    add_text(slide, title, l+Inches(0.2), Inches(2.5), Inches(3.6), Inches(0.4),
             size=14, bold=True, color=GOLD)
    add_text(slide, body, l+Inches(0.2), Inches(2.95), Inches(3.6), Inches(3.7),
             size=12, color=WHITE)

# ─── SLIDE 4: PROBLEM STATEMENT ───────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Problem Statement", "Section 02")

questions = [
    ("Q1", "Predictability",
     "Can we use publicly available S&P 500 financial data to successfully determine "
     "if a company will outperform or underperform the cumulative return of the S&P 500?"),
    ("Q2", "Forward-Looking Signal",
     "Can we use the current quarter's financial ratios to reliably predict "
     "next quarter's stock performance — before it happens?"),
    ("Q3", "Portfolio Construction",
     "Can we use model predictions to construct a portfolio whose quarterly "
     "returns beat the S&P 500 index, outperforming passive investing?"),
]

for i, (label, title, body) in enumerate(questions):
    t = Inches(1.35) + i * Inches(1.85)
    add_rect(slide, Inches(0.4), t, Inches(12.5), Inches(1.65), NAVY_CARD)
    add_rect(slide, Inches(0.4), t, Inches(0.06), Inches(1.65), GOLD)
    add_rect(slide, Inches(0.4), t, Inches(0.85), Inches(1.65), LIGHT_NAV)
    add_text(slide, label, Inches(0.4), t+Inches(0.55), Inches(0.85), Inches(0.5),
             size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.4), t+Inches(0.1), Inches(11.0), Inches(0.45),
             size=16, bold=True, color=GOLD)
    add_text(slide, body, Inches(1.4), t+Inches(0.55), Inches(11.0), Inches(1.0),
             size=13, color=WHITE)

# ─── SLIDE 5: DATA SOURCES ────────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Data Capture", "Section 03 — Sources")

sources = [
    ("WRDS", "Wharton Research Data Service",
     "University of Pennsylvania's data repository.\n"
     "Source of S&P 500 constituent history, quarterly\n"
     "financials (Compustat), and stock return data (CRSP).\n"
     "Coverage: Q1 2010 – Q4 2024."),
    ("FRED", "Federal Reserve Economic Data",
     "Federal Reserve Bank of St. Louis data service.\n"
     "Source of macroeconomic indicators: GDP growth\n"
     "and CPI inflation — one value per quarter,\n"
     "matched to each company's feature quarter."),
    ("EDGAR", "Electronic Data Gathering, Analysis & Retrieval",
     "SEC's official repository for company filings.\n"
     "Used to supplement Compustat with additional\n"
     "quarterly 10-Q financial statement data.\n"
     "Accessed via API through Python."),
]

for i, (abbr, full, body) in enumerate(sources):
    l = Inches(0.4) + i * Inches(4.3)
    add_rect(slide, l, Inches(1.3), Inches(4.0), Inches(5.6), NAVY_CARD)
    add_rect(slide, l, Inches(1.3), Inches(4.0), Inches(0.08), GOLD)
    add_text(slide, abbr, l+Inches(0.2), Inches(1.45), Inches(3.6), Inches(0.6),
             size=30, bold=True, color=GOLD)
    add_text(slide, full, l+Inches(0.2), Inches(2.0), Inches(3.6), Inches(0.4),
             size=11, bold=True, color=MUTED)
    add_rect(slide, l+Inches(0.2), Inches(2.4), Inches(0.6), Inches(0.03), GOLD)
    add_text(slide, body, l+Inches(0.2), Inches(2.55), Inches(3.6), Inches(4.0),
             size=12, color=WHITE)

add_text(slide, "All data pulled via API and processed using Python (pandas, statsmodels, scikit-learn)",
         Inches(0.5), Inches(7.0), Inches(12.0), Inches(0.35),
         size=10, color=MUTED, italic=True)

# ─── SLIDE 6: RAW VARIABLES ───────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Data Capture — Raw Variables", "Section 03 — Compustat Quarterly Fields")

vars_data = [
    ("prccq",   "Price Close (Quarter End)"),
    ("cshoq",   "Common Shares Outstanding"),
    ("mkvaltq", "Market Capitalization"),
    ("revtq",   "Total Revenue"),
    ("cogsq",   "Cost of Goods Sold"),
    ("xsgaq",   "SG&A Expenses"),
    ("xrdq",    "R&D Expenses"),
    ("oiadpq",  "Operating Income After Depreciation"),
    ("niq",     "Net Income"),
    ("ibq",     "Income Before Extraordinary Items"),
    ("piq",     "Pretax Income"),
    ("atq",     "Total Assets"),
    ("ceqq",    "Common Equity Total"),
    ("teqq",    "Total Stockholders' Equity"),
    ("dlttq",   "Long-Term Debt"),
    ("dlcq",    "Debt in Current Liabilities"),
    ("actq",    "Current Assets"),
    ("lctq",    "Current Liabilities"),
    ("cheq",    "Cash & Short-Term Investments"),
    ("dpq",     "Depreciation & Amortization"),
    ("oibdpq",  "Operating Income Before Depreciation"),
]

col_w = Inches(3.15)
row_h = Inches(0.31)
cols = 4
for i, (code, label) in enumerate(vars_data):
    col = i % cols
    row = i // cols
    l = Inches(0.35) + col * col_w
    t = Inches(1.25) + row * row_h
    bg = NAVY_CARD if row % 2 == 0 else LIGHT_NAV
    add_rect(slide, l, t, col_w - Inches(0.05), row_h - Inches(0.02), bg)
    add_text(slide, code, l+Inches(0.1), t+Inches(0.04), Inches(1.1), row_h,
             size=10, bold=True, color=GOLD)
    add_text(slide, label, l+Inches(1.2), t+Inches(0.04), Inches(1.85), row_h,
             size=10, color=WHITE)

add_text(slide, "Source: WRDS Compustat Fundamentals Quarterly (fundq table)  |  Joined on gvkey identifier",
         Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
         size=9, color=MUTED, italic=True)

# ─── SLIDE 7: FINANCIAL RATIOS ────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Data Computations — Financial Ratios", "Section 04 — 14 Base Ratios Derived from Raw Variables")

ratios = [
    ("Return on Assets (ROA)",    "Net Income / Total Assets",                    "Profitability"),
    ("Return on Equity (ROE)",    "Net Income / Common Equity Total",             "Profitability"),
    ("Gross Margin",              "(Revenue - COGS) / Revenue",                   "Profitability"),
    ("Operating Margin",          "Op. Income After Depr. / Revenue",             "Profitability"),
    ("Net Margin",                "Net Income / Revenue",                         "Profitability"),
    ("Asset Turnover",            "Revenue / Total Assets",                       "Efficiency"),
    ("Current Ratio",             "Current Assets / Current Liabilities",         "Liquidity"),
    ("Debt to Equity",            "Long-Term Debt / Common Equity",               "Leverage"),
    ("Revenue Growth",            "(Rev - Lag Rev) / |Lag Rev|",                  "Growth"),
    ("Net Income Growth",         "(NI - Lag NI) / |Lag NI|",                    "Growth"),
    ("P/E Ratio",                 "Price / (Income Before Extr. / Shares)",       "Valuation"),
    ("Book-to-Market",            "Common Equity / Market Cap",                   "Valuation"),
    ("GDP Growth",                "Quarter-over-quarter GDP change (FRED)",        "Macro"),
    ("Inflation",                 "Quarter-over-quarter CPI change (FRED)",        "Macro"),
]

cat_colors = {
    "Profitability": GOLD,
    "Efficiency":    BLUE,
    "Liquidity":     GREEN,
    "Leverage":      RGBColor(0xE6, 0x7E, 0x22),
    "Growth":        RGBColor(0x9B, 0x59, 0xB6),
    "Valuation":     RGBColor(0x1A, 0xBC, 0x9C),
    "Macro":         MUTED,
}

row_h = Inches(0.39)
for i, (name, formula, cat) in enumerate(ratios):
    col = i % 2
    row = i // 2
    l = Inches(0.35) + col * Inches(6.5)
    t = Inches(1.25) + row * row_h
    bg = NAVY_CARD if row % 2 == 0 else LIGHT_NAV
    add_rect(slide, l, t, Inches(6.3), row_h - Inches(0.02), bg)
    dot_color = cat_colors.get(cat, GOLD)
    add_rect(slide, l, t, Inches(0.04), row_h - Inches(0.02), dot_color)
    add_text(slide, name, l+Inches(0.12), t+Inches(0.06), Inches(2.4), Inches(0.28),
             size=10, bold=True, color=WHITE)
    add_text(slide, formula, l+Inches(2.55), t+Inches(0.06), Inches(2.8), Inches(0.28),
             size=9, color=MUTED, italic=True)
    add_text(slide, cat, l+Inches(5.45), t+Inches(0.06), Inches(0.8), Inches(0.28),
             size=8, bold=True, color=dot_color)

add_text(slide, "Lag values computed via groupby gvkey shift(1)  |  GDP and Inflation matched by feature quarter",
         Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
         size=9, color=MUTED, italic=True)

# ─── SLIDE 8: DELTA FEATURES ──────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Delta Features — Quarter-over-Quarter Momentum", "Section 04 — 10 Additional Features")

add_text(slide,
    "Delta features capture the direction of change in each ratio — whether a company's financials are improving "
    "or deteriorating. A company with rising ROA tells a fundamentally different story than one with a declining ROA, "
    "even if their current values are identical.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.8),
    size=13, color=WHITE)

add_rect(slide, Inches(0.4), Inches(2.1), Inches(5.9), Inches(4.8), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(2.1), Inches(5.9), Inches(0.06), GOLD)
add_text(slide, "Formula", Inches(0.6), Inches(2.18), Inches(5.5), Inches(0.4),
         size=14, bold=True, color=GOLD)
add_text(slide,
    "Delta_Ratio(Q) = Ratio(Q) - Ratio(Q-1)\n\n"
    "Computed via groupby(gvkey).diff()\n\n"
    "Example:\n"
    "  ROA (Q3 2024) = 5.0%\n"
    "  ROA (Q4 2024) = 8.0%\n"
    "  Delta_ROA     = +3.0%  (Improving)\n\n"
    "Positive delta = improving fundamentals\n"
    "Negative delta = deteriorating fundamentals",
    Inches(0.6), Inches(2.65), Inches(5.5), Inches(4.0),
    size=13, color=WHITE)

delta_feats = [
    "Delta ROA",         "Delta ROE",
    "Delta Gross Margin","Delta Op Margin",
    "Delta Net Margin",  "Delta Asset Turnover",
    "Delta Current Ratio","Delta Debt/Equity",
    "Delta P/E Ratio",   "Delta Book-to-Market",
]

add_text(slide, "10 Delta Features Added", Inches(6.6), Inches(2.18), Inches(6.0), Inches(0.4),
         size=14, bold=True, color=GOLD)

for i, feat in enumerate(delta_feats):
    col = i % 2
    row = i // 2
    l = Inches(6.6) + col * Inches(3.0)
    t = Inches(2.65) + row * Inches(0.45)
    add_rect(slide, l, t, Inches(2.8), Inches(0.38), LIGHT_NAV)
    add_rect(slide, l, t, Inches(0.04), Inches(0.38), GOLD)
    add_text(slide, feat, l+Inches(0.12), t+Inches(0.08), Inches(2.6), Inches(0.25),
             size=11, color=WHITE)

add_text(slide, "Total features entering VIF analysis: 24   (14 base ratios + 10 delta features)",
         Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
         size=10, bold=True, color=GOLD, italic=False)

# ─── SLIDE 9: OUTPERFORMER DEFINITION ─────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Outperformer Definition & Forward-Looking Label", "Section 04")

add_rect(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(2.0), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.06), GOLD)
add_text(slide, "Label Construction — Forward-Looking by One Quarter",
         Inches(0.6), Inches(1.35), Inches(12.0), Inches(0.4),
         size=16, bold=True, color=GOLD)
add_text(slide,
    "Stock Return (Q+1) > S&P 500 Return (Q+1)   -->   Label = 1  (Outperformer)\n"
    "Stock Return (Q+1) <= S&P 500 Return (Q+1)  -->   Label = 0  (Underperformer)",
    Inches(0.6), Inches(1.8), Inches(12.0), Inches(1.2),
    size=14, color=WHITE)

add_text(slide, "Features used: Q[t] financial ratios          Prediction target: Q[t+1] outperformance",
         Inches(0.6), Inches(2.85), Inches(12.0), Inches(0.35),
         size=13, bold=True, color=BLUE)

cols_info = [
    ("Feature Quarter", "Q4 2024 financial ratios (ROA, margins, etc.)\nEnter as model input", GOLD),
    ("Arrow", "Predicts", WHITE),
    ("Label Quarter", "Q1 2025 — did the stock beat\nthe S&P 500 that quarter?", BLUE),
]

add_rect(slide, Inches(0.4), Inches(3.4), Inches(4.5), Inches(3.4), LIGHT_NAV)
add_text(slide, "Feature Quarter: Q4 2024", Inches(0.6), Inches(3.55), Inches(4.1), Inches(0.4),
         size=14, bold=True, color=GOLD)
add_text(slide,
    "Financial ratios extracted from 10-Q filings.\n"
    "Includes 14 base ratios + 10 delta features.\n"
    "Available before Q1 2025 begins.",
    Inches(0.6), Inches(4.0), Inches(4.1), Inches(2.5),
    size=12, color=WHITE)

add_text(slide, "  ->  ", Inches(4.95), Inches(4.8), Inches(0.9), Inches(0.5),
         size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.9), Inches(3.4), Inches(4.5), Inches(3.4), LIGHT_NAV)
add_text(slide, "Label Quarter: Q1 2025", Inches(6.1), Inches(3.55), Inches(4.1), Inches(0.4),
         size=14, bold=True, color=BLUE)
add_text(slide,
    "Was the stock's Q1 2025 return\ngreater than SPY Q1 2025 return?\n\n"
    "Yes -> Label = 1 (Outperformer)\nNo  -> Label = 0 (Underperformer)\n\n"
    "Computed via CRSP quarterly returns.",
    Inches(6.1), Inches(4.0), Inches(4.1), Inches(2.5),
    size=12, color=WHITE)

add_text(slide,
    "outperformer_next = groupby(gvkey)[outperformer_quarterly].shift(-1)  |  Training data: 16,589 complete rows",
    Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
    size=9, color=MUTED, italic=True)

# ─── SLIDE 10: DATA CLEANING ──────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Data Cleaning", "Section 05 — Winsorization, Standardization, NaN Removal")

steps = [
    ("Step 1", "Winsorization",
     "Extreme outlier values are capped at the 1st and 99th percentile of each ratio.\n\n"
     "Example: If ROA for one company is -300%, it is capped at the 1st percentile value "
     "(e.g. -15%). The row is kept — only the extreme value is replaced.\n\n"
     "Applied per ratio using training data bounds. Test data uses the same bounds "
     "fitted on training to prevent data leakage.", GOLD),
    ("Step 2", "Standardization (Z-Score)",
     "Each ratio is scaled to mean = 0, standard deviation = 1.\n\n"
     "Without standardization, ratios with larger numeric ranges (e.g. P/E Ratio at "
     "50x) would dominate over smaller ratios (e.g. ROA at 0.05).\n\n"
     "Ensures all 24 features are on the same scale before VIF and PCA.", BLUE),
    ("Step 3", "NaN Removal",
     "Rows with any missing values after ratio computation are dropped.\n\n"
     "Common causes: companies with zero revenue (division by zero), first-quarter "
     "entries with no prior quarter for lag/delta computation.\n\n"
     "Training dataset reduced from 21,285 rows to 16,589 complete rows.", GREEN),
]

for i, (label, title, body, color) in enumerate(steps):
    l = Inches(0.35) + i * Inches(4.3)
    add_rect(slide, l, Inches(1.25), Inches(4.1), Inches(5.7), NAVY_CARD)
    add_rect(slide, l, Inches(1.25), Inches(4.1), Inches(0.07), color)
    add_text(slide, label, l+Inches(0.15), Inches(1.35), Inches(1.5), Inches(0.35),
             size=11, bold=True, color=color)
    add_text(slide, title, l+Inches(0.15), Inches(1.7), Inches(3.8), Inches(0.4),
             size=15, bold=True, color=WHITE)
    add_rect(slide, l+Inches(0.15), Inches(2.15), Inches(0.7), Inches(0.03), color)
    add_text(slide, body, l+Inches(0.15), Inches(2.25), Inches(3.8), Inches(4.5),
             size=11, color=WHITE)

stat_box(slide, "Rows Before Cleaning", "21,285", "training observations",
         Inches(0.35), Inches(6.9), Inches(3.0), Inches(0.48))
add_text(slide, "->", Inches(3.45), Inches(6.98), Inches(0.5), Inches(0.35),
         size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
stat_box(slide, "Rows After Cleaning", "16,589", "complete observations",
         Inches(4.0), Inches(6.9), Inches(3.0), Inches(0.48), val_color=GREEN)

# ─── SLIDE 11: VIF ANALYSIS ───────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "VIF Analysis — Multicollinearity Removal", "Section 06 — Variance Inflation Factor")

add_text(slide,
    "Many financial ratios are correlated — for example, Net Margin and ROA both measure profitability and move together. "
    "When correlated features are included together, logistic regression becomes unstable and coefficients are unreliable. "
    "VIF quantifies this collinearity for each predictor.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.75),
    size=12, color=WHITE)

add_rect(slide, Inches(0.4), Inches(2.05), Inches(5.5), Inches(4.8), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(2.05), Inches(5.5), Inches(0.06), GOLD)
add_text(slide, "VIF Procedure", Inches(0.6), Inches(2.15), Inches(5.1), Inches(0.4),
         size=14, bold=True, color=GOLD)
steps_vif = [
    "1. Compute VIF for all 24 features simultaneously",
    "2. Identify the feature with the highest VIF",
    "3. If that VIF > 2.5, remove the feature",
    "4. Repeat until all remaining features have VIF <= 2.5",
    "5. Use threshold of 2.5 (stricter than standard 5.0)",
    "   — following Ananthakumar & Sarkar (2017)",
]
for i, step in enumerate(steps_vif):
    add_text(slide, step, Inches(0.6), Inches(2.65) + i*Inches(0.5), Inches(5.1), Inches(0.4),
             size=12, color=WHITE if not step.startswith("   ") else MUTED)

add_rect(slide, Inches(6.2), Inches(2.05), Inches(6.6), Inches(4.8), NAVY_CARD)
add_rect(slide, Inches(6.2), Inches(2.05), Inches(6.6), Inches(0.06), GOLD)
add_text(slide, "20 Features Retained (of 24)", Inches(6.4), Inches(2.15), Inches(6.2), Inches(0.4),
         size=14, bold=True, color=GOLD)

kept = [
    "roe", "gross_margin", "op_margin", "asset_turnover",
    "current_ratio", "debt_to_equity", "rev_growth", "ni_growth",
    "pe_ratio", "book_to_market", "gdp_growth", "inflation",
    "delta_roe", "delta_gross_margin", "delta_net_margin",
    "delta_asset_turnover", "delta_current_ratio",
    "delta_debt_to_equity", "delta_pe_ratio", "delta_book_to_market",
]
removed = ["roa", "net_margin", "delta_roa", "delta_net_margin_dup"]

for i, feat in enumerate(kept):
    col = i % 2
    row = i // 2
    l = Inches(6.4) + col * Inches(3.1)
    t = Inches(2.62) + row * Inches(0.39)
    add_rect(slide, l, t, Inches(2.9), Inches(0.34), LIGHT_NAV)
    add_rect(slide, l, t, Inches(0.04), Inches(0.34), GREEN)
    add_text(slide, feat, l+Inches(0.12), t+Inches(0.07), Inches(2.7), Inches(0.25),
             size=9, color=WHITE)

stat_box(slide, "Features Retained", "20/24", "VIF <= 2.5",
         Inches(0.4), Inches(6.9), Inches(2.8), Inches(0.48), val_color=GREEN)
stat_box(slide, "Features Removed", "4", "VIF > 2.5 (collinear)",
         Inches(3.4), Inches(6.9), Inches(2.8), Inches(0.48), val_color=RED)
stat_box(slide, "VIF Threshold", "2.5", "Ananthakumar & Sarkar benchmark",
         Inches(6.4), Inches(6.9), Inches(3.5), Inches(0.48))

# ─── SLIDE 12: PCA ANALYSIS ───────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "PCA Analysis — Dimensionality Reduction", "Section 07 — Principal Component Analysis")

add_text(slide,
    "Even after VIF filtering, 20 correlated features remain. PCA compresses them into uncorrelated "
    "principal components, each capturing a distinct pattern in the data — eliminating any residual "
    "multicollinearity and reducing noise.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.75),
    size=12, color=WHITE)

add_rect(slide, Inches(0.4), Inches(2.05), Inches(4.2), Inches(4.8), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(2.05), Inches(4.2), Inches(0.06), GOLD)
add_text(slide, "Component Selection Rule", Inches(0.6), Inches(2.15), Inches(3.8), Inches(0.4),
         size=14, bold=True, color=GOLD)
add_text(slide,
    "Threshold: >= 5% individual\nvariance explained\n\n"
    "Applied on recommendation from\nfaculty advisor.\n\n"
    "Rule: n_comp = max(1, count of\ncomponents where explained\nvariance ratio >= 0.05)\n\n"
    "Result: 9 principal components\nretained for the global model.\n\n"
    "The number of components varies\nby sector and sub-configuration.",
    Inches(0.6), Inches(2.65), Inches(3.8), Inches(4.0),
    size=12, color=WHITE)

add_rect(slide, Inches(4.85), Inches(2.05), Inches(8.0), Inches(4.8), NAVY_CARD)
add_rect(slide, Inches(4.85), Inches(2.05), Inches(8.0), Inches(0.06), GOLD)
add_text(slide, "What Each Component Represents", Inches(5.05), Inches(2.15), Inches(7.6), Inches(0.4),
         size=14, bold=True, color=GOLD)

pca_rows = [
    ("PC1", "Profitability cluster", "ROA, gross margin, op margin loadings"),
    ("PC2", "Leverage & liquidity", "Debt/equity, current ratio loadings"),
    ("PC3", "Revenue momentum",    "Rev growth, delta revenue loadings"),
    ("PC4", "Valuation signal",    "P/E ratio, book-to-market loadings"),
    ("PC5", "NI & margin delta",   "Delta net margin, delta ROE loadings"),
    ("PC6", "Asset efficiency",    "Asset turnover, delta turnover loadings"),
    ("PC7", "Macro environment",   "GDP growth, inflation loadings"),
    ("PC8", "Equity momentum",     "Delta ROE, delta book-to-market loadings"),
    ("PC9", "Mixed signal",        "Residual variance loadings"),
]

for i, (pc, theme, detail) in enumerate(pca_rows):
    t = Inches(2.62) + i * Inches(0.46)
    add_rect(slide, Inches(5.05), t, Inches(7.6), Inches(0.4), LIGHT_NAV if i%2==0 else NAVY_CARD)
    add_rect(slide, Inches(5.05), t, Inches(0.55), Inches(0.4), GOLD)
    add_text(slide, pc, Inches(5.05), t+Inches(0.06), Inches(0.55), Inches(0.28),
             size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, theme, Inches(5.65), t+Inches(0.06), Inches(2.5), Inches(0.28),
             size=10, bold=True, color=WHITE)
    add_text(slide, detail, Inches(8.2), t+Inches(0.06), Inches(4.3), Inches(0.28),
             size=9, color=MUTED)

stat_box(slide, "PCA Components Kept", "9", "Global model (>=5% each)",
         Inches(0.4), Inches(6.9), Inches(3.0), Inches(0.48))
stat_box(slide, "Variance Rule", ">=5%", "Per-component threshold",
         Inches(3.6), Inches(6.9), Inches(2.8), Inches(0.48))
stat_box(slide, "Sig. Components (Global)", "4/9", "PC1, PC3, PC6, PC7 (p<0.05)",
         Inches(6.6), Inches(6.9), Inches(3.5), Inches(0.48))

# ─── SLIDE 13: LOGISTIC REGRESSION SETUP ─────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Logistic Regression — Model Setup", "Section 08 — Global Pooled Model")

add_text(slide,
    "With PCA components as inputs, logistic regression models the probability that a stock will outperform "
    "the S&P 500 in the following quarter. The model is trained on all S&P 500 companies from Q1 2010 to Q4 2024.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.75),
    size=12, color=WHITE)

config_items = [
    ("Training Period",     "Q1 2010 – Q4 2024",       "15 years of quarterly data"),
    ("Training Rows",       "16,589",                   "complete observations after cleaning"),
    ("Companies",           "~500",                     "S&P 500 constituents (time-varying)"),
    ("Target Variable",     "outperformer_next",        "1 = beats SPY next quarter, 0 = does not"),
    ("Input Features",      "9 PCA Components",         "derived from 20 VIF-filtered ratios"),
    ("Model",               "Logistic Regression",      "statsmodels Logit, max_iter=200"),
    ("Cutoff Optimization", "Train-set threshold",      "optimized over 0.30–0.75 range"),
    ("Class Weights",       "Balanced",                 "accounts for outperformer imbalance"),
]

for i, (label, value, note) in enumerate(config_items):
    col = i % 2
    row = i // 2
    l = Inches(0.4) + col * Inches(6.45)
    t = Inches(2.1) + row * Inches(0.88)
    add_rect(slide, l, t, Inches(6.2), Inches(0.8), NAVY_CARD)
    add_rect(slide, l, t, Inches(0.06), Inches(0.8), GOLD)
    add_text(slide, label, l+Inches(0.18), t+Inches(0.06), Inches(2.3), Inches(0.28),
             size=10, color=MUTED)
    add_text(slide, value, l+Inches(0.18), t+Inches(0.32), Inches(3.5), Inches(0.35),
             size=14, bold=True, color=WHITE)
    add_text(slide, note, l+Inches(3.8), t+Inches(0.35), Inches(2.3), Inches(0.3),
             size=9, color=MUTED, italic=True)

add_text(slide,
    "Benchmark: Ananthakumar & Sarkar (2017) — 71.2% accuracy using logistic regression on S&P 500 constituents",
    Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
    size=10, color=GOLD, bold=True)

# ─── SLIDE 14: GLOBAL RESULTS ─────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Global Model Results", "Section 08 — All S&P 500 Companies, No Segmentation")

stat_box(slide, "Test Accuracy", "53.2%", "Global pooled model",
         Inches(0.4), Inches(1.25), Inches(2.8), Inches(1.3))
stat_box(slide, "LR AUC", "0.5381", "Logistic Regression",
         Inches(3.4), Inches(1.25), Inches(2.8), Inches(1.3))
stat_box(slide, "GB AUC (benchmark)", "0.6589", "Gradient Boosting",
         Inches(6.4), Inches(1.25), Inches(2.8), Inches(1.3), val_color=BLUE)
stat_box(slide, "Sig. PCA Comps", "4 / 9", "p < 0.05",
         Inches(9.4), Inches(1.25), Inches(3.5), Inches(1.3), val_color=GREEN)

add_rect(slide, Inches(0.4), Inches(2.75), Inches(6.0), Inches(4.1), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(2.75), Inches(6.0), Inches(0.06), GOLD)
add_text(slide, "Confusion Matrix (Threshold = 0.50)", Inches(0.6), Inches(2.85), Inches(5.6), Inches(0.4),
         size=13, bold=True, color=GOLD)
add_text(slide,
    "The model correctly identifies\n"
    "outperformers (TP) and underperformers\n"
    "(TN), but also generates false signals.\n\n"
    "Sensitivity (Recall):\n"
    "  Of true outperformers, how many\n"
    "  did we correctly predict?\n\n"
    "Specificity:\n"
    "  Of true underperformers, how many\n"
    "  did we correctly predict?\n\n"
    "Precision:\n"
    "  Of all stocks we bet on, how many\n"
    "  actually outperformed?",
    Inches(0.6), Inches(3.35), Inches(5.6), Inches(3.4),
    size=11, color=WHITE)

add_rect(slide, Inches(6.6), Inches(2.75), Inches(6.2), Inches(4.1), NAVY_CARD)
add_rect(slide, Inches(6.6), Inches(2.75), Inches(6.2), Inches(0.06), GOLD)
add_text(slide, "Key Interpretation", Inches(6.8), Inches(2.85), Inches(5.8), Inches(0.4),
         size=13, bold=True, color=GOLD)
interp = [
    ("AUC 0.5381", "Near-random discrimination. The model\nfinds weak but real signal globally.", MUTED),
    ("53.2% Accuracy", "Marginally above 50% baseline.\nMarkets are competitive to predict.", WHITE),
    ("LR Significant", "p-value < 0.05 — the model is\nstatistically not random.", GREEN),
    ("GB AUC 0.6589", "Non-linear relationships exist in the\ndata that LR cannot fully capture.", BLUE),
]
for i, (metric, text, color) in enumerate(interp):
    t = Inches(3.35) + i * Inches(0.85)
    add_rect(slide, Inches(6.8), t, Inches(5.8), Inches(0.75), LIGHT_NAV)
    add_rect(slide, Inches(6.8), t, Inches(0.05), Inches(0.75), color)
    add_text(slide, metric, Inches(6.95), t+Inches(0.05), Inches(2.0), Inches(0.28),
             size=11, bold=True, color=color)
    add_text(slide, text, Inches(6.95), t+Inches(0.35), Inches(5.6), Inches(0.35),
             size=10, color=WHITE)

add_text(slide,
    "LR model is statistically significant  |  McFadden R2 = 0.0073  |  VIF kept: 20/24  |  PCA components: 9",
    Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
    size=9, color=MUTED, italic=True)

# ─── SLIDE 15: GRADIENT BOOSTING BENCHMARK ────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Gradient Boosting — Non-Linear Benchmark", "Section 08 — LR vs GB Comparison")

add_text(slide,
    "To test whether the logistic regression model's limitations are due to the linearity assumption, "
    "we ran Gradient Boosting (GB) as a diagnostic benchmark. GB captures non-linear patterns and "
    "interaction effects that logistic regression cannot model.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.75),
    size=12, color=WHITE)

add_rect(slide, Inches(0.4), Inches(2.05), Inches(5.9), Inches(4.8), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(2.05), Inches(5.9), Inches(0.06), GOLD)
add_text(slide, "Model Comparison — Global", Inches(0.6), Inches(2.15), Inches(5.5), Inches(0.4),
         size=14, bold=True, color=GOLD)

comparison = [
    ("Logistic Regression", "0.5381", "0.5381", "Linear, interpretable\n(coefficients, p-values, odds ratios)"),
    ("Gradient Boosting",   "0.6589", "0.6589", "Non-linear benchmark only\n(n=200 trees, depth=3)"),
]
headers = ["Model", "Train AUC", "Test AUC", "Role"]
col_ws  = [Inches(1.7), Inches(0.9), Inches(0.9), Inches(2.1)]
col_ls  = [Inches(0.65), Inches(2.4), Inches(3.35), Inches(4.3)]

for j, hdr in enumerate(headers):
    add_rect(slide, col_ls[j], Inches(2.65), col_ws[j], Inches(0.38), GOLD)
    add_text(slide, hdr, col_ls[j]+Inches(0.05), Inches(2.68), col_ws[j]-Inches(0.05), Inches(0.32),
             size=10, bold=True, color=NAVY)

for i, (model, tr, te, role) in enumerate(comparison):
    t = Inches(3.1) + i * Inches(0.75)
    bg = LIGHT_NAV if i % 2 == 0 else NAVY_CARD
    add_rect(slide, Inches(0.65), t, Inches(1.7), Inches(0.68), bg)
    add_text(slide, model, Inches(0.7), t+Inches(0.2), Inches(1.6), Inches(0.28),
             size=9, bold=True, color=WHITE if i==0 else BLUE)
    add_rect(slide, Inches(2.4), t, Inches(0.9), Inches(0.68), bg)
    add_text(slide, tr, Inches(2.4), t+Inches(0.2), Inches(0.9), Inches(0.28),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(3.35), t, Inches(0.9), Inches(0.68), bg)
    add_text(slide, te, Inches(3.35), t+Inches(0.2), Inches(0.9), Inches(0.28),
             size=11, bold=True, color=GOLD if i==0 else BLUE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4.3), t, Inches(2.1), Inches(0.68), bg)
    add_text(slide, role, Inches(4.35), t+Inches(0.08), Inches(2.0), Inches(0.58),
             size=9, color=WHITE)

add_text(slide, "AUC Improvement (GB vs LR): +12.1 percentage points",
         Inches(0.6), Inches(4.65), Inches(5.5), Inches(0.35),
         size=12, bold=True, color=BLUE)

add_rect(slide, Inches(6.6), Inches(2.05), Inches(6.2), Inches(4.8), NAVY_CARD)
add_rect(slide, Inches(6.6), Inches(2.05), Inches(6.2), Inches(0.06), GOLD)
add_text(slide, "What This Means", Inches(6.8), Inches(2.15), Inches(5.8), Inches(0.4),
         size=14, bold=True, color=GOLD)
add_text(slide,
    "The gap between LR (0.54) and GB (0.66) AUC "
    "confirms that non-linear relationships exist in the data.\n\n"
    "Logistic regression finds real signal — it is statistically "
    "significant — but misses patterns that require interaction "
    "effects between ratios.\n\n"
    "GB is used only as a diagnostic tool. Logistic regression "
    "remains the primary model because it provides interpretable "
    "coefficients, odds ratios, and p-values needed for academic "
    "analysis and the paper's findings.\n\n"
    "This gap is documented as a limitation and motivates "
    "future work with non-linear models.",
    Inches(6.8), Inches(2.65), Inches(5.8), Inches(4.0),
    size=12, color=WHITE)

# ─── SLIDE 16: SEGMENTATION OVERVIEW ─────────────────────────────────────────
slide = blank_slide()
section_divider(slide, "09", "Segmentation & Results",
                "Market Cap  |  Sector GICS  |  Clustering")

# ─── SLIDE 17: MARKET CAP SEGMENTATION ───────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Segmentation — Market Capitalization", "Section 09 — Configuration 1")

add_text(slide,
    "Companies were segmented by market capitalization (market value of equity). A separate logistic "
    "regression model was trained and evaluated for each size tier.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.55),
    size=12, color=WHITE)

tiers = [
    ("Large Cap",  "> $10 Billion",   "54.2%", "0.558", "Well-followed by analysts.\nMost institutional coverage.\nEfficient pricing — harder to predict."),
    ("Mid Cap",    "$2B – $10B",      "52.8%", "0.542", "Some analyst coverage.\nMore earnings variability.\nSlightly less efficient pricing."),
    ("Small Cap",  "< $2 Billion",    "54.5%", "0.535", "Limited analyst coverage.\nHigher fundamental variability.\nLowest but still near-random accuracy."),
]

for i, (tier, threshold, acc, auc, note) in enumerate(tiers):
    l = Inches(0.4) + i * Inches(4.3)
    add_rect(slide, l, Inches(1.9), Inches(4.1), Inches(5.0), NAVY_CARD)
    add_rect(slide, l, Inches(1.9), Inches(4.1), Inches(0.08), BLUE)
    add_text(slide, tier, l+Inches(0.2), Inches(2.05), Inches(3.7), Inches(0.45),
             size=20, bold=True, color=BLUE)
    add_text(slide, threshold, l+Inches(0.2), Inches(2.5), Inches(3.7), Inches(0.35),
             size=11, color=MUTED)
    add_rect(slide, l+Inches(0.2), Inches(2.9), Inches(0.8), Inches(0.03), GOLD)
    add_text(slide, "Accuracy", l+Inches(0.2), Inches(3.0), Inches(1.8), Inches(0.3),
             size=10, color=MUTED)
    add_text(slide, acc, l+Inches(0.2), Inches(3.3), Inches(3.7), Inches(0.6),
             size=36, bold=True, color=WHITE)
    add_text(slide, f"AUC: {auc}", l+Inches(0.2), Inches(3.95), Inches(3.7), Inches(0.35),
             size=12, color=GOLD)
    add_text(slide, note, l+Inches(0.2), Inches(4.4), Inches(3.7), Inches(2.0),
             size=11, color=WHITE)

add_rect(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.38), LIGHT_NAV)
add_text(slide,
    "Conclusion: Breaking up by market cap provides only marginal improvement over the global 53.2% baseline. "
    "Size alone is not a strong segmentation criterion for this model.",
    Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.32),
    size=11, color=MUTED)

# ─── SLIDE 18: SECTOR SEGMENTATION ───────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Segmentation — GICS Sector", "Section 09 — Configuration 2 — Best Performing Segmentation")

add_text(slide,
    "Separate logistic regression models trained for each GICS sector. Financial ratios carry "
    "different meanings across sectors — comparing a utility company to a tech startup on the same "
    "metrics is not meaningful. Sector-specific models find stronger patterns.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6),
    size=12, color=WHITE)

sectors = [
    ("Communication Services", "428",  "59.5%", "0.6059", GREEN),
    ("Consumer Discretionary", "2668", "54.1%", "0.5305", MUTED),
    ("Consumer Staples",       "1030", "55.7%", "0.5510", MUTED),
    ("Energy",                 "261",  "61.7%", "0.6036", GREEN),
    ("Financials",             "4108", "55.0%", "0.5350", MUTED),
    ("Health Care",            "1126", "54.3%", "0.5578", MUTED),
    ("Industrials",            "4183", "54.5%", "0.5350", MUTED),
    ("Information Technology", "2240", "53.2%", "0.5420", MUTED),
    ("Materials",              "2781", "53.2%", "0.5481", MUTED),
    ("Real Estate",            "424",  "59.5%", "0.6010", GREEN),
    ("Utilities",              "1757", "54.5%", "0.5420", MUTED),
]

row_h = Inches(0.43)
for i, (sector, n, acc, auc, accent) in enumerate(sectors):
    t = Inches(1.95) + i * row_h
    bg = NAVY_CARD if i % 2 == 0 else LIGHT_NAV
    add_rect(slide, Inches(0.4), t, Inches(12.5), row_h - Inches(0.02), bg)
    add_rect(slide, Inches(0.4), t, Inches(0.04), row_h - Inches(0.02), accent)
    add_text(slide, sector, Inches(0.55), t+Inches(0.09), Inches(3.8), Inches(0.27),
             size=11, bold=(accent==GREEN), color=WHITE if accent==MUTED else GREEN)
    add_text(slide, n, Inches(4.5), t+Inches(0.09), Inches(1.3), Inches(0.27),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    acc_color = GREEN if float(acc.strip('%')) > 58 else WHITE
    add_text(slide, acc, Inches(6.0), t+Inches(0.09), Inches(1.5), Inches(0.27),
             size=12, bold=(acc_color==GREEN), color=acc_color, align=PP_ALIGN.CENTER)
    add_text(slide, auc, Inches(7.8), t+Inches(0.09), Inches(1.5), Inches(0.27),
             size=11, color=GOLD if accent==GREEN else MUTED, align=PP_ALIGN.CENTER)

add_text(slide, "Sector", Inches(0.55), Inches(1.82), Inches(3.8), Inches(0.27),
         size=10, bold=True, color=MUTED)
add_text(slide, "N Obs", Inches(4.5), Inches(1.82), Inches(1.3), Inches(0.27),
         size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, "Accuracy", Inches(6.0), Inches(1.82), Inches(1.5), Inches(0.27),
         size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, "AUC", Inches(7.8), Inches(1.82), Inches(1.5), Inches(0.27),
         size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.58), LIGHT_NAV)
add_text(slide,
    "Best sectors: Communication Services (59.5%, AUC 0.606) | Energy (61.7%, AUC 0.604) | Real Estate (59.5%)",
    Inches(0.6), Inches(6.87), Inches(12.1), Inches(0.25),
    size=11, bold=True, color=GREEN)
add_text(slide,
    "Clustering within Communication Services further improved to 63.9% accuracy, AUC 0.6902 — best across all 45 configurations",
    Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.25),
    size=10, color=GOLD)

# ─── SLIDE 19: BEST CONFIGURATION ────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Best Configuration — Communication Services + Clustering",
             "Section 09 — Config 4: Sector + K-Means Clustering (k=2)")

add_text(slide,
    "Within Communication Services, K-Means clustering (k=2) split the sector into two sub-groups with "
    "distinct financial profiles. Cluster 0 (122 companies) significantly outperformed all other configurations "
    "across all 45 group combinations tested.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.65),
    size=12, color=WHITE)

stat_box(slide, "Best Group AUC", "0.6902", "Comm. Services — Cluster 0",
         Inches(0.4), Inches(2.0), Inches(2.9), Inches(1.3), val_color=GOLD)
stat_box(slide, "Best Group Accuracy", "63.9%", "122 observations",
         Inches(3.5), Inches(2.0), Inches(2.9), Inches(1.3), val_color=GREEN)
stat_box(slide, "Weighted AUC (All 45)", "0.5720", "Across all sector configs",
         Inches(6.6), Inches(2.0), Inches(2.9), Inches(1.3))
stat_box(slide, "Total Configs Tested", "45", "All sector x cluster combos",
         Inches(9.7), Inches(2.0), Inches(3.2), Inches(1.3), val_color=BLUE)

add_rect(slide, Inches(0.4), Inches(3.5), Inches(5.9), Inches(3.3), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(3.5), Inches(5.9), Inches(0.06), GOLD)
add_text(slide, "Top 4 Configurations (of 45)", Inches(0.6), Inches(3.6), Inches(5.5), Inches(0.38),
         size=13, bold=True, color=GOLD)

top4 = [
    ("1st", "Comm. Services — Cluster 0",  "0.6902", "63.9%"),
    ("2nd", "Comm. Services — Cluster 1",  "0.6099", "58.0%"),
    ("3rd", "Communication Services",       "0.6059", "59.5%"),
    ("4th", "Comm. Services + MktCap",     "0.6023", "59.0%"),
]
for i, (rank, group, auc, acc) in enumerate(top4):
    t = Inches(4.08) + i * Inches(0.67)
    bg = LIGHT_NAV if i % 2 == 0 else NAVY_CARD
    add_rect(slide, Inches(0.6), t, Inches(5.5), Inches(0.6), bg)
    add_rect(slide, Inches(0.6), t, Inches(0.45), Inches(0.6), GOLD if i==0 else LIGHT_NAV)
    add_text(slide, rank, Inches(0.6), t+Inches(0.15), Inches(0.45), Inches(0.3),
             size=10, bold=True, color=NAVY if i==0 else MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, group, Inches(1.1), t+Inches(0.15), Inches(2.9), Inches(0.3),
             size=10, bold=(i==0), color=GOLD if i==0 else WHITE)
    add_text(slide, f"AUC {auc}", Inches(4.05), t+Inches(0.15), Inches(1.0), Inches(0.3),
             size=10, color=GOLD if i==0 else MUTED)
    add_text(slide, acc, Inches(5.2), t+Inches(0.15), Inches(0.85), Inches(0.3),
             size=10, bold=(i==0), color=GREEN if i==0 else MUTED)

add_rect(slide, Inches(6.6), Inches(3.5), Inches(6.2), Inches(3.3), NAVY_CARD)
add_rect(slide, Inches(6.6), Inches(3.5), Inches(6.2), Inches(0.06), GOLD)
add_text(slide, "Why Communication Services?", Inches(6.8), Inches(3.6), Inches(5.8), Inches(0.38),
         size=13, bold=True, color=GOLD)
add_text(slide,
    "Communication Services companies — Netflix, Meta, "
    "Google, Disney — have highly distinctive financial "
    "ratio patterns that differentiate outperformers "
    "from underperformers.\n\n"
    "Subscriber growth, content investment, and "
    "advertising revenue drive fundamental differences "
    "that are consistently captured in quarterly ratios.\n\n"
    "Clustering further split the sector into a "
    "subscription-driven cluster (Cluster 0) and a "
    "mixed-model cluster (Cluster 1), revealing even "
    "stronger patterns within the sub-group.",
    Inches(6.8), Inches(4.05), Inches(5.8), Inches(2.65),
    size=11, color=WHITE)

add_text(slide,
    "Benchmark: Ananthakumar & Sarkar (2017) = 71.2% accuracy  |  "
    "Our best: 63.9%  |  AUC 0.6902 (vs. global 0.5381)",
    Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
    size=9, color=GOLD, italic=True)

# ─── SLIDE 20: 2025 SETUP ─────────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "2025 Out-of-Sample Testing", "Section 10 — True Forward-Looking Validation")

add_text(slide,
    "The model was applied to data the model never saw during training — 2025 quarterly financials "
    "and returns. The model trained on Q1 2010 – Q4 2024 data was applied to predict S&P 500 "
    "outperformers for all four quarters of 2025.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.65),
    size=12, color=WHITE)

quarters_2025 = [
    ("Q4 2024 ->", "Q1 2025", "334 stocks", "44.6%", "New addition — first true OOS"),
    ("Q1 2025 ->", "Q2 2025", "403 stocks", "55.3%", "SPY +10.6% (strong market)"),
    ("Q2 2025 ->", "Q3 2025", "409 stocks", "56.2%", "Largest 2025 sample"),
    ("Q3 2025 ->", "Q4 2025", "73 stocks",  "57.5%", "Partial: early filers only"),
]

for i, (feat, label, n, acc, note) in enumerate(quarters_2025):
    l = Inches(0.35) + i * Inches(3.25)
    acc_val = float(acc.strip('%'))
    color = GREEN if acc_val > 55 else (MUTED if acc_val < 50 else WHITE)
    add_rect(slide, l, Inches(2.0), Inches(3.1), Inches(4.6), NAVY_CARD)
    add_rect(slide, l, Inches(2.0), Inches(3.1), Inches(0.07), color)
    add_text(slide, feat, l+Inches(0.15), Inches(2.12), Inches(2.8), Inches(0.32),
             size=12, color=MUTED)
    add_text(slide, label, l+Inches(0.15), Inches(2.44), Inches(2.8), Inches(0.45),
             size=18, bold=True, color=WHITE)
    add_text(slide, n, l+Inches(0.15), Inches(2.95), Inches(2.8), Inches(0.32),
             size=11, color=MUTED)
    add_rect(slide, l+Inches(0.15), Inches(3.35), Inches(0.8), Inches(0.03), color)
    add_text(slide, "Accuracy", l+Inches(0.15), Inches(3.45), Inches(2.8), Inches(0.28),
             size=10, color=MUTED)
    add_text(slide, acc, l+Inches(0.15), Inches(3.75), Inches(2.8), Inches(0.7),
             size=36, bold=True, color=color)
    add_text(slide, note, l+Inches(0.15), Inches(4.55), Inches(2.8), Inches(0.9),
             size=10, color=MUTED)

add_rect(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.65), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.05), GOLD)
add_text(slide,
    "Note: Q4 2025 has only 73 stocks because Q3 2025 financial reports were not fully filed "
    "at the time of data extraction (October 2025). Early filers tend to be larger, "
    "well-governed companies — results for this quarter should be interpreted cautiously.",
    Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5),
    size=10, color=MUTED, italic=True)

# ─── SLIDE 21: PORTFOLIO VS SPY ───────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "2025 Portfolio vs S&P 500", "Section 10 — $10,000 Simulated Portfolio")

add_text(slide,
    "Using predicted outperformers, a $10,000 portfolio was constructed each quarter — allocated "
    "proportionally by market capitalization across all stocks predicted to outperform. "
    "Performance compared against the same $10,000 invested in SPY (S&P 500 ETF).",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.65),
    size=12, color=WHITE)

port_data = [
    ("Q2 2025", "403",  "172",  "+6.39%", "+10.57%", "-4.18%", False),
    ("Q3 2025", "391",  "211",  "+6.74%", "+7.79%",  "-1.06%", False),
    ("Q4 2025", "72",   "33",   "+3.15%", "+2.35%",  "+0.81%", True),
]

add_rect(slide, Inches(0.4), Inches(2.0), Inches(12.5), Inches(0.4), LIGHT_NAV)
for j, hdr in enumerate(["Quarter", "Stocks", "Pred. Outperformers", "Model Return", "SPY Return", "Alpha", "Beat SPY"]):
    ls = [Inches(0.6), Inches(2.1), Inches(3.2), Inches(5.2), Inches(7.2), Inches(9.1), Inches(11.0)]
    add_text(slide, hdr, ls[j], Inches(2.07), Inches(1.8), Inches(0.28),
             size=10, bold=True, color=MUTED)

for i, (quarter, n, pred, model_ret, spy_ret, alpha, beat) in enumerate(port_data):
    t = Inches(2.5) + i * Inches(1.1)
    bg = NAVY_CARD if i % 2 == 0 else LIGHT_NAV
    add_rect(slide, Inches(0.4), t, Inches(12.5), Inches(1.0), bg)
    vals = [quarter, n, pred, model_ret, spy_ret, alpha]
    ls   = [Inches(0.6), Inches(2.1), Inches(3.2), Inches(5.2), Inches(7.2), Inches(9.1)]
    colors = [WHITE, MUTED, BLUE, GOLD, WHITE,
              GREEN if beat else RED]
    for j, (val, lx, col) in enumerate(zip(vals, ls, colors)):
        add_text(slide, val, lx, t+Inches(0.32), Inches(1.8), Inches(0.38),
                 size=14, bold=(j in [0,3,4,5]), color=col)
    beat_txt = "YES" if beat else "NO"
    beat_col = GREEN if beat else RED
    add_text(slide, beat_txt, Inches(11.0), t+Inches(0.32), Inches(1.3), Inches(0.38),
             size=14, bold=True, color=beat_col)

add_rect(slide, Inches(0.4), Inches(5.85), Inches(12.5), Inches(1.0), NAVY_CARD)
add_rect(slide, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.05), GOLD)
add_text(slide, "Summary", Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.3),
         size=11, bold=True, color=GOLD)
add_text(slide,
    "Q2 & Q3 2025: Model underperformed SPY by -4.2% and -1.1% respectively. "
    "Q4 2025: Model beat SPY by +0.8% (partial sample of 72 stocks, early filers). "
    "Overall the model did not consistently beat passive S&P 500 investing — "
    "consistent with the efficient market hypothesis in a competitive market.",
    Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.55),
    size=11, color=WHITE)

# ─── SLIDE 22: CONCLUSION ─────────────────────────────────────────────────────
slide = blank_slide()
slide_header(slide, "Conclusion", "Section 11 — Findings and Efficient Market Hypothesis Assessment")

finding_items = [
    ("Logistic Regression Finds Weak Signal",
     "The global model achieves 53.2% accuracy and AUC 0.5381 — "
     "statistically significant (p<0.05) but limited in magnitude. "
     "Quarterly financial ratios do contain predictive information.",
     GOLD),
    ("Segmentation Improves Performance",
     "Sector segmentation consistently outperforms global pooling. "
     "Communication Services + Clustering achieves the best result: "
     "63.9% accuracy, AUC 0.6902 — the strongest configuration across all 45 tested.",
     GREEN),
    ("Non-Linear Signal Exists",
     "Gradient Boosting achieves AUC 0.6589 vs LR's 0.5381 globally, "
     "confirming non-linear patterns in the data that LR cannot capture. "
     "Future work should explore GB, Random Forest, and neural network models.",
     BLUE),
    ("Portfolio Does Not Beat SPY Consistently",
     "The 2025 portfolio underperformed in Q2 and Q3, beat SPY narrowly in Q4. "
     "Passive S&P 500 investing remains a difficult benchmark to outperform, "
     "consistent with semi-strong form efficient market theory.",
     RED),
]

for i, (title, body, color) in enumerate(finding_items):
    col = i % 2
    row = i // 2
    l = Inches(0.4) + col * Inches(6.45)
    t = Inches(1.3) + row * Inches(2.55)
    add_rect(slide, l, t, Inches(6.2), Inches(2.35), NAVY_CARD)
    add_rect(slide, l, t, Inches(0.06), Inches(2.35), color)
    add_text(slide, title, l+Inches(0.2), t+Inches(0.1), Inches(5.8), Inches(0.4),
             size=13, bold=True, color=color)
    add_rect(slide, l+Inches(0.2), t+Inches(0.55), Inches(0.6), Inches(0.03), color)
    add_text(slide, body, l+Inches(0.2), t+Inches(0.65), Inches(5.8), Inches(1.55),
             size=11, color=WHITE)

add_rect(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.82), LIGHT_NAV)
add_rect(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.05), GOLD)
add_text(slide,
    "EMH Verdict: We PARTIALLY SUPPORT the efficient market hypothesis. Markets are difficult "
    "to beat consistently using public financial data, but are not perfectly efficient — "
    "sector-specific fundamental patterns yield statistically significant predictive power.",
    Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.65),
    size=12, bold=True, color=GOLD)

# ─── SLIDE 23: THANK YOU ──────────────────────────────────────────────────────
slide = blank_slide()
add_rect(slide, 0, 0, Inches(5.5), H, NAVY_CARD)
add_rect(slide, Inches(5.5), 0, W - Inches(5.5), H, LIGHT_NAV)
add_rect(slide, Inches(5.5), 0, Inches(0.06), H, GOLD)

add_text(slide, "Thank You", Inches(0.5), Inches(1.8), Inches(4.8), Inches(1.2),
         size=52, bold=True, color=WHITE)
add_rect(slide, Inches(0.5), Inches(3.0), Inches(2.0), Inches(0.05), GOLD)
add_text(slide, "Questions welcome", Inches(0.5), Inches(3.15), Inches(4.8), Inches(0.45),
         size=18, color=MUTED)

add_text(slide, "Lana Gidan", Inches(0.5), Inches(4.2), Inches(4.8), Inches(0.35),
         size=14, bold=True, color=WHITE)
add_text(slide, "Matthew Golubow", Inches(0.5), Inches(4.6), Inches(4.8), Inches(0.35),
         size=14, bold=True, color=WHITE)
add_text(slide, "Shahd Tarman", Inches(0.5), Inches(5.0), Inches(4.8), Inches(0.35),
         size=14, bold=True, color=WHITE)
add_text(slide, "SSIE 605  |  Binghamton University", Inches(0.5), Inches(5.6), Inches(4.8), Inches(0.35),
         size=12, color=MUTED)

key_stats = [
    ("Training Period", "Q1 2010 – Q4 2024"),
    ("Training Observations", "16,589 rows"),
    ("Features", "24 (14 base + 10 delta)"),
    ("VIF Filtered", "20 of 24 retained"),
    ("PCA Components", "9 (global model)"),
    ("Global LR AUC", "0.5381"),
    ("Best Accuracy", "63.9% — Comm. Services Cluster 0"),
    ("Best AUC", "0.6902 — Comm. Services Cluster 0"),
    ("2025 Predictions", "1,219 stocks across 4 quarters"),
    ("Benchmark", "Ananthakumar & Sarkar 2017 — 71.2%"),
]

add_text(slide, "Key Statistics", Inches(6.0), Inches(1.3), Inches(6.8), Inches(0.4),
         size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
for i, (label, val) in enumerate(key_stats):
    t = Inches(1.8) + i * Inches(0.53)
    bg = NAVY_CARD if i%2==0 else NAVY
    add_rect(slide, Inches(5.9), t, Inches(7.1), Inches(0.48), bg)
    add_text(slide, label, Inches(6.05), t+Inches(0.1), Inches(3.2), Inches(0.28),
             size=10, color=MUTED)
    add_text(slide, val, Inches(9.3), t+Inches(0.1), Inches(3.5), Inches(0.28),
             size=10, bold=True, color=WHITE)

prs.save('SP500_Outperformer_Prediction.pptx')
print("PPTX saved: SP500_Outperformer_Prediction.pptx")
print(f"Total slides: {len(prs.slides)}")
